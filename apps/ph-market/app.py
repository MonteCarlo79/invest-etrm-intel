"""Philippines Renewable Energy Investment Advisory — Streamlit app. Port 8510.

Standalone app (does NOT use run_market_app template).
Tabs: Market Structure | Green Energy Auctions | BESS Opportunity |
      Investment Analysis | Investment Advisor | Knowledge Base |
      Data Management | Grid Analysis (PyPSA)
"""
import io
import json
import logging
import os
import sys
import uuid
from datetime import date, timedelta

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")))
from dotenv import load_dotenv
load_dotenv(os.path.join(os.path.dirname(__file__), "..", "..", "config", ".env"))

import anthropic
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import psycopg2
import streamlit as st

st.set_page_config(
    page_title="Philippines RE Investment Advisory",
    page_icon="🇵🇭",
    layout="wide",
    initial_sidebar_state="expanded",
)

from services.ph_knowledge.config import MARKET_CONFIG
from services.ph_knowledge.ingest import run_knowledge_ingest
from services.ph_knowledge.wesm_scraper import WESMPriceScraper, run_wesm_price_scrape
from services.intl_market_common.advanced_retrieval_base import retrieve_for_agent
from services.intl_market_common.expert_memory_base import (
    extract_insights, get_insights, inject_memory, digest_kb_docs,
)

logger = logging.getLogger(__name__)

_ANTHROPIC_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
_client = anthropic.Anthropic(api_key=_ANTHROPIC_KEY)

CFG = MARKET_CONFIG
PREFIX = CFG.table_prefix      # "ph_"
APP_KEY = CFG.app_key          # "ph_market"
CURRENCY = CFG.currency_sym    # "₱"
_USD_PHP = 58.0


# ── DB connection ────────────────────────────────────────────────────────────

@st.cache_resource(ttl=3600)
def _get_conn():
    url = (
        os.environ.get("PGURL")
        or "postgresql://postgres:root@127.0.0.1:5433/marketdata"
    )
    conn = psycopg2.connect(url, connect_timeout=10)
    conn.autocommit = True
    return conn


def _conn():
    conn = _get_conn()
    if conn.closed:
        _get_conn.clear()
        conn = _get_conn()
    return conn


def _query(sql: str, params=None) -> pd.DataFrame:
    return pd.read_sql(sql, _conn(), params=params)


# ── Ensure tables ────────────────────────────────────────────────────────────

def _ensure_tables():
    cur = _conn().cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS marketdata.agent_memory (
            id SERIAL PRIMARY KEY,
            app TEXT NOT NULL DEFAULT 'ph_market',
            category TEXT NOT NULL,
            subject TEXT NOT NULL,
            content TEXT NOT NULL,
            source TEXT DEFAULT 'manual',
            created_at TIMESTAMPTZ DEFAULT NOW(),
            active BOOLEAN DEFAULT TRUE
        )
    """)
    cur.execute("ALTER TABLE marketdata.agent_memory ADD COLUMN IF NOT EXISTS app TEXT DEFAULT 'ph_market'")
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}analyst_sessions (
            session_id TEXT PRIMARY KEY,
            messages   JSONB NOT NULL DEFAULT '[]',
            created_at TIMESTAMPTZ DEFAULT NOW(),
            updated_at TIMESTAMPTZ DEFAULT NOW()
        )
    """)
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}knowledge_docs (
            id              SERIAL PRIMARY KEY,
            source          TEXT NOT NULL,
            doc_type        TEXT NOT NULL,
            title           TEXT,
            url             TEXT UNIQUE,
            published_date  DATE,
            content         TEXT NOT NULL,
            fetched_at      TIMESTAMPTZ NOT NULL DEFAULT NOW(),
            search_vector   TSVECTOR GENERATED ALWAYS AS (
                to_tsvector('english', coalesce(title,'') || ' ' || left(content,100000))
            ) STORED
        )
    """)
    cur.execute(
        f"CREATE INDEX IF NOT EXISTS {PREFIX}knowledge_docs_fts "
        f"ON intl_market.{PREFIX}knowledge_docs USING GIN(search_vector)"
    )
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}expert_insights (
            id SERIAL PRIMARY KEY,
            insight_text TEXT,
            insight_type TEXT,
            confidence TEXT,
            source_session TEXT,
            source_doc_url TEXT,
            active BOOLEAN DEFAULT TRUE,
            validated_at TIMESTAMPTZ DEFAULT NOW()
        )
    """)
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.ph_wesm_prices (
            id           SERIAL PRIMARY KEY,
            trading_date DATE        NOT NULL,
            hour         INTEGER     NOT NULL DEFAULT 0,
            interval_no  INTEGER     NOT NULL DEFAULT 0,
            region       TEXT        NOT NULL,
            node         TEXT,
            price_php_kwh NUMERIC(10,4) NOT NULL,
            price_type   TEXT        NOT NULL DEFAULT 'HSIP',
            fetched_at   TIMESTAMPTZ NOT NULL DEFAULT NOW(),
            CONSTRAINT ph_wesm_prices_uq UNIQUE (trading_date, hour, region, price_type)
        )
    """)
    cur.execute(
        "CREATE INDEX IF NOT EXISTS ph_wesm_prices_date_idx "
        "ON intl_market.ph_wesm_prices (trading_date DESC, region)"
    )
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}report_library (
            id           SERIAL PRIMARY KEY,
            report_name  TEXT NOT NULL,
            frequency    TEXT NOT NULL CHECK (frequency IN ('daily', 'weekly', 'monthly')),
            period       DATE NOT NULL,
            filename     TEXT NOT NULL,
            file_data    BYTEA NOT NULL,
            file_size_kb INTEGER,
            uploaded_at  TIMESTAMPTZ DEFAULT NOW()
        )
    """)
    cur.execute(
        f"CREATE INDEX IF NOT EXISTS {PREFIX}report_library_freq_period "
        f"ON intl_market.{PREFIX}report_library (frequency, period DESC)"
    )
    _conn().commit()


try:
    _ensure_tables()
except Exception as _e:
    logger.warning("_ensure_tables: %s", _e)


# ── Scheduler ────────────────────────────────────────────────────────────────

@st.cache_resource
def _start_scheduler():
    from apscheduler.schedulers.background import BackgroundScheduler

    def _ingest_job():
        try:
            run_knowledge_ingest(verbose=False)
        except Exception as exc:
            logger.error("[ph_scheduler] ingest failed: %s", exc)

    def _digest_job():
        try:
            digest_kb_docs(_ANTHROPIC_KEY, PREFIX, CFG.name, limit=100)
        except Exception as exc:
            logger.error("[ph_scheduler] digest failed: %s", exc)

    def _wesm_price_job():
        try:
            import psycopg2
            conn = psycopg2.connect(
                os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                keepalives=1, keepalives_idle=30,
            )
            conn.autocommit = True
            results = run_wesm_price_scrape(conn, days_back=2)
            conn.close()
            total = sum(v for v in results.values() if v > 0)
            logger.info("[ph_scheduler] WESM price scrape done: %s (%d new rows)", results, total)
        except Exception as exc:
            logger.error("[ph_scheduler] WESM price scrape failed: %s", exc)

    sched = BackgroundScheduler(timezone="Asia/Manila")
    sched.add_job(_ingest_job,      "cron", hour=3,  minute=30, id="ph_ingest",      misfire_grace_time=3600)
    sched.add_job(_digest_job,      "cron", hour=3,  minute=45, id="ph_digest",      misfire_grace_time=3600)
    sched.add_job(_wesm_price_job,  "cron", hour=7,  minute=15, id="ph_wesm_price",  misfire_grace_time=3600)
    sched.start()
    return sched


_start_scheduler()


# ── Cached DB helpers ────────────────────────────────────────────────────────

@st.cache_data(ttl=60)
def _load_memories(app_key: str) -> pd.DataFrame:
    try:
        return _query(
            "SELECT id, category, subject, content, source, created_at "
            "FROM marketdata.agent_memory WHERE app=%s AND active=TRUE ORDER BY created_at DESC",
            (app_key,),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=60)
def _search_knowledge(prefix: str, query: str, limit: int = 10) -> pd.DataFrame:
    try:
        return _query(
            f"SELECT source, doc_type, title, url, published_date, "
            f"left(content,1500) AS snippet, "
            f"ts_rank(search_vector, plainto_tsquery('english',%s)) AS rank "
            f"FROM intl_market.{prefix}knowledge_docs "
            f"WHERE search_vector @@ to_tsquery('english',"
            f"  regexp_replace(plainto_tsquery('english',%s)::text,' & ',' | ','g')) "
            f"ORDER BY rank DESC LIMIT %s",
            (query, query, limit),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=300)
def _knowledge_doc_counts(prefix: str) -> pd.DataFrame:
    try:
        return _query(
            f"SELECT source, doc_type, COUNT(*) AS docs, MAX(fetched_at) AS last_fetch "
            f"FROM intl_market.{prefix}knowledge_docs GROUP BY source, doc_type ORDER BY source"
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=300)
def _ph_table_counts(prefix: str) -> pd.DataFrame:
    rows = []
    for table in [f"{prefix}knowledge_docs", f"{prefix}expert_insights", f"{prefix}analyst_sessions"]:
        try:
            df = _query(f"SELECT COUNT(*) AS n FROM intl_market.{table}")
            rows.append({"Table": f"intl_market.{table}", "Rows": int(df["n"].iloc[0])})
        except Exception:
            rows.append({"Table": f"intl_market.{table}", "Rows": "error"})
    try:
        df2 = _query("SELECT COUNT(*) AS n FROM marketdata.agent_memory WHERE app=%s AND active=TRUE", (APP_KEY,))
        rows.append({"Table": "marketdata.agent_memory", "Rows": int(df2["n"].iloc[0])})
    except Exception:
        rows.append({"Table": "marketdata.agent_memory", "Rows": "error"})
    try:
        df3 = _query("SELECT COUNT(*) AS n FROM intl_market.ph_wesm_prices")
        rows.append({"Table": "intl_market.ph_wesm_prices", "Rows": int(df3["n"].iloc[0])})
    except Exception:
        rows.append({"Table": "intl_market.ph_wesm_prices", "Rows": "error"})
    return pd.DataFrame(rows)


@st.cache_data(ttl=300)
def _wesm_price_history(days: int = 30) -> pd.DataFrame:
    try:
        return _query(
            "SELECT trading_date, region, price_type, "
            "AVG(price_php_kwh) AS avg_price, "
            "MIN(price_php_kwh) AS min_price, "
            "MAX(price_php_kwh) AS max_price, "
            "COUNT(*) AS intervals "
            "FROM intl_market.ph_wesm_prices "
            "WHERE trading_date >= CURRENT_DATE - %s "
            "GROUP BY trading_date, region, price_type "
            "ORDER BY trading_date DESC, region",
            (days,),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=60)
def _wesm_latest_date() -> str:
    try:
        df = _query(
            "SELECT MAX(trading_date) AS latest, COUNT(DISTINCT trading_date) AS n_days "
            "FROM intl_market.ph_wesm_prices"
        )
        if df.empty or df["latest"].iloc[0] is None:
            return "No data yet"
        latest = df["latest"].iloc[0]
        n_days = int(df["n_days"].iloc[0])
        return f"{latest} ({n_days} trading days stored)"
    except Exception:
        return "—"


@st.cache_data(ttl=120)
def _list_library_reports(prefix: str, frequency: str | None = None) -> pd.DataFrame:
    try:
        if frequency:
            return _query(
                f"SELECT id, report_name, frequency, period, filename, file_size_kb, uploaded_at "
                f"FROM intl_market.{prefix}report_library "
                f"WHERE frequency=%s ORDER BY period DESC, uploaded_at DESC",
                (frequency,),
            )
        return _query(
            f"SELECT id, report_name, frequency, period, filename, file_size_kb, uploaded_at "
            f"FROM intl_market.{prefix}report_library ORDER BY period DESC, uploaded_at DESC"
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=600)
def _get_library_report_data(prefix: str, report_id: int) -> bytes | None:
    try:
        cur = _conn().cursor()
        cur.execute(
            f"SELECT file_data FROM intl_market.{prefix}report_library WHERE id=%s",
            (report_id,),
        )
        row = cur.fetchone()
        return bytes(row[0]) if row else None
    except Exception:
        return None


def _save_library_report(report_name: str, frequency: str, period,
                          filename: str, file_data: bytes):
    import psycopg2
    kb = len(file_data) // 1024
    cur = _conn().cursor()
    cur.execute(
        f"INSERT INTO intl_market.{PREFIX}report_library "
        f"(report_name, frequency, period, filename, file_data, file_size_kb) "
        f"VALUES (%s,%s,%s,%s,%s,%s)",
        (report_name, frequency, period, filename, psycopg2.Binary(file_data), kb),
    )
    _list_library_reports.clear()
    _get_library_report_data.clear()


def _delete_library_report(report_id: int):
    cur = _conn().cursor()
    cur.execute(f"DELETE FROM intl_market.{PREFIX}report_library WHERE id=%s", (report_id,))
    _list_library_reports.clear()
    _get_library_report_data.clear()


@st.cache_data(ttl=60)
def _list_recent_sessions(prefix: str, limit: int = 3) -> pd.DataFrame:
    try:
        return _query(
            f"SELECT session_id, jsonb_array_length(messages) AS msg_count, updated_at "
            f"FROM intl_market.{prefix}analyst_sessions "
            f"WHERE jsonb_array_length(messages) > 0 ORDER BY updated_at DESC LIMIT %s",
            (limit,),
        )
    except Exception:
        return pd.DataFrame()


# ── Session persistence ───────────────────────────────────────────────────────

def _save_session(session_id: str, messages: list):
    try:
        cur = _conn().cursor()
        cur.execute(
            f"INSERT INTO intl_market.{PREFIX}analyst_sessions (session_id, messages, updated_at) "
            f"VALUES (%s, %s::jsonb, NOW()) "
            f"ON CONFLICT (session_id) DO UPDATE SET messages=EXCLUDED.messages, updated_at=NOW()",
            (session_id, json.dumps(messages)),
        )
    except Exception as exc:
        logger.debug("_save_session: %s", exc)


def _load_session(session_id: str) -> list:
    try:
        cur = _conn().cursor()
        cur.execute(
            f"SELECT messages FROM intl_market.{PREFIX}analyst_sessions WHERE session_id=%s",
            (session_id,),
        )
        row = cur.fetchone()
        return row[0] if row else []
    except Exception:
        return []


# ── Memory helpers ────────────────────────────────────────────────────────────

def _save_memory(category, subject, content, source="manual"):
    cur = _conn().cursor()
    cur.execute(
        "INSERT INTO marketdata.agent_memory (app, category, subject, content, source) "
        "VALUES (%s,%s,%s,%s,%s)",
        (APP_KEY, category, subject, content, source),
    )
    _load_memories.clear()


def _delete_memory(mem_id: int):
    cur = _conn().cursor()
    cur.execute("UPDATE marketdata.agent_memory SET active=FALSE WHERE id=%s", (mem_id,))


def _get_insight_count() -> int:
    try:
        df = _query(f"SELECT COUNT(*) AS n FROM intl_market.{PREFIX}expert_insights WHERE active=TRUE")
        return int(df.iloc[0]["n"]) if not df.empty else 0
    except Exception:
        return 0


# ── Document ingest helpers ───────────────────────────────────────────────────

def _ingest_url(url: str) -> dict:
    import requests
    from bs4 import BeautifulSoup
    try:
        resp = requests.get(url, timeout=30, headers={"User-Agent": "BESSPlatformBot/1.0"})
        resp.raise_for_status()
    except Exception as exc:
        return {"status": "error", "msg": f"Fetch failed: {exc}"}
    soup = BeautifulSoup(resp.text, "html.parser")
    title_el = soup.find("h1") or soup.find("title")
    title = title_el.get_text(" ", strip=True) if title_el else url
    for tag in soup.find_all(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    content = "\n\n".join(
        el.get_text(" ", strip=True) for el in soup.find_all(["p", "h1", "h2", "h3", "li"])
        if el.get_text(" ", strip=True)
    ) or soup.get_text(" ", strip=True)
    if not content.strip():
        return {"status": "error", "msg": "No text extracted."}
    try:
        cur = _conn().cursor()
        cur.execute(
            f"INSERT INTO intl_market.{PREFIX}knowledge_docs "
            f"(source, doc_type, title, url, published_date, content) VALUES (%s,%s,%s,%s,%s,%s) "
            f"ON CONFLICT (url) DO UPDATE SET content=EXCLUDED.content, title=EXCLUDED.title, fetched_at=NOW()",
            ("upload", "article", title, url, date.today(), content),
        )
        return {"status": "success", "msg": f"Ingested '{title}' ({len(content):,} chars)"}
    except Exception as exc:
        return {"status": "error", "msg": f"DB insert failed: {exc}"}


def _ingest_uploaded_file(filename: str, data: bytes) -> dict:
    ext = filename.rsplit(".", 1)[-1].lower()
    try:
        if ext == "txt":
            content = data.decode("utf-8", errors="replace")
        elif ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            content = "\n\n".join(p.extract_text() for p in reader.pages if p.extract_text())
        elif ext in ("xlsx", "xls"):
            xl = pd.ExcelFile(io.BytesIO(data))
            content = "\n\n".join(
                f"Sheet: {s}\n{xl.parse(s).to_string(index=False)}" for s in xl.sheet_names
            )
        elif ext in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            content = "\n\n".join(slides)
        elif ext in ("docx", "doc"):
            from docx import Document
            doc = Document(io.BytesIO(data))
            content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        else:
            return {"status": "error", "msg": f"Unsupported type: .{ext}"}
    except Exception as exc:
        return {"status": "error", "msg": f"Extraction failed: {exc}"}
    if not content.strip():
        return {"status": "error", "msg": "No text extracted."}
    doc_type = {"pdf": "pdf", "txt": "text", "xlsx": "excel", "xls": "excel",
                "pptx": "presentation", "ppt": "presentation", "docx": "word"}.get(ext, "document")
    try:
        cur = _conn().cursor()
        cur.execute(
            f"INSERT INTO intl_market.{PREFIX}knowledge_docs "
            f"(source, doc_type, title, url, published_date, content) VALUES (%s,%s,%s,%s,%s,%s) "
            f"ON CONFLICT (url) DO UPDATE SET content=EXCLUDED.content, fetched_at=NOW()",
            ("upload", doc_type, filename, f"upload://{filename}", date.today(), content),
        )
        return {"status": "success", "msg": f"Ingested '{filename}' ({len(content):,} chars)"}
    except Exception as exc:
        return {"status": "error", "msg": f"DB insert failed: {exc}"}


# ── IRR model ────────────────────────────────────────────────────────────────

_TECH_PRESETS = {
    "solar":        {"capex_usd_kw": 720,  "cf_pct": 20.0, "om_pct": 1.5, "degradation": 0.005, "life": 25, "label": "Solar PV"},
    "onshore_wind": {"capex_usd_kw": 1550, "cf_pct": 31.0, "om_pct": 2.0, "degradation": 0.000, "life": 25, "label": "Onshore Wind"},
    "bess_2h":      {"capex_usd_kwh": 300, "duration_h": 2, "om_pct": 2.0, "degradation": 0.020, "life": 15, "label": "BESS 2h"},
    "bess_4h":      {"capex_usd_kwh": 280, "duration_h": 4, "om_pct": 2.0, "degradation": 0.020, "life": 15, "label": "BESS 4h"},
    "iress":        {"capex_usd_kw": 1000, "cf_pct": 19.0, "om_pct": 2.0, "degradation": 0.005, "life": 25, "label": "IRESS (Solar+Storage)"},
}


def _compute_irr(cashflows: list) -> float:
    rate = 0.10
    for _ in range(200):
        npv  = sum(cf / (1 + rate) ** t for t, cf in enumerate(cashflows))
        dnpv = sum(-t * cf / (1 + rate) ** (t + 1) for t, cf in enumerate(cashflows))
        if abs(dnpv) < 1e-10:
            break
        rate -= npv / dnpv
        if rate <= -0.999:
            rate = -0.999
    return rate


def _run_irr_model(
    technology: str,
    capacity_mw: float,
    capex_usd_per_kw: float | None = None,
    revenue_php_per_kwh: float | None = None,
    capacity_factor_pct: float | None = None,
    wacc_pct: float = 10.0,
    project_life_yrs: int | None = None,
    leverage_pct: float = 60.0,
    cost_of_debt_pct: float = 7.0,
) -> dict:
    p = _TECH_PRESETS.get(technology, _TECH_PRESETS["solar"])
    is_bess = technology.startswith("bess")

    # CAPEX
    if is_bess:
        duration_h = p.get("duration_h", 2)
        if capex_usd_per_kw:
            capex_usd_kw = capex_usd_per_kw
        else:
            capex_usd_kw = p["capex_usd_kwh"] * duration_h
        energy_mwh = capacity_mw * duration_h
        capex_total_usd = capacity_mw * capex_usd_kw * 1000
    else:
        capex_usd_kw = capex_usd_per_kw or p["capex_usd_kw"]
        capex_total_usd = capacity_mw * capex_usd_kw * 1000
        duration_h = 0

    capex_total_php = capex_total_usd * _USD_PHP
    life = project_life_yrs or p["life"]
    om_annual_php = capex_total_php * (p["om_pct"] / 100)
    degradation = p.get("degradation", 0.0)

    # Revenue
    if is_bess:
        # BESS: revenue from reserves + arbitrage — use PHP/kWh on energy throughput
        rev_php_kwh = revenue_php_per_kwh or 2.50  # default arbitrage spread
        cycles_per_year = 300
        annual_gen_mwh = capacity_mw * duration_h * cycles_per_year
    else:
        cf = (capacity_factor_pct or p["cf_pct"]) / 100
        annual_gen_mwh = capacity_mw * 8760 * cf
        rev_php_kwh = revenue_php_per_kwh or 6.00  # default GET tariff

    annual_rev_php_yr1 = annual_gen_mwh * rev_php_kwh * 1000  # MWh → kWh

    # Cashflows (unlevered)
    cashflows_unlev = [-capex_total_php]
    for yr in range(1, life + 1):
        rev = annual_rev_php_yr1 * (1 - degradation) ** (yr - 1)
        cashflows_unlev.append(rev - om_annual_php)

    unlevered_irr = _compute_irr(cashflows_unlev)
    npv_wacc = sum(cf / (1 + wacc_pct / 100) ** t for t, cf in enumerate(cashflows_unlev))

    # LCOE (PHP/kWh)
    total_gen_kwh = sum(
        annual_gen_mwh * 1000 * (1 - degradation) ** (yr - 1)
        for yr in range(1, life + 1)
    )
    total_costs_pv = sum(
        om_annual_php / (1 + wacc_pct / 100) ** yr for yr in range(1, life + 1)
    ) + capex_total_php
    lcoe = total_costs_pv / max(total_gen_kwh, 1)

    # Levered (equity) IRR
    debt = capex_total_php * (leverage_pct / 100)
    equity = capex_total_php - debt
    debt_service = debt * (cost_of_debt_pct / 100) / (1 - (1 + cost_of_debt_pct / 100) ** (-min(life, 15)))
    cashflows_eq = [-equity]
    for yr in range(1, life + 1):
        rev = annual_rev_php_yr1 * (1 - degradation) ** (yr - 1)
        ds = debt_service if yr <= min(life, 15) else 0
        cashflows_eq.append(rev - om_annual_php - ds)
    equity_irr = _compute_irr(cashflows_eq)

    # 3×3 sensitivity (CAPEX × Revenue)
    sensitivity = []
    for cm in [0.8, 1.0, 1.2]:
        for rm in [0.8, 1.0, 1.2]:
            cfs = [-capex_total_php * cm]
            for yr in range(1, life + 1):
                rev = annual_rev_php_yr1 * rm * (1 - degradation) ** (yr - 1)
                cfs.append(rev - om_annual_php)
            sensitivity.append({
                "capex": f"{cm:.0%}", "revenue": f"{rm:.0%}",
                "unlevered_irr": f"{_compute_irr(cfs) * 100:.1f}%",
            })

    return {
        "technology": p["label"],
        "capacity_mw": capacity_mw,
        "capex_usd_per_kw": round(capex_usd_kw, 0),
        "capex_total_php_m": round(capex_total_php / 1e6, 1),
        "annual_gen_mwh": round(annual_gen_mwh, 0),
        "revenue_php_per_kwh": rev_php_kwh,
        "unlevered_irr_pct": round(unlevered_irr * 100, 1),
        "equity_irr_pct": round(equity_irr * 100, 1),
        "lcoe_php_per_kwh": round(lcoe, 4),
        "npv_at_wacc_php_m": round(npv_wacc / 1e6, 1),
        "sensitivity": sensitivity,
    }


# ── Agent tools ───────────────────────────────────────────────────────────────

_GEAP_DATA = {
    "GEA-1": {"round": "GEA-1", "year": 2022, "technology": "Solar + Wind", "target_mw": None,
               "awarded_mw": 1967, "get_price_php_kwh": 5.80, "notes": "First GEAP round; heavily oversubscribed"},
    "GEA-2": {"round": "GEA-2", "year": 2023, "technology": "Solar + Wind + IRESS",
               "target_mw": 11_000, "awarded_mw": 10_653, "get_price_php_kwh": 6.10,
               "notes": "Largest round; IRESS included for first time"},
    "GEA-3": {"round": "GEA-3", "year": 2024, "technology": "Offshore Wind",
               "target_mw": 2000, "awarded_mw": None, "get_price_php_kwh": None,
               "notes": "Offshore wind dedicated; results pending"},
    "GEA-4": {"round": "GEA-4", "year": 2024, "technology": "Solar + Wind + IRESS + BESS",
               "target_mw": 3441, "awarded_mw": None, "get_price_php_kwh": None,
               "notes": "Active round; BESS included for first time"},
    "GEA-5": {"round": "GEA-5", "year": 2025, "technology": "Offshore Wind",
               "target_mw": 3300, "awarded_mw": None, "get_price_php_kwh": None,
               "notes": "3,300 MW offshore wind; South China Sea focus"},
}

_WESM_PRICES = {
    "Luzon":    {"2025": 4.26, "2030": 5.40, "2040": 6.80, "2060": 8.25, "unit": "PHP/kWh"},
    "Visayas":  {"2025": 4.90, "2030": 6.10, "2040": 7.50, "2060": 9.10, "unit": "PHP/kWh"},
    "Mindanao": {"2025": 3.80, "2030": 4.90, "2040": 6.20, "2060": 7.80, "unit": "PHP/kWh"},
}

_AS_CONTEXT = {
    "framework": "NGCP procures ancillary services under the Grid Code and Market Manual",
    "reserve_types": {
        "Regulating Reserve": "±20 MW response within 1 sec; BESS premium service; ASPA firm contracts",
        "Contingency Reserve": "Restore frequency within 10 min post-contingency; 10-MW minimum bid",
        "Dispatchable Reserve": "30-min response; capacity market basis; manual dispatch by NGCP",
    },
    "reserve_market_start": "January 2024",
    "aspa": "Ancillary Services Procurement Agreement — firm bilateral contract with NGCP for regulating/contingency reserves",
    "key_holder": "GNPOWER (GNPower Kauswagan) won first ASPA for BESS regulating reserves",
    "bess_revenue_stack": [
        "Regulating reserves (highest value, ~PHP 800-1,200/MW/h)",
        "Contingency reserves (~PHP 300-600/MW/h)",
        "WESM energy arbitrage (charge off-peak, discharge peak, ~PHP 1.50-3.00/kWh spread)",
        "IRESS GEAP tariff (if co-located with solar; ~PHP 6.00/kWh GET)",
    ],
    "caper": "IEMOP Capacity and Energy Reserve programme — BESS can register as capacity resource",
}

_MARKET_STRUCTURE = {
    "peak_demand_gw": 19.1,
    "installed_capacity_mw": 29_962,
    "demand_growth_pct": 8.1,
    "gdp_growth_pct": 5.6,
    "grids": {
        "Luzon":    {"share_pct": 73, "capacity_mw": 21_872},
        "Visayas":  {"share_pct": 12, "capacity_mw": 3_595},
        "Mindanao": {"share_pct": 15, "capacity_mw": 4_494},
    },
    "generation_mix": {
        "Coal": 59, "Natural Gas": 12, "Geothermal": 6,
        "Hydro": 8, "Solar": 7, "Wind": 4, "Oil/Others": 4,
    },
    "key_players": [
        "AES Philippines — Masinloc coal + IRESS Alaminos",
        "ACEN (Ayala) — solar, wind, geothermal; RE-heavy",
        "SMC Global Power — coal + gas (Ilijan, San Roque)",
        "Aboitiz Power — multi-fuel; Therma group + SNAP hydro",
        "GNPOWER — coal (Kauswagan) + BESS ancillary services",
        "Total Eren / TotalEnergies — offshore wind development",
        "Vena Energy — solar pipeline, Luzon focus",
    ],
    "routes_to_market": {
        "WESM Spot":        {"range_php_kwh": "4.26–8.25", "approval": "None", "term": "Day-ahead"},
        "PSA with DU/EC":   {"range_php_kwh": "5.19–6.52", "approval": "ERC (1-2yr)", "term": "10-25yr"},
        "Retail/GEOP":      {"range_php_kwh": "4.10–7.97", "approval": "None", "term": "1-5yr"},
        "GEAP (COE-GET)":   {"range_php_kwh": "~6.00 avg", "approval": "DOE", "term": "20yr government-backed"},
        "FiT":              {"range_php_kwh": "8.84–11.14", "approval": "ERC", "term": "FULLY SUBSCRIBED 2019"},
    },
}

_POLICY_SNAPSHOT = {
    "epira": "Electric Power Industry Reform Act (2001) — unbundled generation, transmission, distribution; created WESM, IEMOP, ERC",
    "re_act": "Renewable Energy Act of 2008 (RA 9513) — FiT, RPS, Net Metering, 100% foreign ownership for RE developers",
    "foreign_ownership": {
        "RE development": "100% foreign ownership allowed under RE Act (exception to 40% FDI cap)",
        "NGCP (transmission)": "Max 40% foreign ownership",
        "Distribution utilities": "Max 40% FDI cap applies",
    },
    "pep_2023_2050": {
        "name": "Philippine Energy Plan 2023-2050",
        "re_targets": {"2030": "35% RE share", "2040": "50% RE share", "2050": "net-zero pathway"},
        "capacity_targets_gw": {"2030": 55, "2040": 100, "2050": ">150"},
    },
    "geap_rules": {
        "mechanism": "COE-GET (Certificate of Endorsement — Green Energy Tariff)",
        "price_discovery": "Competitive auction; GET price = goal-seek equity IRR approach",
        "erc_approval": "ERC approves GET before commercial operations",
        "contract_term": "20 years with DU/EC offtaker",
        "eligibility": "Solar, Wind, Biomass, Run-of-River Hydro, Geothermal, IRESS, BESS (GEA-4+)",
    },
    "rps": "Renewable Portfolio Standard — DUs must source minimum % from RE (1% initially, rising to 2.52% by 2030)",
    "net_metering": "Prosumers can sell excess RE back to grid at avoided cost",
    "key_risks": [
        "Grid congestion and curtailment (especially Mindanao MRU dispatch priority)",
        "ERC approval delays for PSA contracts (1-2 year timeline)",
        "Permitting: Environmental Compliance Certificate (ECC) 6-18 months",
        "Transmission interconnection — NGCP TDP backlog; connection agreements",
        "Foreign land ownership restriction (workaround: long-term lease)",
    ],
}


def _dispatch_tool(name: str, inputs: dict) -> str:
    try:
        if name == "search_knowledge_base":
            try:
                return retrieve_for_agent(inputs["query"], _ANTHROPIC_KEY, CFG,
                                          sources=inputs.get("sources") or None, top_k=6)
            except Exception:
                results = _search_knowledge(PREFIX, inputs["query"], limit=6)
                if results.empty:
                    return "No matching knowledge documents found."
                return "\n\n---\n\n".join(
                    f"[{r['source']}] {r['title']} ({r['published_date']})\n{r['snippet']}"
                    for _, r in results.iterrows()
                )

        elif name == "get_geap_data":
            rnd = inputs.get("round")
            if rnd and rnd in _GEAP_DATA:
                return json.dumps(_GEAP_DATA[rnd], indent=2)
            return json.dumps(list(_GEAP_DATA.values()), indent=2)

        elif name == "get_wesm_price_context":
            grid = inputs.get("grid")
            if grid and grid in _WESM_PRICES:
                return json.dumps({grid: _WESM_PRICES[grid]}, indent=2)
            return json.dumps({
                "prices": _WESM_PRICES,
                "routes_to_market": _MARKET_STRUCTURE["routes_to_market"],
                "source": "AFRY AIMR 2024Q4",
            }, indent=2)

        elif name == "get_ancillary_services_context":
            return json.dumps(_AS_CONTEXT, indent=2)

        elif name == "estimate_re_irr":
            result = _run_irr_model(
                technology=inputs.get("technology", "solar"),
                capacity_mw=float(inputs.get("capacity_mw", 50)),
                capex_usd_per_kw=inputs.get("capex_usd_per_kw"),
                revenue_php_per_kwh=inputs.get("revenue_php_per_kwh"),
                capacity_factor_pct=inputs.get("capacity_factor_pct"),
                wacc_pct=float(inputs.get("wacc_pct", 10.0)),
                project_life_yrs=inputs.get("project_life_yrs"),
                leverage_pct=float(inputs.get("leverage_pct", 60.0)),
                cost_of_debt_pct=float(inputs.get("cost_of_debt_pct", 7.0)),
            )
            return json.dumps(result, indent=2)

        elif name == "get_market_structure":
            topic = inputs.get("topic", "").lower()
            if "mix" in topic or "generation" in topic:
                return json.dumps(_MARKET_STRUCTURE["generation_mix"], indent=2)
            if "player" in topic or "developer" in topic:
                return json.dumps(_MARKET_STRUCTURE["key_players"], indent=2)
            if "route" in topic or "market" in topic:
                return json.dumps(_MARKET_STRUCTURE["routes_to_market"], indent=2)
            return json.dumps(_MARKET_STRUCTURE, indent=2)

        elif name == "get_policy_snapshot":
            return json.dumps(_POLICY_SNAPSHOT, indent=2)

    except Exception as exc:
        return f"Tool error: {exc}"
    return "Unknown tool"


_TOOLS = [
    {
        "name": "search_knowledge_base",
        "description": "Semantic search (HyDE + FTS + rerank) over Philippines market reports and documents.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string"},
                "sources": {"type": "array", "items": {"type": "string"}},
            },
            "required": ["query"],
        },
    },
    {
        "name": "get_geap_data",
        "description": "Returns Green Energy Auction Program (GEAP) data — rounds GEA-1 through GEA-5, target vs awarded capacity, GET pricing, technologies.",
        "input_schema": {
            "type": "object",
            "properties": {"round": {"type": "string"}},
            "required": [],
        },
    },
    {
        "name": "get_wesm_price_context",
        "description": "WESM spot price projections and routes-to-market price ranges (PHP/kWh) by grid (Luzon/Visayas/Mindanao) from AFRY AIMR 2024Q4.",
        "input_schema": {
            "type": "object",
            "properties": {"grid": {"type": "string"}},
            "required": [],
        },
    },
    {
        "name": "get_ancillary_services_context",
        "description": "Philippines ancillary services market — NGCP reserve types, ASPA procurement, reserve market mechanics, BESS revenue stack.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "estimate_re_irr",
        "description": (
            "Parametric IRR/LCOE model for Philippine RE projects (Solar, Onshore Wind, BESS, IRESS). "
            "Returns unlevered IRR, equity IRR, LCOE, NPV, sensitivity table."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "technology":           {"type": "string", "enum": ["solar", "onshore_wind", "bess_2h", "bess_4h", "iress"]},
                "capacity_mw":          {"type": "number"},
                "capex_usd_per_kw":     {"type": "number"},
                "revenue_php_per_kwh":  {"type": "number"},
                "capacity_factor_pct":  {"type": "number"},
                "wacc_pct":             {"type": "number"},
                "project_life_yrs":     {"type": "integer"},
                "leverage_pct":         {"type": "number"},
                "cost_of_debt_pct":     {"type": "number"},
            },
            "required": ["technology", "capacity_mw"],
        },
    },
    {
        "name": "get_market_structure",
        "description": "Philippines power market structure — installed capacity, demand, generation mix, key players, market participants.",
        "input_schema": {
            "type": "object",
            "properties": {"topic": {"type": "string"}},
            "required": [],
        },
    },
    {
        "name": "get_policy_snapshot",
        "description": "Philippines energy policy and regulatory landscape — EPIRA, RE Act, foreign ownership rules, DOE PEP 2023-2050, GEAP rules.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
]


def _build_system(query: str = "") -> str:
    base = """\
You are a senior Philippines Renewable Power Investment Expert at a global infrastructure fund.

GROUNDING RULE: For specific current prices, project data, and recent developments → use your tools.
For regulatory framework, market mechanics, and historical context → use embedded knowledge below.

MARKET CONTEXT:
- WESM: Wholesale Electricity Spot Market, operated by IEMOP/PEMC
- Three grids: Luzon (73% of capacity), Visayas (12%), Mindanao (15%)
- Peak demand: 19.1 GW (2024); demand growth 8.1%/yr; GDP growth 5.6%
- Generation mix: coal 56-63%, geothermal 6% (highest utilisation at 66% CF), renewables 23%
- RE target: 35% by 2030, 50% by 2040 (DOE Philippine Energy Plan 2023-2050)

ROUTES TO MARKET:
- WESM Spot: 4.26–8.25 PHP/kWh (2025–2060); wind has priority dispatch
- PSA with DU/EC: 5.19–6.52 PHP/kWh; ERC-regulated; 1-2yr approval timeline
- Retail/GEOP: 4.10–7.97 PHP/kWh; no ERC approval; short contracts
- GEAP: ~6.00 PHP/kWh avg for wind (GEA-1/2); 20-year COE-GET; government-backed
- FiT: 8.84–11.14 PHP/kWh; FULLY SUBSCRIBED since 2019

BESS OPPORTUNITY:
- Reserve market started January 2024 (regulating, contingency, dispatchable)
- NGCP procures via ASPA (firm contracts) + real-time reserve market
- BESS excels at regulating reserves (fast response premium)
- Revenue stack: regulating reserves + contingency reserves + WESM arbitrage + IRESS GEAP tariff

FOREIGN OWNERSHIP:
- RE Act allows 100% foreign ownership for RE project developers (exception to 40% FDI cap)
- NGCP (transmission) is 40% foreign-owned max
- Distribution: 40% FDI cap applies

KEY RISKS:
- Grid congestion / curtailment especially Mindanao MRU dispatch
- ERC PSA approval 1-2 year timeline
- Environmental Compliance Certificate (ECC) 6-18 months
- NGCP transmission interconnection backlog
- Foreign land ownership restriction (long-term lease workaround)
"""
    if query:
        try:
            insights = get_insights(query, PREFIX, limit=5)
            mem_block = inject_memory(insights, CFG.name)
            if mem_block:
                base += f"\n\n{mem_block}"
        except Exception:
            pass
    mems = _load_memories(APP_KEY)
    if not mems.empty:
        lines = "\n".join(f"- [{r.category}] {r.subject}: {r.content}" for r in mems.itertuples())
        base += f"\n\n## Analyst notes from prior sessions:\n{lines}"
    return base


def _run_agent_turn(messages: list, system: str) -> tuple[str, list, list]:
    tool_events = []
    while True:
        resp = _client.messages.create(
            model="claude-sonnet-4-6", max_tokens=4096,
            system=system, tools=_TOOLS, messages=messages,
        )
        messages = messages + [{"role": "assistant", "content": resp.content}]
        if resp.stop_reason == "end_turn":
            text = next((b.text for b in resp.content if hasattr(b, "text")), "")
            return text, messages, tool_events
        tool_results = []
        for block in resp.content:
            if block.type == "tool_use":
                result_str = _dispatch_tool(block.name, block.input)
                tool_events.append({"tool": block.name, "result": result_str[:200]})
                tool_results.append({"type": "tool_result", "tool_use_id": block.id, "content": result_str})
        messages = messages + [{"role": "user", "content": tool_results}]


# ── UI ────────────────────────────────────────────────────────────────────────

with st.sidebar:
    st.title("🇵🇭 Philippines RE Investment")
    st.caption("Investment advisory · WESM · GEAP · BESS")
    st.divider()
    n_insights = _get_insight_count()
    st.metric("Expert Insights", n_insights)
    try:
        kb_ct = _knowledge_doc_counts(PREFIX)
        total_docs = int(kb_ct["docs"].sum()) if not kb_ct.empty else 0
        st.metric("KB Documents", total_docs)
    except Exception:
        pass
    st.divider()
    st.caption("Port 8510 · Asia/Manila · ap-southeast-1")

(
    tab_mkt, tab_geap, tab_bess, tab_irr,
    tab_advisor, tab_kb, tab_mgmt, tab_pypsa, tab_library,
) = st.tabs([
    "Market Structure", "Green Energy Auctions", "BESS Opportunity",
    "Investment Analysis", "Investment Advisor",
    "Knowledge Base", "Data Management", "Grid Analysis", "Library",
])


# ═══════════════════════════════════════════════════════════════
# Tab 1 — Market Structure
# ═══════════════════════════════════════════════════════════════
with tab_mkt:
    st.header("Philippines Power Market Structure")
    st.caption("Source: DOE, IEMOP, NGCP — 2024 data")

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Peak Demand",       "19.1 GW")
    c2.metric("Installed Capacity", "29,962 MW")
    c3.metric("Demand Growth",     "8.1% / yr")
    c4.metric("GDP Growth",        "5.6% / yr")

    st.divider()
    col_grid, col_mix = st.columns(2)

    with col_grid:
        st.subheader("Grid Capacity Split")
        grids_df = pd.DataFrame([
            {"Grid": k, "Capacity (MW)": v["capacity_mw"], "Share (%)": v["share_pct"]}
            for k, v in _MARKET_STRUCTURE["grids"].items()
        ])
        fig_grid = px.bar(
            grids_df, x="Grid", y="Capacity (MW)", color="Grid",
            text="Share (%)", color_discrete_sequence=["#1f77b4", "#ff7f0e", "#2ca02c"],
        )
        fig_grid.update_traces(texttemplate="%{text}%", textposition="outside")
        fig_grid.update_layout(height=320, showlegend=False, margin=dict(l=0, r=0, t=0, b=0))
        st.plotly_chart(fig_grid, use_container_width=True)

    with col_mix:
        st.subheader("Generation Mix (2024)")
        mix_df = pd.DataFrame([
            {"Technology": k, "Share (%)": v}
            for k, v in _MARKET_STRUCTURE["generation_mix"].items()
        ])
        fig_mix = px.pie(
            mix_df, names="Technology", values="Share (%)",
            color_discrete_sequence=px.colors.qualitative.Set2,
        )
        fig_mix.update_traces(textposition="inside", textinfo="percent+label")
        fig_mix.update_layout(height=320, showlegend=False, margin=dict(l=0, r=0, t=10, b=0))
        st.plotly_chart(fig_mix, use_container_width=True)

    st.subheader("Routes to Market")
    rtm = _MARKET_STRUCTURE["routes_to_market"]
    rtm_df = pd.DataFrame([
        {"Route": k, "Price Range (PHP/kWh)": v["range_php_kwh"],
         "Regulatory Approval": v["approval"], "Contract Term": v["term"]}
        for k, v in rtm.items()
    ])
    st.dataframe(rtm_df, use_container_width=True, hide_index=True)

    st.subheader("Key Market Participants")
    for player in _MARKET_STRUCTURE["key_players"]:
        st.markdown(f"- {player}")


# ═══════════════════════════════════════════════════════════════
# Tab 2 — Green Energy Auctions
# ═══════════════════════════════════════════════════════════════
with tab_geap:
    st.header("Green Energy Auction Program (GEAP)")
    st.caption("DOE-administered competitive auctions · COE-GET mechanism · 20-year government-backed tariff")

    ga1, ga2 = st.columns([3, 2])
    with ga1:
        st.subheader("Rounds Overview")
        geap_df = pd.DataFrame(_GEAP_DATA.values())
        geap_show = geap_df[["round", "year", "technology", "target_mw", "awarded_mw", "get_price_php_kwh", "notes"]].copy()
        geap_show.columns = ["Round", "Year", "Technology", "Target MW", "Awarded MW", "GET Price (PHP/kWh)", "Notes"]
        st.dataframe(geap_show, use_container_width=True, hide_index=True)

    with ga2:
        st.subheader("Target vs Awarded (MW)")
        plot_df = pd.DataFrame([
            {"Round": k, "Type": "Awarded", "MW": v["awarded_mw"]}
            for k, v in _GEAP_DATA.items() if v["awarded_mw"]
        ] + [
            {"Round": k, "Type": "Target", "MW": v["target_mw"]}
            for k, v in _GEAP_DATA.items() if v["target_mw"]
        ])
        if not plot_df.empty:
            fig_geap = px.bar(
                plot_df, x="Round", y="MW", color="Type", barmode="group",
                color_discrete_map={"Target": "#aec7e8", "Awarded": "#1f77b4"},
            )
            fig_geap.update_layout(height=300, margin=dict(l=0, r=0, t=0, b=0))
            st.plotly_chart(fig_geap, use_container_width=True)

    st.divider()
    st.subheader("GET Price Comparison (Awarded Rounds)")
    price_df = pd.DataFrame([
        {"Round": k, "GET Price (PHP/kWh)": v["get_price_php_kwh"]}
        for k, v in _GEAP_DATA.items() if v["get_price_php_kwh"]
    ])
    if not price_df.empty:
        fig_price = px.bar(
            price_df, x="Round", y="GET Price (PHP/kWh)",
            text="GET Price (PHP/kWh)", color_discrete_sequence=["#2ca02c"],
        )
        fig_price.update_traces(texttemplate="%{text:.2f}", textposition="outside")
        fig_price.update_layout(height=280, margin=dict(l=0, r=0, t=0, b=0))
        st.plotly_chart(fig_price, use_container_width=True)

    st.subheader("COE-GET Mechanism")
    st.markdown("""
**How GET pricing works:**
1. DOE issues Notice of Auction specifying technology, capacity target, and ceiling GET price
2. Developers submit bids (GET price = price they will accept for 20 years)
3. Lowest-price bids win; DOE issues Certificate of Endorsement (COE)
4. ERC approves GET before commercial operations
5. Offtakers (distribution utilities / electric cooperatives) must purchase output at GET price

**GET price = goal-seek equity IRR:** Developers back-calculate the GET price needed to achieve target equity IRR (typically 12-15%), given CAPEX, leverage, O&M, and capacity factor assumptions.

**IRESS inclusion (GEA-2+):** Intermittent RE with Storage Systems — solar or wind co-located with BESS; higher GET price allowed (~PHP 6.50-7.00/kWh) to account for BESS CAPEX.
""")


# ═══════════════════════════════════════════════════════════════
# Tab 3 — BESS Opportunity
# ═══════════════════════════════════════════════════════════════
with tab_bess:
    st.header("BESS Opportunity — Philippines")
    st.caption("Reserve market started January 2024 · NGCP AS framework · IRESS GEAP eligibility")

    b1, b2 = st.columns(2)
    with b1:
        st.subheader("NGCP Ancillary Services Framework")
        as_df = pd.DataFrame([
            {"Reserve Type": k, "Details": v}
            for k, v in _AS_CONTEXT["reserve_types"].items()
        ])
        st.dataframe(as_df, use_container_width=True, hide_index=True)

        st.subheader("ASPA (Firm Contracts)")
        st.info(
            "**Ancillary Services Procurement Agreement** — bilateral contracts between NGCP and "
            "qualified service providers. BESS assets are ideally suited due to fast response. "
            f"Reserve market commenced: **{_AS_CONTEXT['reserve_market_start']}**"
        )

    with b2:
        st.subheader("BESS Revenue Stack")
        rev_items = _AS_CONTEXT["bess_revenue_stack"]
        for i, item in enumerate(rev_items, 1):
            st.markdown(f"**{i}.** {item}")

        st.subheader("Indicative Revenue Estimates")
        rev_data = pd.DataFrame([
            {"Revenue Stream": "Regulating Reserves (ASPA)", "PHP/MW/hr": "800 – 1,200", "Certainty": "High (firm contract)"},
            {"Revenue Stream": "Contingency Reserves", "PHP/MW/hr": "300 – 600", "Certainty": "Medium (auction)"},
            {"Revenue Stream": "WESM Arbitrage (spread)", "PHP/kWh": "1.50 – 3.00", "Certainty": "Variable"},
            {"Revenue Stream": "IRESS GEAP Tariff", "PHP/kWh": "~6.50 (GEA-4)", "Certainty": "High (20yr GET)"},
        ])
        st.dataframe(rev_data, use_container_width=True, hide_index=True)

    st.divider()
    st.subheader("IEMOP CAPER Programme")
    st.markdown("""
The **Capacity and Energy Reserve (CAPER)** programme allows BESS to register as a capacity resource in WESM, enabling:
- **Capacity payments** during peak periods
- **Energy delivery** obligation when dispatched
- Stack with ancillary services for maximum revenue

BESS investment case in the Philippines is strengthened by the combination of GEAP (IRESS), reserve market (ASPA), and WESM arbitrage — making the Philippines one of Southeast Asia's most attractive BESS markets.
""")


# ═══════════════════════════════════════════════════════════════
# Tab 4 — Investment Analysis
# ═══════════════════════════════════════════════════════════════
with tab_irr:
    st.header("Investment Analysis — IRR / LCOE Calculator")
    st.caption("Parametric model · All values in PHP unless noted · USD/PHP = 58")

    irr_c1, irr_c2 = st.columns([1, 1])
    with irr_c1:
        tech_options = {v["label"]: k for k, v in _TECH_PRESETS.items()}
        selected_label = st.selectbox("Technology", list(tech_options.keys()))
        tech_key = tech_options[selected_label]
        p = _TECH_PRESETS[tech_key]

        is_bess = tech_key.startswith("bess")
        cap_mw = st.number_input("Project Capacity (MW)", min_value=1.0, max_value=5000.0,
                                  value=100.0, step=10.0)

        if is_bess:
            default_capex = p.get("capex_usd_kwh", 300) * p.get("duration_h", 2)
            capex_val = st.number_input("CAPEX (USD/kW-AC)", min_value=100.0, max_value=3000.0,
                                         value=float(default_capex), step=50.0)
            rev_val = st.number_input("Revenue assumption (PHP/kWh throughput)",
                                       min_value=0.5, max_value=10.0, value=2.50, step=0.25)
            cf_val = None
        else:
            capex_val = st.number_input("CAPEX (USD/kW)", min_value=100.0, max_value=5000.0,
                                         value=float(p["capex_usd_kw"]), step=50.0)
            cf_val = st.number_input("Capacity Factor (%)", min_value=5.0, max_value=60.0,
                                      value=float(p["cf_pct"]), step=1.0)
            rev_val = st.number_input("Revenue (PHP/kWh) — GET or WESM",
                                       min_value=2.0, max_value=15.0, value=6.00, step=0.25)

        with st.expander("Advanced parameters"):
            wacc_val      = st.number_input("WACC (%)",           min_value=5.0, max_value=20.0, value=10.0, step=0.5)
            life_val      = st.number_input("Project Life (yrs)", min_value=5,   max_value=30,   value=int(p["life"]), step=1)
            leverage_val  = st.number_input("Leverage (%)",        min_value=0.0, max_value=80.0, value=60.0, step=5.0)
            debt_rate_val = st.number_input("Cost of Debt (%)",    min_value=3.0, max_value=15.0, value=7.0, step=0.5)

    with irr_c2:
        if st.button("Calculate IRR", type="primary"):
            result = _run_irr_model(
                technology=tech_key, capacity_mw=cap_mw,
                capex_usd_per_kw=capex_val, revenue_php_per_kwh=rev_val,
                capacity_factor_pct=cf_val, wacc_pct=wacc_val,
                project_life_yrs=int(life_val), leverage_pct=leverage_val,
                cost_of_debt_pct=debt_rate_val,
            )
            st.session_state["ph_irr_result"] = result
            st.rerun()

        res = st.session_state.get("ph_irr_result")
        if res:
            r1, r2 = st.columns(2)
            r1.metric("Unlevered IRR", f"{res['unlevered_irr_pct']:.1f}%")
            r2.metric("Equity IRR",    f"{res['equity_irr_pct']:.1f}%")
            r3, r4 = st.columns(2)
            r3.metric("LCOE", f"₱{res['lcoe_php_per_kwh']:.3f}/kWh")
            r4.metric("NPV @ WACC", f"₱{res['npv_at_wacc_php_m']:.0f}M")

            st.caption(
                f"CAPEX: USD {res['capex_usd_per_kw']:,.0f}/kW · "
                f"Total CAPEX: ₱{res['capex_total_php_m']:.0f}M · "
                f"Annual Gen: {res['annual_gen_mwh']:,.0f} MWh"
            )

            st.subheader("Sensitivity (Unlevered IRR) — CAPEX × Revenue")
            sens = res.get("sensitivity", [])
            if sens:
                sens_df = pd.DataFrame(sens)
                pivot = sens_df.pivot(index="capex", columns="revenue", values="unlevered_irr")
                st.dataframe(pivot, use_container_width=True)


# ═══════════════════════════════════════════════════════════════
# Tab 5 — Investment Advisor
# ═══════════════════════════════════════════════════════════════
with tab_advisor:
    st.header("Investment Advisor — Philippines RE")
    n_ins = _get_insight_count()
    st.caption(
        f"Senior Philippines RE Investment Expert · 7 tools · "
        f"Expert memory: {n_ins} accumulated insights · Knowledge base grounded"
    )

    # Session init
    if "ph_adv_session_id" not in st.session_state:
        st.session_state["ph_adv_session_id"] = str(uuid.uuid4())
    if "ph_adv_history" not in st.session_state:
        st.session_state["ph_adv_history"] = []

    # Resume previous session
    if not st.session_state["ph_adv_history"]:
        recent = _list_recent_sessions(PREFIX)
        if not recent.empty:
            with st.expander("Resume a previous conversation?", expanded=False):
                for _, srow in recent.iterrows():
                    lbl = (
                        f"{srow['session_id'][:8]}… — "
                        f"{srow['updated_at'].strftime('%Y-%m-%d %H:%M')} — "
                        f"{int(srow['msg_count'])} messages"
                    )
                    if st.button(lbl, key=f"ph_resume_{srow['session_id']}"):
                        st.session_state["ph_adv_session_id"] = srow["session_id"]
                        st.session_state["ph_adv_history"] = _load_session(srow["session_id"])
                        st.rerun()

    # Quick-start questions
    if not st.session_state["ph_adv_history"]:
        st.markdown("**Quick-start questions:**")
        qs_col1, qs_col2 = st.columns(2)
        quick_questions = [
            "What is the investment case for BESS in the Philippines in 2024?",
            "Compare GEAP vs PSA vs WESM merchant routes for a 100MW solar project",
            "Model IRR for a 50MW / 2h BESS targeting NGCP regulating reserves",
            "What are the key risks for a foreign developer entering the Philippines RE market?",
        ]
        for i, qq in enumerate(quick_questions):
            col = qs_col1 if i % 2 == 0 else qs_col2
            if col.button(qq, key=f"ph_qq_{i}"):
                st.session_state["ph_adv_history"].append({"role": "user", "content": qq})
                st.rerun()

    # Chat history
    for msg in st.session_state["ph_adv_history"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # Chat input
    user_input = st.chat_input("Ask about Philippines RE investment — GEAP, BESS, IRR, policy, market structure…")
    if user_input:
        st.session_state["ph_adv_history"].append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.markdown(user_input)
        with st.chat_message("assistant"):
            with st.spinner("Analysing…"):
                api_msgs = [{"role": m["role"], "content": m["content"]}
                            for m in st.session_state["ph_adv_history"]]
                try:
                    reply, _, tool_events = _run_agent_turn(api_msgs, _build_system(user_input))
                except Exception as err:
                    reply = f"API error: {err}. Please try again."
                    tool_events = []
            st.markdown(reply)
            if tool_events:
                with st.expander(f"Tools used ({len(tool_events)})", expanded=False):
                    for ev in tool_events:
                        st.caption(f"**{ev['tool']}** → {ev['result'][:120]}…")

        st.session_state["ph_adv_history"].append({"role": "assistant", "content": reply})
        try:
            _save_session(st.session_state["ph_adv_session_id"], st.session_state["ph_adv_history"])
        except Exception:
            pass
        try:
            n_new = extract_insights(user_input, reply, _ANTHROPIC_KEY, PREFIX, CFG.name)
            if n_new:
                st.toast(f"Stored {n_new} expert insight(s)")
        except Exception:
            pass

    if st.session_state.get("ph_adv_history"):
        if st.button("Clear conversation", key="ph_clear_adv"):
            st.session_state["ph_adv_history"] = []
            st.session_state["ph_adv_session_id"] = str(uuid.uuid4())
            st.rerun()


# ═══════════════════════════════════════════════════════════════
# Tab 6 — Knowledge Base
# ═══════════════════════════════════════════════════════════════
with tab_kb:
    st.header("Philippines Knowledge Base")
    st.info("Auto-updated daily · Ingest at **03:30 MNL** · Digest at **03:45 MNL**", icon="🔄")

    kb_c1, kb_c2 = st.columns([2, 1])
    with kb_c1:
        kb_counts = _knowledge_doc_counts(PREFIX)
        if kb_counts.empty:
            st.info("Knowledge base is empty. Click 'Auto-ingest Local Reports' to seed it.")
        else:
            st.dataframe(kb_counts, use_container_width=True, hide_index=True)
    with kb_c2:
        if st.button("Refresh KB stats", key="ph_kb_refresh"):
            _knowledge_doc_counts.clear()
            st.rerun()
        if st.button("Auto-ingest Local Reports", type="primary", key="ph_kb_ingest_local"):
            with st.spinner("Ingesting from data/market-fundamentals-ph/…"):
                try:
                    results = run_knowledge_ingest(only=["local_reports"], verbose=False)
                    st.success(f"Done — {results.get('local_reports', 0)} new docs")
                except Exception as exc:
                    st.error(f"Failed: {exc}")
            _knowledge_doc_counts.clear()
            st.rerun()
        if st.button("Ingest All Sources", key="ph_kb_ingest_all"):
            with st.spinner("Fetching from all sources (DOE, IEMOP, local)…"):
                try:
                    results = run_knowledge_ingest(verbose=False)
                    total = sum(results.values())
                    st.success(f"Done — {total} new docs: {results}")
                except Exception as exc:
                    st.error(f"Failed: {exc}")
            _knowledge_doc_counts.clear()
            st.rerun()

    st.divider()
    st.subheader("Search Knowledge Base")
    kb_query = st.text_input("Search query", placeholder="e.g. BESS ancillary services NGCP reserves", key="ph_kb_query")
    if kb_query:
        _search_knowledge.clear()
        results_df = _search_knowledge(PREFIX, kb_query, limit=10)
        if results_df.empty:
            st.info("No results found.")
        else:
            for _, row in results_df.iterrows():
                with st.expander(
                    f"[{row['source']}] {row['title'] or 'Untitled'} — {row['published_date']}",
                    expanded=False,
                ):
                    if row.get("url"):
                        st.markdown(f"[Source]({row['url']})")
                    st.text(row["snippet"])

    st.divider()
    st.subheader("Upload Documents")
    up1, up2 = st.tabs(["Upload Files", "Fetch from URL"])
    with up1:
        uploaded_files = st.file_uploader(
            "PDF, Excel, PPTX, Word, TXT",
            type=["pdf", "xlsx", "xls", "pptx", "ppt", "docx", "doc", "txt"],
            accept_multiple_files=True, key="ph_kb_upload",
        )
        if uploaded_files and st.button("Ingest uploaded files", type="primary", key="ph_upload_btn"):
            prog = st.progress(0)
            ok, errs = [], []
            for i, f in enumerate(uploaded_files):
                res = _ingest_uploaded_file(f.name, f.read())
                (ok if res["status"] == "success" else errs).append((f.name, res))
                prog.progress((i + 1) / len(uploaded_files))
            prog.empty()
            if ok:
                st.success(f"Ingested {len(ok)} file(s).")
            for fname, r in errs:
                st.error(f"{fname}: {r['msg']}")
            _knowledge_doc_counts.clear()
            st.rerun()
    with up2:
        fetch_url = st.text_input("Article URL", key="ph_fetch_url")
        if st.button("Fetch and ingest", type="primary", key="ph_fetch_btn") and fetch_url:
            with st.spinner("Fetching…"):
                res = _ingest_url(fetch_url.strip())
            if res["status"] == "success":
                st.success(res["msg"])
                _knowledge_doc_counts.clear()
            else:
                st.error(res["msg"])


# ═══════════════════════════════════════════════════════════════
# Tab 7 — Data Management
# ═══════════════════════════════════════════════════════════════
with tab_mgmt:
    st.header("Data Management")

    dm_c1, dm_c2 = st.columns(2)
    with dm_c1:
        st.subheader("Table Coverage")
        if st.button("Refresh counts", key="ph_refresh_counts"):
            _ph_table_counts.clear()
        counts_df = _ph_table_counts(PREFIX)
        st.dataframe(counts_df, use_container_width=True, hide_index=True)

    with dm_c2:
        st.subheader("Scheduler Status")
        try:
            sched = _start_scheduler()
            jobs = sched.get_jobs()
            if jobs:
                next_run = min((j.next_run_time for j in jobs if j.next_run_time), default=None)
                next_str = next_run.strftime("%Y-%m-%d %H:%M MNL") if next_run else "—"
                st.success(f"Running · Next job: **{next_str}**")
                for j in jobs:
                    nrt = j.next_run_time.strftime("%H:%M") if j.next_run_time else "—"
                    st.caption(f"`{j.id}` — next: {nrt}")
            else:
                st.warning("Scheduler has no active jobs.")
        except Exception as exc:
            st.error(f"Scheduler error: {exc}")

    st.divider()

    # ── WESM Price Data ───────────────────────────────────────────────────────
    st.subheader("WESM Spot Price Data")
    st.caption(
        "Daily settlement interval prices for Luzon, Visayas, and Mindanao reference "
        "trading nodes scraped from IEMOP. Scheduler runs at **07:15 MNL** daily."
    )

    wesm_stat_c1, wesm_stat_c2, wesm_stat_c3 = st.columns(3)
    with wesm_stat_c1:
        latest_str = _wesm_latest_date()
        st.metric("Latest trading date", latest_str.split(" ")[0] if latest_str != "—" else "—")
    with wesm_stat_c2:
        try:
            n_rows = _query("SELECT COUNT(*) AS n FROM intl_market.ph_wesm_prices")
            st.metric("Price rows stored", int(n_rows["n"].iloc[0]))
        except Exception:
            st.metric("Price rows stored", "—")
    with wesm_stat_c3:
        try:
            n_days_df = _query("SELECT COUNT(DISTINCT trading_date) AS n FROM intl_market.ph_wesm_prices")
            st.metric("Trading days", int(n_days_df["n"].iloc[0]))
        except Exception:
            st.metric("Trading days", "—")

    # Manual trigger
    wesm_btn_c1, wesm_btn_c2 = st.columns(2)
    with wesm_btn_c1:
        if st.button("Scrape Today's WESM Prices", type="primary", key="ph_wesm_scrape_now"):
            with st.spinner("Fetching prices from IEMOP…"):
                try:
                    import psycopg2
                    _wconn = psycopg2.connect(
                        os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                        keepalives=1, keepalives_idle=30,
                    )
                    _wconn.autocommit = True
                    results = run_wesm_price_scrape(_wconn, days_back=1)
                    _wconn.close()
                    _wesm_latest_date.clear()
                    total = sum(v for v in results.values() if v > 0)
                    if total > 0:
                        st.success(f"Fetched {total} new price records: {results}")
                    else:
                        st.warning(
                            f"Scraper ran but found 0 new records ({results}). "
                            "IEMOP may not have published yesterday's data yet — "
                            "prices are typically available by 09:00 MNL."
                        )
                except Exception as exc:
                    st.error(f"WESM price scrape failed: {exc}")
    with wesm_btn_c2:
        backfill_days = st.number_input(
            "Backfill days", min_value=1, max_value=30, value=7, key="ph_wesm_backfill_days"
        )
        if st.button("Backfill Price History", key="ph_wesm_backfill"):
            with st.spinner(f"Backfilling {backfill_days} days…"):
                try:
                    import psycopg2
                    _wconn = psycopg2.connect(
                        os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                        keepalives=1, keepalives_idle=30,
                    )
                    _wconn.autocommit = True
                    results = run_wesm_price_scrape(_wconn, days_back=int(backfill_days))
                    _wconn.close()
                    _wesm_latest_date.clear()
                    total = sum(v for v in results.values() if v > 0)
                    st.success(f"Backfill complete: {total} new rows across {len(results)} dates.")
                    with st.expander("Per-date results"):
                        for dt, n in sorted(results.items(), reverse=True):
                            st.caption(f"{dt}: {n} rows")
                except Exception as exc:
                    st.error(f"Backfill failed: {exc}")

    # WESM price chart
    wesm_df = _wesm_price_history(days=30)
    if not wesm_df.empty:
        st.markdown("**Last 30 days — Daily average spot price (PHP/kWh)**")
        try:
            import plotly.express as px
            fig = px.line(
                wesm_df[wesm_df["avg_price"].notna()],
                x="trading_date", y="avg_price", color="region",
                labels={"trading_date": "Date", "avg_price": "PHP/kWh", "region": "Grid"},
                color_discrete_map={"Luzon": "#1f77b4", "Visayas": "#ff7f0e", "Mindanao": "#2ca02c"},
            )
            fig.update_layout(height=300, margin=dict(t=20, b=20))
            st.plotly_chart(fig, use_container_width=True)
        except Exception:
            st.dataframe(wesm_df[["trading_date", "region", "avg_price"]],
                         use_container_width=True, hide_index=True)
    else:
        st.info("No WESM price data yet. Click 'Scrape Today's WESM Prices' above to fetch.")

    st.divider()

    # ── WESM / IEMOP Market Reports ───────────────────────────────────────────
    st.subheader("WESM Market Reports (IEMOP)")
    st.caption(
        "Scrapes IEMOP market bulletins, advisories, and notices into the Knowledge Base. "
        "Runs as part of the nightly KB ingest at **03:30 MNL**."
    )
    wesm_rep_c1, wesm_rep_c2 = st.columns(2)
    with wesm_rep_c1:
        try:
            n_wesm_docs = _query(
                f"SELECT COUNT(*) AS n FROM intl_market.{PREFIX}knowledge_docs "
                "WHERE source='wesm_iemop'"
            )
            st.metric("WESM docs in KB", int(n_wesm_docs["n"].iloc[0]))
        except Exception:
            st.metric("WESM docs in KB", "—")
    with wesm_rep_c2:
        try:
            latest_doc = _query(
                f"SELECT MAX(fetched_at) AS latest FROM intl_market.{PREFIX}knowledge_docs "
                "WHERE source='wesm_iemop'"
            )
            lt = latest_doc["latest"].iloc[0]
            st.metric("Last fetched", str(lt)[:16] if lt is not None else "Never")
        except Exception:
            st.metric("Last fetched", "—")

    if st.button("Fetch WESM Reports Now", key="ph_wesm_reports_now"):
        with st.spinner("Scraping IEMOP reports (may take 30–60 s)…"):
            try:
                from services.ph_knowledge.ingest import run_knowledge_ingest
                results = run_knowledge_ingest(only=["wesm_reports"], verbose=False)
                n = results.get("wesm_reports", 0)
                st.success(f"Done — {n} new IEMOP documents added to Knowledge Base.")
            except Exception as exc:
                st.error(f"Report scrape failed: {exc}")

    st.divider()
    st.subheader("Knowledge Base Digest → Expert Memory")
    st.caption("Extracts durable insights from KB documents. Runs at 03:45 MNL nightly.")
    if st.button("Digest KB into Expert Memory", type="primary", key="ph_digest_kb"):
        with st.spinner("Extracting insights (1-2 min)…"):
            try:
                n_dk = digest_kb_docs(_ANTHROPIC_KEY, PREFIX, CFG.name, limit=200)
                st.success(f"Extracted {n_dk} new insights.")
            except Exception as exc:
                st.error(f"KB digest failed: {exc}")

    st.divider()
    st.subheader("Expert Memory")
    try:
        insights_df = _query(
            f"SELECT id, insight_type, confidence, insight_text, validated_at "
            f"FROM intl_market.{PREFIX}expert_insights WHERE active=TRUE ORDER BY validated_at DESC LIMIT 50"
        )
        if insights_df.empty:
            st.info("No expert insights yet. Digest the KB above to populate.")
        else:
            st.caption(f"{len(insights_df)} insights shown (max 50)")
            for _, row in insights_df.iterrows():
                with st.expander(
                    f"[{row['insight_type']}] [{row['confidence']}] {str(row['insight_text'])[:80]}…",
                    expanded=False,
                ):
                    st.write(row["insight_text"])
                    st.caption(f"Logged: {row['validated_at']}")
    except Exception as exc:
        st.error(f"Could not load insights: {exc}")

    st.divider()
    st.subheader("Agent Memory")
    mems = _load_memories(APP_KEY)
    if mems.empty:
        st.info("No agent memories saved yet.")
    else:
        for _, row in mems.iterrows():
            c1, c2 = st.columns([10, 1])
            with c1:
                st.markdown(f"**[{row['category']}]** {row['subject']}: {row['content']}")
                st.caption(f"{row['source']} · {row['created_at']}")
            with c2:
                if st.button("🗑", key=f"ph_del_mem_{row['id']}"):
                    _delete_memory(row["id"])
                    st.rerun()

    st.divider()
    st.subheader("Add Memory Manually")
    with st.form("ph_add_mem"):
        cat = st.selectbox("Category",
                           ["market_view", "methodology", "investment_thesis", "asset_note", "red_flag"])
        subj = st.text_input("Subject (≤8 words)")
        cont = st.text_area("Content (one sentence)")
        if st.form_submit_button("Save"):
            _save_memory(cat, subj, cont, source="manual")
            st.success("Saved")


# ═══════════════════════════════════════════════════════════════
# Tab 8 — Grid Analysis (PyPSA)
# ═══════════════════════════════════════════════════════════════
with tab_pypsa:
    st.header("Grid Analysis — PyPSA Power Flow")
    st.info(
        "Upload NGCP network data (buses, lines, generators CSV) to enable power flow analysis. "
        "Models inter-island capacity constraints and RE integration scenarios.",
        icon="⚡",
    )

    col_py1, col_py2, col_py3 = st.columns(3)
    buses_file = col_py1.file_uploader("Buses CSV",      type=["csv"], key="ph_pypsa_buses")
    lines_file  = col_py2.file_uploader("Lines CSV",      type=["csv"], key="ph_pypsa_lines")
    gens_file   = col_py3.file_uploader("Generators CSV", type=["csv"], key="ph_pypsa_gens")

    if buses_file and lines_file:
        try:
            import pypsa
            import tempfile, os as _os

            buses_df = pd.read_csv(buses_file)
            lines_df = pd.read_csv(lines_file)

            n = pypsa.Network()
            for _, row in buses_df.iterrows():
                n.add("Bus", row["name"], v_nom=row.get("v_nom", 230))
            for _, row in lines_df.iterrows():
                n.add("Line", row.get("name", f"L{_}"),
                      bus0=row["bus0"], bus1=row["bus1"], x=row.get("x", 0.1), s_nom=row.get("s_nom", 200))
            if gens_file:
                gens_df = pd.read_csv(gens_file)
                for _, row in gens_df.iterrows():
                    n.add("Generator", row["name"], bus=row["bus"],
                          p_nom=row.get("p_nom", 100), marginal_cost=row.get("marginal_cost", 50))

            st.success(f"Network loaded: {len(n.buses)} buses, {len(n.lines)} lines, {len(n.generators)} generators")

            if st.button("Run DC Power Flow", type="primary", key="ph_pypsa_pf"):
                with st.spinner("Running LOPF…"):
                    n.lopf(pyomo=False, solver_name="highs")

                st.subheader("Generator Dispatch")
                disp = n.generators_t.p.T.reset_index()
                disp.columns = ["Generator"] + [f"t{i}" for i in range(len(disp.columns) - 1)]
                st.dataframe(disp, use_container_width=True, hide_index=True)

                st.subheader("Line Loading")
                if not n.lines_t.p0.empty:
                    loading = (n.lines_t.p0.abs() / n.lines.s_nom).T.reset_index()
                    loading.columns = ["Line"] + [f"t{i}" for i in range(len(loading.columns) - 1)]
                    st.dataframe(loading, use_container_width=True, hide_index=True)

        except ImportError:
            st.error("PyPSA not installed in this environment. Add `pypsa` to the Dockerfile.")
        except Exception as exc:
            st.error(f"PyPSA error: {exc}")
    else:
        st.markdown("""
**Expected CSV formats:**

| buses.csv | lines.csv | generators.csv |
|-----------|-----------|----------------|
| name, v_nom | name, bus0, bus1, x, s_nom | name, bus, p_nom, marginal_cost |

**Philippines grid notes:**
- Three isolated grids: Luzon, Visayas, Mindanao
- Luzon-Visayas HVDC link: 440 MW (LEYTE–LUZON)
- Mindanao grid: largely islanded; MRU has dispatch priority
- NGCP TDP 2024-2050 has full transmission project list

Source: NGCP Transmission Development Plan (TDP), IEMOP CAPER
""")


# ═══════════════════════════════════════════════════════════════
# Tab 9 — Library
# ═══════════════════════════════════════════════════════════════
with tab_library:
    st.header("Report Library")
    st.caption("Save and retrieve daily, weekly, and monthly market reports (PDF)")

    # ── Upload ──
    with st.expander("Upload New Report", expanded=True):
        ul_c1, ul_c2, ul_c3 = st.columns(3)
        with ul_c1:
            lib_name = st.text_input("Report Name",
                                      placeholder="e.g. DOE Market Report May 2026",
                                      key="ph_lib_name")
        with ul_c2:
            lib_freq = st.selectbox("Frequency", ["daily", "weekly", "monthly"],
                                     key="ph_lib_freq")
        with ul_c3:
            lib_period = st.date_input("Report Period", value=date.today(),
                                        key="ph_lib_period")
        lib_file = st.file_uploader("PDF File", type=["pdf"], key="ph_lib_upload")
        if st.button("Save to Library", type="primary", key="ph_lib_save_btn",
                     disabled=not (lib_file and lib_name)):
            try:
                _save_library_report(lib_name.strip(), lib_freq, lib_period,
                                      lib_file.name, lib_file.read())
                st.success(f"Saved '{lib_name}' ({lib_freq}) to library.")
                st.rerun()
            except Exception as exc:
                st.error(f"Failed: {exc}")

    st.divider()

    # ── Browse by frequency ──
    lib_freq_tabs = st.tabs(["All Reports", "Daily", "Weekly", "Monthly"])
    lib_freq_filters = [None, "daily", "weekly", "monthly"]

    for lib_ftab, lib_filt in zip(lib_freq_tabs, lib_freq_filters):
        with lib_ftab:
            lib_df = _list_library_reports(PREFIX, lib_filt)
            if lib_df.empty:
                st.info("No reports in this category yet. Use 'Upload New Report' above.")
                continue

            st.caption(f"{len(lib_df)} report(s)")
            filt_key = lib_filt or "all"

            # Summary table
            display_df = lib_df[["report_name", "frequency", "period",
                                  "filename", "file_size_kb", "uploaded_at"]].copy()
            display_df.columns = ["Report Name", "Frequency", "Period",
                                   "Filename", "Size (KB)", "Uploaded At"]
            st.dataframe(display_df, use_container_width=True, hide_index=True)

            st.divider()
            st.subheader("Download / Delete")
            options = {
                f"{r['report_name']}  —  {str(r['period'])}  [{r['frequency']}]": int(r["id"])
                for _, r in lib_df.iterrows()
            }
            sel_label = st.selectbox("Select report", list(options.keys()),
                                      key=f"ph_lib_sel_{filt_key}")
            if sel_label:
                sel_id = options[sel_label]
                sel_row = lib_df[lib_df["id"] == sel_id].iloc[0]
                col_dl, col_del = st.columns([3, 1])
                with col_dl:
                    report_bytes = _get_library_report_data(PREFIX, sel_id)
                    if report_bytes:
                        st.download_button(
                            label=f"Download  {sel_row['filename']}",
                            data=report_bytes,
                            file_name=sel_row["filename"],
                            mime="application/pdf",
                            key=f"ph_lib_dl_{sel_id}_{filt_key}",
                        )
                    else:
                        st.warning("Report data not found in database.")
                with col_del:
                    if st.button("Delete", key=f"ph_lib_del_{sel_id}_{filt_key}",
                                  type="secondary"):
                        _delete_library_report(sel_id)
                        st.success("Deleted.")
                        st.rerun()
