"""Poland Power Market Investment Advisory — Streamlit app. Port 8511.

Standalone app (does NOT use run_market_app template).
Tabs: Market Structure | Balancing & AS Markets | BESS Opportunity |
      Investment Analysis | Investment Advisor | Knowledge Base |
      Data Management | Grid Analysis (PyPSA)
"""
import io
import json
import logging
import os
import sys
import uuid
from datetime import date, timedelta

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..")))
from dotenv import load_dotenv
load_dotenv(os.path.join(os.path.dirname(__file__), "..", "..", "config", ".env"))

import anthropic
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import psycopg2
import streamlit as st

st.set_page_config(
    page_title="Poland Power Investment Advisory",
    page_icon="🇵🇱",
    layout="wide",
    initial_sidebar_state="expanded",
)

from services.po_knowledge.config import MARKET_CONFIG
from services.po_knowledge.ingest import run_knowledge_ingest
from services.po_knowledge.entso_scraper import (
    ENTSOEPriceScraper,
    run_entso_price_scrape,
    run_po_doc_backfill,
    _PO_DOC_CONNECTOR_MAP,
    get_as_revenue_estimate,
)
from services.intl_market_common.advanced_retrieval_base import retrieve_for_agent
from services.intl_market_common.expert_memory_base import (
    extract_insights, get_insights, inject_memory, digest_kb_docs,
)
from services.intl_market_common.export_helpers import export_pdf, export_pptx, export_docx
from services.intl_market_common.audio_ingest import transcribe_and_contextualize, store_voice_memo

logger = logging.getLogger(__name__)

_ANTHROPIC_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
_client = anthropic.Anthropic(api_key=_ANTHROPIC_KEY)

CFG = MARKET_CONFIG
PREFIX = CFG.table_prefix      # "po_"
APP_KEY = CFG.app_key          # "po_market"
CURRENCY = CFG.currency_sym    # "zł"
_EUR_PLN = 4.25
_USD_PLN = 3.95


# ── DB connection ────────────────────────────────────────────────────────────

@st.cache_resource(ttl=3600)
def _get_conn():
    url = (
        os.environ.get("PGURL")
        or "postgresql://postgres:root@127.0.0.1:5433/marketdata"
    )
    conn = psycopg2.connect(url, connect_timeout=10)
    conn.autocommit = True
    return conn


def _conn():
    conn = _get_conn()
    if conn.closed:
        _get_conn.clear()
        conn = _get_conn()
    return conn


def _query(sql: str, params=None) -> pd.DataFrame:
    return pd.read_sql(sql, _conn(), params=params)


# ── Ensure tables ────────────────────────────────────────────────────────────

def _ensure_tables():
    cur = _conn().cursor()
    cur.execute("ALTER TABLE marketdata.agent_memory ADD COLUMN IF NOT EXISTS app TEXT DEFAULT 'po_market'")
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}analyst_sessions (
            session_id TEXT PRIMARY KEY,
            messages   JSONB NOT NULL DEFAULT '[]',
            created_at TIMESTAMPTZ DEFAULT NOW(),
            updated_at TIMESTAMPTZ DEFAULT NOW()
        )
    """)
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}knowledge_docs (
            id              SERIAL PRIMARY KEY,
            source          TEXT NOT NULL,
            doc_type        TEXT NOT NULL,
            title           TEXT,
            url             TEXT UNIQUE,
            published_date  DATE,
            content         TEXT NOT NULL,
            fetched_at      TIMESTAMPTZ NOT NULL DEFAULT NOW(),
            search_vector   TSVECTOR GENERATED ALWAYS AS (
                to_tsvector('english', coalesce(title,'') || ' ' || left(content,100000))
            ) STORED
        )
    """)
    cur.execute(
        f"CREATE INDEX IF NOT EXISTS {PREFIX}knowledge_docs_fts "
        f"ON intl_market.{PREFIX}knowledge_docs USING GIN(search_vector)"
    )
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}expert_insights (
            id SERIAL PRIMARY KEY,
            insight_text TEXT,
            insight_type TEXT,
            confidence TEXT,
            source_session TEXT,
            source_doc_url TEXT,
            active BOOLEAN DEFAULT TRUE,
            validated_at TIMESTAMPTZ DEFAULT NOW()
        )
    """)
    cur.execute(f"""
        CREATE TABLE IF NOT EXISTS intl_market.{PREFIX}report_library (
            id           SERIAL PRIMARY KEY,
            report_name  TEXT NOT NULL,
            frequency    TEXT NOT NULL CHECK (frequency IN ('daily', 'weekly', 'monthly')),
            period       DATE NOT NULL,
            filename     TEXT NOT NULL,
            file_data    BYTEA NOT NULL,
            file_size_kb INTEGER,
            uploaded_at  TIMESTAMPTZ DEFAULT NOW()
        )
    """)
    cur.execute(
        f"CREATE INDEX IF NOT EXISTS {PREFIX}report_library_freq_period "
        f"ON intl_market.{PREFIX}report_library (frequency, period DESC)"
    )
    cur.execute("""
        CREATE TABLE IF NOT EXISTS intl_market.po_day_ahead_prices (
            id            SERIAL PRIMARY KEY,
            trading_date  DATE          NOT NULL,
            hour          INTEGER       NOT NULL,
            price_pln_mwh NUMERIC(10,4),
            price_eur_mwh NUMERIC(10,4),
            source        TEXT          NOT NULL DEFAULT 'energy_charts',
            fetched_at    TIMESTAMPTZ   NOT NULL DEFAULT NOW(),
            CONSTRAINT po_day_ahead_prices_uq UNIQUE (trading_date, hour, source)
        )
    """)
    cur.execute(
        "CREATE INDEX IF NOT EXISTS po_day_ahead_prices_date_idx "
        "ON intl_market.po_day_ahead_prices (trading_date DESC)"
    )
    # Ancillary service weekly auction prices (FCR, aFRR)
    cur.execute("""
        CREATE TABLE IF NOT EXISTS intl_market.po_as_prices (
            id                SERIAL PRIMARY KEY,
            week_start        DATE NOT NULL,
            market_type       TEXT NOT NULL,
            price_pln_mw_week NUMERIC(12,2),
            accepted_mw       NUMERIC(10,2),
            source            TEXT DEFAULT 'pse',
            fetched_at        TIMESTAMPTZ DEFAULT now(),
            UNIQUE (week_start, market_type)
        )
    """)
    cur.execute(
        "CREATE INDEX IF NOT EXISTS po_as_prices_week_idx "
        "ON intl_market.po_as_prices (week_start DESC)"
    )
    # Rynek Mocy (Capacity Market) annual auction results
    cur.execute("""
        CREATE TABLE IF NOT EXISTS intl_market.po_capacity_market (
            id              SERIAL PRIMARY KEY,
            delivery_year   INT  NOT NULL,
            auction_date    DATE,
            price_pln_mw_yr NUMERIC(12,2),
            accepted_mw     NUMERIC(10,2),
            source          TEXT DEFAULT 'tge',
            fetched_at      TIMESTAMPTZ DEFAULT now(),
            UNIQUE (delivery_year)
        )
    """)
    _conn().commit()


try:
    _ensure_tables()
except Exception as _e:
    logger.warning("_ensure_tables: %s", _e)


# ── Scheduler ────────────────────────────────────────────────────────────────

@st.cache_resource
def _start_scheduler():
    from apscheduler.schedulers.background import BackgroundScheduler

    def _ingest_job():
        try:
            run_knowledge_ingest(verbose=False)
        except Exception as exc:
            logger.error("[po_scheduler] ingest: %s", exc)

    def _digest_job():
        try:
            digest_kb_docs(_ANTHROPIC_KEY, PREFIX, CFG.name, limit=100)
        except Exception as exc:
            logger.error("[po_scheduler] digest: %s", exc)

    def _price_job():
        try:
            import psycopg2
            conn = psycopg2.connect(
                os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                keepalives=1, keepalives_idle=30,
            )
            conn.autocommit = True
            results = run_entso_price_scrape(conn, days_back=2)
            conn.close()
            total = sum(v for v in results.values() if v > 0)
            logger.info("[po_scheduler] price scrape done: %s (%d new rows)", results, total)
        except Exception as exc:
            logger.error("[po_scheduler] price scrape failed: %s", exc)

    def _docs_job():
        """Weekly: scrape PSE, TGE, URE, and ENTSO-E document sources."""
        try:
            run_knowledge_ingest(only=[
                "pse_pl", "pse_grid", "pse_afrr",
                "tge_reports", "ure_regulatory", "entsoe_publications",
            ], verbose=False)
            logger.info("[po_scheduler] docs job done")
        except Exception as exc:
            logger.error("[po_scheduler] docs job failed: %s", exc)

    def _as_scrape_job():
        """Scheduled: scrape FCR and aFRR prices from PSE (Tuesdays 06:05 CET)."""
        try:
            import psycopg2 as _psycopg2
            _c = _psycopg2.connect(
                os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                keepalives=1, keepalives_idle=30,
            )
            _c.autocommit = True
            from services.po_knowledge.entso_scraper import (
                scrape_po_fcr_prices, scrape_po_afrr_prices,
            )
            n_fcr  = scrape_po_fcr_prices(_c, weeks_back=4)
            n_afrr = scrape_po_afrr_prices(_c, weeks_back=4)
            _c.close()
            logger.info("[po_scheduler] po_as_prices: FCR=%d rows, aFRR=%d rows", n_fcr, n_afrr)
        except Exception as exc:
            logger.error("[po_scheduler] po_as_prices failed: %s", exc)

    def _cap_market_job():
        """Scheduled: scrape Rynek Mocy results from TGE (1st of month 05:10 CET)."""
        try:
            import psycopg2 as _psycopg2
            _c = _psycopg2.connect(
                os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                keepalives=1, keepalives_idle=30,
            )
            _c.autocommit = True
            from services.po_knowledge.entso_scraper import scrape_po_capacity_market
            n = scrape_po_capacity_market(_c)
            _c.close()
            logger.info("[po_scheduler] po_capacity_market: %d rows", n)
        except Exception as exc:
            logger.error("[po_scheduler] po_capacity_market failed: %s", exc)

    sched = BackgroundScheduler(timezone="Europe/Warsaw")
    sched.add_job(_ingest_job,     "cron", hour=3, minute=30, id="po_ingest",     misfire_grace_time=3600)
    sched.add_job(_digest_job,     "cron", hour=3, minute=45, id="po_digest",     misfire_grace_time=3600)
    sched.add_job(_price_job,      "cron", hour=7, minute=15, id="po_price",      misfire_grace_time=3600)
    sched.add_job(_docs_job,       "cron", day_of_week="mon", hour=4, minute=5,
                  id="po_docs",       misfire_grace_time=7200)
    sched.add_job(_as_scrape_job,  "cron", day_of_week="tue", hour=6, minute=5,
                  id="po_as_prices",  misfire_grace_time=3600)
    sched.add_job(_cap_market_job, "cron", day=1, hour=5, minute=10,
                  id="po_cap_market", misfire_grace_time=3600)
    sched.start()
    return sched


_start_scheduler()


# ── Cached helpers ────────────────────────────────────────────────────────────

@st.cache_data(ttl=60)
def _load_memories(app_key: str) -> pd.DataFrame:
    try:
        return _query(
            "SELECT id, category, subject, content, source, created_at "
            "FROM marketdata.agent_memory WHERE app=%s AND active=TRUE ORDER BY created_at DESC",
            (app_key,),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=60)
def _search_knowledge(prefix: str, query: str, limit: int = 10) -> pd.DataFrame:
    try:
        return _query(
            f"SELECT source, doc_type, title, url, published_date, "
            f"left(content,1500) AS snippet, "
            f"ts_rank(search_vector, plainto_tsquery('english',%s)) AS rank "
            f"FROM intl_market.{prefix}knowledge_docs "
            f"WHERE search_vector @@ to_tsquery('english',"
            f"  regexp_replace(plainto_tsquery('english',%s)::text,' & ',' | ','g')) "
            f"ORDER BY rank DESC LIMIT %s",
            (query, query, limit),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=300)
def _knowledge_doc_counts(prefix: str) -> pd.DataFrame:
    try:
        return _query(
            f"SELECT source, doc_type, COUNT(*) AS docs, MAX(fetched_at) AS last_fetch "
            f"FROM intl_market.{prefix}knowledge_docs GROUP BY source, doc_type ORDER BY source"
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=300)
def _po_table_counts(prefix: str) -> pd.DataFrame:
    rows = []
    for table in [f"{prefix}knowledge_docs", f"{prefix}expert_insights", f"{prefix}analyst_sessions"]:
        try:
            df = _query(f"SELECT COUNT(*) AS n FROM intl_market.{table}")
            rows.append({"Table": f"intl_market.{table}", "Rows": int(df["n"].iloc[0])})
        except Exception:
            rows.append({"Table": f"intl_market.{table}", "Rows": "error"})
    try:
        df2 = _query("SELECT COUNT(*) AS n FROM marketdata.agent_memory WHERE app=%s AND active=TRUE", (APP_KEY,))
        rows.append({"Table": "marketdata.agent_memory", "Rows": int(df2["n"].iloc[0])})
    except Exception:
        rows.append({"Table": "marketdata.agent_memory", "Rows": "error"})
    try:
        df3 = _query("SELECT COUNT(*) AS n FROM intl_market.po_day_ahead_prices")
        rows.append({"Table": "intl_market.po_day_ahead_prices", "Rows": int(df3["n"].iloc[0])})
    except Exception:
        rows.append({"Table": "intl_market.po_day_ahead_prices", "Rows": "error"})
    return pd.DataFrame(rows)


@st.cache_data(ttl=60)
def _po_latest_price_date() -> str:
    try:
        df = _query("SELECT MAX(trading_date) AS latest FROM intl_market.po_day_ahead_prices")
        val = df["latest"].iloc[0]
        return str(val) if val is not None else "—"
    except Exception:
        return "—"


@st.cache_data(ttl=60)
def _po_price_history(days: int = 30) -> pd.DataFrame:
    try:
        since = date.today() - timedelta(days=days)
        return _query(
            "SELECT trading_date, "
            "  AVG(price_pln_mwh) AS avg_pln, "
            "  MIN(price_pln_mwh) AS min_pln, "
            "  MAX(price_pln_mwh) AS max_pln, "
            "  AVG(price_eur_mwh) AS avg_eur "
            "FROM intl_market.po_day_ahead_prices "
            "WHERE trading_date >= %s "
            "GROUP BY trading_date ORDER BY trading_date",
            (since,),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=120)
def _list_library_reports(prefix: str, frequency: str | None = None) -> pd.DataFrame:
    try:
        if frequency:
            return _query(
                f"SELECT id, report_name, frequency, period, filename, file_size_kb, uploaded_at "
                f"FROM intl_market.{prefix}report_library "
                f"WHERE frequency=%s ORDER BY period DESC, uploaded_at DESC",
                (frequency,),
            )
        return _query(
            f"SELECT id, report_name, frequency, period, filename, file_size_kb, uploaded_at "
            f"FROM intl_market.{prefix}report_library ORDER BY period DESC, uploaded_at DESC"
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=600)
def _get_library_report_data(prefix: str, report_id: int) -> bytes | None:
    try:
        cur = _conn().cursor()
        cur.execute(
            f"SELECT file_data FROM intl_market.{prefix}report_library WHERE id=%s",
            (report_id,),
        )
        row = cur.fetchone()
        return bytes(row[0]) if row else None
    except Exception:
        return None


def _save_library_report(report_name: str, frequency: str, period,
                          filename: str, file_data: bytes):
    import psycopg2
    kb = len(file_data) // 1024
    cur = _conn().cursor()
    cur.execute(
        f"INSERT INTO intl_market.{PREFIX}report_library "
        f"(report_name, frequency, period, filename, file_data, file_size_kb) "
        f"VALUES (%s,%s,%s,%s,%s,%s)",
        (report_name, frequency, period, filename, psycopg2.Binary(file_data), kb),
    )
    _list_library_reports.clear()
    _get_library_report_data.clear()


def _delete_library_report(report_id: int):
    cur = _conn().cursor()
    cur.execute(f"DELETE FROM intl_market.{PREFIX}report_library WHERE id=%s", (report_id,))
    _list_library_reports.clear()
    _get_library_report_data.clear()


@st.cache_data(ttl=60)
def _list_recent_sessions(prefix: str, limit: int = 3) -> pd.DataFrame:
    try:
        return _query(
            f"SELECT session_id, jsonb_array_length(messages) AS msg_count, updated_at "
            f"FROM intl_market.{prefix}analyst_sessions "
            f"WHERE jsonb_array_length(messages) > 0 ORDER BY updated_at DESC LIMIT %s",
            (limit,),
        )
    except Exception:
        return pd.DataFrame()


# ── Session / memory helpers ──────────────────────────────────────────────────

def _save_session(session_id: str, messages: list):
    try:
        cur = _conn().cursor()
        cur.execute(
            f"INSERT INTO intl_market.{PREFIX}analyst_sessions (session_id, messages, updated_at) "
            f"VALUES (%s, %s::jsonb, NOW()) "
            f"ON CONFLICT (session_id) DO UPDATE SET messages=EXCLUDED.messages, updated_at=NOW()",
            (session_id, json.dumps(messages)),
        )
    except Exception as exc:
        logger.debug("_save_session: %s", exc)


def _load_session(session_id: str) -> list:
    try:
        cur = _conn().cursor()
        cur.execute(
            f"SELECT messages FROM intl_market.{PREFIX}analyst_sessions WHERE session_id=%s",
            (session_id,),
        )
        row = cur.fetchone()
        return row[0] if row else []
    except Exception:
        return []


def _save_memory(category, subject, content, source="manual"):
    cur = _conn().cursor()
    cur.execute(
        "INSERT INTO marketdata.agent_memory (app, category, subject, content, source) "
        "VALUES (%s,%s,%s,%s,%s)",
        (APP_KEY, category, subject, content, source),
    )
    _load_memories.clear()


def _delete_memory(mem_id: int):
    cur = _conn().cursor()
    cur.execute("UPDATE marketdata.agent_memory SET active=FALSE WHERE id=%s", (mem_id,))


def _get_insight_count() -> int:
    try:
        df = _query(f"SELECT COUNT(*) AS n FROM intl_market.{PREFIX}expert_insights WHERE active=TRUE")
        return int(df.iloc[0]["n"]) if not df.empty else 0
    except Exception:
        return 0


# ── Ingest helpers ────────────────────────────────────────────────────────────

def _ingest_url(url: str) -> dict:
    import requests
    from bs4 import BeautifulSoup
    try:
        resp = requests.get(url, timeout=30, headers={"User-Agent": "BESSPlatformBot/1.0"})
        resp.raise_for_status()
    except Exception as exc:
        return {"status": "error", "msg": f"Fetch failed: {exc}"}
    soup = BeautifulSoup(resp.text, "html.parser")
    title_el = soup.find("h1") or soup.find("title")
    title = title_el.get_text(" ", strip=True) if title_el else url
    for tag in soup.find_all(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    content = "\n\n".join(
        el.get_text(" ", strip=True) for el in soup.find_all(["p", "h1", "h2", "h3", "li"])
        if el.get_text(" ", strip=True)
    ) or soup.get_text(" ", strip=True)
    if not content.strip():
        return {"status": "error", "msg": "No text extracted."}
    try:
        cur = _conn().cursor()
        cur.execute(
            f"INSERT INTO intl_market.{PREFIX}knowledge_docs "
            f"(source, doc_type, title, url, published_date, content) VALUES (%s,%s,%s,%s,%s,%s) "
            f"ON CONFLICT (url) DO UPDATE SET content=EXCLUDED.content, title=EXCLUDED.title, fetched_at=NOW()",
            ("upload", "article", title, url, date.today(), content),
        )
        return {"status": "success", "msg": f"Ingested '{title}' ({len(content):,} chars)"}
    except Exception as exc:
        return {"status": "error", "msg": f"DB insert failed: {exc}"}


def _ingest_uploaded_file(filename: str, data: bytes) -> dict:
    ext = filename.rsplit(".", 1)[-1].lower()
    try:
        if ext == "txt":
            content = data.decode("utf-8", errors="replace")
        elif ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            content = "\n\n".join(p.extract_text() for p in reader.pages if p.extract_text())
        elif ext in ("xlsx", "xls"):
            xl = pd.ExcelFile(io.BytesIO(data))
            content = "\n\n".join(
                f"Sheet: {s}\n{xl.parse(s).to_string(index=False)}" for s in xl.sheet_names
            )
        elif ext in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            content = "\n\n".join(slides)
        elif ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(data))
            content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == "doc":
            content = ""
            try:
                import subprocess, tempfile, os as _os
                with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as _f:
                    _f.write(data)
                    _tmppath = _f.name
                _r = subprocess.run(["antiword", _tmppath], capture_output=True, text=True, timeout=30)
                _os.unlink(_tmppath)
                if _r.returncode == 0:
                    content = _r.stdout.strip()
            except Exception:
                pass
            if not content:
                try:
                    from docx import Document as _Doc
                    _doc = _Doc(io.BytesIO(data))
                    content = "\n".join(p.text for p in _doc.paragraphs if p.text.strip())
                except Exception:
                    pass
            if not content:
                return {"status": "error", "msg": (
                    "Could not extract text from .doc file. "
                    "Please convert to .docx or .pdf and re-upload."
                )}
        else:
            return {"status": "error", "msg": f"Unsupported type: .{ext}"}
    except Exception as exc:
        return {"status": "error", "msg": f"Extraction failed: {exc}"}
    if not content.strip():
        return {"status": "error", "msg": "No text extracted."}
    doc_type = {"pdf": "pdf", "txt": "text", "xlsx": "excel", "xls": "excel",
                "pptx": "presentation", "ppt": "presentation", "docx": "word"}.get(ext, "document")
    try:
        cur = _conn().cursor()
        cur.execute(
            f"INSERT INTO intl_market.{PREFIX}knowledge_docs "
            f"(source, doc_type, title, url, published_date, content) VALUES (%s,%s,%s,%s,%s,%s) "
            f"ON CONFLICT (url) DO UPDATE SET content=EXCLUDED.content, fetched_at=NOW()",
            ("upload", doc_type, filename, f"upload://{filename}", date.today(), content),
        )
        return {"status": "success", "msg": f"Ingested '{filename}' ({len(content):,} chars)"}
    except Exception as exc:
        return {"status": "error", "msg": f"DB insert failed: {exc}"}


# ── IRR model ────────────────────────────────────────────────────────────────

_TECH_PRESETS_PO = {
    "bess_1h":      {"capex_eur_kwh": 280, "duration_h": 1, "om_pct": 2.0, "degradation": 0.020, "life": 15, "label": "BESS 1h"},
    "bess_2h":      {"capex_eur_kwh": 260, "duration_h": 2, "om_pct": 2.0, "degradation": 0.020, "life": 15, "label": "BESS 2h"},
    "bess_4h":      {"capex_eur_kwh": 240, "duration_h": 4, "om_pct": 2.0, "degradation": 0.020, "life": 15, "label": "BESS 4h"},
    "solar":        {"capex_eur_kw": 650,  "cf_pct": 11.5, "om_pct": 1.5, "degradation": 0.005, "life": 25, "label": "Solar PV"},
    "onshore_wind": {"capex_eur_kw": 1400, "cf_pct": 27.0, "om_pct": 2.0, "degradation": 0.000, "life": 25, "label": "Onshore Wind"},
}


def _compute_irr(cashflows: list) -> float:
    rate = 0.10
    for _ in range(200):
        npv  = sum(cf / (1 + rate) ** t for t, cf in enumerate(cashflows))
        dnpv = sum(-t * cf / (1 + rate) ** (t + 1) for t, cf in enumerate(cashflows))
        if abs(dnpv) < 1e-10:
            break
        rate -= npv / dnpv
        if rate <= -0.999:
            rate = -0.999
    return rate


def _run_irr_model_po(
    technology: str,
    capacity_mw: float,
    capex_eur_per_kw: float | None = None,
    revenue_pln_per_mw_yr: float | None = None,
    capacity_factor_pct: float | None = None,
    wacc_pct: float = 9.0,
    project_life_yrs: int | None = None,
    leverage_pct: float = 60.0,
    cost_of_debt_pct: float = 5.5,
) -> dict:
    p = _TECH_PRESETS_PO.get(technology, _TECH_PRESETS_PO["bess_2h"])
    is_bess = technology.startswith("bess")

    if is_bess:
        duration_h = p["duration_h"]
        capex_eur_kw = (capex_eur_per_kw or p["capex_eur_kwh"] * duration_h)
        capex_total_pln = capacity_mw * capex_eur_kw * 1000 * _EUR_PLN
        # Revenue: FCR + aFRR expressed as PLN/MW/yr
        annual_rev_pln = capacity_mw * (revenue_pln_per_mw_yr or 250_000)  # ~250k PLN/MW/yr default
    else:
        capex_eur_kw = capex_eur_per_kw or p["capex_eur_kw"]
        capex_total_pln = capacity_mw * capex_eur_kw * 1000 * _EUR_PLN
        cf = (capacity_factor_pct or p["cf_pct"]) / 100
        annual_gen_mwh = capacity_mw * 8760 * cf
        # Revenue: TGE day-ahead price ~330 PLN/MWh (2024 avg)
        rev_pln_mwh = revenue_pln_per_mw_yr / 8760 if revenue_pln_per_mw_yr else 330
        annual_rev_pln = annual_gen_mwh * rev_pln_mwh

    life = project_life_yrs or p["life"]
    om_annual_pln = capex_total_pln * (p["om_pct"] / 100)
    degradation = p.get("degradation", 0.0)

    cashflows_unlev = [-capex_total_pln]
    for yr in range(1, life + 1):
        rev = annual_rev_pln * (1 - degradation) ** (yr - 1)
        cashflows_unlev.append(rev - om_annual_pln)

    unlevered_irr = _compute_irr(cashflows_unlev)
    npv_wacc = sum(cf / (1 + wacc_pct / 100) ** t for t, cf in enumerate(cashflows_unlev))

    # Levered IRR
    debt = capex_total_pln * (leverage_pct / 100)
    equity = capex_total_pln - debt
    debt_service = debt * (cost_of_debt_pct / 100) / (1 - (1 + cost_of_debt_pct / 100) ** (-min(life, 15)))
    cashflows_eq = [-equity]
    for yr in range(1, life + 1):
        rev = annual_rev_pln * (1 - degradation) ** (yr - 1)
        ds = debt_service if yr <= min(life, 15) else 0
        cashflows_eq.append(rev - om_annual_pln - ds)
    equity_irr = _compute_irr(cashflows_eq)

    # Sensitivity
    sensitivity = []
    for cm in [0.8, 1.0, 1.2]:
        for rm in [0.8, 1.0, 1.2]:
            cfs = [-capex_total_pln * cm]
            for yr in range(1, life + 1):
                rev = annual_rev_pln * rm * (1 - degradation) ** (yr - 1)
                cfs.append(rev - om_annual_pln)
            sensitivity.append({
                "capex": f"{cm:.0%}", "revenue": f"{rm:.0%}",
                "unlevered_irr": f"{_compute_irr(cfs) * 100:.1f}%",
            })

    return {
        "technology": p["label"],
        "capacity_mw": capacity_mw,
        "capex_eur_per_kw": round(capex_eur_kw, 0),
        "capex_total_pln_m": round(capex_total_pln / 1e6, 1),
        "annual_rev_pln_m": round(annual_rev_pln / 1e6, 1),
        "unlevered_irr_pct": round(unlevered_irr * 100, 1),
        "equity_irr_pct": round(equity_irr * 100, 1),
        "npv_at_wacc_pln_m": round(npv_wacc / 1e6, 1),
        "sensitivity": sensitivity,
    }


# ── Embedded market data ──────────────────────────────────────────────────────

_MARKET_STRUCTURE_PO = {
    "installed_capacity_gw": 65,
    "peak_demand_gw": 27,
    "generation_mix_2024": {
        "Coal": 38, "Lignite": 14, "Natural Gas": 8, "Nuclear": 0,
        "Solar": 15, "Wind Onshore": 13, "Hydro": 4, "Other": 8,
    },
    "re_targets": {"2030": "50% RES share (PEP2040)", "2040": "80% RES share"},
    "key_players": [
        "PGE — largest utility; coal + pumped hydro; transitioning to RE",
        "Tauron — coal-heavy; RE ambitions (wind, solar)",
        "Enea — coal + biomass; RE build-out",
        "DTEK / international investors — utility-scale solar+storage",
        "RWE, Statkraft, Total Eren — BESS + wind developers",
        "Orlen — oil major transitioning to offshore wind",
    ],
    "tge_prices_pln_mwh": {
        "2023_avg": 550,
        "2024_avg": 330,
        "2025_fwd": 280,
        "note": "Sharp fall from 2022 energy crisis peak; Aurora forecasts ~200-250 PLN/MWh long-run",
    },
}

_AS_CONTEXT_PO = {
    "fcr": {
        "name": "FCR — Frequency Containment Reserve (Primary)",
        "entso_e_framework": "Symmetric product; activation when freq deviates ±200 mHz from 50 Hz",
        "procurement": "Monthly auctions; European integrated FCR market since 2021",
        "bess_advantage": "BESS provides very fast response (<1s); no fuel cost; high availability",
        "typical_price_eur_mw_week": "2–8 EUR/MW/week (variable; can spike during scarcity)",
        "bess_revenue_pln_mw_yr": "~60,000 – 180,000 PLN/MW/yr (FCR alone)",
    },
    "afrr": {
        "name": "aFRR — Automatic Frequency Restoration Reserve (Secondary)",
        "procurement": "Weekly/monthly auctions by PSE; symmetric and asymmetric products",
        "bess_advantage": "BESS provides both upward and downward aFRR; no fuel cost per activation",
        "capacity_price": "Significant capacity payment component (~20-60 EUR/MW/h depending on auction)",
        "energy_settlement": "Additional activation energy settlement at real-time balancing price",
        "bess_revenue_pln_mw_yr": "~80,000 – 200,000 PLN/MW/yr (aFRR alone)",
    },
    "mfrr": {
        "name": "mFRR — Manual Frequency Restoration Reserve (Tertiary)",
        "procurement": "PSE direct contracting; less relevant for BESS (slower product, 30-min response)",
        "bess_relevance": "Lower priority for BESS vs FCR/aFRR",
    },
    "capacity_market": {
        "name": "Rynek Mocy — Polish Capacity Market",
        "auctions": "T-4 (4 years ahead) and T-1 (1 year ahead) auctions",
        "bess_eligibility": "BESS eligible since 2021; must demonstrate ≥1h duration for T-4, ≥2h preferred",
        "derating": "BESS derated based on available energy duration vs. peak obligation period (4h)",
        "typical_price_pln_kw_yr": "150–250 PLN/kW/yr",
        "bess_revenue_pln_mw_yr": "~150,000 – 250,000 PLN/MW/yr",
    },
    "bess_revenue_stack": [
        "FCR (Primary Reserve): ~60-180k PLN/MW/yr",
        "aFRR (Secondary Reserve): ~80-200k PLN/MW/yr",
        "mFRR (Tertiary): lower priority for BESS",
        "Capacity Market (Rynek Mocy): ~150-250k PLN/MW/yr",
        "Energy arbitrage (TGE day-ahead spread): ~30-80k PLN/MW/yr",
        "Total stacked: ~320-710k PLN/MW/yr (FCR+aFRR+capacity market)",
    ],
}

_POLICY_SNAPSHOT_PO = {
    "pep2040": {
        "name": "Polish Energy Policy to 2040 (PEP2040)",
        "key_targets": {
            "2030": "23% RES in gross final energy consumption; 32 GW+ offshore wind by 2040",
            "offshore_wind": "Phase 1 auctions (5.9 GW) underway; OWF Baltica projects (Orlen/PGE)",
            "nuclear": "First Polish nuclear plant planned ~2033-2035 (6 GW total by 2043)",
        },
    },
    "res_act": "Act on Renewable Energy Sources (RES Act) — support mechanisms: CfD-style auctions (OZE auctions), prosumer net billing",
    "oze_auctions": {
        "mechanism": "Competitive CfD auctions for new RE capacity; 15-year price guarantee",
        "2024": "Active auctions for solar, onshore wind, offshore wind; BESS not yet standalone eligible",
        "bess_path": "BESS expected in future auction rounds or via capacity market + ancillary services",
    },
    "foreign_investment": {
        "restrictions": "No FDI restrictions for RE/storage in Poland; EU Single Market applies",
        "grid_connection": "URE (energy regulator) and DSO/PSE approval required; typical 2-4yr timeline for large BESS",
        "land": "Agricultural land conversion requires Ministry approval; industrial land preferred",
    },
    "key_risks": [
        "Grid connection capacity constraints (distribution network saturation in SW/central Poland)",
        "TGE power price decline (from 550 PLN/MWh in 2023 to ~280-330 PLN/MWh 2024-2025) reduces arbitrage revenue",
        "FCR/aFRR price volatility — European integrated market adds competition from German/Nordic BESS",
        "Capacity market cannibalisation — large BESS pipeline may suppress Rynek Mocy prices by 2027-2030",
        "Permitting risk: zoning and environmental approvals 18-36 months",
    ],
}

_AURORA_KEY_FINDINGS = {
    "source": "Aurora Energy Research Q1/Q2 2026 Poland Power & Renewables Market Forecast",
    "power_prices": {
        "2026_avg_pln_mwh": "~280-310",
        "2030_avg_pln_mwh": "~200-240",
        "2040_avg_pln_mwh": "~160-200 (with high RE penetration)",
        "driver": "Rapid solar/wind build-out suppresses baseload prices; nuclear addition moderates decline",
    },
    "bess_economics": {
        "fcr_afrr_pln_mw_yr_2026": "~200,000 – 350,000",
        "capacity_market_pln_kw_yr": "150-200 (declining as BESS pipeline grows)",
        "optimal_duration": "1-2h for FCR/aFRR; 2-4h for energy arbitrage + capacity market",
        "irr_range_pct": "8-14% unlevered depending on revenue stack and CAPEX",
    },
    "flexible_market_summary_apr26": "Apr 2026 flexible energy market showed sustained FCR prices; aFRR symmetric products attracted new BESS registrations",
}


def _dispatch_tool_po(name: str, inputs: dict) -> str:
    try:
        if name == "search_knowledge_base":
            try:
                return retrieve_for_agent(inputs["query"], _ANTHROPIC_KEY, CFG,
                                          sources=inputs.get("sources") or None, top_k=6)
            except Exception:
                results = _search_knowledge(PREFIX, inputs["query"], limit=6)
                if results.empty:
                    return "No matching knowledge documents found."
                return "\n\n---\n\n".join(
                    f"[{r['source']}] {r['title']} ({r['published_date']})\n{r['snippet']}"
                    for _, r in results.iterrows()
                )

        elif name == "get_aurora_forecast_data":
            topic = inputs.get("topic", "").lower()
            if "price" in topic or "power" in topic:
                return json.dumps(_AURORA_KEY_FINDINGS["power_prices"], indent=2)
            if "bess" in topic or "storage" in topic or "economics" in topic:
                return json.dumps(_AURORA_KEY_FINDINGS["bess_economics"], indent=2)
            return json.dumps(_AURORA_KEY_FINDINGS, indent=2)

        elif name == "get_balancing_market_context":
            return json.dumps(_AS_CONTEXT_PO, indent=2)

        elif name == "get_capacity_market_context":
            return json.dumps(_AS_CONTEXT_PO["capacity_market"], indent=2)

        elif name == "estimate_bess_irr":
            result = _run_irr_model_po(
                technology=inputs.get("technology", "bess_2h"),
                capacity_mw=float(inputs.get("capacity_mw", 50)),
                capex_eur_per_kw=inputs.get("capex_eur_per_kw"),
                revenue_pln_per_mw_yr=inputs.get("revenue_pln_per_mw_yr"),
                wacc_pct=float(inputs.get("wacc_pct", 9.0)),
                project_life_yrs=inputs.get("project_life_yrs"),
                leverage_pct=float(inputs.get("leverage_pct", 60.0)),
                cost_of_debt_pct=float(inputs.get("cost_of_debt_pct", 5.5)),
            )
            return json.dumps(result, indent=2)

        elif name == "get_market_structure":
            return json.dumps(_MARKET_STRUCTURE_PO, indent=2)

        elif name == "get_policy_snapshot":
            return json.dumps(_POLICY_SNAPSHOT_PO, indent=2)

        elif name == "get_entso_price_data":
            try:
                days = int(inputs.get("days", 30))
                df = _po_price_history(days=days)
                if df.empty:
                    return (
                        "No day-ahead price data available yet. "
                        "Use 'Scrape Prices Now' in the Data Management tab to fetch historical prices."
                    )
                latest = df["trading_date"].max()
                avg_pln = df["avg_pln"].mean()
                min_pln = df["min_pln"].min()
                max_pln = df["max_pln"].max()
                summary = {
                    "period": f"Last {days} days",
                    "latest_trading_date": str(latest),
                    "avg_da_price_pln_mwh": round(float(avg_pln), 1),
                    "min_da_price_pln_mwh": round(float(min_pln), 1),
                    "max_da_price_pln_mwh": round(float(max_pln), 1),
                    "avg_da_price_eur_mwh": round(float(avg_pln / _EUR_PLN), 1),
                    "source": "PSE.pl day-ahead market (TGE RDN)",
                    "days_of_data": len(df),
                    "note": "Prices from PSE.pl CSV export. ENTSO-E API available with ENTSOE_API_KEY env var.",
                }
                # Include recent daily series
                recent = df.tail(10)[["trading_date", "avg_pln", "avg_eur"]].copy()
                recent["trading_date"] = recent["trading_date"].astype(str)
                recent["avg_pln"] = recent["avg_pln"].round(1)
                recent["avg_eur"] = recent["avg_eur"].round(1)
                summary["recent_daily_avg"] = recent.to_dict("records")
                return json.dumps(summary, indent=2)
            except Exception as exc:
                return f"Price data query failed: {exc}"

        elif name == "list_knowledge_docs":
            limit = int(inputs.get("limit") or 30)
            try:
                df = _query(
                    f"SELECT source, doc_type, title, url, published_date::text, fetched_at::text "
                    f"FROM intl_market.{PREFIX}knowledge_docs "
                    f"ORDER BY fetched_at DESC LIMIT %s",
                    (limit,),
                )
                if df.empty:
                    return "Knowledge base is empty — no documents have been ingested yet."
                total = _query(f"SELECT COUNT(*) AS n FROM intl_market.{PREFIX}knowledge_docs").iloc[0]["n"]
                header = f"Knowledge base contains {total} documents total. Most recent {len(df)}:\n\n"
                rows = []
                for _, r in df.iterrows():
                    rows.append(
                        f"- [{r['source']}] [{r['doc_type']}] {r['title'] or 'Untitled'} "
                        f"(uploaded: {r['fetched_at'][:16]})"
                    )
                return header + "\n".join(rows)
            except Exception as exc:
                return f"DB error listing docs: {exc}"

    except Exception as exc:
        return f"Tool error: {exc}"
    return "Unknown tool"


_TOOLS_PO = [
    {
        "name": "search_knowledge_base",
        "description": (
            "Semantic search (HyDE + FTS + rerank) over ALL documents in the knowledge base. "
            "This includes: user-uploaded files (PDFs, Word, Excel, PPTX), voice memo transcripts "
            "from expert interviews, meeting minutes, investigation reports, Aurora Energy Research "
            "forecasts, and any document the user has shared today or in previous sessions. "
            "ALWAYS call this tool first for any user question — especially when the user mentions "
            "'uploaded', 'shared', 'sent', 'today', 'meeting minutes', 'report', or any specific document."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string"},
                "sources": {"type": "array", "items": {"type": "string"}},
            },
            "required": ["query"],
        },
    },
    {
        "name": "get_aurora_forecast_data",
        "description": "Returns key projections from Aurora Energy Research Q1/Q2 2026 Poland Power & Renewables Market Forecast — power prices, BESS economics.",
        "input_schema": {
            "type": "object",
            "properties": {"topic": {"type": "string"}},
            "required": [],
        },
    },
    {
        "name": "get_balancing_market_context",
        "description": "PSE balancing market (Rynek Bilansujący), FCR/aFRR/mFRR ancillary services, BESS participation and revenue stack in Poland.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "get_capacity_market_context",
        "description": "Polish capacity market (Rynek Mocy) — T-4/T-1 auction mechanics, BESS eligibility, duration requirements, derating, pricing.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "estimate_bess_irr",
        "description": "Parametric IRR model for a Polish BESS project (PLN, FCR+aFRR revenue stack). Returns unlevered IRR, equity IRR, NPV, sensitivity.",
        "input_schema": {
            "type": "object",
            "properties": {
                "technology":            {"type": "string", "enum": ["bess_1h", "bess_2h", "bess_4h", "solar", "onshore_wind"]},
                "capacity_mw":           {"type": "number"},
                "capex_eur_per_kw":      {"type": "number"},
                "revenue_pln_per_mw_yr": {"type": "number"},
                "wacc_pct":              {"type": "number"},
                "project_life_yrs":      {"type": "integer"},
                "leverage_pct":          {"type": "number"},
                "cost_of_debt_pct":      {"type": "number"},
            },
            "required": ["technology", "capacity_mw"],
        },
    },
    {
        "name": "get_market_structure",
        "description": "Poland power market structure — installed capacity, generation mix, TGE prices, key developers and utilities.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "get_policy_snapshot",
        "description": "Poland energy policy — PEP2040, RES Act, OZE auctions, offshore wind programme, foreign investment rules, key risks.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "get_entso_price_data",
        "description": (
            "Returns Polish day-ahead electricity price statistics from the PSE.pl database "
            "(TGE RDN prices in PLN/MWh and EUR/MWh). Includes recent daily average/min/max "
            "prices from the live scraped dataset. Use this for up-to-date price trend analysis, "
            "arbitrage revenue estimation, and any question about recent Polish power prices."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "days": {"type": "integer", "description": "Number of historical days to summarise (default 30)"},
            },
            "required": [],
        },
    },
    {
        "name": "list_knowledge_docs",
        "description": (
            "List the most recently added documents in the knowledge base, showing title, source, "
            "type, and upload date. Use this to verify exactly which documents are available — "
            "especially when the user says they uploaded something and search_knowledge_base returns "
            "nothing. Call this first to confirm what is actually in the database."
        ),
        "input_schema": {
            "type": "object",
            "properties": {"limit": {"type": "integer", "description": "Max documents to list (default 30)"}},
            "required": [],
        },
    },
]


def _build_system_po(query: str = "") -> str:
    base = """\
You are a senior Poland Power Market Investment Expert at a global infrastructure fund, specialising in BESS and renewable energy.

MANDATORY KB RULE — FOLLOW THIS EVERY TIME:
1. Call `search_knowledge_base` as your FIRST action for EVERY user question, no exceptions.
2. The knowledge base contains ALL user-uploaded documents: PDFs, Excel, Word, PPTX, voice memo \
transcripts from expert interviews, meeting minutes, investigation reports, and any file the user \
has shared in this or previous sessions.
3. If the user asks about something they "uploaded", "shared", "sent", or mentions "today's document", \
"meeting minutes", "investigation", or any named report → search the KB immediately. Never claim \
documents are missing before searching.
4. If the KB search returns relevant content, prioritise it over your embedded knowledge.
5. Only fall back to embedded knowledge for foundational market mechanics when KB returns nothing relevant.

MARKET CONTEXT:
- Market operator: TGE (Towarowa Giełda Energii) — day-ahead (RDN) + intraday (RDT)
- TSO and balancing: PSE (Polskie Sieci Elektroenergetyczne) — Rynek Bilansujący (RB)
- Regulator: URE (Urząd Regulacji Energetyki)
- Installed capacity: ~65 GW; peak demand ~27 GW
- Generation: still coal/lignite dominant but rapid solar build-out (15% share 2024)
- TGE day-ahead: ~330 PLN/MWh (2024 avg); down sharply from 2022-23 energy crisis

ANCILLARY SERVICES (ENTSO-E framework):
- FCR (Primary): symmetric, monthly European auctions; BESS ideal (fast response, high availability)
- aFRR (Secondary): weekly/monthly PSE auctions; symmetric + asymmetric products; capacity + energy payments
- mFRR (Tertiary): PSE direct contracting; less valuable for BESS
- Revenue guidance: FCR+aFRR combined ~140,000-380,000 PLN/MW/yr (2024 observed)

CAPACITY MARKET (RYNEK MOCY):
- T-4 and T-1 competitive auctions; 15-year obligation periods
- BESS eligible since 2021; ~150-250 PLN/kW/yr clearing price (2023-2024)
- Duration requirement: ≥1h for T-4; ≥2h preferred for full derating benefit
- Peak obligation: 4h; BESS derated proportionally for shorter durations

RE POLICY (PEP2040):
- 50% RES in gross final energy by 2030
- Offshore wind Baltic Sea: Phase 1 auctions (OWF Baltica 2&3, ~5.9 GW); Orlen + PGE joint venture
- OZE CfD auctions: 15-year price guarantee; solar + onshore wind active rounds
- First nuclear plant ~2033 (APR-1400 technology); 6 GW total by 2043

CURRENCY: PLN (Polish Zloty); EUR/PLN ≈ 4.25; USD/PLN ≈ 3.95
"""
    if query:
        # Auto-inject KB retrieval so context is present even before the agent calls the tool
        try:
            kb_context = retrieve_for_agent(query, _ANTHROPIC_KEY, CFG)
            if kb_context:
                base += f"\n\n## Pre-fetched knowledge base context for this query:\n{kb_context}"
        except Exception:
            pass
        try:
            insights = get_insights(query, PREFIX, limit=5)
            mem_block = inject_memory(insights, CFG.name)
            if mem_block:
                base += f"\n\n{mem_block}"
        except Exception:
            pass
    mems = _load_memories(APP_KEY)
    if not mems.empty:
        lines = "\n".join(f"- [{r.category}] {r.subject}: {r.content}" for r in mems.itertuples())
        base += f"\n\n## Analyst notes from prior sessions:\n{lines}"
    return base


def _run_agent_turn(messages: list, system: str) -> tuple[str, list, list]:
    tool_events = []
    while True:
        resp = _client.messages.create(
            model="claude-sonnet-4-6", max_tokens=4096,
            system=system, tools=_TOOLS_PO, messages=messages,
        )
        messages = messages + [{"role": "assistant", "content": resp.content}]
        if resp.stop_reason == "end_turn":
            text = next((b.text for b in resp.content if hasattr(b, "text")), "")
            return text, messages, tool_events
        tool_results = []
        for block in resp.content:
            if block.type == "tool_use":
                result_str = _dispatch_tool_po(block.name, block.input) or "(tool returned no output)"
                tool_events.append({"tool": block.name, "result": result_str[:200]})
                tool_results.append({"type": "tool_result", "tool_use_id": block.id, "content": result_str})
        if not tool_results:
            # stop_reason was not end_turn but no tool_use blocks — return any text
            text = next((b.text for b in resp.content if hasattr(b, "text")), "")
            return text, messages, tool_events
        messages = messages + [{"role": "user", "content": tool_results}]


# ── BESS Dispatch & Calibration ───────────────────────────────────────────────

def _run_bess_dispatch_po(
    power_mw: float,
    duration_h: float,
    roundtrip_eff: float,
    price_col: str = "price_pln_mwh",
) -> pd.DataFrame:
    """Run LP perfect-forecast BESS dispatch against Polish day-ahead prices.

    Args:
        power_mw: Arbitrage-slice power rating (MW) — caller passes power × arb_pct/100
        duration_h: Battery duration in hours (energy = power_mw × duration_h MWh)
        roundtrip_eff: Round-trip efficiency (e.g. 0.85)
        price_col: 'price_pln_mwh' or 'price_eur_mwh'

    Returns DataFrame with columns:
        trading_date, pf_profit_pln, naive_profit_pln, options_value_pln,
        charge_mwh, discharge_mwh
    """
    from services.bess_map.optimisation_engine import optimise_day

    prices_df = _query(
        "SELECT trading_date, hour, price_pln_mwh, price_eur_mwh "
        "FROM intl_market.po_day_ahead_prices "
        "ORDER BY trading_date, hour"
    )
    if prices_df.empty:
        return pd.DataFrame(columns=[
            "trading_date", "pf_profit_pln", "naive_profit_pln",
            "options_value_pln", "charge_mwh", "discharge_mwh",
        ])

    # Keep only complete 24-hour days
    day_counts = prices_df.groupby("trading_date")["hour"].count()
    complete_days = day_counts[day_counts == 24].index
    prices_df = prices_df[prices_df["trading_date"].isin(complete_days)]

    rows = []
    for day, grp in prices_df.groupby("trading_date"):
        grp = grp.sort_values("hour")
        prices_arr = grp[price_col].to_numpy(dtype=float)  # PLN/MWh (or EUR/MWh)

        # LP dispatch — prices in PLN/MWh → profit in PLN directly (MW × PLN/MWh × 1h = PLN)
        res = optimise_day(prices_arr, power_mw, duration_h, roundtrip_eff)
        pf_profit = res.profit if res.status == "Optimal" else 0.0

        # Naive: charge at cheapest hour, discharge at most expensive hour (1 cycle)
        min_h, max_h = int(np.argmin(prices_arr)), int(np.argmax(prices_arr))
        eta_c = np.sqrt(roundtrip_eff)
        eta_d = np.sqrt(roundtrip_eff)
        energy_mwh = power_mw * duration_h
        if max_h > min_h:
            naive_profit = (
                prices_arr[max_h] * eta_d * energy_mwh
                - prices_arr[min_h] / eta_c * energy_mwh
            )
        else:
            naive_profit = 0.0

        options_value = max(pf_profit - max(naive_profit, 0.0), 0.0)

        charge_mwh    = float(np.sum(res.charge_mw))    if res.status == "Optimal" else 0.0
        discharge_mwh = float(np.sum(res.discharge_mw)) if res.status == "Optimal" else 0.0

        rows.append({
            "trading_date":     day,
            "pf_profit_pln":    pf_profit,
            "naive_profit_pln": naive_profit,
            "options_value_pln": options_value,
            "charge_mwh":       charge_mwh,
            "discharge_mwh":    discharge_mwh,
        })

    return pd.DataFrame(rows)


def _calibrate_po_strip_params(
    conn,
    peak_start_h: int = 8,
    peak_end_h: int = 20,
    window_days: int = 90,
) -> dict:
    """Calibrate Kirk-Margrabe inputs from po_day_ahead_prices history.

    Args:
        conn: psycopg2 connection (unused directly — uses _query)
        peak_start_h: First peak hour (inclusive), default 8
        peak_end_h: Last peak hour (exclusive), default 20
        window_days: Look-back window in days

    Returns dict:
        peak_forward_pln, offpeak_forward_pln,
        peak_vol, offpeak_vol,
        n_days  (data coverage)
    """
    from datetime import date, timedelta

    cutoff = (date.today() - timedelta(days=window_days)).isoformat()
    df = _query(
        "SELECT trading_date, hour, price_pln_mwh "
        "FROM intl_market.po_day_ahead_prices "
        "WHERE trading_date >= %s AND price_pln_mwh IS NOT NULL "
        "ORDER BY trading_date, hour",
        params=(cutoff,),
    )

    _default = {
        "peak_forward_pln": 0.0, "offpeak_forward_pln": 0.0,
        "peak_vol": 0.30, "offpeak_vol": 0.30, "n_days": 0,
    }

    if df.empty:
        return _default

    is_peak = df["hour"].between(peak_start_h, peak_end_h - 1)
    peak_df    = df[is_peak].groupby("trading_date")["price_pln_mwh"].mean()
    offpeak_df = df[~is_peak].groupby("trading_date")["price_pln_mwh"].mean()

    # Align to common dates
    common = peak_df.index.intersection(offpeak_df.index)
    if len(common) < 5:
        return _default

    peak_series    = peak_df.loc[common].sort_index()
    offpeak_series = offpeak_df.loc[common].sort_index()

    peak_fwd    = float(peak_series.mean())
    offpeak_fwd = float(offpeak_series.mean())

    def _annualised_vol(series: pd.Series) -> float:
        log_ret = np.log(series.values[1:] / np.maximum(series.values[:-1], 1e-6))
        return float(np.std(log_ret) * np.sqrt(252)) if len(log_ret) > 1 else 0.30

    return {
        "peak_forward_pln":    peak_fwd,
        "offpeak_forward_pln": offpeak_fwd,
        "peak_vol":    _annualised_vol(peak_series),
        "offpeak_vol": _annualised_vol(offpeak_series),
        "n_days":      len(common),
    }


# ── UI ────────────────────────────────────────────────────────────────────────

with st.sidebar:
    st.title("🇵🇱 Poland Power Investment")
    st.caption("Investment advisory · FCR/aFRR · Rynek Mocy · BESS")
    st.divider()
    n_insights = _get_insight_count()
    st.metric("Expert Insights", n_insights)
    try:
        kb_ct = _knowledge_doc_counts(PREFIX)
        total_docs = int(kb_ct["docs"].sum()) if not kb_ct.empty else 0
        st.metric("KB Documents", total_docs)
    except Exception:
        pass
    st.divider()
    st.caption("Port 8511 · Europe/Warsaw · ap-southeast-1")

(
    tab_mkt, tab_as, tab_bess, tab_irr,
    tab_advisor, tab_kb, tab_mgmt, tab_pypsa, tab_library,
) = st.tabs([
    "Market Structure", "Balancing & AS Markets", "BESS Opportunity",
    "Investment Analysis", "Investment Advisor",
    "Knowledge Base", "Data Management", "Grid Analysis", "Library",
])


# ═══════════════════════════════════════════════════════════════
# Tab 1 — Market Structure
# ═══════════════════════════════════════════════════════════════
with tab_mkt:
    st.header("Poland Power Market Structure")
    st.caption("Source: PSE, TGE, URE, Aurora Energy Research Q2 2026")

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Installed Capacity", "~65 GW")
    c2.metric("Peak Demand",        "~27 GW")
    c3.metric("TGE DA Price (2024)", "~330 PLN/MWh")
    c4.metric("Solar Share",        "~15%")

    col_mix, col_price = st.columns(2)
    with col_mix:
        st.subheader("Generation Mix (2024)")
        mix_df = pd.DataFrame([
            {"Technology": k, "Share (%)": v}
            for k, v in _MARKET_STRUCTURE_PO["generation_mix_2024"].items()
        ])
        fig_mix = px.pie(
            mix_df, names="Technology", values="Share (%)",
            color_discrete_sequence=px.colors.qualitative.Set2,
        )
        fig_mix.update_traces(textposition="inside", textinfo="percent+label")
        fig_mix.update_layout(height=320, showlegend=False, margin=dict(l=0, r=0, t=10, b=0))
        st.plotly_chart(fig_mix, use_container_width=True)

    with col_price:
        st.subheader("TGE Power Price Trend (PLN/MWh)")
        price_df = pd.DataFrame([
            {"Year": "2022 (crisis)", "PLN/MWh": 800},
            {"Year": "2023 avg",      "PLN/MWh": 550},
            {"Year": "2024 avg",      "PLN/MWh": 330},
            {"Year": "2025 fwd",      "PLN/MWh": 280},
            {"Year": "2030 Aurora",   "PLN/MWh": 220},
            {"Year": "2040 Aurora",   "PLN/MWh": 180},
        ])
        fig_price = px.bar(price_df, x="Year", y="PLN/MWh",
                           text="PLN/MWh", color_discrete_sequence=["#1f77b4"])
        fig_price.update_traces(textposition="outside")
        fig_price.update_layout(height=320, margin=dict(l=0, r=0, t=0, b=0))
        st.plotly_chart(fig_price, use_container_width=True)

    st.subheader("Key Market Participants")
    for p in _MARKET_STRUCTURE_PO["key_players"]:
        st.markdown(f"- {p}")

    st.subheader("RE Targets (PEP2040)")
    st.info(
        "**2030:** 50% RES in gross final energy consumption  \n"
        "**2040:** 80% RES share  \n"
        "**Offshore wind:** 5.9 GW Phase 1 (OWF Baltica 2&3 by Orlen/PGE); 18 GW by 2040  \n"
        "**Nuclear:** First plant ~2033; 6 GW total by 2043"
    )


# ═══════════════════════════════════════════════════════════════
# Tab 2 — Balancing & AS Markets
# ═══════════════════════════════════════════════════════════════
with tab_as:
    st.header("Balancing & Ancillary Services Markets — Poland")
    st.caption("FCR (Primary) · aFRR (Secondary) · mFRR (Tertiary) · Rynek Bilansujący")

    as1, as2 = st.columns(2)
    with as1:
        st.subheader("FCR — Primary Control Reserve")
        fcr = _AS_CONTEXT_PO["fcr"]
        st.markdown(f"""
**Framework:** {fcr['entso_e_framework']}
**Procurement:** {fcr['procurement']}
**BESS advantage:** {fcr['bess_advantage']}
**Typical price:** {fcr['typical_price_eur_mw_week']}
**Indicative BESS revenue:** {fcr['bess_revenue_pln_mw_yr']}
""")

        st.subheader("mFRR — Tertiary")
        mfrr = _AS_CONTEXT_PO["mfrr"]
        st.caption(f"{mfrr['name']} — {mfrr['bess_relevance']}")

    with as2:
        st.subheader("aFRR — Secondary Reserve")
        afrr = _AS_CONTEXT_PO["afrr"]
        st.markdown(f"""
**Procurement:** {afrr['procurement']}
**BESS advantage:** {afrr['bess_advantage']}
**Capacity price:** {afrr['capacity_price']}
**Energy settlement:** {afrr['energy_settlement']}
**Indicative BESS revenue:** {afrr['bess_revenue_pln_mw_yr']}
""")

    st.divider()
    st.subheader("Capacity Market — Rynek Mocy")
    cm = _AS_CONTEXT_PO["capacity_market"]
    col_cm1, col_cm2 = st.columns(2)
    with col_cm1:
        st.markdown(f"""
**Auctions:** {cm['auctions']}
**BESS eligibility:** {cm['bess_eligibility']}
**Derating:** {cm['derating']}
**Typical price:** {cm['typical_price_pln_kw_yr']}
**Indicative revenue:** {cm['bess_revenue_pln_mw_yr']}
""")
    with col_cm2:
        st.subheader("BESS Revenue Stack Summary")
        rev_df = pd.DataFrame([
            {"Stream": "FCR (Primary)", "PLN/MW/yr": "60-180k"},
            {"Stream": "aFRR (Secondary)", "PLN/MW/yr": "80-200k"},
            {"Stream": "Capacity Market (Rynek Mocy)", "PLN/MW/yr": "150-250k"},
            {"Stream": "Energy Arbitrage (TGE spread)", "PLN/MW/yr": "30-80k"},
            {"Stream": "TOTAL (FCR+aFRR+CM)", "PLN/MW/yr": "290-630k"},
        ])
        st.dataframe(rev_df, use_container_width=True, hide_index=True)


# ═══════════════════════════════════════════════════════════════
# Tab 3 — BESS Opportunity
# ═══════════════════════════════════════════════════════════════
with tab_bess:
    st.header("BESS Investment Opportunity — Poland")
    st.caption("Aurora Q1/Q2 2026 · FCR+aFRR+Rynek Mocy revenue stack · EU taxonomy eligible")

    b1, b2 = st.columns([3, 2])
    with b1:
        st.subheader("Why Poland for BESS?")
        st.markdown("""
1. **Deep FCR/aFRR revenue pool** — European integrated FCR market + domestic aFRR auctions
2. **Capacity market (Rynek Mocy)** — BESS eligible; 15-year capacity obligation provides revenue visibility
3. **Rapid RE build-out** → growing flexibility need; intraday price volatility increasing
4. **Grid investment** — PSE 2040 grid modernisation plan includes >35 GW new transmission
5. **EU Taxonomy** compliance — BESS explicitly listed as sustainable finance eligible
6. **No FDI restrictions** — EU Single Market; no foreign ownership limits on storage

**Key risk:** FCR/aFRR price compression as BESS fleet grows and European market integrates further.
""")

    with b2:
        st.subheader("Aurora BESS Economics (2026)")
        aurora_bess = _AURORA_KEY_FINDINGS["bess_economics"]
        for k, v in aurora_bess.items():
            if k != "irr_range_pct":
                st.markdown(f"**{k.replace('_', ' ').title()}:** {v}")
        st.metric("IRR Range (Aurora)", aurora_bess["irr_range_pct"])

    st.divider()
    st.subheader("Optimal BESS Configuration for Poland")
    config_df = pd.DataFrame([
        {"Duration": "1h", "Best for": "FCR only", "Revenue stream": "FCR (primary reserves)", "Note": "Minimum viable for FCR"},
        {"Duration": "2h", "Best for": "FCR + aFRR + partial CM", "Revenue stream": "FCR + aFRR + Rynek Mocy (50% derating)", "Note": "Sweet spot by Aurora 2026"},
        {"Duration": "4h", "Best for": "Full revenue stack", "Revenue stream": "FCR + aFRR + Rynek Mocy (full credit) + arbitrage", "Note": "Higher CAPEX but more CM revenue"},
    ])
    st.dataframe(config_df, use_container_width=True, hide_index=True)

    # ── Subsection A: Perfect-Forecast Dispatch P&L ────────────────────────
    st.divider()
    st.subheader("BESS P&L Analysis — Perfect-Forecast Dispatch")
    st.caption(
        "LP optimal dispatch on arbitrage slice · Compares to naive 1-cycle · "
        "AS revenue stacked from DB average auction prices"
    )

    da_c1, da_c2, da_c3, da_c4 = st.columns(4)
    with da_c1:
        da_power = st.number_input("Power (MW)", min_value=1.0, max_value=1000.0,
                                    value=50.0, step=10.0, key="po_da_power")
    with da_c2:
        da_dur   = st.number_input("Duration (h)", min_value=0.5, max_value=8.0,
                                    value=2.0, step=0.5, key="po_da_dur")
    with da_c3:
        da_eff   = st.number_input("Efficiency (%)", min_value=50.0, max_value=100.0,
                                    value=85.0, step=1.0, key="po_da_eff") / 100.0
    with da_c4:
        da_pcol  = st.selectbox("Price", ["price_pln_mwh", "price_eur_mwh"],
                                 format_func=lambda x: "PLN/MWh" if "pln" in x else "EUR/MWh",
                                 key="po_da_pcol")

    al_c1, al_c2, al_c3 = st.columns(3)
    with al_c1:
        fcr_pct  = st.number_input("FCR allocation (%)", 0.0, 100.0, 20.0, 5.0, key="po_fcr_pct")
    with al_c2:
        afrr_pct = st.number_input("aFRR allocation (%)", 0.0, 100.0, 20.0, 5.0, key="po_afrr_pct")
    with al_c3:
        arb_pct  = 100.0 - fcr_pct - afrr_pct
        st.metric("Arbitrage allocation (%)", f"{arb_pct:.0f}")
        if arb_pct < 0:
            st.error("FCR + aFRR > 100% — reduce allocations")

    if st.button("Run Dispatch Model", type="primary", key="po_run_dispatch",
                 disabled=(arb_pct < 0)):
        with st.spinner("Running LP dispatch…"):
            arb_mw = da_power * arb_pct / 100.0
            dispatch_df = _run_bess_dispatch_po(arb_mw, da_dur, da_eff, da_pcol)

            as_rev = get_as_revenue_estimate(_conn(), da_power, fcr_pct, afrr_pct)

            pf_annual   = float(dispatch_df["pf_profit_pln"].sum())   if not dispatch_df.empty else 0.0
            opts_annual = float(dispatch_df["options_value_pln"].sum()) if not dispatch_df.empty else 0.0
            total_rev   = pf_annual + as_rev["total_pln_yr"]

            st.session_state["po_dispatch_results"] = {
                "arb_pln_yr":      pf_annual,
                "fcr_pln_yr":      as_rev["fcr_pln_yr"],
                "afrr_pln_yr":     as_rev["afrr_pln_yr"],
                "capacity_pln_yr": as_rev["capacity_pln_yr"],
                "total_pln_yr":    total_rev,
                "options_pln_yr":  opts_annual,
                "fcr_pct":         fcr_pct,
                "afrr_pct":        afrr_pct,
                "arb_pct":         arb_pct,
                "df":              dispatch_df,
            }
            st.rerun()

    disp = st.session_state.get("po_dispatch_results")
    if disp is not None:
        m1, m2, m3, m4, m5, m6 = st.columns(6)
        m1.metric("Total Annual Revenue",
                  f"zł{disp['total_pln_yr']/1e6:.2f}M")
        m2.metric("Arbitrage P&L",
                  f"zł{disp['arb_pln_yr']/1e6:.2f}M")
        m3.metric("FCR Revenue",
                  f"zł{disp['fcr_pln_yr']/1e6:.2f}M",
                  f"{disp['fcr_pct']:.0f}% capacity")
        m4.metric("aFRR Revenue",
                  f"zł{disp['afrr_pln_yr']/1e6:.2f}M",
                  f"{disp['afrr_pct']:.0f}% capacity")
        m5.metric("Rynek Mocy",
                  f"zł{disp['capacity_pln_yr']/1e6:.2f}M")
        m6.metric("Options Value",
                  f"zł{disp['options_pln_yr']/1e6:.2f}M",
                  help="PF dispatch premium over naive 1-cycle dispatch")

        df_disp = disp["df"]
        if not df_disp.empty:
            # Chart 1: Daily P&L line
            fig1 = px.line(
                df_disp, x="trading_date", y="pf_profit_pln",
                title="Daily Arbitrage P&L (PLN)",
                labels={"trading_date": "Date", "pf_profit_pln": "Profit (PLN)"},
            )
            st.plotly_chart(fig1, use_container_width=True)

            # Chart 2: Monthly stacked bar (arb + AS layers)
            df2 = df_disp.copy()
            df2["month"] = pd.to_datetime(df2["trading_date"]).dt.to_period("M").astype(str)
            monthly_arb = df2.groupby("month")["pf_profit_pln"].sum().reset_index()
            monthly_arb["FCR"]        = disp["fcr_pln_yr"]        / 12 if disp["fcr_pln_yr"] else 0
            monthly_arb["aFRR"]       = disp["afrr_pln_yr"]       / 12 if disp["afrr_pln_yr"] else 0
            monthly_arb["Rynek Mocy"] = disp["capacity_pln_yr"]   / 12 if disp["capacity_pln_yr"] else 0
            monthly_arb = monthly_arb.rename(columns={"pf_profit_pln": "Arbitrage"})
            fig2 = px.bar(
                monthly_arb.melt(id_vars="month",
                                  value_vars=["Arbitrage", "FCR", "aFRR", "Rynek Mocy"]),
                x="month", y="value", color="variable",
                title="Monthly Revenue Stack (PLN)",
                labels={"month": "Month", "value": "Revenue (PLN)", "variable": "Source"},
                barmode="stack",
            )
            st.plotly_chart(fig2, use_container_width=True)

            # Chart 3: Dispatch profile for selected date
            sel_date = st.selectbox(
                "Dispatch profile date",
                options=sorted(df_disp["trading_date"].unique()),
                key="po_disp_date",
            )
            day_prices = _query(
                "SELECT hour, price_pln_mwh FROM intl_market.po_day_ahead_prices "
                "WHERE trading_date = %s ORDER BY hour",
                params=(str(sel_date),),
            )
            if not day_prices.empty:
                from services.bess_map.optimisation_engine import optimise_day
                arb_mw_sel = da_power * arb_pct / 100.0
                res = optimise_day(
                    day_prices["price_pln_mwh"].to_numpy(dtype=float),
                    arb_mw_sel, da_dur, da_eff,
                )
                fig3 = go.Figure()
                hours = list(range(24))
                fig3.add_bar(x=hours, y=list(-res.charge_mw),
                              name="Charge (MW)", marker_color="steelblue")
                fig3.add_bar(x=hours, y=list(res.discharge_mw),
                              name="Discharge (MW)", marker_color="coral")
                fig3.add_scatter(x=hours, y=day_prices["price_pln_mwh"].tolist(),
                                  name="Price (PLN/MWh)", yaxis="y2",
                                  line=dict(color="gold", width=2))
                fig3.update_layout(
                    title=f"Dispatch Profile — {sel_date}",
                    barmode="relative",
                    yaxis=dict(title="Power (MW)"),
                    yaxis2=dict(title="Price (PLN/MWh)", overlaying="y", side="right"),
                    legend=dict(orientation="h"),
                )
                st.plotly_chart(fig3, use_container_width=True)

    # ── Subsection B: Kirk-Margrabe Strip Valuation ────────────────────────
    st.divider()
    st.subheader("BESS Spread Option Strip Valuation (Kirk-Margrabe)")
    st.caption(
        "Treats BESS as a strip of N daily peak/offpeak spread call options · "
        "Calibrated from last 90 days of TGE day-ahead prices"
    )

    km_c1, km_c2 = st.columns(2)
    with km_c1:
        km_peak_start = st.slider("Peak hours start", 6, 12, 8, key="po_km_pk_start")
        km_peak_end   = st.slider("Peak hours end",   14, 22, 20, key="po_km_pk_end")
        km_om_cost    = st.number_input("O&M cost / strike K (PLN/MWh)",
                                         0.0, 200.0, 20.0, 5.0, key="po_km_om")
        km_horizon    = st.number_input("Valuation horizon (days)",
                                         30, 730, 365, 30, key="po_km_horizon")
    with km_c2:
        km_corr = st.slider("Peak/offpeak correlation", 0.0, 1.0, 0.85, 0.05,
                              key="po_km_corr")
        _prev = st.session_state.get("po_dispatch_results", {})
        _prev_arb_pct = _prev.get("arb_pct", 100) if _prev else 100
        _prev_power   = st.session_state.get("po_da_power", 50.0) or 50.0
        _prev_dur     = st.session_state.get("po_da_dur", 2.0) or 2.0
        _prev_eff_pct = st.session_state.get("po_da_eff", 85.0) or 85.0  # widget stores % (85.0), not fraction
        km_power = st.number_input("Power (MW)", 1.0, 1000.0,
                                    float(_prev_arb_pct / 100 * _prev_power),
                                    10.0, key="po_km_power")
        km_dur   = st.number_input("Duration (h)", 0.5, 8.0,
                                    float(_prev_dur), 0.5, key="po_km_dur")
        km_eff   = st.number_input("Efficiency (%)", 50.0, 100.0,
                                    float(_prev_eff_pct),
                                    1.0, key="po_km_eff") / 100.0

    if st.button("Value Strip", type="primary", key="po_km_run"):
        with st.spinner("Calibrating from price history and pricing strip…"):
            params = _calibrate_po_strip_params(
                _conn(), km_peak_start, km_peak_end, window_days=90
            )

            if params["peak_forward_pln"] == 0.0:
                st.warning("Insufficient price history for calibration (< 5 days). "
                           "Scrape more day-ahead prices first.")
            else:
                from libs.decision_models.bess_spread_call_strip import _run as _km_run

                km_result = _km_run(
                    asset_code="PO-BESS",
                    as_of_date=str(pd.Timestamp.today().date()),
                    n_days_remaining=int(km_horizon),
                    peak_forward_yuan=params["peak_forward_pln"],
                    offpeak_forward_yuan=params["offpeak_forward_pln"],
                    peak_vol=params["peak_vol"],
                    offpeak_vol=params["offpeak_vol"],
                    peak_offpeak_corr=km_corr,
                    roundtrip_eff=km_eff,
                    power_mw=km_power,
                    duration_h=km_dur,
                    om_cost_yuan_per_mwh=km_om_cost,
                )

                # Display metrics (output fields use "_yuan" suffix but values are PLN)
                sv    = km_result["strip_value_yuan"]
                iv    = km_result["intrinsic_value_yuan"]
                tv    = km_result["time_value_yuan"]
                mon   = km_result["moneyness_pct"]
                delta = km_result["delta_yuan_per_yuan"]
                vega  = km_result["vega_yuan_per_vol_point"]

                km1, km2, km3, km4, km5, km6 = st.columns(6)
                km1.metric("Strip Value",     f"zł{sv/1e6:.2f}M")
                km2.metric("Intrinsic Value", f"zł{iv/1e6:.2f}M")
                km3.metric("Time Value",      f"zł{tv/1e6:.2f}M")
                km4.metric(
                    "Moneyness",
                    f"{mon:+.1f}%",
                    delta="ITM" if mon > 0 else "OTM",
                    delta_color="normal" if mon > 0 else "inverse",
                )
                km5.metric("Delta", f"{delta:.3f}")
                km6.metric("Vega",  f"zł{vega/1e3:.1f}K / 1% vol")

                with st.expander("Calibration details"):
                    col_a, col_b = st.columns(2)
                    col_a.markdown(
                        f"**Peak forward:** zł{params['peak_forward_pln']:.1f}/MWh  \n"
                        f"**Peak vol:** {params['peak_vol']*100:.1f}%  \n"
                        f"**Data window:** {params['n_days']} days"
                    )
                    col_b.markdown(
                        f"**Offpeak forward:** zł{params['offpeak_forward_pln']:.1f}/MWh  \n"
                        f"**Offpeak vol:** {params['offpeak_vol']*100:.1f}%  \n"
                        f"**Spread vol:** {km_result['spread_vol_used']*100:.1f}%"
                    )


# ═══════════════════════════════════════════════════════════════
# Tab 4 — Investment Analysis
# ═══════════════════════════════════════════════════════════════
with tab_irr:
    st.header("Investment Analysis — BESS IRR Calculator")
    st.caption("Parametric model · PLN · EUR/PLN = 4.25 · FCR+aFRR+Capacity Market revenue stack")

    irr_c1, irr_c2 = st.columns([1, 1])
    with irr_c1:
        tech_opts = {v["label"]: k for k, v in _TECH_PRESETS_PO.items()}
        sel_lbl = st.selectbox("Technology", list(tech_opts.keys()), key="po_tech")
        tech_key = tech_opts[sel_lbl]
        p = _TECH_PRESETS_PO[tech_key]

        is_bess = tech_key.startswith("bess")
        cap_mw = st.number_input("Project Capacity (MW)", min_value=1.0, max_value=2000.0,
                                  value=50.0, step=10.0, key="po_cap")

        if is_bess:
            dur = p["duration_h"]
            default_capex = p["capex_eur_kwh"] * dur
            capex_val = st.number_input("CAPEX (EUR/kW-AC)", min_value=100.0, max_value=2000.0,
                                         value=float(default_capex), step=50.0, key="po_capex")
            st.caption("Revenue = FCR + aFRR + Rynek Mocy combined (PLN/MW/yr)")
            # Load from dispatch model if available
            _disp = st.session_state.get("po_dispatch_results")
            if _disp is not None:
                _total_pln_yr = _disp["total_pln_yr"]
                _total_mw_yr  = _total_pln_yr / da_power if da_power else _total_pln_yr
                if st.button(
                    f"📥 Load from dispatch model  (zł{_total_pln_yr/1e6:.2f}M/yr total)",
                    key="po_load_dispatch",
                ):
                    st.session_state["po_irr_rev_override"] = _total_mw_yr
                with st.expander("Revenue breakdown", expanded=False):
                    _total = _disp["total_pln_yr"]
                    for _label, _key in [
                        ("Arbitrage (PF dispatch)", "arb_pln_yr"),
                        ("FCR", "fcr_pln_yr"),
                        ("aFRR", "afrr_pln_yr"),
                        ("Rynek Mocy", "capacity_pln_yr"),
                    ]:
                        _val = _disp.get(_key, 0.0)
                        _pct = _val / _total * 100 if _total else 0
                        st.markdown(f"**{_label}:** zł{_val/1e6:.2f}M &nbsp;&nbsp; `{_pct:.0f}%`")
                    st.markdown(f"---  \n**Total:** zł{_total/1e6:.2f}M")
            _rev_default = st.session_state.pop("po_irr_rev_override", 300_000.0)
            rev_val = st.number_input("Combined Revenue (PLN/MW/yr)",
                                       min_value=0.0, max_value=2_000_000.0,
                                       value=float(_rev_default), step=10_000.0, key="po_rev")
            cf_val = None
        else:
            capex_val = st.number_input("CAPEX (EUR/kW)", min_value=200.0, max_value=3000.0,
                                         value=float(p["capex_eur_kw"]), step=50.0, key="po_capex_re")
            cf_val = st.number_input("Capacity Factor (%)", min_value=5.0, max_value=60.0,
                                      value=float(p["cf_pct"]), step=1.0, key="po_cf")
            rev_val = st.number_input("TGE price assumption (PLN/MWh) → as PLN/MW/yr",
                                       min_value=50_000.0, max_value=600_000.0,
                                       value=280_000.0, step=10_000.0, key="po_rev_re")

        with st.expander("Advanced parameters"):
            wacc_val      = st.number_input("WACC (%)",           min_value=5.0,  max_value=20.0, value=9.0, step=0.5, key="po_wacc")
            life_val      = st.number_input("Project Life (yrs)", min_value=5,    max_value=30,   value=int(p["life"]), step=1, key="po_life")
            leverage_val  = st.number_input("Leverage (%)",        min_value=0.0,  max_value=80.0, value=60.0, step=5.0, key="po_lev")
            debt_rate_val = st.number_input("Cost of Debt (%)",    min_value=2.0,  max_value=12.0, value=5.5, step=0.5, key="po_debt")

    with irr_c2:
        if st.button("Calculate IRR", type="primary", key="po_calc_irr"):
            result = _run_irr_model_po(
                technology=tech_key, capacity_mw=cap_mw,
                capex_eur_per_kw=capex_val,
                revenue_pln_per_mw_yr=rev_val if is_bess else None,
                capacity_factor_pct=cf_val, wacc_pct=wacc_val,
                project_life_yrs=int(life_val), leverage_pct=leverage_val,
                cost_of_debt_pct=debt_rate_val,
            )
            st.session_state["po_irr_result"] = result
            st.rerun()

        res = st.session_state.get("po_irr_result")
        if res:
            r1, r2 = st.columns(2)
            r1.metric("Unlevered IRR", f"{res['unlevered_irr_pct']:.1f}%")
            r2.metric("Equity IRR",    f"{res['equity_irr_pct']:.1f}%")
            r3, r4 = st.columns(2)
            r3.metric("Total CAPEX",   f"zł{res['capex_total_pln_m']:.0f}M")
            r4.metric("NPV @ WACC",    f"zł{res['npv_at_wacc_pln_m']:.0f}M")
            st.caption(f"Annual revenue: zł{res['annual_rev_pln_m']:.1f}M")

            st.subheader("Sensitivity (Unlevered IRR) — CAPEX × Revenue")
            sens = res.get("sensitivity", [])
            if sens:
                sens_df = pd.DataFrame(sens)
                pivot = sens_df.pivot(index="capex", columns="revenue", values="unlevered_irr")
                st.dataframe(pivot, use_container_width=True)


# ═══════════════════════════════════════════════════════════════
# Tab 5 — Investment Advisor
# ═══════════════════════════════════════════════════════════════
with tab_advisor:
    st.header("Investment Advisor — Poland Power")
    n_ins = _get_insight_count()
    st.caption(
        f"Senior Poland Power Investment Expert · 7 tools · "
        f"Expert memory: {n_ins} insights · Aurora Research grounded"
    )

    if "po_adv_session_id" not in st.session_state:
        st.session_state["po_adv_session_id"] = str(uuid.uuid4())
    if "po_adv_history" not in st.session_state:
        st.session_state["po_adv_history"] = []

    if not st.session_state["po_adv_history"]:
        recent = _list_recent_sessions(PREFIX)
        if not recent.empty:
            with st.expander("Resume a previous conversation?", expanded=False):
                for _, srow in recent.iterrows():
                    lbl = (f"{srow['session_id'][:8]}… — "
                           f"{srow['updated_at'].strftime('%Y-%m-%d %H:%M')} — "
                           f"{int(srow['msg_count'])} messages")
                    if st.button(lbl, key=f"po_resume_{srow['session_id']}"):
                        st.session_state["po_adv_session_id"] = srow["session_id"]
                        st.session_state["po_adv_history"] = _load_session(srow["session_id"])
                        st.rerun()

    if not st.session_state["po_adv_history"]:
        st.markdown("**Quick-start questions:**")
        qq_c1, qq_c2 = st.columns(2)
        quick_qs = [
            "What is the investment case for BESS in Poland in 2026?",
            "Model IRR for a 50MW / 2h BESS targeting FCR + aFRR + Rynek Mocy",
            "How have Aurora's power price forecasts for Poland changed in Q1 vs Q2 2026?",
            "What are the key risks for a foreign developer entering the Polish BESS market?",
        ]
        for i, qq in enumerate(quick_qs):
            col = qq_c1 if i % 2 == 0 else qq_c2
            if col.button(qq, key=f"po_qq_{i}"):
                st.session_state["po_adv_history"].append({"role": "user", "content": qq})
                st.rerun()

    for msg in st.session_state["po_adv_history"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    user_input = st.chat_input("Ask about Poland power markets — FCR, aFRR, Rynek Mocy, BESS IRR, Aurora forecasts…")
    if user_input:
        st.session_state["po_adv_history"].append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.markdown(user_input)
        with st.chat_message("assistant"):
            with st.spinner("Analysing…"):
                api_msgs = [{"role": m["role"], "content": m["content"]}
                            for m in st.session_state["po_adv_history"]]
                try:
                    reply, _, tool_events = _run_agent_turn(api_msgs, _build_system_po(user_input))
                except Exception as err:
                    reply = f"API error: {err}. Please try again."
                    tool_events = []
            st.markdown(reply)
            if tool_events:
                with st.expander(f"Tools used ({len(tool_events)})", expanded=False):
                    for ev in tool_events:
                        st.caption(f"**{ev['tool']}** → {ev['result'][:120]}…")

        st.session_state["po_adv_history"].append({"role": "assistant", "content": reply})
        try:
            _save_session(st.session_state["po_adv_session_id"], st.session_state["po_adv_history"])
        except Exception:
            pass
        try:
            n_new = extract_insights(user_input, reply, _ANTHROPIC_KEY, PREFIX, CFG.name)
            if n_new:
                st.toast(f"Stored {n_new} expert insight(s)")
        except Exception:
            pass

    if st.session_state.get("po_adv_history"):
        st.divider()
        exp_col1, exp_col2, exp_col3, exp_col4 = st.columns([3, 2, 2, 2])
        with exp_col1:
            if st.button("Clear conversation", key="po_clear_adv"):
                st.session_state["po_adv_history"] = []
                st.session_state["po_adv_session_id"] = str(uuid.uuid4())
                st.rerun()
        # Export buttons — only show when there is at least one assistant reply
        has_answers = any(m["role"] == "assistant" for m in st.session_state["po_adv_history"])
        if has_answers:
            _exp_title = f"Poland Power Investment Advisory — {date.today()}"
            with exp_col2:
                try:
                    _pdf_bytes = export_pdf(st.session_state["po_adv_history"], _exp_title, CFG.name)
                    st.download_button(
                        "📄 Export PDF", _pdf_bytes,
                        file_name=f"po_advisory_{date.today()}.pdf",
                        mime="application/pdf", key="po_exp_pdf",
                    )
                except Exception as _e:
                    st.error(f"PDF failed: {_e}")
            with exp_col3:
                try:
                    _pptx_bytes = export_pptx(st.session_state["po_adv_history"], _exp_title, CFG.name)
                    st.download_button(
                        "📊 Export PPT", _pptx_bytes,
                        file_name=f"po_advisory_{date.today()}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="po_exp_pptx",
                    )
                except Exception as _e:
                    st.error(f"PPT failed: {_e}")
            with exp_col4:
                try:
                    _docx_bytes = export_docx(st.session_state["po_adv_history"], _exp_title, CFG.name)
                    st.download_button(
                        "📝 Export Word", _docx_bytes,
                        file_name=f"po_advisory_{date.today()}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key="po_exp_docx",
                    )
                except Exception as _e:
                    st.error(f"Word failed: {_e}")


# ═══════════════════════════════════════════════════════════════
# Tab 6 — Knowledge Base
# ═══════════════════════════════════════════════════════════════
with tab_kb:
    st.header("Poland Knowledge Base")
    st.info("Aurora Energy Research Q1/Q2 2026 PDFs + Excel · Auto-updated daily at **03:30 WAW**", icon="🔄")

    kb_c1, kb_c2 = st.columns([2, 1])
    with kb_c1:
        kb_counts = _knowledge_doc_counts(PREFIX)
        if kb_counts.empty:
            st.info("Knowledge base is empty. Click 'Auto-ingest Local Reports' to seed from Aurora data.")
        else:
            st.dataframe(kb_counts, use_container_width=True, hide_index=True)
    with kb_c2:
        if st.button("Refresh KB stats", key="po_kb_refresh"):
            _knowledge_doc_counts.clear()
            st.rerun()
        if st.button("Auto-ingest Local Reports", type="primary", key="po_kb_ingest_local"):
            with st.spinner("Ingesting from data/market-fundamentals-po/…"):
                try:
                    results = run_knowledge_ingest(only=["local_reports"], verbose=False)
                    st.success(f"Done — {results.get('local_reports', 0)} new docs")
                except Exception as exc:
                    st.error(f"Failed: {exc}")
            _knowledge_doc_counts.clear()
            st.rerun()
        if st.button("Ingest All Sources", key="po_kb_ingest_all"):
            with st.spinner("Fetching from all sources…"):
                try:
                    results = run_knowledge_ingest(verbose=False)
                    total = sum(results.values())
                    st.success(f"Done — {total} new docs: {results}")
                except Exception as exc:
                    st.error(f"Failed: {exc}")
            _knowledge_doc_counts.clear()
            st.rerun()

    st.divider()
    st.subheader("Search Knowledge Base")
    kb_query = st.text_input("Search query", placeholder="e.g. Aurora FCR aFRR Poland BESS revenue", key="po_kb_query")
    if kb_query:
        _search_knowledge.clear()
        results_df = _search_knowledge(PREFIX, kb_query, limit=10)
        if results_df.empty:
            st.info("No results found.")
        else:
            for _, row in results_df.iterrows():
                with st.expander(
                    f"[{row['source']}] {row['title'] or 'Untitled'} — {row['published_date']}",
                    expanded=False,
                ):
                    if row.get("url"):
                        st.markdown(f"[Source]({row['url']})")
                    st.text(row["snippet"])

    st.divider()
    st.subheader("Upload Documents")
    up1, up2 = st.tabs(["Upload Files", "Fetch from URL"])
    with up1:
        uploaded_files = st.file_uploader(
            "PDF, Excel, PPTX, Word, TXT",
            type=["pdf", "xlsx", "xls", "pptx", "ppt", "docx", "doc", "txt"],
            accept_multiple_files=True, key="po_kb_upload",
        )
        if uploaded_files and st.button("Ingest uploaded files", type="primary", key="po_upload_btn"):
            prog = st.progress(0)
            ok, errs = [], []
            for i, f in enumerate(uploaded_files):
                res = _ingest_uploaded_file(f.name, f.read())
                (ok if res["status"] == "success" else errs).append((f.name, res))
                prog.progress((i + 1) / len(uploaded_files))
            prog.empty()
            if ok:
                st.success(f"Ingested {len(ok)} file(s).")
            for fname, r in errs:
                st.error(f"{fname}: {r['msg']}")
            _knowledge_doc_counts.clear()
            st.rerun()
    with up2:
        fetch_url = st.text_input("Article URL", key="po_fetch_url")
        if st.button("Fetch and ingest", type="primary", key="po_fetch_btn") and fetch_url:
            with st.spinner("Fetching…"):
                res = _ingest_url(fetch_url.strip())
            if res["status"] == "success":
                st.success(res["msg"])
                _knowledge_doc_counts.clear()
            else:
                st.error(res["msg"])

    st.divider()
    st.subheader("Voice Memo Interviews")
    st.caption(
        "Upload iPhone voice memos (.m4a) or other audio recordings from domain expert meetings. "
        "Whisper transcribes the audio, then Claude fixes domain terminology and extracts key points."
    )

    _openai_key = os.environ.get("OPENAI_API_KEY", "")
    if not _openai_key:
        st.warning("OPENAI_API_KEY not set in environment. Enter it below for this session only.")
        _openai_key = st.text_input("OpenAI API Key", type="password", key="po_oai_key_input",
                                    help="Required for Whisper audio transcription")

    _audio_file = st.file_uploader(
        "Upload audio recording",
        type=["m4a", "mp3", "mp4", "wav", "webm", "mpeg", "mpga"],
        key="po_audio_upload",
    )
    _speaker_ctx = st.text_input(
        "Speaker / meeting context (optional)",
        placeholder="e.g. Interview with PSE grid operator about FCR capacity requirements 2025",
        key="po_audio_ctx",
    )

    if _audio_file and _openai_key:
        _audio_bytes = _audio_file.read()
        _size_mb = len(_audio_bytes) / (1024 * 1024)
        st.caption(f"File: **{_audio_file.name}** · {_size_mb:.1f} MB")
        if _size_mb > 100:
            st.error("File too large (>100 MB). Please trim the recording before uploading.")
        else:
            if st.button("Transcribe & Contextualize", type="primary", key="po_transcribe_btn"):
                st.session_state.pop("po_audio_result", None)
                with st.status("Processing audio…", expanded=True) as _status:
                    st.write(f"Transcribing with Whisper ({'chunked' if _size_mb > 24 else 'single pass'})…")
                    try:
                        _result = transcribe_and_contextualize(
                            _audio_bytes, _audio_file.name,
                            _openai_key, _ANTHROPIC_KEY,
                            CFG.name, _speaker_ctx,
                        )
                        st.write("Contextualizing with domain knowledge…")
                        st.session_state["po_audio_result"] = _result
                        st.session_state["po_audio_filename"] = _audio_file.name
                        _status.update(label="Done!", state="complete")
                    except Exception as _err:
                        _status.update(label=f"Failed: {_err}", state="error")
                        st.error(str(_err))

    if st.session_state.get("po_audio_result"):
        _res = st.session_state["po_audio_result"]
        _fname = st.session_state.get("po_audio_filename", "recording")
        st.divider()
        with st.expander("Raw Whisper transcript", expanded=False):
            st.text_area("Raw", _res["raw_transcript"], height=200, disabled=True, key="po_raw_transcript")
        st.subheader("Cleaned transcript")
        st.markdown(_res["polished_content"])
        _memo_title = st.text_input(
            "Document title (edit if needed)",
            value=_res["title"], key="po_memo_title",
        )
        if st.button("Add to Knowledge Base", type="primary", key="po_memo_add"):
            try:
                store_voice_memo(_res, _fname, _conn(), PREFIX, custom_title=_memo_title)
                st.success(f"Added **{_memo_title}** to knowledge base.")
                st.session_state.pop("po_audio_result", None)
                _knowledge_doc_counts.clear()
            except Exception as _exc:
                st.error(f"DB insert failed: {_exc}")


# ═══════════════════════════════════════════════════════════════
# Tab 7 — Data Management
# ═══════════════════════════════════════════════════════════════
with tab_mgmt:
    st.header("Data Management")

    dm_c1, dm_c2 = st.columns(2)
    with dm_c1:
        st.subheader("Table Coverage")
        if st.button("Refresh counts", key="po_refresh_counts"):
            _po_table_counts.clear()
        counts_df = _po_table_counts(PREFIX)
        st.dataframe(counts_df, use_container_width=True, hide_index=True)
    with dm_c2:
        st.subheader("Scheduler Status")
        try:
            sched = _start_scheduler()
            jobs = sched.get_jobs()
            if jobs:
                next_run = min((j.next_run_time for j in jobs if j.next_run_time), default=None)
                next_str = next_run.strftime("%Y-%m-%d %H:%M WAW") if next_run else "—"
                st.success(f"Running · Next job: **{next_str}**")
                for j in jobs:
                    nrt = j.next_run_time.strftime("%H:%M") if j.next_run_time else "—"
                    st.caption(f"`{j.id}` — next: {nrt}")
        except Exception as exc:
            st.error(f"Scheduler error: {exc}")

    st.divider()

    # ── Day-Ahead Price Data ──────────────────────────────────────────────────
    st.subheader("TGE Day-Ahead Price Data")
    st.caption(
        "Hourly day-ahead prices (EUR/MWh → PLN/MWh) from **energy-charts.info** "
        "(Fraunhofer ISE, EPEX SPOT Poland, free, no auth). "
        "Optional ENTSO-E API fallback (set **ENTSOE_API_KEY** env var). "
        "Scheduler runs at **07:15 WAW** daily."
    )

    price_stat_c1, price_stat_c2, price_stat_c3 = st.columns(3)
    with price_stat_c1:
        latest_price_date = _po_latest_price_date()
        st.metric("Latest trading date", latest_price_date.split(" ")[0] if latest_price_date != "—" else "—")
    with price_stat_c2:
        try:
            n_price_rows = _query("SELECT COUNT(*) AS n FROM intl_market.po_day_ahead_prices")
            st.metric("Price rows stored", int(n_price_rows["n"].iloc[0]))
        except Exception:
            st.metric("Price rows stored", "—")
    with price_stat_c3:
        try:
            n_days_df = _query("SELECT COUNT(DISTINCT trading_date) AS n FROM intl_market.po_day_ahead_prices")
            st.metric("Trading days", int(n_days_df["n"].iloc[0]))
        except Exception:
            st.metric("Trading days", "—")

    price_btn_c1, price_btn_c2 = st.columns(2)
    with price_btn_c1:
        if st.button("Scrape Prices Now", type="primary", key="po_price_scrape_now"):
            with st.spinner("Fetching prices from PSE.pl…"):
                try:
                    import psycopg2
                    _pconn = psycopg2.connect(
                        os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                        keepalives=1, keepalives_idle=30,
                    )
                    _pconn.autocommit = True
                    results = run_entso_price_scrape(_pconn, days_back=1)
                    _pconn.close()
                    _po_latest_price_date.clear()
                    _po_price_history.clear()
                    total = sum(v for v in results.values() if v > 0)
                    if total > 0:
                        st.success(f"Fetched {total} new price records: {results}")
                    else:
                        st.info(
                            f"All records already stored or no new data ({results}). "
                            "PSE.pl publishes day-ahead prices with ~1-day lag. "
                            "Verify the latest date shown above."
                        )
                except Exception as exc:
                    st.error(f"Price scrape failed: {exc}")
    with price_btn_c2:
        backfill_price_days = st.number_input(
            "Backfill days", min_value=1, max_value=365, value=30, key="po_price_backfill_days"
        )
        if st.button("Backfill Price History", key="po_price_backfill"):
            with st.spinner(f"Backfilling {backfill_price_days} days of day-ahead prices…"):
                try:
                    import psycopg2
                    _pconn = psycopg2.connect(
                        os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                        keepalives=1, keepalives_idle=30,
                    )
                    _pconn.autocommit = True
                    results = run_entso_price_scrape(_pconn, days_back=int(backfill_price_days))
                    _pconn.close()
                    _po_latest_price_date.clear()
                    _po_price_history.clear()
                    total = sum(v for v in results.values() if v > 0)
                    st.success(f"Backfill complete: {total} new rows across {len(results)} dates.")
                    with st.expander("Per-date results"):
                        for dt, n in sorted(results.items(), reverse=True):
                            st.caption(f"{dt}: {n} rows")
                except Exception as exc:
                    st.error(f"Backfill failed: {exc}")

    # Price chart
    price_df = _po_price_history(days=30)
    if not price_df.empty:
        st.markdown("**Last 30 days — Daily average day-ahead price**")
        try:
            fig_p = px.line(
                price_df[price_df["avg_pln"].notna()],
                x="trading_date", y="avg_pln",
                labels={"trading_date": "Date", "avg_pln": "PLN/MWh"},
                color_discrete_sequence=["#1f77b4"],
            )
            fig_p.update_layout(height=300, margin=dict(t=20, b=20))
            st.plotly_chart(fig_p, use_container_width=True)
        except Exception:
            st.dataframe(price_df[["trading_date", "avg_pln", "avg_eur"]],
                         use_container_width=True, hide_index=True)
    else:
        st.info("No day-ahead price data yet. Click 'Scrape Prices Now' above to fetch.")

    st.caption(
        "Source: **energy-charts.info** (Fraunhofer ISE, CC BY 4.0). "
        "Optional: set `ENTSOE_API_KEY` env var for direct ENTSO-E Transparency Platform API access."
    )

    st.divider()

    # ── Polish Market Document Sources ────────────────────────────────────────
    st.subheader("Polish Market Document Sources")
    st.caption(
        "PSE balancing/grid/aFRR reports · TGE market reports · URE regulatory publications · "
        "ENTSO-E news and publications. Scheduled every **Monday 04:05 WAW**."
    )

    _PO_DOC_SOURCE_LABELS = [
        ("local_reports",       "Local Reports (Aurora PDFs/Excel)"),
        ("pse_pl",              "PSE Balancing Reports"),
        ("pse_grid",            "PSE Grid/Transmission Reports"),
        ("pse_afrr",            "PSE aFRR/FCR Tender Results"),
        ("tge_reports",         "TGE Market Reports"),
        ("ure_regulatory",      "URE Regulatory Publications"),
        ("entsoe_publications",  "ENTSO-E News & Publications"),
    ]

    # Status table
    try:
        src_rows = []
        for src_key, src_label in _PO_DOC_SOURCE_LABELS:
            row_df = _query(
                f"SELECT COUNT(*) AS n, MAX(fetched_at) AS last_fetched "
                f"FROM intl_market.{PREFIX}knowledge_docs WHERE source=%s",
                (src_key,),
            )
            n_docs = int(row_df["n"].iloc[0]) if not row_df.empty else 0
            last_ft = row_df["last_fetched"].iloc[0] if not row_df.empty else None
            last_str = str(last_ft)[:16] if last_ft is not None else "Never"
            src_rows.append({"Source": src_label, "Docs": n_docs, "Last Fetched": last_str})
        st.dataframe(pd.DataFrame(src_rows), use_container_width=True, hide_index=True)
    except Exception as exc:
        st.warning(f"Could not load source stats: {exc}")

    doc_btn_c1, doc_btn_c2 = st.columns(2)
    with doc_btn_c1:
        if st.button("Auto-ingest Local Reports", type="primary", key="po_ingest_local_dm"):
            with st.spinner("Ingesting from data/market-fundamentals-po/…"):
                try:
                    results = run_knowledge_ingest(only=["local_reports"], verbose=False)
                    st.success(f"Done — {results.get('local_reports', 0)} new docs")
                    _knowledge_doc_counts.clear()
                except Exception as exc:
                    st.error(f"Failed: {exc}")
    with doc_btn_c2:
        if st.button("Fetch All Online Sources Now", key="po_docs_fetch_all"):
            with st.spinner("Scraping PSE/TGE/URE/ENTSO-E (may take 2–5 min)…"):
                try:
                    online_keys = [s for s, _ in _PO_DOC_SOURCE_LABELS if s != "local_reports"]
                    results = run_knowledge_ingest(only=online_keys, verbose=False)
                    total = sum(v for v in results.values() if isinstance(v, int) and v > 0)
                    st.success(f"Done — {total} new documents added.")
                    with st.expander("Per-source results"):
                        for src_key, src_label in _PO_DOC_SOURCE_LABELS:
                            if src_key in results:
                                st.caption(f"{src_label}: {results.get(src_key, 0)} new docs")
                    _knowledge_doc_counts.clear()
                except Exception as exc:
                    st.error(f"Fetch failed: {exc}")

    st.markdown("**Historical Backfill**")
    st.caption("Paginate deeply through listing pages to retrieve documents since a chosen start date.")

    po_bf_c1, po_bf_c2 = st.columns(2)
    with po_bf_c1:
        po_backfill_start = st.date_input(
            "Backfill from date", value=date(2023, 1, 1), key="po_bf_start"
        )
        po_backfill_sources = st.multiselect(
            "Sources to backfill",
            options=[s for s, _ in _PO_DOC_SOURCE_LABELS if s != "local_reports"],
            format_func=lambda s: dict(_PO_DOC_SOURCE_LABELS).get(s, s),
            default=[s for s, _ in _PO_DOC_SOURCE_LABELS if s != "local_reports"],
            key="po_bf_sources",
        )
    with po_bf_c2:
        st.markdown("")
        st.markdown("")
        if st.button("Run Historical Backfill", type="primary", key="po_bf_run"):
            if not po_backfill_sources:
                st.warning("Select at least one source.")
            else:
                with st.status(
                    f"Backfilling {len(po_backfill_sources)} sources since {po_backfill_start}…",
                    expanded=True,
                ) as bf_status:
                    try:
                        import psycopg2
                        _bf_conn = psycopg2.connect(
                            os.environ.get("PGURL", "postgresql://postgres:root@127.0.0.1:5433/marketdata"),
                            keepalives=1, keepalives_idle=30,
                        )
                        _bf_conn.autocommit = False
                        bf_results: dict = {}
                        for src_key in po_backfill_sources:
                            src_label = dict(_PO_DOC_SOURCE_LABELS).get(src_key, src_key)
                            st.write(f"Fetching **{src_label}** …")
                            per = run_po_doc_backfill(
                                _bf_conn, [src_key], po_backfill_start, PREFIX
                            )
                            n_src = per.get(src_key, 0)
                            bf_results[src_key] = n_src
                            st.write(f"  → {n_src} new docs")
                        _bf_conn.close()
                        total_bf = sum(v for v in bf_results.values() if v > 0)
                        bf_status.update(
                            label=f"Backfill complete — {total_bf} new documents added.",
                            state="complete",
                        )
                        with st.expander("Per-source results", expanded=True):
                            for src_key, src_label in _PO_DOC_SOURCE_LABELS:
                                if src_key in bf_results:
                                    st.caption(f"{src_label}: {bf_results[src_key]} new docs")
                        _knowledge_doc_counts.clear()
                    except Exception as exc:
                        st.error(f"Backfill failed: {exc}")

    st.divider()
    st.subheader("Knowledge Base Digest → Expert Memory")
    st.caption("Extracts durable insights from KB documents. Runs at 03:45 WAW nightly.")
    if st.button("Digest KB into Expert Memory", type="primary", key="po_digest_kb"):
        with st.spinner("Extracting insights (1-2 min)…"):
            try:
                n_dk = digest_kb_docs(_ANTHROPIC_KEY, PREFIX, CFG.name, limit=200)
                st.success(f"Extracted {n_dk} new insights.")
            except Exception as exc:
                st.error(f"KB digest failed: {exc}")

    st.divider()
    st.subheader("Expert Memory")
    try:
        insights_df = _query(
            f"SELECT id, insight_type, confidence, insight_text, validated_at "
            f"FROM intl_market.{PREFIX}expert_insights WHERE active=TRUE ORDER BY validated_at DESC LIMIT 50"
        )
        if insights_df.empty:
            st.info("No expert insights yet. Digest the KB above.")
        else:
            for _, row in insights_df.iterrows():
                with st.expander(
                    f"[{row['insight_type']}] [{row['confidence']}] {str(row['insight_text'])[:80]}…",
                    expanded=False,
                ):
                    st.write(row["insight_text"])
                    st.caption(f"Logged: {row['validated_at']}")
    except Exception as exc:
        st.error(f"Could not load insights: {exc}")

    st.divider()
    st.subheader("Agent Memory")
    mems = _load_memories(APP_KEY)
    if mems.empty:
        st.info("No agent memories saved yet.")
    else:
        for _, row in mems.iterrows():
            c1, c2 = st.columns([10, 1])
            with c1:
                st.markdown(f"**[{row['category']}]** {row['subject']}: {row['content']}")
                st.caption(f"{row['source']} · {row['created_at']}")
            with c2:
                if st.button("🗑", key=f"po_del_mem_{row['id']}"):
                    _delete_memory(row["id"])
                    st.rerun()

    st.divider()
    st.subheader("Add Memory Manually")
    with st.form("po_add_mem"):
        cat = st.selectbox("Category",
                           ["market_view", "methodology", "investment_thesis", "asset_note", "red_flag"])
        subj = st.text_input("Subject (≤8 words)")
        cont = st.text_area("Content (one sentence)")
        if st.form_submit_button("Save"):
            _save_memory(cat, subj, cont, source="manual")
            st.success("Saved")

    # ── Ancillary Service Prices ───────────────────────────────────────────
    st.divider()
    st.subheader("Ancillary Service Market Prices")
    st.caption(
        "FCR and aFRR weekly clearing prices from PSE reporting API · "
        "Rynek Mocy annual auction results from TGE · "
        "Scheduler: FCR/aFRR every Tuesday 06:05, Capacity Market 1st of month 05:10 (WAW)"
    )

    try:
        as_status = _query(
            "SELECT market_type, COUNT(*) as weeks, "
            "MAX(week_start) as latest_week, "
            "AVG(price_pln_mw_week) as avg_price "
            "FROM intl_market.po_as_prices "
            "GROUP BY market_type"
        )
    except Exception:
        as_status = pd.DataFrame()
    try:
        cap_status = _query(
            "SELECT delivery_year, price_pln_mw_yr, auction_date "
            "FROM intl_market.po_capacity_market "
            "ORDER BY delivery_year DESC LIMIT 1"
        )
    except Exception:
        cap_status = pd.DataFrame()

    sc1, sc2, sc3 = st.columns(3)
    _fcr_row = as_status[as_status["market_type"] == "FCR"]        if not as_status.empty else pd.DataFrame()
    _afr_row = as_status[as_status["market_type"] == "aFRR_capacity"] if not as_status.empty else pd.DataFrame()

    with sc1:
        if not _fcr_row.empty:
            st.metric("FCR",
                       f"zł{float(_fcr_row['avg_price'].iloc[0]):,.0f}/MW/wk",
                       f"{int(_fcr_row['weeks'].iloc[0])} weeks · latest {_fcr_row['latest_week'].iloc[0]}")
        else:
            st.metric("FCR", "No data")
    with sc2:
        if not _afr_row.empty:
            st.metric("aFRR capacity",
                       f"zł{float(_afr_row['avg_price'].iloc[0]):,.0f}/MW/wk",
                       f"{int(_afr_row['weeks'].iloc[0])} weeks · latest {_afr_row['latest_week'].iloc[0]}")
        else:
            st.metric("aFRR capacity", "No data")
    with sc3:
        if not cap_status.empty:
            st.metric("Rynek Mocy",
                       f"zł{float(cap_status['price_pln_mw_yr'].iloc[0]):,.0f}/MW/yr",
                       f"{int(cap_status['delivery_year'].iloc[0])} delivery year")
        else:
            st.metric("Rynek Mocy", "No data")

    bt1, bt2, bt3 = st.columns(3)
    with bt1:
        if st.button("Scrape FCR prices", key="po_scrape_fcr"):
            with st.spinner("Fetching FCR auction results from PSE…"):
                from services.po_knowledge.entso_scraper import scrape_po_fcr_prices
                n = scrape_po_fcr_prices(_conn(), weeks_back=52)
            st.success(f"FCR: {n} new rows inserted")
    with bt2:
        if st.button("Scrape aFRR prices", key="po_scrape_afrr"):
            with st.spinner("Fetching aFRR auction results from PSE…"):
                from services.po_knowledge.entso_scraper import scrape_po_afrr_prices
                n = scrape_po_afrr_prices(_conn(), weeks_back=52)
            st.success(f"aFRR: {n} new rows inserted")
    with bt3:
        if st.button("Scrape Capacity Market", key="po_scrape_cap"):
            with st.spinner("Fetching Rynek Mocy results from TGE…"):
                from services.po_knowledge.entso_scraper import scrape_po_capacity_market
                n = scrape_po_capacity_market(_conn())
            st.success(f"Rynek Mocy: {n} rows upserted")

    # 52-week AS price chart
    try:
        as_history = _query(
            "SELECT week_start, market_type, price_pln_mw_week "
            "FROM intl_market.po_as_prices "
            "WHERE week_start >= CURRENT_DATE - INTERVAL '52 weeks' "
            "AND price_pln_mw_week IS NOT NULL "
            "ORDER BY week_start"
        )
    except Exception:
        as_history = pd.DataFrame()
    if not as_history.empty:
        fig_as = px.line(
            as_history, x="week_start", y="price_pln_mw_week", color="market_type",
            title="FCR & aFRR Weekly Clearing Prices (PLN/MW/week)",
            labels={"week_start": "Week", "price_pln_mw_week": "PLN/MW/week",
                     "market_type": "Market"},
        )
        st.plotly_chart(fig_as, use_container_width=True)
    else:
        st.info("No AS price history yet. Click 'Scrape FCR prices' or 'Scrape aFRR prices' above.")

    # ── AS Backfill ────────────────────────────────────────────────────────
    st.subheader("AS Data Backfill")
    bf_c1, bf_c2, bf_c3 = st.columns(3)
    with bf_c1:
        bf_start = st.date_input("Backfill from week", key="po_as_bf_start",
                                  value=pd.Timestamp.today() - pd.Timedelta(weeks=104))
    with bf_c2:
        bf_type = st.selectbox("Market", ["FCR", "aFRR", "Both"], key="po_as_bf_type")
    with bf_c3:
        st.write("")
        st.write("")
        if st.button("Run Backfill", key="po_as_bf_run"):
            weeks_back = max(1, (pd.Timestamp.today().date() - bf_start).days // 7)
            with st.spinner(f"Backfilling {bf_type} for {weeks_back} weeks…"):
                from services.po_knowledge.entso_scraper import (
                    scrape_po_fcr_prices, scrape_po_afrr_prices,
                )
                total = 0
                if bf_type in ("FCR", "Both"):
                    total += scrape_po_fcr_prices(_conn(), weeks_back=weeks_back)
                if bf_type in ("aFRR", "Both"):
                    total += scrape_po_afrr_prices(_conn(), weeks_back=weeks_back)
            st.success(f"Backfill complete: {total} new rows inserted")


# ═══════════════════════════════════════════════════════════════
# Tab 8 — Grid Analysis (PyPSA)
# ═══════════════════════════════════════════════════════════════
with tab_pypsa:
    st.header("Grid Analysis — PyPSA Zonal Model")
    st.info(
        "Upload PSE grid data or use simplified 4-zone Poland model (N/S/E/W) "
        "to run zonal power flow and analyse congestion patterns affecting BESS dispatch.",
        icon="⚡",
    )

    pypsa_mode = st.radio("Model type", ["4-Zone Poland (built-in)", "Upload CSV files"], horizontal=True)

    if pypsa_mode == "4-Zone Poland (built-in)":
        st.subheader("Simplified 4-Zone Poland Network")
        st.caption("Approximate zonal prices and transmission from Aurora Q2 2026 data")

        if st.button("Build and run 4-zone model", type="primary", key="po_pypsa_run"):
            try:
                import pypsa
                n = pypsa.Network()
                n.set_snapshots(pd.date_range("2026-01-01", periods=24, freq="h"))

                # Zones
                for zone, price in [("North", 290), ("South", 310), ("East", 285), ("West", 305)]:
                    n.add("Bus", zone, v_nom=400)
                    n.add("Generator", f"{zone}_gen", bus=zone, p_nom=5000,
                          marginal_cost=price, p_max_pu=0.8)
                    n.add("Load", f"{zone}_load", bus=zone,
                          p_set=pd.Series([2000] * 24, index=n.snapshots))

                # Transmission
                for b0, b1, cap in [("North","South",2000),("North","East",1500),
                                     ("South","West",2500),("East","West",1000)]:
                    n.add("Line", f"{b0}-{b1}", bus0=b0, bus1=b1, x=0.05, s_nom=cap)

                n.lopf(pyomo=False, solver_name="highs")

                st.success("LOPF solved successfully")
                col1, col2 = st.columns(2)
                with col1:
                    st.subheader("Zonal Dispatch (MW)")
                    disp = n.generators_t.p.mean().reset_index()
                    disp.columns = ["Zone_Gen", "Avg MW"]
                    st.dataframe(disp, use_container_width=True, hide_index=True)
                with col2:
                    st.subheader("Line Loading (%)")
                    if not n.lines_t.p0.empty:
                        loading = (n.lines_t.p0.abs() / n.lines.s_nom * 100).mean().reset_index()
                        loading.columns = ["Line", "Avg Loading %"]
                        st.dataframe(loading, use_container_width=True, hide_index=True)
            except ImportError:
                st.error("PyPSA not installed. Add `pypsa` to the Dockerfile.")
            except Exception as exc:
                st.error(f"PyPSA error: {exc}")

    else:
        col1, col2, col3 = st.columns(3)
        buses_file = col1.file_uploader("Buses CSV",      type=["csv"], key="po_buses")
        lines_file  = col2.file_uploader("Lines CSV",      type=["csv"], key="po_lines")
        gens_file   = col3.file_uploader("Generators CSV", type=["csv"], key="po_gens")

        if buses_file and lines_file:
            try:
                import pypsa
                buses_df = pd.read_csv(buses_file)
                lines_df = pd.read_csv(lines_file)
                n = pypsa.Network()
                for _, row in buses_df.iterrows():
                    n.add("Bus", row["name"], v_nom=row.get("v_nom", 400))
                for _, row in lines_df.iterrows():
                    n.add("Line", row.get("name", "L"), bus0=row["bus0"], bus1=row["bus1"],
                          x=row.get("x", 0.05), s_nom=row.get("s_nom", 500))
                if gens_file:
                    gens_df = pd.read_csv(gens_file)
                    for _, row in gens_df.iterrows():
                        n.add("Generator", row["name"], bus=row["bus"],
                              p_nom=row.get("p_nom", 100), marginal_cost=row.get("marginal_cost", 50))
                st.success(f"Network loaded: {len(n.buses)} buses, {len(n.lines)} lines")
            except ImportError:
                st.error("PyPSA not installed.")
            except Exception as exc:
                st.error(f"Network build error: {exc}")
        else:
            st.markdown("""
**Data sources for Poland grid model:**
- PSE publishes annual grid data (generation units, transmission topology)
- ENTSO-E Transparency Platform: cross-border flows, installed capacity by zone
- Aurora Excel data (Q1/Q2 2026): zonal prices, capacity by technology
""")


# ═══════════════════════════════════════════════════════════════
# Tab 9 — Library
# ═══════════════════════════════════════════════════════════════
with tab_library:
    st.header("Report Library")
    st.caption("Save and retrieve daily, weekly, and monthly market reports (PDF)")

    # ── Upload ──
    with st.expander("Upload New Report", expanded=True):
        ul_c1, ul_c2, ul_c3 = st.columns(3)
        with ul_c1:
            lib_name = st.text_input("Report Name",
                                      placeholder="e.g. Aurora Monthly Flexible Summary Apr 2026",
                                      key="po_lib_name")
        with ul_c2:
            lib_freq = st.selectbox("Frequency", ["daily", "weekly", "monthly"],
                                     key="po_lib_freq")
        with ul_c3:
            lib_period = st.date_input("Report Period", value=date.today(),
                                        key="po_lib_period")
        lib_file = st.file_uploader("PDF File", type=["pdf"], key="po_lib_upload")
        if st.button("Save to Library", type="primary", key="po_lib_save_btn",
                     disabled=not (lib_file and lib_name)):
            try:
                _save_library_report(lib_name.strip(), lib_freq, lib_period,
                                      lib_file.name, lib_file.read())
                st.success(f"Saved '{lib_name}' ({lib_freq}) to library.")
                st.rerun()
            except Exception as exc:
                st.error(f"Failed: {exc}")

    st.divider()

    # ── Browse by frequency ──
    lib_freq_tabs = st.tabs(["All Reports", "Daily", "Weekly", "Monthly"])
    lib_freq_filters = [None, "daily", "weekly", "monthly"]

    for lib_ftab, lib_filt in zip(lib_freq_tabs, lib_freq_filters):
        with lib_ftab:
            lib_df = _list_library_reports(PREFIX, lib_filt)
            if lib_df.empty:
                st.info("No reports in this category yet. Use 'Upload New Report' above.")
                continue

            st.caption(f"{len(lib_df)} report(s)")
            filt_key = lib_filt or "all"

            # Summary table
            display_df = lib_df[["report_name", "frequency", "period",
                                  "filename", "file_size_kb", "uploaded_at"]].copy()
            display_df.columns = ["Report Name", "Frequency", "Period",
                                   "Filename", "Size (KB)", "Uploaded At"]
            st.dataframe(display_df, use_container_width=True, hide_index=True)

            st.divider()
            st.subheader("Download / Delete")
            options = {
                f"{r['report_name']}  —  {str(r['period'])}  [{r['frequency']}]": int(r["id"])
                for _, r in lib_df.iterrows()
            }
            sel_label = st.selectbox("Select report", list(options.keys()),
                                      key=f"po_lib_sel_{filt_key}")
            if sel_label:
                sel_id = options[sel_label]
                sel_row = lib_df[lib_df["id"] == sel_id].iloc[0]
                col_dl, col_del = st.columns([3, 1])
                with col_dl:
                    report_bytes = _get_library_report_data(PREFIX, sel_id)
                    if report_bytes:
                        st.download_button(
                            label=f"Download  {sel_row['filename']}",
                            data=report_bytes,
                            file_name=sel_row["filename"],
                            mime="application/pdf",
                            key=f"po_lib_dl_{sel_id}_{filt_key}",
                        )
                    else:
                        st.warning("Report data not found in database.")
                with col_del:
                    if st.button("Delete", key=f"po_lib_del_{sel_id}_{filt_key}",
                                  type="secondary"):
                        _delete_library_report(sel_id)
                        st.success("Deleted.")
                        st.rerun()
