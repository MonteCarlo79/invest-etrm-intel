"""GB Market Intelligence — Streamlit app.

Two Claude agents:
  - Strategist: GB market fundamentals (system price, EPEX, NIV, ancillary)
  - Quant: BESS investment economics (benchmarking index, leaderboard, IRR)

Port 8508 | ALB slug: gb-market | Memory app key: gb_market
"""
import json
import logging
import os
import sys
import uuid

logger = logging.getLogger(__name__)
from datetime import date, datetime, timedelta

import anthropic
import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import psycopg2
import streamlit as st
from dotenv import load_dotenv

load_dotenv(os.path.join(os.path.dirname(__file__), "..", "..", "config", ".env"))

# Ensure repo root is on sys.path so all services.* packages are importable
_REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", ".."))
if _REPO_ROOT not in sys.path:
    sys.path.insert(0, _REPO_ROOT)

# ---------------------------------------------------------------------------
# Page config
# ---------------------------------------------------------------------------

st.set_page_config(
    page_title="GB Market Intelligence",
    page_icon="🇬🇧",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------
# DB connection
# ---------------------------------------------------------------------------

@st.cache_resource(ttl=3600)
def _get_conn():
    url = (
        os.environ.get("PGURL")
        or os.environ.get("DATABASE_URL")
        or "postgresql://postgres:root@127.0.0.1:5433/marketdata"
    )
    print("[DB] connecting", flush=True)
    conn = psycopg2.connect(url, connect_timeout=10)
    conn.autocommit = True
    print("[DB] connected", flush=True)
    return conn


def _conn():
    conn = _get_conn()
    if conn.closed:
        _get_conn.clear()
        conn = _get_conn()
    return conn


def _query(sql: str, params=None) -> pd.DataFrame:
    return pd.read_sql(sql, _conn(), params=params)


# ---------------------------------------------------------------------------
# Agent memory helpers
# ---------------------------------------------------------------------------

_APP_KEY = "gb_market"
_ANTHROPIC_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
_client = anthropic.Anthropic(api_key=_ANTHROPIC_KEY)


def _ensure_memory_table():
    cur = _conn().cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS marketdata.agent_memory (
            id SERIAL PRIMARY KEY,
            app TEXT NOT NULL DEFAULT 'gb_market',
            category TEXT NOT NULL,
            subject TEXT NOT NULL,
            content TEXT NOT NULL,
            source TEXT DEFAULT 'manual',
            created_at TIMESTAMPTZ DEFAULT NOW(),
            active BOOLEAN DEFAULT TRUE
        )
    """)
    cur.execute("ALTER TABLE marketdata.agent_memory ADD COLUMN IF NOT EXISTS app TEXT DEFAULT 'gb_market'")
    _conn().commit()


def _ensure_ingestion_log_table():
    cur = _conn().cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS intl_market.gb_ingestion_log (
            id SERIAL PRIMARY KEY,
            run_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
            trigger TEXT NOT NULL,
            date_from DATE,
            date_to DATE,
            status TEXT NOT NULL,
            rows_ingested JSONB,
            error_msg TEXT,
            duration_seconds NUMERIC
        )
    """)
    _conn().commit()


def _ensure_knowledge_table():
    cur = _conn().cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS intl_market.gb_knowledge_docs (
            id              SERIAL PRIMARY KEY,
            source          TEXT NOT NULL,
            doc_type        TEXT NOT NULL,
            title           TEXT,
            url             TEXT UNIQUE,
            published_date  DATE,
            content         TEXT NOT NULL,
            fetched_at      TIMESTAMPTZ NOT NULL DEFAULT NOW(),
            search_vector   TSVECTOR GENERATED ALWAYS AS (
                                to_tsvector('english',
                                    coalesce(title, '') || ' ' || left(content, 100000))
                            ) STORED
        )
    """)
    cur.execute(
        "CREATE INDEX IF NOT EXISTS gb_knowledge_docs_fts "
        "ON intl_market.gb_knowledge_docs USING GIN(search_vector)"
    )
    _conn().commit()


@st.cache_data(ttl=60)
def _search_knowledge(query: str, sources: list | None = None, limit: int = 8) -> pd.DataFrame:
    try:
        src_clause = "AND source = ANY(%s)" if sources else ""

        def _run(sql, params):
            return _query(sql, params)

        # OR-based FTS: replace & with | so ANY word match returns results;
        # still rank by plainto_tsquery (AND-based) so best-matching docs rank highest.
        fts_params: list = [query, query]
        if sources:
            fts_params.append(sources)
        fts_params.append(limit)
        df = _run(
            "SELECT source, doc_type, title, url, published_date, "
            "left(content, 1500) AS snippet, "
            "ts_rank(search_vector, plainto_tsquery('english', %s)) AS rank "
            "FROM intl_market.gb_knowledge_docs "
            "WHERE search_vector @@ to_tsquery('english', "
            "  regexp_replace(plainto_tsquery('english', %s)::text, ' & ', ' | ', 'g')"
            f") {src_clause} "
            "ORDER BY rank DESC LIMIT %s",
            fts_params,
        )
        if not df.empty:
            return df

        # Fallback: ILIKE title search (catches proper nouns / non-English terms)
        like_q = "%" + query.strip().replace("%", "").replace("_", "") + "%"
        ilike_params: list = [like_q]
        if sources:
            ilike_params.append(sources)
        ilike_params.append(limit)
        return _run(
            "SELECT source, doc_type, title, url, published_date, "
            "left(content, 1500) AS snippet, 0.0::float AS rank "
            "FROM intl_market.gb_knowledge_docs "
            f"WHERE title ILIKE %s {src_clause} "
            "ORDER BY published_date DESC NULLS LAST LIMIT %s",
            ilike_params,
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=300)
def _knowledge_doc_counts() -> pd.DataFrame:
    try:
        return _query(
            "SELECT source, doc_type, COUNT(*) AS docs, MAX(fetched_at) AS last_fetch "
            "FROM intl_market.gb_knowledge_docs "
            "GROUP BY source, doc_type ORDER BY source, doc_type"
        )
    except Exception:
        return pd.DataFrame()


# ---------------------------------------------------------------------------
# Document upload helpers (PDF / Excel / Word / PPTX / TXT text extraction)
# ---------------------------------------------------------------------------

def _extract_text_from_upload(filename: str, data: bytes) -> str:
    """Extract plain text from an uploaded file. Returns extracted text or raises."""
    import io
    ext = filename.rsplit(".", 1)[-1].lower()

    if ext == "txt":
        return data.decode("utf-8", errors="replace")

    if ext == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
        return "\n\n".join(pages)

    if ext in ("xlsx", "xls"):
        xl = pd.ExcelFile(io.BytesIO(data))
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            parts.append(f"Sheet: {sheet}\n{df.to_string(index=False)}")
        return "\n\n".join(parts)

    if ext in ("docx", "doc"):
        from docx import Document
        doc = Document(io.BytesIO(data))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                paras.append("  ".join(cell.text for cell in row.cells if cell.text.strip()))
        return "\n".join(paras)

    if ext in ("pptx", "ppt"):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        texts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                texts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)

    raise ValueError(f"Unsupported file type: .{ext}")


_UPLOAD_DOC_TYPES = {
    "pdf": "pdf", "txt": "text",
    "xlsx": "excel", "xls": "excel",
    "docx": "word", "doc": "word",
    "pptx": "pptx", "ppt": "pptx",
}


def _ingest_url(url: str) -> dict:
    """Fetch a public web URL, extract text, upsert into gb_knowledge_docs."""
    import requests
    from bs4 import BeautifulSoup
    try:
        resp = requests.get(
            url,
            timeout=30,
            headers={"User-Agent": "Mozilla/5.0 (compatible; BESSPlatformBot/1.0)"},
            allow_redirects=True,
        )
    except Exception as exc:
        return {"status": "error", "msg": f"Fetch failed: {exc}"}
    if resp.status_code != 200:
        return {"status": "error", "msg": f"HTTP {resp.status_code} from {url}"}
    if "/sign-in" in resp.url or "/login" in resp.url:
        return {"status": "error", "msg": "Page requires login — download as PDF and upload instead."}

    soup = BeautifulSoup(resp.text, "html.parser")
    # Extract title
    title_el = soup.find("h1") or soup.find("title")
    title = title_el.get_text(" ", strip=True) if title_el else url

    # Strip boilerplate
    for tag in soup.find_all(["script", "style", "nav", "header", "footer",
                               "aside", "form", "noscript", "iframe", "button"]):
        tag.decompose()
    content_parts = []
    for el in soup.find_all(["p", "h1", "h2", "h3", "h4", "li", "blockquote"]):
        t = el.get_text(" ", strip=True)
        if t:
            content_parts.append(t)
    content = "\n\n".join(content_parts) if content_parts else soup.get_text(" ", strip=True)

    if not content.strip():
        return {"status": "error", "msg": "No text could be extracted from the page."}

    try:
        cur = _conn().cursor()
        cur.execute(
            "INSERT INTO intl_market.gb_knowledge_docs "
            "(source, doc_type, title, url, published_date, content) "
            "VALUES (%s, %s, %s, %s, %s, %s) "
            "ON CONFLICT (url) DO UPDATE SET "
            "  content = EXCLUDED.content, title = EXCLUDED.title, fetched_at = NOW()",
            ("upload", "article", title, url, date.today(), content),
        )
        return {"status": "success", "msg": f"Ingested '{title}' ({len(content):,} chars)"}
    except Exception as exc:
        return {"status": "error", "msg": f"DB insert failed: {exc}"}


def _ingest_uploaded_file(filename: str, data: bytes) -> dict:
    """Extract text, upsert into gb_knowledge_docs. Returns {status, msg}."""
    try:
        content = _extract_text_from_upload(filename, data)
    except Exception as exc:
        return {"status": "error", "msg": f"Text extraction failed: {exc}"}

    if not content.strip():
        return {"status": "error", "msg": "No text could be extracted from this file."}

    ext = filename.rsplit(".", 1)[-1].lower()
    doc_type = _UPLOAD_DOC_TYPES.get(ext, "document")
    url_key = f"upload://{filename}"

    try:
        cur = _conn().cursor()
        cur.execute(
            "INSERT INTO intl_market.gb_knowledge_docs "
            "(source, doc_type, title, url, published_date, content) "
            "VALUES (%s, %s, %s, %s, %s, %s) "
            "ON CONFLICT (url) DO UPDATE SET "
            "  content = EXCLUDED.content, "
            "  fetched_at = NOW()",
            ("upload", doc_type, filename, url_key, date.today(), content),
        )
        return {"status": "success", "msg": f"Ingested '{filename}' ({len(content):,} chars)"}
    except Exception as exc:
        return {"status": "error", "msg": f"DB insert failed: {exc}"}


def _log_ingestion_run(trigger: str, date_from, date_to, status: str,
                        rows: dict | None, error_msg: str | None, duration: float):
    try:
        cur = _conn().cursor()
        cur.execute(
            "INSERT INTO intl_market.gb_ingestion_log "
            "(trigger, date_from, date_to, status, rows_ingested, error_msg, duration_seconds) "
            "VALUES (%s, %s, %s, %s, %s, %s, %s)",
            (trigger, date_from, date_to, status,
             json.dumps(rows) if rows else None, error_msg, round(duration, 1)),
        )
        _conn().commit()
    except Exception:
        pass  # never crash the scheduler thread due to logging failure


@st.cache_data(ttl=30)
def _get_ingestion_logs(limit: int = 20) -> pd.DataFrame:
    try:
        return _query(
            "SELECT id, run_at AT TIME ZONE 'Asia/Singapore' AS run_at_sgt, "
            "trigger, date_from, date_to, status, rows_ingested, error_msg, duration_seconds "
            "FROM intl_market.gb_ingestion_log "
            "ORDER BY run_at DESC LIMIT %s",
            (limit,),
        )
    except Exception:
        return pd.DataFrame()


def _run_ingestion_job(date_from, date_to, trigger: str = "manual") -> dict:
    """Run full GB ingestion; return result dict. Called by scheduler and UI button."""
    import io, time
    from contextlib import redirect_stdout

    from services.modo_energy.gb_ingestion import run_gb_backfill

    t0 = time.time()
    buf = io.StringIO()
    try:
        with redirect_stdout(buf):
            run_gb_backfill(date_from, date_to)
        duration = time.time() - t0
        _log_ingestion_run(trigger, date_from, date_to, "success", None, None, duration)
        return {"status": "success", "log": buf.getvalue(), "duration": duration}
    except Exception as exc:
        duration = time.time() - t0
        _log_ingestion_run(trigger, date_from, date_to, "error", None, str(exc), duration)
        return {"status": "error", "error": str(exc), "log": buf.getvalue(), "duration": duration}


def _run_knowledge_ingest_job(only: list[str] | None = None, trigger: str = "manual") -> dict:
    """Run knowledge base ingestion; return result dict."""
    import time
    t0 = time.time()
    try:
        from services.gb_knowledge.ingest import run_knowledge_ingest
        results = run_knowledge_ingest(only=only, verbose=False)
        duration = time.time() - t0
        total = sum(results.values())
        return {"status": "success", "results": results, "total": total, "duration": duration}
    except Exception as exc:
        duration = time.time() - t0
        return {"status": "error", "error": str(exc), "duration": duration}


@st.cache_resource
def _start_scheduler():
    """Start APScheduler background scheduler (runs once per process via cache_resource)."""
    from apscheduler.schedulers.background import BackgroundScheduler

    def _daily_market_job():
        import importlib.util, pathlib
        yesterday = date.today() - timedelta(days=1)
        _run_ingestion_job(yesterday, yesterday, trigger="scheduled")
        _table_counts.clear()
        _get_ingestion_logs.clear()
        # Ingest fuel mix from NESO CKAN
        try:
            _fm_path = pathlib.Path(__file__).with_name("fuel_mix_ingest.py")
            _fm_spec = importlib.util.spec_from_file_location("fuel_mix_ingest", _fm_path)
            _fm_mod  = importlib.util.module_from_spec(_fm_spec)
            _fm_spec.loader.exec_module(_fm_mod)
            n = _fm_mod.ingest_fuel_mix(yesterday, _conn())
            logger.info("Fuel mix ingest: %d rows for %s", n, yesterday)
        except Exception as _fm_exc:
            logger.warning("Fuel mix ingest failed: %s", _fm_exc)

    def _daily_knowledge_job():
        _run_knowledge_ingest_job(trigger="scheduled")

    def _daily_report_job():
        import importlib.util, pathlib
        today     = date.today()
        yesterday = today - timedelta(days=1)

        # Ensure today's market-data ingestion succeeded before reporting.
        # If the 03:00 job ran and logged success, skip; otherwise run it now.
        ingestion_done = False
        try:
            cur = _conn().cursor()
            cur.execute(
                "SELECT COUNT(*) FROM intl_market.gb_ingestion_log "
                "WHERE trigger IN ('scheduled', 'report_triggered') "
                "AND status = 'success' "
                "AND date_from = %s AND run_at::date = %s",
                (yesterday, today),
            )
            ingestion_done = cur.fetchone()[0] > 0
        except Exception as _chk_exc:
            logger.warning("Report job: could not check ingestion log: %s", _chk_exc)

        if not ingestion_done:
            logger.info("Report job: ingestion not complete, running now")
            _run_ingestion_job(yesterday, yesterday, trigger="report_triggered")
            _table_counts.clear()
            _get_ingestion_logs.clear()

        _rpt_path = pathlib.Path(__file__).with_name("daily_report.py")
        _spec = importlib.util.spec_from_file_location("daily_report", _rpt_path)
        _mod  = importlib.util.module_from_spec(_spec)
        _spec.loader.exec_module(_mod)

        # Generate PDF once; send via email and WeCom
        _rpt_date = _mod._get_latest_data_date(_mod._get_conn())
        pdf_bytes, ai_commentary = _mod.generate_report_pdf(_rpt_date)

        # Email
        try:
            _mod.send_daily_report_email(pdf_bytes, _rpt_date, ai_commentary=ai_commentary)
            logger.info("Daily report emailed for %s (%d bytes)", _rpt_date, len(pdf_bytes))
        except Exception as _email_exc:
            logger.error("Daily report email failed: %s", _email_exc)

        # WeCom (optional — only if webhook URL is configured)
        _wecom_url = os.environ.get("WECOM_WEBHOOK_URL", "")
        if _wecom_url:
            try:
                _mod.send_daily_report_wecom(pdf_bytes, _rpt_date,
                                             webhook_url=_wecom_url,
                                             ai_commentary=ai_commentary)
                logger.info("Daily report sent to WeCom for %s", _rpt_date)
            except Exception as _wc_exc:
                logger.error("Daily report WeCom send failed: %s", _wc_exc)

    def _pricing_batch_job():
        import importlib.util, pathlib
        yesterday = date.today() - timedelta(days=1)
        try:
            _pb_path = pathlib.Path(__file__).with_name("pricing_batch.py")
            _pb_spec = importlib.util.spec_from_file_location("pricing_batch", _pb_path)
            _pb_mod  = importlib.util.module_from_spec(_pb_spec)
            _pb_spec.loader.exec_module(_pb_mod)
            result = _pb_mod.run_pricing_batch(yesterday, _conn())
            logger.info("Pricing batch: %s", result)
        except Exception as _pb_exc:
            logger.error("Pricing batch failed: %s", _pb_exc)

    def _modo_ai_job():
        """Distill daily GB BESS intelligence from Modo Energy's AI agent."""
        try:
            from services.gb_knowledge.modo_ai import ModoAIConnector
            from services.gb_knowledge.base import get_db_conn, ensure_table, upsert_doc
            conn = get_db_conn()
            ensure_table(conn)
            connector = ModoAIConnector()
            n = connector.run(conn)
            conn.close()
            logger.info("Modo AI distillation: %d new docs inserted", n)
        except Exception as _ma_exc:
            logger.error("Modo AI distillation failed: %s", _ma_exc)

    scheduler = BackgroundScheduler(timezone="Asia/Singapore")
    scheduler.add_job(_daily_market_job, "cron", hour=3, minute=0,
                      id="gb_daily_market", misfire_grace_time=3600)
    scheduler.add_job(_daily_knowledge_job, "cron", hour=3, minute=30,
                      id="gb_daily_knowledge", misfire_grace_time=3600)
    scheduler.add_job(_modo_ai_job, "cron", hour=4, minute=0,
                      id="gb_modo_ai_distill", misfire_grace_time=3600)

    def _kb_digest_job():
        """Digest unprocessed KB docs into structured expert insights."""
        try:
            from services.gb_knowledge.expert_memory import digest_kb_docs
            n = digest_kb_docs(_ANTHROPIC_KEY, limit=100)
            logger.info("KB digest: %d new insights extracted", n)
        except Exception as _exc:
            logger.error("KB digest failed: %s", _exc)

    scheduler.add_job(_kb_digest_job, "cron", hour=3, minute=45,
                      id="gb_kb_digest", misfire_grace_time=3600)
    scheduler.add_job(_pricing_batch_job, "cron", hour=4, minute=30,
                      id="gb_pricing_batch", misfire_grace_time=3600)
    scheduler.add_job(_daily_report_job, "cron", hour=6, minute=0,
                      id="gb_daily_report", misfire_grace_time=3600)
    scheduler.start()
    return scheduler


@st.cache_data(ttl=60)
def _load_memories(app_key: str) -> pd.DataFrame:
    try:
        return _query(
            "SELECT id, category, subject, content, source, created_at "
            "FROM marketdata.agent_memory WHERE app = %s AND active = TRUE "
            "ORDER BY created_at DESC",
            (app_key,),
        )
    except Exception:
        return pd.DataFrame()


def _save_memory(category: str, subject: str, content: str, source: str = "manual"):
    cur = _conn().cursor()
    cur.execute(
        "INSERT INTO marketdata.agent_memory (app, category, subject, content, source) "
        "VALUES (%s, %s, %s, %s, %s)",
        (_APP_KEY, category, subject, content, source),
    )
    _conn().commit()
    _load_memories.clear()


def _delete_memory(mem_id: int):
    cur = _conn().cursor()
    cur.execute("UPDATE marketdata.agent_memory SET active = FALSE WHERE id = %s", (mem_id,))
    _conn().commit()


# ── Session persistence (Strategist chat history survives page reload) ────────

def _ensure_sessions_table():
    cur = _conn().cursor()
    cur.execute("""
        CREATE TABLE IF NOT EXISTS intl_market.gb_analyst_sessions (
            session_id TEXT PRIMARY KEY,
            messages   JSONB NOT NULL DEFAULT '[]',
            created_at TIMESTAMPTZ DEFAULT NOW(),
            updated_at TIMESTAMPTZ DEFAULT NOW()
        )
    """)
    _conn().commit()


def _save_session(session_id: str, messages: list):
    try:
        _ensure_sessions_table()
        cur = _conn().cursor()
        cur.execute(
            "INSERT INTO intl_market.gb_analyst_sessions (session_id, messages, updated_at) "
            "VALUES (%s, %s::jsonb, NOW()) "
            "ON CONFLICT (session_id) DO UPDATE "
            "  SET messages = EXCLUDED.messages, updated_at = NOW()",
            (session_id, json.dumps(messages)),
        )
        _conn().commit()
    except Exception as _exc:
        logger.debug("_save_session failed: %s", _exc)


def _load_session(session_id: str) -> list:
    try:
        _ensure_sessions_table()
        cur = _conn().cursor()
        cur.execute(
            "SELECT messages FROM intl_market.gb_analyst_sessions WHERE session_id = %s",
            (session_id,),
        )
        row = cur.fetchone()
        return row[0] if row else []
    except Exception:
        return []


@st.cache_data(ttl=30)
def _list_recent_sessions(limit: int = 3) -> pd.DataFrame:
    try:
        return _query(
            "SELECT session_id, jsonb_array_length(messages) AS msg_count, updated_at "
            "FROM intl_market.gb_analyst_sessions "
            "WHERE jsonb_array_length(messages) > 0 "
            "ORDER BY updated_at DESC LIMIT %s",
            (limit,),
        )
    except Exception:
        return pd.DataFrame()
    _load_memories.clear()


def _generate_interview_questions() -> list[dict]:
    """
    Review current insight pool + KB coverage, identify knowledge gaps,
    return up to 5 targeted questions for the user to answer.
    """
    summary = _query(
        "SELECT insight_type, confidence, COUNT(*) AS n "
        "FROM intl_market.gb_expert_insights WHERE active = TRUE "
        "GROUP BY insight_type, confidence ORDER BY n DESC"
    )
    kb_cov = _query(
        "SELECT source, COUNT(*) AS n "
        "FROM intl_market.gb_knowledge_docs GROUP BY source ORDER BY n DESC"
    )
    sample = _query(
        "SELECT insight_text, insight_type "
        "FROM intl_market.gb_expert_insights WHERE active = TRUE "
        "ORDER BY id DESC LIMIT 15"
    )

    ctx_lines = ["Current expert insight pool:"]
    if not summary.empty:
        for _, r in summary.iterrows():
            ctx_lines.append(f"  {r['insight_type']} ({r['confidence']}): {int(r['n'])} insights")
    else:
        ctx_lines.append("  (empty)")

    ctx_lines.append("\nKnowledge base coverage:")
    if not kb_cov.empty:
        for _, r in kb_cov.iterrows():
            ctx_lines.append(f"  {r['source']}: {int(r['n'])} docs")

    ctx_lines.append("\nSample of already-known insights (do NOT duplicate):")
    if not sample.empty:
        for _, r in sample.iterrows():
            ctx_lines.append(f"  [{r['insight_type']}] {str(r['insight_text'])[:120]}")

    system = """\
You are the GB BESS market strategist agent auditing your own knowledge base to find gaps.
Identify the 5 most valuable areas where knowledge is THIN, UNCERTAIN, or MISSING.
Generate one precise expert interview question per gap — something only a practitioner
with hands-on GB BESS experience can answer from their own observation.

Prioritise gaps in these areas (in order):
1. Specific operational nuances of top-performing BESS assets (dispatch strategies, SoC management)
2. Counterintuitive market patterns the expert has personally observed (NIV chasing, BM dynamics)
3. Upcoming or recent regulatory changes with concrete operational implications
4. Revenue stack combinations that differentiate the top 10% of assets from the median
5. Grid constraint or locational patterns that affect BESS dispatch or revenue

Do NOT generate questions that are already answered in the sample insights above.
Do NOT generate generic textbook questions about UK power markets.

Respond ONLY with valid JSON:
{"questions": [{"question": "...", "topic": "market_structure|regulation|operations|bess_economics|grid_services", "why_asking": "one sentence on what knowledge gap this fills"}]}
"""
    try:
        resp = _client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=900,
            system=system,
            messages=[{"role": "user", "content": "\n".join(ctx_lines)}],
        )
        raw = resp.content[0].text.strip()
        if raw.startswith("```"):
            raw = raw.split("```", 2)[1]
            if raw.startswith("json"):
                raw = raw[4:]
        return json.loads(raw).get("questions", [])[:5]
    except Exception as exc:
        logger.warning("Gap analysis failed: %s", exc)
        return []


def _store_interview_answer(question: str, answer: str, topic: str) -> None:
    """Store a user's expert interview answer as a high-confidence insight."""
    insight_text = f"[Expert interview] Q: {question[:150]} | A: {answer}"
    cur = _conn().cursor()
    cur.execute(
        "INSERT INTO intl_market.gb_expert_insights "
        "(insight_text, insight_type, confidence, source_session) "
        "VALUES (%s, %s, 'high', %s)",
        (insight_text[:1000], topic, date.today().isoformat()),
    )
    _conn().commit()


def _extract_memories(user_msg: str, agent_reply: str) -> list[dict]:
    resp = _client.messages.create(
        model="claude-haiku-4-5",
        max_tokens=512,
        system=(
            "Extract memorable analyst preferences or domain facts from this GB power market "
            "conversation. Return a JSON array of objects with keys: "
            "category (one of: market_view, methodology, asset_note, investment_thesis, red_flag), "
            "subject (short title ≤8 words), content (one sentence). "
            "Only extract genuinely reusable insights — not ephemeral data points. "
            "Return [] if nothing is worth remembering."
        ),
        messages=[{"role": "user", "content": f"User: {user_msg}\n\nAgent: {agent_reply[:1500]}"}],
    )
    raw = next((b.text for b in resp.content if hasattr(b, "text")), "[]")
    start, end = raw.find("["), raw.rfind("]")
    if start == -1:
        return []
    try:
        return json.loads(raw[start:end + 1])
    except json.JSONDecodeError:
        return []


# ---------------------------------------------------------------------------
# System prompts
# ---------------------------------------------------------------------------

_GB_STRATEGIST_BASE_SYSTEM = """You are the GB Market Strategist, an expert in Great Britain electricity markets.

GROUNDING RULE: Answer only from data returned by your tools in this conversation. Never state specific prices, volumes, system conditions, or market events from your training data. If you have not called a tool yet, call one before answering factual questions.

DOMAIN CONTEXT:
- Settlement periods: 30-min half-hourly blocks, SP1–SP48 each day
- System price: the Balancing Mechanism clearing price (can be negative)
- NIV (Net Imbalance Volume): positive = system long (oversupplied), negative = system short
- EPEX DA: day-ahead auction clearing prices in GBP/MWh
- Ancillary services: Dynamic Containment (DC), Dynamic Moderation (DM), Dynamic Regulation (DR)
  - Each service has Low (L) and High (H) directions, with EFA block auctions (EFA 1–6)
  - Clearing prices in GBP/MW/h (capacity payment)
- Key BESS revenue streams: BM (balancing mechanism), CM (capacity market), frequency_response (DC/DM/DR)

ANALYTICAL FRAMEWORK:
- For price/market questions → call get_system_price or get_epex_prices
- For market tightness/balance → call get_system_price (NIV direction indicates system balance)
- For ancillary market questions → call get_ancillary_results
- For combined market overview → call get_market_summary
- For BESS asset performance / which asset is most profitable → call get_bess_leaderboard
- For BESS revenue trends by market stream → call get_bess_revenue_index
- For asset-specific data (owner, operator, capacity, location) → call get_bess_assets
- For market context, regulation, policy → call search_knowledge_base
"""


def _build_strategist_system(query: str = "") -> str:
    base = _GB_STRATEGIST_BASE_SYSTEM

    # Inject structured expert insights (HyDE-retrieved, confidence-ranked)
    if query:
        try:
            from services.gb_knowledge.expert_memory import get_gb_insights, inject_gb_memory
            insights = get_gb_insights(query, limit=5)
            mem_block = inject_gb_memory(insights)
            if mem_block:
                base += f"\n\n{mem_block}"
        except Exception:
            pass

    # Inject flat analyst notes from prior sessions
    mems = _load_memories(_APP_KEY)
    if not mems.empty:
        mem_lines = "\n".join(
            f"- [{r.category}] {r.subject}: {r.content}" for r in mems.itertuples()
        )
        base += f"\n\n## Analyst notes from prior sessions:\n{mem_lines}"

    return base


_GB_QUANT_BASE_SYSTEM = """You are the GB Quant, an expert in BESS investment economics for the GB market.

GROUNDING RULE: Answer only from data returned by your tools in this conversation. Never state specific revenue figures, IRR estimates, or market benchmarks from training data. Always fetch data before answering.

DOMAIN CONTEXT:
- GB BESS revenues are expressed in £/MW/day or £/MW/month (capacity-normalised)
- Revenue streams: BM (balancing mechanism), CM (capacity market), frequency_response, wholesale (EPEX arbitrage)
- Modo Energy BESS index = industry-average revenue across all GB BESS assets
- Leaderboard = per-asset, per-settlement-period, per-market revenue breakdown
- Duration matters: longer duration (2h+) captures more BM/wholesale; shorter (0.5h) is better for DC
- IRR methodology: unlevered, 15-year life, O&M 2% of capex/yr, degradation 2%/yr revenue reduction

ANALYTICAL FRAMEWORK:
- For revenue trend questions → call get_bess_daily_index or get_bess_monthly_index
- For asset comparison → call get_leaderboard
- For market landscape → call get_asset_database
- For investment return → call estimate_irr (note this is a parametric estimate, not a full LP model)
"""


def _build_quant_system() -> str:
    mems = _load_memories(_APP_KEY)
    if mems.empty:
        return _GB_QUANT_BASE_SYSTEM
    mem_lines = "\n".join(
        f"- [{r.category}] {r.subject}: {r.content}" for r in mems.itertuples()
    )
    return _GB_QUANT_BASE_SYSTEM + f"\n\n## Analyst memory from prior sessions:\n{mem_lines}"


# ---------------------------------------------------------------------------
# Strategist tools
# ---------------------------------------------------------------------------

_STRATEGIST_TOOLS = [
    {
        "name": "get_system_price",
        "description": "Half-hourly GB system price (£/MWh) and NIV (MW) for a date range.",
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string", "description": "ISO date e.g. '2026-05-01'"},
                "end_date":   {"type": "string", "description": "ISO date e.g. '2026-05-10'"},
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "get_epex_prices",
        "description": "EPEX day-ahead half-hourly prices (GBP/MWh) including daily baseload/peak/offpeak.",
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string"},
                "end_date":   {"type": "string"},
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "get_ancillary_results",
        "description": "DX ancillary service clearing prices (GBP/MW/h) and cleared volumes (MW) by service and EFA block.",
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string"},
                "end_date":   {"type": "string"},
                "services": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Service codes e.g. ['DCL','DRL','DCH','DRH','DMH','DML']. Leave empty for all.",
                },
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "get_market_summary",
        "description": "Daily summary combining avg system price, EPEX baseload, spread (system-EPEX), avg NIV.",
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string"},
                "end_date":   {"type": "string"},
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "search_knowledge_base",
        "description": (
            "Semantic full-text search over GB energy market knowledge: articles, reports, "
            "market notices, regulatory changes, and commentary from Elexon, ENTSO-E, "
            "Timera Energy, Modo Energy, and Meteologica. "
            "Use this to answer questions about market context, policy changes, or research."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Keywords or natural language question to search for.",
                },
                "sources": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Filter by source. Options: elexon, entso_e, timera, modo, meteologica. "
                        "Leave empty to search all sources."
                    ),
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "get_bess_leaderboard",
        "description": (
            "Asset-level GB BESS revenue leaderboard from Modo Energy data. "
            "Returns total revenue (£) and normalised £/MW/day per asset for a date range. "
            "Use to answer questions about which assets are most profitable, revenue rankings, "
            "or performance by market stream."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string", "description": "ISO date e.g. '2025-01-01'"},
                "end_date":   {"type": "string", "description": "ISO date e.g. '2025-12-31'"},
                "market": {
                    "type": "string",
                    "description": "Filter by market stream e.g. 'bm','frequency_response','wholesale'. Leave empty for total across all markets.",
                },
                "top_n": {"type": "integer", "description": "Number of top assets to return (default 20)."},
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "get_bess_revenue_index",
        "description": (
            "GB BESS industry-average revenue index (£/MW/day or £/MW/month) from Modo Energy. "
            "Shows market-wide revenue trends by stream (BM, frequency response, wholesale, etc.). "
            "Use for trend analysis, seasonal patterns, or market-level benchmarks."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string", "description": "ISO date e.g. '2025-01-01'"},
                "end_date":   {"type": "string", "description": "ISO date e.g. '2025-12-31'"},
                "granularity": {
                    "type": "string",
                    "enum": ["daily", "monthly"],
                    "description": "daily = £/MW/day per settlement date; monthly = £/MW/month. Default: monthly.",
                },
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "get_bess_assets",
        "description": (
            "GB BESS asset register from Modo Energy. "
            "Returns asset names, rated power (MW), energy capacity (MWh), owner, operator, "
            "developer, location, commissioning date. "
            "Use to answer questions about specific assets, fleet composition, or capacity data."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "min_power_mw": {"type": "number", "description": "Minimum rated power filter (MW)."},
                "owner":        {"type": "string", "description": "Filter by owner name (partial match)."},
                "operator":     {"type": "string", "description": "Filter by operator name (partial match)."},
            },
            "required": [],
        },
    },
]


def _dispatch_strategist(name: str, inputs: dict) -> str:
    try:
        if name == "search_knowledge_base":
            # Advanced retrieval: HyDE query expansion → OR-based FTS → re-ranking
            try:
                from services.gb_knowledge.advanced_retrieval import retrieve_for_gb_agent
                return retrieve_for_gb_agent(
                    query=inputs["query"],
                    api_key=_ANTHROPIC_KEY,
                    sources=inputs.get("sources") or None,
                    top_k=6,
                )
            except Exception as exc:
                # Fallback to plain FTS if advanced retrieval fails
                logger.warning("Advanced retrieval failed, falling back to FTS: %s", exc)
            results = _search_knowledge(
                inputs["query"],
                sources=inputs.get("sources") or None,
                limit=6,
            )
            if results.empty:
                return "No matching knowledge documents found."
            out = []
            for _, row in results.iterrows():
                out.append(
                    f"[{row['source']} / {row['doc_type']}] {row['title']} "
                    f"({row['published_date']})\n{row['snippet']}"
                )
            return "\n\n---\n\n".join(out)

        elif name == "get_system_price":
            df = _query(
                "SELECT sp.date, sp.settlement_period, sp.system_price, n.niv "
                "FROM intl_market.gb_system_price sp "
                "LEFT JOIN intl_market.gb_niv n "
                "  ON sp.date = n.date AND sp.settlement_period = n.settlement_period "
                "WHERE sp.date BETWEEN %s AND %s ORDER BY sp.date, sp.settlement_period",
                (inputs["start_date"], inputs["end_date"]),
            )
            if df.empty:
                return "No system price data for the requested period."
            summary = df.groupby("date").agg(
                avg_system_price=("system_price", "mean"),
                min_system_price=("system_price", "min"),
                max_system_price=("system_price", "max"),
                avg_niv=("niv", "mean"),
            ).round(2).reset_index()
            return summary.to_json(orient="records", date_format="iso")

        elif name == "get_epex_prices":
            df = _query(
                "SELECT delivery_date, settlement_period, price, volume, "
                "daily_baseload, daily_peakload, daily_offpeak "
                "FROM intl_market.gb_epex_da_hh "
                "WHERE delivery_date BETWEEN %s AND %s ORDER BY delivery_date, settlement_period",
                (inputs["start_date"], inputs["end_date"]),
            )
            if df.empty:
                return "No EPEX DA data for the requested period."
            daily = df.groupby("delivery_date").agg(
                daily_baseload=("daily_baseload", "first"),
                daily_peakload=("daily_peakload", "first"),
                daily_offpeak=("daily_offpeak", "first"),
                avg_price=("price", "mean"),
            ).round(2).reset_index()
            return daily.to_json(orient="records", date_format="iso")

        elif name == "get_ancillary_results":
            services = inputs.get("services") or []
            if services:
                placeholders = ",".join(["%s"] * len(services))
                df = _query(
                    f"SELECT efa_date, efa, service, clearing_price, cleared_volume, service_type "
                    f"FROM intl_market.gb_dx_results "
                    f"WHERE efa_date BETWEEN %s AND %s AND service IN ({placeholders}) "
                    f"ORDER BY efa_date, efa, service",
                    (inputs["start_date"], inputs["end_date"], *services),
                )
            else:
                df = _query(
                    "SELECT efa_date, efa, service, clearing_price, cleared_volume, service_type "
                    "FROM intl_market.gb_dx_results "
                    "WHERE efa_date BETWEEN %s AND %s ORDER BY efa_date, efa, service",
                    (inputs["start_date"], inputs["end_date"]),
                )
            if df.empty:
                return "No DX ancillary data for the requested period."
            summary = df.groupby(["service"]).agg(
                avg_clearing_price=("clearing_price", "mean"),
                avg_cleared_volume=("cleared_volume", "mean"),
                min_price=("clearing_price", "min"),
                max_price=("clearing_price", "max"),
            ).round(2).reset_index()
            return summary.to_json(orient="records")

        elif name == "get_market_summary":
            sp = _query(
                "SELECT date, AVG(system_price) AS avg_system_price "
                "FROM intl_market.gb_system_price "
                "WHERE date BETWEEN %s AND %s GROUP BY date ORDER BY date",
                (inputs["start_date"], inputs["end_date"]),
            )
            epex = _query(
                "SELECT delivery_date AS date, MAX(daily_baseload) AS epex_baseload "
                "FROM intl_market.gb_epex_da_hh "
                "WHERE delivery_date BETWEEN %s AND %s GROUP BY delivery_date ORDER BY delivery_date",
                (inputs["start_date"], inputs["end_date"]),
            )
            niv = _query(
                "SELECT date, AVG(niv) AS avg_niv "
                "FROM intl_market.gb_niv "
                "WHERE date BETWEEN %s AND %s GROUP BY date ORDER BY date",
                (inputs["start_date"], inputs["end_date"]),
            )
            merged = sp.merge(epex, on="date", how="outer").merge(niv, on="date", how="outer")
            merged["spread_sys_epex"] = (merged["avg_system_price"] - merged["epex_baseload"]).round(2)
            merged = merged.round(2)
            if merged.empty:
                return "No data available for the requested period."
            return merged.to_json(orient="records", date_format="iso")

        elif name == "get_bess_leaderboard":
            market = inputs.get("market")
            top_n  = inputs.get("top_n", 20)
            if market:
                df = _query(
                    "WITH lb AS ("
                    "  SELECT asset, SUM(revenue) AS total_revenue, "
                    "    AVG(revspermw)*48 AS avg_rev_per_mw_day "
                    "  FROM intl_market.gb_bess_leaderboard "
                    "  WHERE settlement_date BETWEEN %s AND %s AND market = %s "
                    "  GROUP BY asset ORDER BY total_revenue DESC LIMIT %s"
                    "), "
                    "op AS (SELECT DISTINCT ON (asset) asset, value AS operator "
                    "       FROM intl_market.gb_bess_assets WHERE history_table='operator' "
                    "       ORDER BY asset, valid_from DESC), "
                    "ow AS (SELECT DISTINCT ON (asset) asset, value AS owner "
                    "       FROM intl_market.gb_bess_assets WHERE history_table='owner' "
                    "       ORDER BY asset, valid_from DESC) "
                    "SELECT lb.asset, ow.owner, op.operator, lb.total_revenue, lb.avg_rev_per_mw_day "
                    "FROM lb LEFT JOIN op ON op.asset=lb.asset LEFT JOIN ow ON ow.asset=lb.asset",
                    (inputs["start_date"], inputs["end_date"], market, top_n),
                )
            else:
                df = _query(
                    "WITH lb AS ("
                    "  SELECT asset, SUM(revenue) AS total_revenue, "
                    "    AVG(revspermw)*48 AS avg_rev_per_mw_day "
                    "  FROM intl_market.gb_bess_leaderboard "
                    "  WHERE settlement_date BETWEEN %s AND %s "
                    "  GROUP BY asset ORDER BY total_revenue DESC LIMIT %s"
                    "), "
                    "op AS (SELECT DISTINCT ON (asset) asset, value AS operator "
                    "       FROM intl_market.gb_bess_assets WHERE history_table='operator' "
                    "       ORDER BY asset, valid_from DESC), "
                    "ow AS (SELECT DISTINCT ON (asset) asset, value AS owner "
                    "       FROM intl_market.gb_bess_assets WHERE history_table='owner' "
                    "       ORDER BY asset, valid_from DESC) "
                    "SELECT lb.asset, ow.owner, op.operator, lb.total_revenue, lb.avg_rev_per_mw_day "
                    "FROM lb LEFT JOIN op ON op.asset=lb.asset LEFT JOIN ow ON ow.asset=lb.asset",
                    (inputs["start_date"], inputs["end_date"], top_n),
                )
            if df.empty:
                return "No leaderboard data for the requested period."
            return df.round(2).to_json(orient="records")

        elif name == "get_bess_revenue_index":
            granularity = inputs.get("granularity", "monthly")
            if granularity == "daily":
                df = _query(
                    "SELECT settlement_date, market, revenue_permw, revenue_permwh "
                    "FROM intl_market.gb_bess_daily_index "
                    "WHERE settlement_date BETWEEN %s AND %s AND duration='*' "
                    "ORDER BY settlement_date, market",
                    (inputs["start_date"], inputs["end_date"]),
                )
            else:
                month_from = inputs["start_date"][:7] + "-01"
                month_to   = inputs["end_date"][:7]   + "-01"
                df = _query(
                    "SELECT month, market, revenue_permw, revenue_permwh "
                    "FROM intl_market.gb_bess_monthly_index "
                    "WHERE month BETWEEN %s AND %s AND duration='*' "
                    "ORDER BY month, market",
                    (month_from, month_to),
                )
            if df.empty:
                return "No BESS revenue index data for the requested period."
            return df.round(2).to_json(orient="records", date_format="iso")

        elif name == "get_bess_assets":
            conditions = ["history_table = 'rated_power'"]
            params: list = []
            if inputs.get("min_power_mw"):
                conditions.append("CAST(value AS NUMERIC) >= %s")
                params.append(inputs["min_power_mw"])
            where = " AND ".join(conditions)
            df = _query(
                f"WITH rp AS ("
                f"  SELECT DISTINCT ON (asset) asset, CAST(value AS NUMERIC) AS rated_power_mw, "
                f"  commissioning_date, gsp, developer, manufacturer, latitude, longitude "
                f"  FROM intl_market.gb_bess_assets WHERE {where} ORDER BY asset, valid_from DESC"
                f"), "
                f"ec AS (SELECT DISTINCT ON (asset) asset, CAST(value AS NUMERIC) AS energy_capacity_mwh "
                f"       FROM intl_market.gb_bess_assets WHERE history_table='energy_capacity' "
                f"       ORDER BY asset, valid_from DESC), "
                f"op AS (SELECT DISTINCT ON (asset) asset, value AS operator "
                f"       FROM intl_market.gb_bess_assets WHERE history_table='operator' "
                f"       ORDER BY asset, valid_from DESC), "
                f"ow AS (SELECT DISTINCT ON (asset) asset, value AS owner "
                f"       FROM intl_market.gb_bess_assets WHERE history_table='owner' "
                f"       ORDER BY asset, valid_from DESC) "
                f"SELECT rp.asset, ow.owner, op.operator, rp.developer, rp.manufacturer, "
                f"  rp.rated_power_mw, ec.energy_capacity_mwh, rp.commissioning_date, rp.gsp "
                f"FROM rp LEFT JOIN ec ON ec.asset=rp.asset "
                f"LEFT JOIN op ON op.asset=rp.asset LEFT JOIN ow ON ow.asset=rp.asset",
                params or None,
            )
            # Apply owner/operator post-filters
            if inputs.get("owner") and not df.empty:
                df = df[df["owner"].str.contains(inputs["owner"], case=False, na=False)]
            if inputs.get("operator") and not df.empty:
                df = df[df["operator"].str.contains(inputs["operator"], case=False, na=False)]
            if df.empty:
                return "No assets found matching the criteria."
            return (
                f"Total: {len(df)} assets, "
                f"Total rated power: {df['rated_power_mw'].sum():.0f} MW\n"
                + df.to_json(orient="records", date_format="iso")
            )

    except Exception as e:
        return f"Error: {e}"
    return "Unknown tool"


# ---------------------------------------------------------------------------
# Quant tools
# ---------------------------------------------------------------------------

_QUANT_TOOLS = [
    {
        "name": "get_bess_daily_index",
        "description": "Daily GB BESS industry-average revenue (£/MW/day, £/MWh/day) by market stream.",
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string"},
                "end_date":   {"type": "string"},
                "markets": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Filter by market e.g. ['bm','cm','frequency_response']. Leave empty for all.",
                },
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "get_bess_monthly_index",
        "description": "Monthly GB BESS industry-average revenue (£/MW/month) by market stream.",
        "input_schema": {
            "type": "object",
            "properties": {
                "month_from": {"type": "string", "description": "YYYY-MM e.g. '2024-01'"},
                "month_to":   {"type": "string", "description": "YYYY-MM e.g. '2026-05'"},
                "markets": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Filter by market. Leave empty for all.",
                },
            },
            "required": ["month_from", "month_to"],
        },
    },
    {
        "name": "get_leaderboard",
        "description": "Asset-level BESS performance leaderboard (revenue and £/MW) for a date range.",
        "input_schema": {
            "type": "object",
            "properties": {
                "start_date": {"type": "string"},
                "end_date":   {"type": "string"},
                "market": {
                    "type": "string",
                    "description": "Filter by specific market e.g. 'dml','dcl','bm'. Leave empty for all.",
                },
                "top_n": {"type": "integer", "description": "Number of top assets to return (default 20)."},
            },
            "required": ["start_date", "end_date"],
        },
    },
    {
        "name": "get_asset_database",
        "description": "GB BESS asset registry: power, capacity, location, developer, GSP.",
        "input_schema": {
            "type": "object",
            "properties": {
                "min_power_mw": {"type": "number", "description": "Minimum rated power filter."},
                "developer":    {"type": "string", "description": "Filter by developer name (partial match)."},
            },
            "required": [],
        },
    },
    {
        "name": "estimate_irr",
        "description": (
            "Parametric unlevered IRR for a GB BESS project. "
            "Uses the Modo monthly index (last 12 available months) as the revenue proxy. "
            "Returns IRR + sensitivity table across capex and revenue scenarios."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "power_mw":         {"type": "number", "description": "BESS rated power in MW."},
                "duration_h":       {"type": "number", "description": "Storage duration in hours."},
                "capex_gbp_per_kw": {"type": "number", "description": "Total capex in £/kW."},
                "opex_pct_capex":   {"type": "number", "description": "Annual O&M as % of capex (default 2.0)."},
                "project_life_yrs": {"type": "integer", "description": "Project life in years (default 15)."},
            },
            "required": ["power_mw", "duration_h", "capex_gbp_per_kw"],
        },
    },
]


def _compute_irr(cashflows: list[float]) -> float:
    """Newton-Raphson IRR from a list of cashflows starting with capex (negative)."""
    rate = 0.1
    for _ in range(100):
        npv = sum(cf / (1 + rate) ** t for t, cf in enumerate(cashflows))
        dnpv = sum(-t * cf / (1 + rate) ** (t + 1) for t, cf in enumerate(cashflows))
        if abs(dnpv) < 1e-10:
            break
        rate -= npv / dnpv
        if rate <= -1:
            rate = -0.999
    return rate


def _dispatch_quant(name: str, inputs: dict) -> str:
    try:
        if name == "get_bess_daily_index":
            markets = inputs.get("markets") or []
            if markets:
                ph = ",".join(["%s"] * len(markets))
                df = _query(
                    f"SELECT settlement_date, market, revenue_permw, revenue_permwh "
                    f"FROM intl_market.gb_bess_daily_index "
                    f"WHERE settlement_date BETWEEN %s AND %s AND market IN ({ph}) "
                    f"ORDER BY settlement_date, market",
                    (inputs["start_date"], inputs["end_date"], *markets),
                )
            else:
                df = _query(
                    "SELECT settlement_date, market, revenue_permw, revenue_permwh "
                    "FROM intl_market.gb_bess_daily_index "
                    "WHERE settlement_date BETWEEN %s AND %s AND duration = '*' "
                    "ORDER BY settlement_date, market",
                    (inputs["start_date"], inputs["end_date"]),
                )
            if df.empty:
                return "No daily index data for the requested period."
            return df.to_json(orient="records", date_format="iso")

        elif name == "get_bess_monthly_index":
            month_from = inputs["month_from"] + "-01"
            month_to   = inputs["month_to"]   + "-01"
            markets = inputs.get("markets") or []
            if markets:
                ph = ",".join(["%s"] * len(markets))
                df = _query(
                    f"SELECT month, market, revenue_permw, revenue_permwh "
                    f"FROM intl_market.gb_bess_monthly_index "
                    f"WHERE month BETWEEN %s AND %s AND market IN ({ph}) "
                    f"ORDER BY month, market",
                    (month_from, month_to, *markets),
                )
            else:
                df = _query(
                    "SELECT month, market, revenue_permw, revenue_permwh "
                    "FROM intl_market.gb_bess_monthly_index "
                    "WHERE month BETWEEN %s AND %s AND duration = '*' "
                    "ORDER BY month, market",
                    (month_from, month_to),
                )
            if df.empty:
                return "No monthly index data for the requested period."
            return df.to_json(orient="records", date_format="iso")

        elif name == "get_leaderboard":
            market = inputs.get("market")
            top_n  = inputs.get("top_n", 20)
            if market:
                df = _query(
                    "SELECT asset, market, SUM(revenue) AS total_revenue, "
                    "AVG(revspermw) AS avg_revspermw, AVG(rated_power) AS rated_power_mw "
                    "FROM intl_market.gb_bess_leaderboard "
                    "WHERE settlement_date BETWEEN %s AND %s AND market = %s "
                    "GROUP BY asset, market ORDER BY total_revenue DESC LIMIT %s",
                    (inputs["start_date"], inputs["end_date"], market, top_n),
                )
            else:
                df = _query(
                    "SELECT asset, SUM(revenue) AS total_revenue, "
                    "AVG(revspermw) AS avg_revspermw, AVG(rated_power) AS rated_power_mw "
                    "FROM intl_market.gb_bess_leaderboard "
                    "WHERE settlement_date BETWEEN %s AND %s "
                    "GROUP BY asset ORDER BY total_revenue DESC LIMIT %s",
                    (inputs["start_date"], inputs["end_date"], top_n),
                )
            if df.empty:
                return "No leaderboard data for the requested period."
            return df.round(2).to_json(orient="records")

        elif name == "get_asset_database":
            # Get unique assets with their rated power (history_table = 'rated_power')
            conditions = ["history_table = 'rated_power'"]
            params: list = []
            min_mw = inputs.get("min_power_mw")
            developer = inputs.get("developer")
            if min_mw is not None:
                conditions.append("CAST(value AS NUMERIC) >= %s")
                params.append(min_mw)
            if developer:
                conditions.append("developer ILIKE %s")
                params.append(f"%{developer}%")
            where = " AND ".join(conditions)
            df = _query(
                f"SELECT DISTINCT ON (asset) asset, "
                f"CAST(value AS NUMERIC) AS rated_power_mw, latitude, longitude, "
                f"gsp, developer, manufacturer, commissioning_date, dno, is_co_located "
                f"FROM intl_market.gb_bess_assets WHERE {where} "
                f"ORDER BY asset, valid_from DESC",
                params or None,
            )
            if df.empty:
                return "No assets found matching the criteria."
            return (
                f"Total: {len(df)} assets, "
                f"Total rated power: {df['rated_power_mw'].sum():.0f} MW\n"
                + df.head(50).to_json(orient="records", date_format="iso")
            )

        elif name == "estimate_irr":
            power_mw   = float(inputs["power_mw"])
            duration_h = float(inputs["duration_h"])
            capex_per_kw = float(inputs["capex_gbp_per_kw"])
            opex_pct   = float(inputs.get("opex_pct_capex", 2.0)) / 100
            life_yrs   = int(inputs.get("project_life_yrs", 15))

            # Fetch last 12 months of monthly index, all markets combined
            df = _query(
                "SELECT month, SUM(revenue_permw) AS total_revpermw "
                "FROM intl_market.gb_bess_monthly_index "
                "WHERE duration = '*' "
                "GROUP BY month ORDER BY month DESC LIMIT 12"
            )
            if df.empty:
                return "No monthly index data available to estimate IRR."

            avg_monthly_rev_per_mw = df["total_revpermw"].mean()
            annual_rev_per_mw = avg_monthly_rev_per_mw * 12

            capex_total = power_mw * capex_per_kw * 1000  # £
            opex_annual = capex_total * opex_pct

            cashflows = [-capex_total]
            for yr in range(1, life_yrs + 1):
                degrad = (1 - 0.02) ** (yr - 1)
                rev = power_mw * annual_rev_per_mw * degrad
                cashflows.append(rev - opex_annual)

            irr = _compute_irr(cashflows)
            npv_10 = sum(cf / 1.10 ** t for t, cf in enumerate(cashflows))

            # Sensitivity: ±20% capex, ±20% revenue
            sens = []
            for capex_mult in [0.8, 1.0, 1.2]:
                for rev_mult in [0.8, 1.0, 1.2]:
                    c0 = -capex_total * capex_mult
                    cfs = [c0] + [
                        power_mw * annual_rev_per_mw * rev_mult * (1 - 0.02) ** (yr - 1) - opex_annual
                        for yr in range(1, life_yrs + 1)
                    ]
                    sens.append({
                        "capex_factor": f"{capex_mult:.0%}",
                        "revenue_factor": f"{rev_mult:.0%}",
                        "irr": f"{_compute_irr(cfs) * 100:.1f}%",
                    })

            return json.dumps({
                "inputs": {
                    "power_mw": power_mw, "duration_h": duration_h,
                    "capex_gbp_per_kw": capex_per_kw,
                    "avg_monthly_rev_per_mw_gbp": round(avg_monthly_rev_per_mw, 0),
                    "annual_rev_per_mw_gbp": round(annual_rev_per_mw, 0),
                    "capex_total_gbp": round(capex_total, 0),
                    "opex_annual_gbp": round(opex_annual, 0),
                    "project_life_yrs": life_yrs,
                },
                "result": {
                    "unlevered_irr": f"{irr * 100:.1f}%",
                    "npv_at_10pct_gbp": round(npv_10, 0),
                    "revenue_source": f"Modo monthly index avg (last {len(df)} months)",
                },
                "sensitivity": sens,
            }, indent=2)

    except Exception as e:
        return f"Error: {e}"
    return "Unknown tool"


# ---------------------------------------------------------------------------
# Agent turn loop
# ---------------------------------------------------------------------------

def _run_agent_turn(messages: list, system: str, tools: list, dispatch_fn) -> tuple[str, list, list]:
    tool_events: list[dict] = []
    while True:
        resp = _client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=4096,
            system=system,
            tools=tools,
            messages=messages,
        )
        messages = messages + [{"role": "assistant", "content": resp.content}]

        if resp.stop_reason == "end_turn":
            text = next((b.text for b in resp.content if hasattr(b, "text")), "")
            return text, messages, tool_events

        tool_results = []
        for block in resp.content:
            if block.type == "tool_use":
                result_str = dispatch_fn(block.name, block.input)
                tool_events.append({"tool": block.name, "result": result_str[:200]})
                tool_results.append({
                    "type": "tool_result",
                    "tool_use_id": block.id,
                    "content": result_str,
                })
        messages = messages + [{"role": "user", "content": tool_results}]


# ---------------------------------------------------------------------------
# Data query helpers for visualisation tabs
# ---------------------------------------------------------------------------

@st.cache_data(ttl=300)
def _get_system_price_range(start: str, end: str) -> pd.DataFrame:
    return _query(
        "SELECT sp.date, sp.settlement_period, sp.system_price, n.niv "
        "FROM intl_market.gb_system_price sp "
        "LEFT JOIN intl_market.gb_niv n "
        "  ON sp.date = n.date AND sp.settlement_period = n.settlement_period "
        "WHERE sp.date BETWEEN %s AND %s ORDER BY sp.date, sp.settlement_period",
        (start, end),
    )


@st.cache_data(ttl=3600)
def _get_epex_range(start: str, end: str) -> pd.DataFrame:
    return _query(
        "SELECT delivery_date, settlement_period, price, daily_baseload, daily_peakload, daily_offpeak "
        "FROM intl_market.gb_epex_da_hh WHERE delivery_date BETWEEN %s AND %s "
        "ORDER BY delivery_date, settlement_period",
        (start, end),
    )


@st.cache_data(ttl=3600)
def _get_system_price_daily(start: str, end: str) -> pd.DataFrame:
    """Daily avg system price + avg NIV — fast overview query (30 rows vs 1,440)."""
    return _query(
        "SELECT sp.date, AVG(sp.system_price) AS avg_system_price, "
        "AVG(n.niv) AS avg_niv "
        "FROM intl_market.gb_system_price sp "
        "LEFT JOIN intl_market.gb_niv n "
        "  ON sp.date = n.date AND sp.settlement_period = n.settlement_period "
        "WHERE sp.date BETWEEN %s AND %s GROUP BY sp.date ORDER BY sp.date",
        (start, end),
    )


@st.cache_data(ttl=3600)
def _get_epex_overview(start: str, end: str) -> pd.DataFrame:
    """Daily baseload/peakload/offpeak only — fast overview query (30 rows vs 1,440)."""
    return _query(
        "SELECT delivery_date, MAX(daily_baseload) AS daily_baseload, "
        "MAX(daily_peakload) AS daily_peakload, MAX(daily_offpeak) AS daily_offpeak "
        "FROM intl_market.gb_epex_da_hh WHERE delivery_date BETWEEN %s AND %s "
        "GROUP BY delivery_date ORDER BY delivery_date",
        (start, end),
    )


@st.cache_data(ttl=3600)
def _get_system_price_hourly(start: str, end: str) -> pd.DataFrame:
    """Hourly system price — avg of the two half-hourly SPs within each clock-hour."""
    return _query(
        "SELECT date, "
        "  (FLOOR((settlement_period - 1) / 2))::int AS hour_of_day, "
        "  AVG(system_price) AS system_price "
        "FROM intl_market.gb_system_price "
        "WHERE date BETWEEN %s AND %s "
        "GROUP BY date, hour_of_day "
        "ORDER BY date, hour_of_day",
        (start, end),
    )


@st.cache_data(ttl=3600)
def _get_epex_hourly(start: str, end: str) -> pd.DataFrame:
    """Hourly EPEX DA price — avg of the two half-hourly SPs within each clock-hour."""
    return _query(
        "SELECT delivery_date, "
        "  (FLOOR((settlement_period - 1) / 2))::int AS hour_of_day, "
        "  AVG(price) AS price "
        "FROM intl_market.gb_epex_da_hh "
        "WHERE delivery_date BETWEEN %s AND %s "
        "GROUP BY delivery_date, hour_of_day "
        "ORDER BY delivery_date, hour_of_day",
        (start, end),
    )


@st.cache_data(ttl=300)
def _get_pricing_missing_dates(start: str, end: str) -> list[str]:
    """Return ISO date strings in [start, end] with no rows in gb_pricing_results."""
    from datetime import date as _d, timedelta as _td
    have = set()
    try:
        df = _query(
            "SELECT DISTINCT settlement_date FROM intl_market.gb_pricing_results "
            "WHERE settlement_date BETWEEN %s AND %s",
            (start, end),
        )
        have = {str(r) for r in df["settlement_date"]}
    except Exception:
        pass
    s, e = _d.fromisoformat(start), _d.fromisoformat(end)
    return [str(s + _td(days=i)) for i in range((e - s).days + 1)
            if str(s + _td(days=i)) not in have]


@st.cache_data(ttl=300)
def _get_fuel_mix_missing_dates(start: str, end: str) -> list[str]:
    """Return ISO date strings in [start, end] with no rows in gb_fuel_mix."""
    from datetime import date as _d, timedelta as _td
    have = set()
    try:
        df = _query(
            "SELECT DISTINCT settlement_date FROM intl_market.gb_fuel_mix "
            "WHERE settlement_date BETWEEN %s AND %s",
            (start, end),
        )
        have = {str(r) for r in df["settlement_date"]}
    except Exception:
        pass
    s, e = _d.fromisoformat(start), _d.fromisoformat(end)
    return [str(s + _td(days=i)) for i in range((e - s).days + 1)
            if str(s + _td(days=i)) not in have]


@st.cache_data(ttl=300)
def _get_dx_range(start: str, end: str) -> pd.DataFrame:
    return _query(
        "SELECT efa_date, efa, service, clearing_price, cleared_volume "
        "FROM intl_market.gb_dx_results WHERE efa_date BETWEEN %s AND %s "
        "ORDER BY efa_date, efa, service",
        (start, end),
    )


@st.cache_data(ttl=3600, show_spinner=False)
def _get_fuel_mix_daily(start: str, end: str) -> pd.DataFrame:
    """Daily average MW by fuel from gb_fuel_mix."""
    try:
        return _query(
            "SELECT settlement_date, "
            "  AVG(gas_mw) AS gas, AVG(coal_mw) AS coal, "
            "  AVG(nuclear_mw) AS nuclear, "
            "  AVG(COALESCE(wind_mw,0) + COALESCE(wind_emb_mw,0)) AS wind, "
            "  AVG(solar_mw) AS solar, AVG(hydro_mw) AS hydro, "
            "  AVG(imports_mw) AS imports, AVG(biomass_mw) AS biomass, "
            "  AVG(storage_mw) AS storage "
            "FROM intl_market.gb_fuel_mix "
            "WHERE settlement_date BETWEEN %s AND %s "
            "GROUP BY settlement_date ORDER BY settlement_date",
            (start, end),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=3600, show_spinner=False)
def _get_bidding_space_hh(start: str, end: str) -> pd.DataFrame:
    """Half-hourly bidding space (thermal + flexible) from gb_fuel_mix."""
    try:
        return _query(
            "SELECT settlement_date, settlement_period, "
            "  (COALESCE(generation_mw,0) "
            "   - COALESCE(nuclear_mw,0) - COALESCE(wind_mw,0) "
            "   - COALESCE(wind_emb_mw,0) - COALESCE(solar_mw,0) "
            "   - COALESCE(hydro_mw,0) - COALESCE(imports_mw,0)) AS bidding_space_mw "
            "FROM intl_market.gb_fuel_mix "
            "WHERE settlement_date BETWEEN %s AND %s "
            "ORDER BY settlement_date, settlement_period",
            (start, end),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=3600, show_spinner=False)
def _get_fuel_mix_daily_dates(dates: tuple) -> pd.DataFrame:
    """Daily average MW by fuel for specific dates (tuple of ISO date strings)."""
    if not dates:
        return pd.DataFrame()
    placeholders = ",".join(["%s"] * len(dates))
    try:
        return _query(
            "SELECT settlement_date, "
            "  AVG(gas_mw) AS gas, AVG(coal_mw) AS coal, "
            "  AVG(nuclear_mw) AS nuclear, "
            "  AVG(COALESCE(wind_mw,0) + COALESCE(wind_emb_mw,0)) AS wind, "
            "  AVG(solar_mw) AS solar, AVG(hydro_mw) AS hydro, "
            "  AVG(imports_mw) AS imports, AVG(biomass_mw) AS biomass, "
            "  AVG(storage_mw) AS storage "
            "FROM intl_market.gb_fuel_mix "
            f"WHERE settlement_date IN ({placeholders}) "
            "GROUP BY settlement_date ORDER BY settlement_date",
            dates,
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=3600, show_spinner=False)
def _get_bidding_space_hh_dates(dates: tuple) -> pd.DataFrame:
    """Half-hourly bidding space for specific dates (tuple of ISO date strings)."""
    if not dates:
        return pd.DataFrame()
    placeholders = ",".join(["%s"] * len(dates))
    try:
        return _query(
            "SELECT settlement_date, settlement_period, "
            "  (COALESCE(generation_mw,0) "
            "   - COALESCE(nuclear_mw,0) - COALESCE(wind_mw,0) "
            "   - COALESCE(wind_emb_mw,0) - COALESCE(solar_mw,0) "
            "   - COALESCE(hydro_mw,0) - COALESCE(imports_mw,0)) AS bidding_space_mw "
            "FROM intl_market.gb_fuel_mix "
            f"WHERE settlement_date IN ({placeholders}) "
            "ORDER BY settlement_date, settlement_period",
            dates,
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=3600, show_spinner=False)
def _get_epex_hh_dates(dates: tuple) -> pd.DataFrame:
    """Half-hourly EPEX DA prices for specific delivery dates (tuple of ISO date strings)."""
    if not dates:
        return pd.DataFrame()
    placeholders = ",".join(["%s"] * len(dates))
    try:
        return _query(
            "SELECT delivery_date, settlement_period, price "
            "FROM intl_market.gb_epex_da_hh "
            f"WHERE delivery_date IN ({placeholders}) "
            "ORDER BY delivery_date, settlement_period",
            dates,
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=1800, show_spinner=False)
def _get_pricing_table(start: str, end: str, top_n: int = 20) -> pd.DataFrame:
    """Top N assets by wholesale revenue with pricing model results."""
    try:
        return _query(
            "WITH wholesale AS ( "
            "  SELECT asset, SUM(revenue) AS wholesale_gbp, "
            "    AVG(rated_power) AS power_mw "
            "  FROM intl_market.gb_bess_leaderboard "
            "  WHERE settlement_date BETWEEN %s AND %s AND market = 'wholesale' "
            "  GROUP BY asset ORDER BY wholesale_gbp DESC LIMIT %s "
            "), "
            "owner_op AS ( "
            "  SELECT DISTINCT ON (asset) asset, "
            "    MAX(CASE WHEN history_table='owner'    THEN value END) OVER (PARTITION BY asset) AS owner, "
            "    MAX(CASE WHEN history_table='operator' THEN value END) OVER (PARTITION BY asset) AS operator "
            "  FROM intl_market.gb_bess_assets "
            "  WHERE history_table IN ('owner','operator') "
            "), "
            "pricing AS ( "
            "  SELECT asset_name, "
            "    AVG(options_value_gbp_per_mw)   AS options_val_per_mw, "
            "    AVG(pf_actual_da_pnl_gbp)       AS pf_actual_da, "
            "    AVG(COALESCE(pf_actual_sp_pnl_gbp, pf_actual_da_pnl_gbp)) AS pf_actual, "
            "    AVG(duration_h)                 AS duration_h, "
            "    MAX(settlement_date)            AS latest_batch_date, "
            "    AVG(COALESCE("
            "      (SELECT SUM(-LEAST(elem::numeric, 0)) * 0.5 "
            "       FROM jsonb_array_elements_text("
            "         COALESCE(pf_actual_sp_dispatch_48, pf_actual_dispatch_48)) AS elem), "
            "      0)) AS avg_charged_mwh "
            "  FROM intl_market.gb_pricing_results "
            "  WHERE settlement_date BETWEEN %s AND %s "
            "  GROUP BY asset_name "
            ") "
            "SELECT ROW_NUMBER() OVER (ORDER BY w.wholesale_gbp DESC) AS rank, "
            "  w.asset, oo.owner, oo.operator, "
            "  p.duration_h, w.power_mw, "
            "  w.wholesale_gbp, "
            "  p.options_val_per_mw, p.pf_actual, p.avg_charged_mwh, "
            "  p.latest_batch_date "
            "FROM wholesale w "
            "LEFT JOIN owner_op oo ON oo.asset = w.asset "
            "LEFT JOIN pricing p   ON p.asset_name = w.asset "
            "ORDER BY w.wholesale_gbp DESC",
            (start, end, top_n, start, end),
        )
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=1800, show_spinner=False)
def _get_dispatch_comparison(asset: str, settlement_date: str) -> pd.DataFrame:
    """PF dispatch arrays + prices for an asset on a given date."""
    try:
        df = _query(
            "SELECT pf_actual_dispatch_48, pf_forecast_dispatch_48, "
            "actual_epex_48, forecast_epex_48, "
            "COALESCE(pf_actual_sp_dispatch_48, pf_actual_dispatch_48) AS pf_sp_dispatch_48, "
            "COALESCE(actual_sp_48, actual_epex_48) AS actual_sp_48 "
            "FROM intl_market.gb_pricing_results "
            "WHERE asset_name = %s AND settlement_date = %s",
            (asset, settlement_date),
        )
        return df
    except Exception:
        return pd.DataFrame()


@st.cache_data(ttl=300)
def _get_daily_index_range(start: str, end: str) -> pd.DataFrame:
    return _query(
        "SELECT settlement_date, market, revenue_permw, revenue_permwh "
        "FROM intl_market.gb_bess_daily_index WHERE settlement_date BETWEEN %s AND %s AND duration = '*' "
        "ORDER BY settlement_date, market",
        (start, end),
    )


@st.cache_data(ttl=300)
def _get_monthly_index_range(start: str, end: str) -> pd.DataFrame:
    return _query(
        "SELECT month, market, revenue_permw "
        "FROM intl_market.gb_bess_monthly_index WHERE month BETWEEN %s AND %s AND duration = '*' "
        "ORDER BY month, market",
        (start, end),
    )


@st.cache_data(ttl=300)
def _get_asset_revenue_map(start: str, end: str, market: str) -> pd.DataFrame:
    """Per-asset revenue per MW for the map colour scale."""
    return _query(
        "SELECT asset, "
        "  SUM(revenue) / NULLIF(SUM(rated_power), 0) AS rev_per_mw, "
        "  SUM(revenue) AS total_revenue "
        "FROM intl_market.gb_bess_leaderboard "
        "WHERE settlement_date BETWEEN %s AND %s AND market = %s "
        "GROUP BY asset",
        (start, end, market),
    )


@st.cache_data(ttl=300)
def _get_leaderboard_range(start: str, end: str, top_n: int = 20) -> pd.DataFrame:
    # Aggregate in two passes: first sum across markets per SP, then aggregate across SPs.
    # This avoids the broken market='total' filter (that market name doesn't exist).
    df = _query(
        "WITH per_sp AS ( "
        "  SELECT settlement_date, settlement_period, asset, "
        "    SUM(revenue)     AS total_rev, "
        "    SUM(revspermw)   AS total_revspermw, "
        "    SUM(revspermwh)  AS total_revspermwh, "
        "    AVG(rated_power) AS rated_power, "
        "    SUM(CASE WHEN market='wholesale'          THEN revenue ELSE 0 END) AS wholesale, "
        "    SUM(CASE WHEN market='frequency_response' THEN revenue ELSE 0 END) AS freq_response, "
        "    SUM(CASE WHEN market='bm'                 THEN revenue ELSE 0 END) AS bm, "
        "    SUM(CASE WHEN market='imbalance'          THEN revenue ELSE 0 END) AS imbalance, "
        "    SUM(CASE WHEN market='reserve'            THEN revenue ELSE 0 END) AS reserve "
        "  FROM intl_market.gb_bess_leaderboard "
        "  WHERE settlement_date BETWEEN %s AND %s "
        "  GROUP BY settlement_date, settlement_period, asset "
        "), "
        "lb AS ( "
        "  SELECT asset, "
        "    SUM(total_rev)        AS total_revenue, "
        "    SUM(wholesale)        AS wholesale, "
        "    SUM(freq_response)    AS freq_response, "
        "    SUM(bm)               AS bm, "
        "    SUM(imbalance)        AS imbalance, "
        "    SUM(reserve)          AS reserve, "
        "    AVG(total_revspermw)  AS avg_revspermw, "
        "    AVG(total_revspermwh) AS avg_revspermwh, "
        "    AVG(rated_power)      AS rated_power_mw "
        "  FROM per_sp "
        "  GROUP BY asset ORDER BY total_revenue DESC LIMIT %s "
        "), "
        "assets AS ( "
        "  SELECT DISTINCT ON (asset) asset, developer, integrator "
        "  FROM intl_market.gb_bess_assets "
        "  WHERE history_table = 'rated_power' "
        "  ORDER BY asset, valid_from DESC "
        "), "
        "op AS ( "
        "  SELECT DISTINCT ON (asset) asset, value AS operator "
        "  FROM intl_market.gb_bess_assets WHERE history_table = 'operator' "
        "  ORDER BY asset, valid_from DESC "
        "), "
        "ow AS ( "
        "  SELECT DISTINCT ON (asset) asset, value AS owner "
        "  FROM intl_market.gb_bess_assets WHERE history_table = 'owner' "
        "  ORDER BY asset, valid_from DESC "
        "), "
        "ec AS ( "
        "  SELECT DISTINCT ON (asset) asset, CAST(value AS NUMERIC) AS energy_capacity_mwh "
        "  FROM intl_market.gb_bess_assets WHERE history_table = 'energy_capacity' "
        "  ORDER BY asset, valid_from DESC "
        ") "
        "SELECT lb.asset, ow.owner, op.operator, a.integrator, "
        "  lb.total_revenue, lb.wholesale, lb.freq_response, lb.bm, lb.imbalance, lb.reserve, "
        "  lb.avg_revspermw, lb.avg_revspermwh, lb.rated_power_mw, ec.energy_capacity_mwh "
        "FROM lb "
        "LEFT JOIN assets a ON a.asset = lb.asset "
        "LEFT JOIN op ON op.asset = lb.asset "
        "LEFT JOIN ow ON ow.asset = lb.asset "
        "LEFT JOIN ec ON ec.asset = lb.asset",
        (start, end, top_n),
    )
    return _reclassify_manufacturers(df) if not df.empty else df


# Companies that are BESS manufacturers, not operators — Modo sometimes populates
# the integrator field with these; we reclassify them at query time.
_KNOWN_MANUFACTURERS = {
    "BYD", "CATL", "Samsung SDI", "LG Energy Solution",
}

# Knowledge base source guide — descriptions and suggested topics per source
_KB_SOURCE_GUIDE = [
    {
        "source": "elexon",
        "label": "Elexon",
        "description": "Balancing Mechanism operator. Publishes system warnings, network transmission notices (NOTOs), and market notices affecting GB grid operations.",
        "topics": [
            "System price high/low alerts",
            "Network transmission outage notices",
            "Balancing mechanism instructions",
            "Grid constraint events",
            "Market suspension notices",
        ],
    },
    {
        "source": "entso_e",
        "label": "National Grid ESO",
        "description": "GB electricity system operator. Covers flexibility market design, demand side response, capacity market, ESO transformation, and Future Energy Scenarios.",
        "topics": [
            "Flexibility market roadmap",
            "Dynamic containment and moderation",
            "Stability pathfinder contracts",
            "Capacity market auctions",
            "Net zero system pathways",
        ],
    },
    {
        "source": "timera",
        "label": "Timera Energy",
        "description": "Independent energy analysis. In-depth BESS revenue stack analysis, battery storage market outlook, merchant risk, and European storage trends.",
        "topics": [
            "BESS revenue stack decomposition",
            "Battery storage investment returns",
            "Frequency response market outlook",
            "Merchant vs contracted risk",
            "European battery storage market",
        ],
    },
    {
        "source": "modo",
        "label": "Modo Energy",
        "description": "BESS-focused market intelligence. Asset performance benchmarks, operator strategies, dispatch analytics, and storage market commentary.",
        "topics": [
            "BESS asset performance benchmarks",
            "Battery operator revenue strategies",
            "Frequency response optimisation",
            "Wholesale trading vs ancillary",
            "Storage asset dispatch analysis",
        ],
    },
    {
        "source": "meteologica",
        "label": "Meteologica",
        "description": "Energy forecasting specialist. Wind and solar generation forecasts, demand forecasting methodology, and grid balancing implications.",
        "topics": [
            "Wind generation forecast accuracy",
            "Solar PV output forecasting",
            "Demand forecasting methodology",
            "Forecast uncertainty and imbalance",
            "Renewable intermittency impacts",
        ],
    },
]


def _reclassify_manufacturers(df: pd.DataFrame) -> pd.DataFrame:
    """Move integrator values that are manufacturers → manufacturer column."""
    if "integrator" not in df.columns:
        return df
    is_mfr = df["integrator"].isin(_KNOWN_MANUFACTURERS)
    if "manufacturer" in df.columns:
        df.loc[is_mfr & df["manufacturer"].isna(), "manufacturer"] = df.loc[
            is_mfr & df["manufacturer"].isna(), "integrator"
        ]
    df.loc[is_mfr, "integrator"] = None
    return df


@st.cache_data(ttl=3600)
def _get_assets() -> pd.DataFrame:
    df = _query(
        "WITH rp AS ( "
        "  SELECT DISTINCT ON (asset) asset, "
        "  CAST(value AS NUMERIC) AS rated_power_mw, latitude, longitude, "
        "  gsp, developer, integrator, manufacturer, commissioning_date, dno, "
        "  is_co_located, co_located_type "
        "  FROM intl_market.gb_bess_assets WHERE history_table = 'rated_power' "
        "  ORDER BY asset, valid_from DESC "
        "), "
        "op AS ( "
        "  SELECT DISTINCT ON (asset) asset, value AS operator "
        "  FROM intl_market.gb_bess_assets WHERE history_table = 'operator' "
        "  ORDER BY asset, valid_from DESC "
        "), "
        "ow AS ( "
        "  SELECT DISTINCT ON (asset) asset, value AS owner "
        "  FROM intl_market.gb_bess_assets WHERE history_table = 'owner' "
        "  ORDER BY asset, valid_from DESC "
        ") "
        "SELECT rp.*, op.operator, ow.owner "
        "FROM rp "
        "LEFT JOIN op ON op.asset = rp.asset "
        "LEFT JOIN ow ON ow.asset = rp.asset"
    )
    return _reclassify_manufacturers(df)


@st.cache_data(ttl=300)
def _table_counts() -> pd.DataFrame:
    # (table_name, date_column) — None means no time-series date
    _COVERAGE = [
        ("gb_system_price",      "date"),
        ("gb_niv",               "date"),
        ("gb_epex_da_hh",        "delivery_date"),
        ("gb_dx_results",        "efa_date"),
        ("gb_bess_daily_index",  "settlement_date"),
        ("gb_bess_monthly_index","month"),
        ("gb_bess_leaderboard",  "settlement_date"),
        ("gb_bess_assets",       None),
        ("gb_knowledge_docs",    None),
    ]
    rows = []
    for table, date_col in _COVERAGE:
        try:
            if date_col:
                df = _query(
                    f"SELECT COUNT(*) AS n, MIN({date_col})::text AS min_d, "
                    f"MAX({date_col})::text AS max_d FROM intl_market.{table}"
                )
                row = df.iloc[0]
                rows.append({
                    "Table": table,
                    "Rows": int(row["n"]),
                    "From": row["min_d"],
                    "To": row["max_d"],
                })
            else:
                df = _query(f"SELECT COUNT(*) AS n FROM intl_market.{table}")
                rows.append({
                    "Table": table,
                    "Rows": int(df["n"].iloc[0]),
                    "From": None,
                    "To": None,
                })
        except Exception:
            rows.append({"Table": table, "Rows": "error", "From": None, "To": None})
    return pd.DataFrame(rows)


# ---------------------------------------------------------------------------
# App initialisation (runs once per process via cache_resource / explicit call)
# ---------------------------------------------------------------------------

_start_scheduler()  # no-op after first call (cache_resource)
print("[INIT] done _start_scheduler", flush=True)

# ---------------------------------------------------------------------------
# Sidebar
# ---------------------------------------------------------------------------

with st.sidebar:
    st.title("🇬🇧 GB Market")
    st.caption("Powered by Modo Energy API")
    today = date.today()
    default_start = today - timedelta(days=30)

    st.subheader("Date Range")
    d_start = st.date_input("From", value=default_start, key="d_start")
    d_end   = st.date_input("To",   value=today,         key="d_end")
    date_start = d_start.isoformat()
    date_end   = d_end.isoformat()

    st.divider()
    st.caption("v1 · ap-southeast-1")


# ---------------------------------------------------------------------------
# Tabs
# ---------------------------------------------------------------------------

tab_overview, tab_ancillary, tab_bess, tab_pricing, tab_map, tab_strategist, tab_quant, tab_knowledge, tab_mgmt = st.tabs([
    "Market Overview", "Ancillary Markets", "BESS Benchmarking",
    "Pricing", "Asset Map", "Strategist", "Quant", "Knowledge Base", "Data Management",
])

# ---- Market Overview -------------------------------------------------------
with tab_overview:
    st.header("GB Market Overview")

    col1, col2 = st.columns(2)

    with col1:
        st.subheader("System Price (SBP/SSP) + EPEX DA (£/MWh)")
        sp_df = _get_system_price_daily(date_start, date_end)   # kept for NIV in col2
        _sp_hourly   = _get_system_price_hourly(date_start, date_end)
        _epex_hourly = _get_epex_hourly(date_start, date_end)
        if _sp_hourly.empty and _epex_hourly.empty:
            st.info("No system price data. Run a backfill in Data Management.")
        else:
            # Build datetime index: date + hour
            if not _sp_hourly.empty:
                _sp_hourly["ts"] = (
                    pd.to_datetime(_sp_hourly["date"].astype(str))
                    + pd.to_timedelta(_sp_hourly["hour_of_day"].astype(int), unit="h")
                )
            if not _epex_hourly.empty:
                _epex_hourly["ts"] = (
                    pd.to_datetime(_epex_hourly["delivery_date"].astype(str))
                    + pd.to_timedelta(_epex_hourly["hour_of_day"].astype(int), unit="h")
                )
            fig = go.Figure()
            if not _sp_hourly.empty:
                fig.add_trace(go.Scatter(
                    x=_sp_hourly["ts"], y=_sp_hourly["system_price"],
                    mode="lines", name="System Price (hourly)",
                    line=dict(color="#d62728", width=1),
                ))
            if not _epex_hourly.empty:
                fig.add_trace(go.Scatter(
                    x=_epex_hourly["ts"], y=_epex_hourly["price"],
                    mode="lines", name="EPEX DA (hourly)",
                    line=dict(color="#2ca02c", width=1.5),
                ))
            fig.add_hline(y=0, line_dash="dash", line_color="gray", line_width=1)
            fig.update_layout(
                margin=dict(l=0, r=0, t=0, b=0), height=300,
                legend=dict(orientation="h", yanchor="bottom", y=1.01, x=0),
                yaxis_title="£/MWh",
            )
            st.plotly_chart(fig, use_container_width=True)

    with col2:
        st.subheader("Net Imbalance Volume (MW, daily avg)")
        if sp_df.empty:
            st.info("No NIV data.")
        else:
            niv_df = sp_df.dropna(subset=["avg_niv"])
            if not niv_df.empty:
                niv_df = niv_df.copy()
                niv_df["colour"] = niv_df["avg_niv"].apply(lambda x: "short" if x < 0 else "long")
                fig2 = px.bar(niv_df, x="date", y="avg_niv", color="colour",
                              color_discrete_map={"short": "#d62728", "long": "#2ca02c"},
                              labels={"avg_niv": "MW", "date": ""})
                fig2.update_layout(margin=dict(l=0, r=0, t=0, b=0), height=300, showlegend=False)
                st.plotly_chart(fig2, use_container_width=True)
            else:
                st.info("No NIV data in range.")

    st.caption(
        "**SBP** (System Buy Price) = price paid by short parties; "
        "**SSP** (System Sell Price) = price received by long parties. "
        "SBP > SSP when the system is short (Elexon buys energy to balance). "
        "**EPEX DA** = day-ahead auction clearing price settled at gate closure (11:00 D-1)."
    )

    st.subheader("EPEX Day-Ahead Prices — Heatmap (£/MWh)")
    epex_df = _get_epex_range(date_start, date_end)
    if epex_df.empty:
        st.info("No EPEX DA data.")
    else:
        pivot = epex_df.pivot_table(index="settlement_period", columns="delivery_date",
                                    values="price", aggfunc="mean")
        fig3 = go.Figure(go.Heatmap(
            z=pivot.values,
            x=[str(c) for c in pivot.columns],
            y=pivot.index.tolist(),
            colorscale="RdYlGn_r",
            colorbar=dict(title="£/MWh"),
        ))
        fig3.update_layout(
            height=400,
            xaxis_title="Delivery date",
            yaxis_title="Settlement period",
            margin=dict(l=40, r=0, t=0, b=40),
        )
        st.plotly_chart(fig3, use_container_width=True)

        # Daily summary table
        daily_epex = epex_df.groupby("delivery_date").agg(
            Baseload=("daily_baseload", "first"),
            Peak=("daily_peakload", "first"),
            Offpeak=("daily_offpeak", "first"),
        ).round(2).reset_index()
        daily_epex.columns = ["Date", "Baseload (£/MWh)", "Peak (£/MWh)", "Offpeak (£/MWh)"]
        st.dataframe(daily_epex.tail(14).sort_values("Date", ascending=False),
                     use_container_width=True, hide_index=True)
        st.caption(
            "**EPEX product definitions** — "
            "**Baseload**: simple 24-hour average (all 48 SPs). "
            "**Peak**: Mon–Fri 08:00–20:00 only (SP 17–40); no peak product on weekends or public holidays. "
            "**Offpeak**: Mon–Fri 00:00–08:00 and 20:00–24:00 (SP 1–16, 41–48); "
            "weekends and holidays are treated as offpeak-only days. "
            "⚠️ In the current GB market, peak prices are frequently *lower* than offpeak — "
            "midday solar (duck curve effect) suppresses SP 20–36 prices while the evening gas ramp "
            "(SP 33–38, 16:30–19:00) and overnight gas baseload (SP 1–8) keep offpeak prices elevated."
        )

    # ---- Fuel Mix, Bidding Space & EPEX DA (shared date selector) ----------
    st.markdown("---")
    _date_opts = [(date.today() - timedelta(days=i)).isoformat() for i in range(1, 61)]
    _default_sel = _date_opts[:7]
    _sel_dates_raw = st.multiselect(
        "Select dates — Fuel Mix / Bidding Space / EPEX DA",
        options=_date_opts,
        default=_default_sel,
        key="shared_fm_dates",
    )
    _sel_dates = tuple(sorted(_sel_dates_raw, reverse=True))  # newest first; hashable for cache

    _FUEL_COLORS = {
        "gas":     "#f4a460",
        "coal":    "#444444",
        "nuclear": "#7b2d8b",
        "wind":    "#2ca02c",
        "solar":   "#ffdd44",
        "hydro":   "#1f77b4",
        "imports": "#17becf",
        "biomass": "#8c564b",
        "storage": "#e377c2",
    }

    st.subheader("GB Generation Fuel Mix (Daily Average MW)")
    if not _sel_dates:
        st.info("Select at least one date above.")
    else:
        fm_daily = _get_fuel_mix_daily_dates(_sel_dates)
        if fm_daily.empty:
            st.info("No fuel mix data for the selected dates. Run a backfill in Data Management.")
        else:
            fig_fm = go.Figure()
            for fuel in ["coal", "gas", "nuclear", "biomass", "hydro", "imports", "storage", "solar", "wind"]:
                if fuel in fm_daily.columns:
                    fig_fm.add_trace(go.Bar(
                        x=fm_daily["settlement_date"],
                        y=fm_daily[fuel],
                        name=fuel.capitalize(),
                        marker_color=_FUEL_COLORS.get(fuel, "#aaa"),
                    ))
            fig_fm.update_layout(
                barmode="stack",
                height=320,
                margin=dict(l=0, r=0, t=0, b=0),
                xaxis_title="",
                yaxis_title="Avg MW",
                legend=dict(orientation="h", yanchor="bottom", y=1.01, x=0),
            )
            st.plotly_chart(fig_fm, use_container_width=True)

    st.subheader("Bidding Space — Half-Hourly (MW, Thermal + Flexible)")
    if not _sel_dates:
        st.info("Select at least one date above.")
    else:
        bs_df = _get_bidding_space_hh_dates(_sel_dates)
        if bs_df.empty:
            st.info("No bidding space data for the selected dates. Awaiting fuel mix ingestion.")
        else:
            fig_bs = go.Figure()
            for d_val, grp in bs_df.groupby("settlement_date"):
                fig_bs.add_trace(go.Scatter(
                    x=grp["settlement_period"],
                    y=grp["bidding_space_mw"],
                    mode="lines",
                    name=str(d_val),
                    line=dict(width=1.5),
                ))
            fig_bs.update_layout(
                height=280,
                margin=dict(l=0, r=0, t=0, b=0),
                xaxis=dict(title="Settlement Period (1=00:00, 48=23:30)", tickmode="linear", dtick=4),
                yaxis_title="MW",
                legend=dict(orientation="h", yanchor="bottom", y=1.01, x=0),
            )
            st.plotly_chart(fig_bs, use_container_width=True)

    st.subheader("EPEX DA Price — Half-Hourly (£/MWh)")
    st.caption("Each line = one delivery date. Same dates as above.")
    if not _sel_dates:
        st.info("Select at least one date above.")
    else:
        _epex_hh_df = _get_epex_hh_dates(_sel_dates)
        if _epex_hh_df.empty:
            st.info("No EPEX DA data for the selected dates.")
        else:
            fig_bs_epex = go.Figure()
            for d_val, grp in _epex_hh_df.groupby("delivery_date"):
                fig_bs_epex.add_trace(go.Scatter(
                    x=grp["settlement_period"],
                    y=grp["price"],
                    mode="lines",
                    name=str(d_val),
                    line=dict(width=1.5),
                ))
            fig_bs_epex.update_layout(
                height=280,
                margin=dict(l=0, r=0, t=0, b=0),
                xaxis=dict(title="Settlement Period (1=00:00, 48=23:30)", tickmode="linear", dtick=4),
                yaxis_title="£/MWh",
                legend=dict(orientation="h", yanchor="bottom", y=1.01, x=0),
            )
            st.plotly_chart(fig_bs_epex, use_container_width=True)

    st.subheader("Bidding Space vs EPEX DA Price")
    st.caption("Each point = one settlement period across selected dates. Colour = date.")
    if not _sel_dates:
        st.info("Select at least one date above.")
    else:
        _bs_sc = _get_bidding_space_hh_dates(_sel_dates)
        _ep_sc = _get_epex_hh_dates(_sel_dates)
        if _bs_sc.empty or _ep_sc.empty:
            st.info("No overlapping bidding space / EPEX DA data for the selected dates.")
        else:
            _ep_sc = _ep_sc.rename(columns={"delivery_date": "settlement_date"})
            _ep_sc["settlement_date"] = pd.to_datetime(_ep_sc["settlement_date"]).dt.date
            _bs_sc["settlement_date"] = pd.to_datetime(_bs_sc["settlement_date"]).dt.date
            _scatter_df = _bs_sc.merge(_ep_sc, on=["settlement_date", "settlement_period"])
            if _scatter_df.empty:
                st.info("No overlapping data between bidding space and EPEX DA for these dates.")
            else:
                _scatter_df["date_str"] = _scatter_df["settlement_date"].astype(str)
                fig_scat = px.scatter(
                    _scatter_df,
                    x="bidding_space_mw",
                    y="price",
                    color="date_str",
                    opacity=0.65,
                    trendline="ols",
                    labels={
                        "bidding_space_mw": "Bidding Space (MW)",
                        "price": "EPEX DA (£/MWh)",
                        "date_str": "Date",
                    },
                )
                fig_scat.update_layout(
                    height=340,
                    margin=dict(l=0, r=0, t=10, b=0),
                    xaxis_title="Bidding Space (MW)",
                    yaxis_title="EPEX DA Price (£/MWh)",
                    legend=dict(orientation="h", yanchor="bottom", y=1.02, x=0),
                )
                st.plotly_chart(fig_scat, use_container_width=True)

# ---- Ancillary Markets -----------------------------------------------------
with tab_ancillary:
    st.header("Ancillary Markets — Dynamic Services (DX)")

    dx_df = _get_dx_range(date_start, date_end)
    if dx_df.empty:
        st.info("No DX results data in range.")
    else:
        services_available = sorted(dx_df["service"].unique().tolist())
        selected_services = st.multiselect(
            "Services", services_available,
            default=services_available[:6] if len(services_available) >= 6 else services_available,
        )
        if selected_services:
            filtered = dx_df[dx_df["service"].isin(selected_services)]

            col1, col2 = st.columns(2)
            with col1:
                st.subheader("Clearing Price (£/MW/h)")
                fig = px.line(
                    filtered.groupby(["efa_date", "service"])["clearing_price"].mean().reset_index(),
                    x="efa_date", y="clearing_price", color="service",
                    labels={"clearing_price": "£/MW/h", "efa_date": ""},
                )
                fig.update_layout(height=300, margin=dict(l=0, r=0, t=0, b=0))
                st.plotly_chart(fig, use_container_width=True)

            with col2:
                st.subheader("Cleared Volume (MW)")
                fig2 = px.line(
                    filtered.groupby(["efa_date", "service"])["cleared_volume"].sum().reset_index(),
                    x="efa_date", y="cleared_volume", color="service",
                    labels={"cleared_volume": "MW", "efa_date": ""},
                )
                fig2.update_layout(height=300, margin=dict(l=0, r=0, t=0, b=0))
                st.plotly_chart(fig2, use_container_width=True)

            st.subheader("Summary Statistics by Service")
            summary = filtered.groupby("service").agg(
                avg_price=("clearing_price", "mean"),
                max_price=("clearing_price", "max"),
                min_price=("clearing_price", "min"),
                avg_volume=("cleared_volume", "mean"),
            ).round(2).reset_index()
            summary.columns = ["Service", "Avg Price (£/MW/h)", "Max", "Min", "Avg Volume (MW)"]
            st.dataframe(summary, use_container_width=True, hide_index=True)

# ---- BESS Benchmarking -----------------------------------------------------
with tab_bess:
    st.header("BESS Benchmarking Index")

    # Tab-level date range (independent of sidebar)
    bc1, bc2 = st.columns(2)
    bess_from = bc1.date_input("From", value=date.today() - timedelta(days=90), key="bess_from")
    bess_to   = bc2.date_input("To",   value=date.today() - timedelta(days=1),  key="bess_to")
    _bess_start = bess_from.isoformat()
    _bess_end   = bess_to.isoformat()

    col1, col2 = st.columns([2, 1])

    with col1:
        st.subheader("Daily Revenue Index (£/MW/day) by Market")
        daily_idx = _get_daily_index_range(_bess_start, _bess_end)
        if daily_idx.empty:
            st.info("No daily index data.")
        else:
            _daily_components = daily_idx[daily_idx["market"] != "total"]
            _daily_total = daily_idx[daily_idx["market"] == "total"].sort_values("settlement_date")
            fig = px.bar(
                _daily_components,
                x="settlement_date", y="revenue_permw", color="market",
                labels={"revenue_permw": "£/MW/day", "settlement_date": ""},
                barmode="relative",
            )
            if not _daily_total.empty:
                fig.add_scatter(
                    x=_daily_total["settlement_date"], y=_daily_total["revenue_permw"],
                    mode="lines", name="total",
                    line=dict(color="black", width=2, dash="dash"),
                )
            fig.update_layout(height=350, margin=dict(l=0, r=0, t=0, b=0))
            st.plotly_chart(fig, use_container_width=True)

    with col2:
        st.subheader("Market Avg")
        if not daily_idx.empty:
            # revenue_permw / revenue_permwh in gb_bess_daily_index are already per-day
            # values (daily aggregates from Modo), NOT per settlement-period — so no ×48.
            mkt_avg = (
                daily_idx[daily_idx["market"] != "total"]
                .groupby("market")[["revenue_permw", "revenue_permwh"]]
                .mean()
                .reset_index()
            )
            mkt_avg["avg_permw_day"]   = mkt_avg["revenue_permw"].round(2)
            mkt_avg["avg_permwh_year"] = (mkt_avg["revenue_permwh"] * 365).round(0)
            mkt_avg = mkt_avg[["market", "avg_permw_day", "avg_permwh_year"]].sort_values("avg_permw_day", ascending=False)
            # Total row — use the pre-aggregated 'total' market rows from Modo
            _tot = daily_idx[daily_idx["market"] == "total"][["revenue_permw", "revenue_permwh"]].mean()
            if not _tot.isna().all():
                total_row = pd.DataFrame([{
                    "market": "Total",
                    "avg_permw_day":   round(float(_tot["revenue_permw"]), 2),
                    "avg_permwh_year": round(float(_tot["revenue_permwh"]) * 365, 0),
                }])
                mkt_avg = pd.concat([mkt_avg, total_row], ignore_index=True)
            mkt_avg.columns = ["Market", "Avg £/MW/day", "Avg £/MWh/year"]
            st.dataframe(mkt_avg, use_container_width=True, hide_index=True)

    st.subheader("Monthly Revenue Index (£/MW/month)")
    _month_start = bess_from.replace(day=1).isoformat()
    monthly_idx = _get_monthly_index_range(_month_start, _bess_end)
    if monthly_idx.empty:
        st.info("No monthly index data.")
    else:
        fig2 = px.bar(
            monthly_idx[monthly_idx["market"] != "total"],
            x="month", y="revenue_permw", color="market",
            labels={"revenue_permw": "£/MW/month", "month": ""},
            barmode="stack",
        )
        fig2.update_layout(height=320, margin=dict(l=0, r=0, t=0, b=0))
        st.plotly_chart(fig2, use_container_width=True)

    st.subheader(f"Asset Leaderboard — Top 20, Sorted by £/MWh/year")
    leader_df = _get_leaderboard_range(_bess_start, _bess_end)
    if leader_df.empty:
        st.info("No leaderboard data.")
    else:
        for col in ["total_revenue", "wholesale", "freq_response", "bm", "imbalance", "reserve",
                    "avg_revspermw", "avg_revspermwh", "rated_power_mw", "energy_capacity_mwh"]:
            if col in leader_df.columns:
                leader_df[col] = pd.to_numeric(leader_df[col], errors="coerce")
        for col in ["total_revenue", "wholesale", "freq_response", "bm", "imbalance", "reserve"]:
            if col in leader_df.columns:
                leader_df[col] = leader_df[col].round(0)
        leader_df["avg_revspermw"] = (leader_df["avg_revspermw"] * 48).round(2)
        leader_df["avg_revspermwh"] = (leader_df["avg_revspermwh"] * 48 * 365).round(0)
        leader_df["rated_power_mw"] = leader_df["rated_power_mw"].round(1)
        # Compute duration label from energy_capacity / rated_power
        def _duration_label(row):
            e = row.get("energy_capacity_mwh")
            p = row.get("rated_power_mw")
            if pd.isna(e) or pd.isna(p) or p == 0:
                return None
            h = e / p
            # Snap to well-known C-rate buckets within ±15%
            for hours, label in [(0.5, "0.5h (2C)"), (1.0, "1h (1C)"), (1.5, "1.5h"),
                                  (2.0, "2h (0.5C)"), (4.0, "4h (0.25C)")]:
                if abs(h - hours) / hours < 0.15:
                    return label
            return f"{h:.1f}h"
        leader_df["duration"] = leader_df.apply(_duration_label, axis=1)
        # Sort by £/MWh/year (capacity-normalised, descending)
        leader_df = leader_df.sort_values("avg_revspermwh", ascending=False, na_position="last").reset_index(drop=True)
        # Reorder: asset identity first, then financials
        cols_order = ["asset", "owner", "operator", "integrator", "duration",
                      "total_revenue", "wholesale", "freq_response", "bm", "imbalance", "reserve",
                      "avg_revspermw", "avg_revspermwh", "rated_power_mw"]
        leader_df = leader_df[[c for c in cols_order if c in leader_df.columns]]
        leader_df.columns = [
            {"asset": "Asset", "owner": "Owner", "operator": "Operator",
             "integrator": "Integrator", "duration": "Duration",
             "total_revenue": "Total (£)", "wholesale": "Wholesale (£)",
             "freq_response": "Freq Response (£)", "bm": "BM (£)",
             "imbalance": "Imbalance (£)", "reserve": "Reserve (£)",
             "avg_revspermw": "£/MW/day", "avg_revspermwh": "£/MWh/year",
             "rated_power_mw": "Rated Power (MW)"}.get(c, c)
            for c in leader_df.columns
        ]
        st.dataframe(leader_df, use_container_width=True, hide_index=True)

# ---- Pricing ---------------------------------------------------------------
with tab_pricing:
    st.header("Pricing Models — GB BESS")
    st.caption(
        "Options: Kirk/Margrabe spread call strip (365-day horizon) · "
        "PF Actual: perfect-foresight dispatch on EPEX DA prices · "
        "PF Forecast: PF on OLS+fundamentals price forecast · "
        "Batch runs nightly at 04:30 SGT"
    )

    pr_col1, pr_col2 = st.columns(2)
    pricing_from = pr_col1.date_input("From", value=date.today() - timedelta(days=30), key="pr_from")
    pricing_to   = pr_col2.date_input("To",   value=date.today() - timedelta(days=1),  key="pr_to")
    _pr_start = pricing_from.isoformat()
    _pr_end   = pricing_to.isoformat()

    # ── Section 1: Top 20 Pricing Table ──────────────────────────────────────
    st.subheader("Top 20 Wholesale Revenue — Pricing Models")
    pr_df = _get_pricing_table(_pr_start, _pr_end)

    if pr_df.empty:
        st.info(
            "No pricing data available. "
            "Either the nightly batch (04:30 SGT) hasn't run yet or no EPEX DA prices exist "
            "for this date range. You can trigger a manual run in Data Management."
        )
    else:
        for col in ["wholesale_gbp", "options_val_per_mw", "pf_actual", "avg_charged_mwh",
                    "power_mw", "duration_h"]:
            if col in pr_df.columns:
                pr_df[col] = pd.to_numeric(pr_df[col], errors="coerce")

        # Normalise to annualised £/MWh of installed energy capacity
        n_days = max(1, (pricing_to - pricing_from).days + 1)
        e_cap = (pr_df["power_mw"] * pr_df["duration_h"]).replace(0, np.nan)
        pr_df["wholesale_per_mwh"]    = pr_df["wholesale_gbp"] / e_cap / n_days * 365
        pr_df["options_per_mwh"]      = pr_df["options_val_per_mw"] / pr_df["duration_h"]
        pr_df["pf_actual_per_mwh"]    = pr_df["pf_actual"] / e_cap * 365
        # Derived PF actual DA metrics
        _charged = pr_df["avg_charged_mwh"].replace(0, np.nan)
        pr_df["price_spread_gbp_mwh"] = pr_df["pf_actual"] / _charged
        pr_df["cycles_per_day"]       = _charged / e_cap

        def _dur_label(h):
            if pd.isna(h):
                return "—"
            for hrs, lbl in [(0.5, "0.5h"), (1.0, "1h"), (2.0, "2h"), (4.0, "4h")]:
                if abs(h - hrs) / hrs < 0.15:
                    return lbl
            return f"{h:.1f}h"

        pr_df["Duration"] = pr_df["duration_h"].apply(_dur_label)

        disp_cols = {
            "rank":                 "#",
            "asset":                "Asset",
            "owner":                "Owner",
            "operator":             "Operator",
            "Duration":             "Duration",
            "power_mw":             "Rated MW",
            "wholesale_per_mwh":    "Wholesale (£/MWh/yr)",
            "options_per_mwh":      "Options (£/MWh/yr)",
            "pf_actual_per_mwh":    "PF SP (£/MWh/yr)",
            "price_spread_gbp_mwh": "PF Price Spread (£/MWh)",
            "cycles_per_day":       "Cycles/Day",
        }
        show_cols = [c for c in disp_cols if c in pr_df.columns or c == "Duration"]
        pr_show = pr_df[[c for c in show_cols if c in pr_df.columns]].rename(columns=disp_cols)

        fmts = {
            "Rated MW":                 "{:.0f}",
            "Wholesale (£/MWh/yr)":     "{:,.0f}",
            "Options (£/MWh/yr)":       "{:,.0f}",
            "PF SP (£/MWh/yr)":         "{:,.0f}",
            "PF Price Spread (£/MWh)":  "{:.1f}",
            "Cycles/Day":               "{:.2f}",
        }
        # Warn if PF values look stale (zeros from pre-v32 batch runs)
        _pf_col = pr_df["pf_actual"].dropna()
        if len(_pf_col) > 0 and (_pf_col.abs() < 1.0).all():
            st.warning(
                "PF Actual DA values are near zero — the pricing batch rows in this date range "
                "were likely computed before the v32 EPEX DA fix. "
                "Go to **Data Management → Pricing Batch** and re-run for this date range to refresh."
            )
        st.caption(
            f"Revenue annualised per MWh installed capacity ({n_days}-day window × 365/n). "
            "Options = Kirk/Margrabe scaled by model Cycles/Day (fixes 1-cycle undercount). "
            "PF SP = perfect-foresight dispatch on half-hourly system prices (intraday proxy). "
            "Price Spread = PF P&L ÷ avg daily charged MWh. "
            "Cycles/Day = avg daily charged MWh ÷ energy capacity."
        )
        st.dataframe(pr_show.style.format(fmts, na_rep="—"), use_container_width=True, hide_index=True)

    # ── Section 2: PF Dispatch Comparison Chart ──────────────────────────────
    st.subheader("PF Dispatch Comparison")
    if not pr_df.empty and "asset" in pr_df.columns:
        asset_list = pr_df["asset"].dropna().tolist()
        sel_asset = st.selectbox("Select asset", asset_list, key="pr_asset_sel")

        # Date picker for dispatch chart
        dispatch_dates = pd.date_range(end=pricing_to, periods=14).strftime("%Y-%m-%d").tolist()
        sel_disp_date = st.selectbox("Date", dispatch_dates[::-1], key="pr_disp_date")

        disp_df = _get_dispatch_comparison(sel_asset, sel_disp_date)
        if disp_df.empty:
            st.info(f"No pricing batch results for {sel_asset} on {sel_disp_date}.")
        else:
            import json as _json
            row = disp_df.iloc[0]

            def _parse_json_col(val):
                if val is None:
                    return []
                try:
                    if isinstance(val, str):
                        return _json.loads(val)
                    return list(val)
                except Exception:
                    return []

            pf_actual_disp   = _parse_json_col(row.get("pf_actual_dispatch_48"))
            pf_sp_disp       = _parse_json_col(row.get("pf_sp_dispatch_48"))
            pf_forecast_disp = _parse_json_col(row.get("pf_forecast_dispatch_48"))
            actual_epex      = _parse_json_col(row.get("actual_epex_48"))
            actual_sp        = _parse_json_col(row.get("actual_sp_48"))
            forecast_epex    = _parse_json_col(row.get("forecast_epex_48"))
            sps = list(range(1, 49))

            _sp_nonzero   = any(v != 0 for v in actual_sp)   if actual_sp   else False
            _epex_nonzero = any(v != 0 for v in actual_epex) if actual_epex else False
            if actual_epex and not _epex_nonzero and not _sp_nonzero:
                st.warning(
                    "Prices for this date are all zero — re-run the pricing batch for this "
                    "date range in Data Management to refresh."
                )

            fig_disp = go.Figure()
            # Dispatch traces — primary Y-axis (left, MW)
            if pf_sp_disp:
                fig_disp.add_trace(go.Scatter(
                    x=sps, y=pf_sp_disp, mode="lines",
                    name="PF Dispatch (System Price)",
                    line=dict(color="#1f77b4", width=2),
                    yaxis="y1",
                ))
            if pf_actual_disp:
                fig_disp.add_trace(go.Scatter(
                    x=sps, y=pf_actual_disp, mode="lines",
                    name="PF Dispatch (EPEX DA ref)",
                    line=dict(color="#aec7e8", width=1.5, dash="dash"),
                    yaxis="y1",
                ))
            if pf_forecast_disp:
                fig_disp.add_trace(go.Scatter(
                    x=sps, y=pf_forecast_disp, mode="lines",
                    name="PF Dispatch (OLS Forecast)",
                    line=dict(color="#ff7f0e", width=1.5, dash="dash"),
                    yaxis="y1",
                ))
            # Price traces — secondary Y-axis (right, £/MWh)
            if actual_sp and _sp_nonzero:
                fig_disp.add_trace(go.Scatter(
                    x=sps, y=actual_sp, mode="lines",
                    name="System Price",
                    line=dict(color="#d62728", width=2),
                    yaxis="y2",
                ))
            if actual_epex and _epex_nonzero:
                fig_disp.add_trace(go.Scatter(
                    x=sps, y=actual_epex, mode="lines",
                    name="EPEX DA Price",
                    line=dict(color="#2ca02c", width=1.5, dash="dot"),
                    yaxis="y2",
                ))
            if forecast_epex:
                _fc_nonzero = any(v != 0 for v in forecast_epex)
                if _fc_nonzero:
                    fig_disp.add_trace(go.Scatter(
                        x=sps, y=forecast_epex, mode="lines",
                        name="OLS Forecast Price",
                        line=dict(color="#9467bd", width=1.5, dash="dot"),
                        yaxis="y2",
                    ))
            fig_disp.add_hline(y=0, line_dash="dot", line_color="grey", line_width=1)
            fig_disp.update_layout(
                height=350,
                margin=dict(l=0, r=60, t=10, b=0),
                xaxis=dict(title="Settlement Period", tickmode="linear", dtick=4),
                yaxis=dict(title="MW (+ discharge, − charge)"),
                yaxis2=dict(
                    title="£/MWh",
                    overlaying="y",
                    side="right",
                    showgrid=False,
                ),
                legend=dict(orientation="h", yanchor="bottom", y=1.02, x=0),
            )
            st.plotly_chart(fig_disp, use_container_width=True)
            st.caption(
                "Dispatch (left axis): model-computed PF optimal dispatch — not physical. "
                "Price (right axis, dotted): EPEX DA realised price and OLS forecast. "
                "Dispatch shape follows the price spread — charge at low-price SPs, discharge at high."
            )

            # Price comparison table
            if actual_epex or forecast_epex:
                price_df = pd.DataFrame({
                    "SP": sps,
                    "EPEX DA (£/MWh)":      [round(v, 2) for v in actual_epex]   if actual_epex   else ["—"] * 48,
                    "OLS Forecast (£/MWh)": [round(v, 2) for v in forecast_epex] if forecast_epex else ["—"] * 48,
                })
                with st.expander("Half-hourly price table", expanded=False):
                    st.dataframe(price_df, use_container_width=True, hide_index=True)
    else:
        st.info("Select a date range above to load pricing data first.")

    st.caption("SBP/SSP and EPEX DA price charts are available in the **Market Overview** tab.")


# ---- Asset Map -------------------------------------------------------------
_REV_COLOR_SCALE = [[0, "#d73027"], [0.5, "#fee08b"], [1, "#1a9850"]]  # red → yellow → green
_MAP_MARKETS = ["total", "wholesale", "frequency_response", "bm", "imbalance", "reserve"]

with tab_map:
    st.header("GB BESS Asset Map")

    assets_df = _get_assets()
    if assets_df.empty:
        st.info("No asset data. Run a backfill.")
    else:
        # KPIs
        col1, col2, col3 = st.columns(3)
        col1.metric("Total Assets", len(assets_df))
        col2.metric("Total Capacity", f"{assets_df['rated_power_mw'].sum():.0f} MW")
        col3.metric("Transmission-connected",
                    int(assets_df["transmission_connected"].sum()) if "transmission_connected" in assets_df.columns else 0)

        # Revenue colour controls
        st.markdown("---")
        mc1, mc2, mc3, mc4 = st.columns([2, 2, 2, 2])
        map_from = mc1.date_input("From", value=date.today() - timedelta(days=90), key="map_from")
        map_to   = mc2.date_input("To",   value=date.today() - timedelta(days=1),  key="map_to")
        map_mkt  = mc3.selectbox("Colour by revenue source", _MAP_MARKETS, key="map_mkt")
        map_cat  = mc4.radio("Identify assets by", ["Owner", "Operator", "Integrator"], horizontal=True, key="map_cat")
        _cat_col   = {"Owner": "owner", "Operator": "operator", "Integrator": "integrator"}[map_cat]
        _cat_label = map_cat

        rev_df = _get_asset_revenue_map(map_from.isoformat(), map_to.isoformat(), map_mkt)
        map_df = assets_df.dropna(subset=["latitude", "longitude", "rated_power_mw"])

        if not rev_df.empty:
            map_df = map_df.merge(rev_df[["asset", "rev_per_mw"]], on="asset", how="left")
        else:
            map_df["rev_per_mw"] = float("nan")

        has_rev = "rev_per_mw" in map_df.columns and map_df["rev_per_mw"].notna().any()

        col_map, col_rank = st.columns([3, 1])

        with col_map:
            _hover_common = {
                "rated_power_mw": ":.0f", "gsp": True,
                "owner": True, "operator": True, "integrator": True,
                "commissioning_date": True,
                "latitude": False, "longitude": False,
            }
            if has_rev:
                fig_map = px.scatter_mapbox(
                    map_df,
                    lat="latitude", lon="longitude",
                    size="rated_power_mw",
                    color="rev_per_mw",
                    color_continuous_scale=_REV_COLOR_SCALE,
                    hover_name="asset",
                    hover_data={**_hover_common, "rev_per_mw": ":.2f"},
                    labels={"rev_per_mw": f"£/MW ({map_mkt})", "rated_power_mw": "MW",
                            "owner": "Owner", "operator": "Operator", "integrator": "Integrator"},
                    zoom=5, center={"lat": 53.5, "lon": -1.5},
                    mapbox_style="open-street-map", height=580,
                )
            else:
                fig_map = px.scatter_mapbox(
                    map_df,
                    lat="latitude", lon="longitude",
                    size="rated_power_mw",
                    color=_cat_col,
                    hover_name="asset",
                    hover_data=_hover_common,
                    labels={"owner": "Owner", "operator": "Operator", "integrator": "Integrator",
                            "rated_power_mw": "MW"},
                    zoom=5, center={"lat": 53.5, "lon": -1.5},
                    mapbox_style="open-street-map", height=580,
                )
            fig_map.update_layout(margin=dict(l=0, r=0, t=0, b=0))
            st.plotly_chart(fig_map, use_container_width=True)

        with col_rank:
            if has_rev:
                rank_df = (
                    map_df[["asset", "rev_per_mw"]].dropna()
                    .sort_values("rev_per_mw", ascending=True)
                    .tail(30)
                )
                fig_rank = px.bar(
                    rank_df,
                    x="rev_per_mw", y="asset",
                    orientation="h",
                    color="rev_per_mw",
                    color_continuous_scale=_REV_COLOR_SCALE,
                    labels={"rev_per_mw": f"£/MW", "asset": ""},
                    title=f"Top assets — {map_mkt}",
                )
                fig_rank.update_layout(
                    coloraxis_showscale=False,
                    margin=dict(l=0, r=10, t=30, b=0),
                    height=580,
                    yaxis=dict(tickfont=dict(size=9)),
                )
                st.plotly_chart(fig_rank, use_container_width=True)
            else:
                # Show total installed MW grouped by the selected category
                cat_df = (
                    map_df[["rated_power_mw", _cat_col]].dropna(subset=[_cat_col])
                    .groupby(_cat_col, as_index=False)["rated_power_mw"].sum()
                    .sort_values("rated_power_mw", ascending=True)
                )
                if not cat_df.empty:
                    fig_cat = px.bar(
                        cat_df,
                        x="rated_power_mw", y=_cat_col,
                        orientation="h",
                        labels={"rated_power_mw": "MW", _cat_col: ""},
                        title=f"Capacity by {map_cat}",
                    )
                    fig_cat.update_layout(
                        margin=dict(l=0, r=10, t=30, b=0),
                        height=580,
                        yaxis=dict(tickfont=dict(size=9)),
                    )
                    st.plotly_chart(fig_cat, use_container_width=True)
                else:
                    st.info(f"No {map_cat.lower()} data available.")

        # Table
        with st.expander("Asset details"):
            detail_cols = ["asset", "rated_power_mw", "owner", "operator", "integrator",
                           "manufacturer", "gsp", "commissioning_date", "dno"]
            show_df = assets_df[[c for c in detail_cols if c in assets_df.columns]].copy()
            show_df.rename(columns={
                "asset": "Asset", "rated_power_mw": "Power (MW)",
                "owner": "Owner", "operator": "Operator", "integrator": "Integrator",
                "manufacturer": "Manufacturer", "gsp": "GSP",
                "commissioning_date": "Commissioned", "dno": "DNO",
            }, inplace=True)
            st.dataframe(show_df, use_container_width=True, hide_index=True)

# ---- Knowledge Base --------------------------------------------------------
with tab_knowledge:
    st.header("GB Market Knowledge Base")
    st.info(
        "**Auto-updated daily** · Market data ingested at **03:00 SGT** · "
        "Knowledge base (articles & reports) ingested at **03:30 SGT**",
        icon="🔄",
    )

    # Coverage table + ingest controls
    kc1, kc2 = st.columns([2, 1])
    with kc1:
        kb_counts = _knowledge_doc_counts()
        if kb_counts.empty:
            st.info("Knowledge base is empty. Run ingestion from Data Management tab.")
        else:
            st.dataframe(kb_counts, use_container_width=True, hide_index=True)

    with kc2:
        if st.button("Refresh knowledge stats"):
            _knowledge_doc_counts.clear()
            st.rerun()
        if st.button("Run knowledge ingest now", type="primary"):
            with st.spinner("Fetching from all knowledge sources…"):
                result = _run_knowledge_ingest_job(trigger="manual")
            if result["status"] == "success":
                st.success(f"Done — {result['total']} new docs in {result['duration']:.1f}s")
                if result.get("results"):
                    st.json(result["results"])
            else:
                st.error(f"Failed: {result['error']}")
            _knowledge_doc_counts.clear()
            st.rerun()

    # ---- Source topic guide ------------------------------------------------
    st.divider()
    st.subheader("What's in the Knowledge Base")
    st.caption("Click a topic to pre-fill the search below.")

    guide_cols = st.columns(len(_KB_SOURCE_GUIDE))
    for col, src in zip(guide_cols, _KB_SOURCE_GUIDE):
        with col:
            st.markdown(f"**{src['label']}**")
            st.caption(src["description"])
            for topic in src["topics"]:
                if st.button(topic, key=f"kb_topic_{src['source']}_{topic}"):
                    st.session_state["kb_query"] = topic
                    # Don't auto-filter by source — search all sources so
                    # cross-source hits aren't silently dropped
                    st.session_state["kb_source_filter"] = []
                    st.rerun()

    # ---- Search ------------------------------------------------------------
    st.divider()
    st.subheader("Search Knowledge Base")
    kb_query = st.text_input(
        "Search query",
        placeholder="e.g. BESS frequency response market trends 2025",
        key="kb_query",
    )
    _KB_SOURCE_DISPLAY = {
        "elexon": "Elexon",
        "entso_e": "National Grid ESO",
        "timera": "Timera Energy",
        "modo": "Modo Energy",
        "meteologica": "Meteologica",
        "upload": "Uploaded Documents",
    }
    kb_sources = st.multiselect(
        "Filter by source (all if empty)",
        options=list(_KB_SOURCE_DISPLAY.keys()),
        format_func=lambda x: _KB_SOURCE_DISPLAY.get(x, x),
        key="kb_source_filter",
    )
    if kb_query:
        _search_knowledge.clear()
        results_df = _search_knowledge(kb_query, sources=kb_sources or None, limit=10)
        if results_df.empty:
            st.info("No results found.")
        else:
            for _, row in results_df.iterrows():
                with st.expander(
                    f"[{row['source']}] {row['title'] or 'Untitled'}  —  {row['published_date']}",
                    expanded=False,
                ):
                    if row.get("url"):
                        st.markdown(f"[View source]({row['url']})")
                    st.text(row["snippet"])


# ---- Strategist Agent ------------------------------------------------------
with tab_strategist:
    st.header("Strategist — GB Market Analysis")

    # Show accumulated insight pool count
    _insight_pool_df = _query(
        "SELECT COUNT(*) AS n FROM intl_market.gb_expert_insights WHERE active = TRUE"
    )
    _n_insights = int(_insight_pool_df.iloc[0]["n"]) if not _insight_pool_df.empty else 0
    st.caption(
        f"Grounded on DB data only · Memory persists across sessions · "
        f"Expert memory: {_n_insights} accumulated insights from KB + conversations"
    )

    # ── Knowledge Gap Interview ──────────────────────────────────────────────
    with st.expander("Teach the Agent — Knowledge Gap Interview", expanded=False):
        # Initialise all interview state keys
        for _k, _v in [
            ("interview_questions", []),
            ("interview_idx", 0),
            ("interview_answers", 0),
            ("interview_modo_queried", False),
            ("interview_modo_results", {}),   # {question: answer_or_None}
            ("interview_pending_qs", []),     # questions Modo couldn't answer
        ]:
            if _k not in st.session_state:
                st.session_state[_k] = _v

        _iq  = st.session_state["interview_questions"]
        _ii  = st.session_state["interview_idx"]
        _pqs = st.session_state["interview_pending_qs"]   # user Q&A queue

        # ── Stage 0: no questions yet ──────────────────────────────────────
        if not _iq:
            st.markdown(
                "The agent audits its knowledge base, identifies gaps, then tries "
                "**Modo AI first** to fill them — saving your time for questions "
                "only you can answer from direct experience."
            )
            if st.button("Generate Knowledge Gap Questions", key="gen_interview"):
                with st.spinner("Auditing knowledge base and identifying gaps…"):
                    _new_qs = _generate_interview_questions()
                if _new_qs:
                    st.session_state["interview_questions"]    = _new_qs
                    st.session_state["interview_idx"]          = 0
                    st.session_state["interview_answers"]      = 0
                    st.session_state["interview_modo_queried"] = False
                    st.session_state["interview_modo_results"] = {}
                    st.session_state["interview_pending_qs"]   = []
                    st.rerun()
                else:
                    st.error("Could not generate questions — check API key or try again.")

        # ── Stage 1: questions ready, Modo not yet queried ─────────────────
        elif not st.session_state["interview_modo_queried"]:
            st.markdown("**Generated knowledge gap questions:**")
            for _qi, _qo in enumerate(_iq):
                st.markdown(f"{_qi+1}. **[{_qo['topic']}]** {_qo['question']}")
                st.caption(f"   *{_qo.get('why_asking','')}*")

            st.divider()
            _col_m, _col_u = st.columns(2)
            with _col_m:
                if st.button(
                    "Query Modo AI First (recommended)",
                    key="interview_modo_query", type="primary",
                ):
                    with st.spinner(
                        "Querying Modo AI for each gap question… (~2–4 min, "
                        "uses Modo credits only for unanswered gaps)"
                    ):
                        try:
                            from services.gb_knowledge.modo_ai import distill_gap_questions
                            from services.gb_knowledge.expert_memory import digest_kb_docs
                            _qs_text = [_qo["question"] for _qo in _iq]
                            _mres = distill_gap_questions(_qs_text)
                            # Digest the new docs into expert insights
                            digest_kb_docs(_ANTHROPIC_KEY, limit=len(_iq) + 5)
                        except Exception as _me:
                            st.error(f"Modo query failed: {_me}")
                            _mres = {}
                    st.session_state["interview_modo_results"] = _mres
                    # Pending = questions Modo couldn't answer
                    st.session_state["interview_pending_qs"] = [
                        _qo for _qo in _iq if not _mres.get(_qo["question"])
                    ]
                    st.session_state["interview_modo_queried"] = True
                    st.session_state["interview_idx"] = 0
                    st.rerun()

            with _col_u:
                if st.button("Answer Yourself (skip Modo)", key="interview_skip_modo"):
                    st.session_state["interview_pending_qs"]   = list(_iq)
                    st.session_state["interview_modo_queried"] = True
                    st.session_state["interview_idx"]          = 0
                    st.rerun()

        # ── Stage 2: show Modo results + user Q&A for unanswered ──────────
        elif _ii >= len(_pqs):
            # All pending questions done — show summary
            _mres   = st.session_state["interview_modo_results"]
            _n_modo = sum(1 for v in _mres.values() if v)
            _n_user = st.session_state["interview_answers"]
            if _n_modo:
                st.success(
                    f"Modo AI answered **{_n_modo}** gap question(s) — insights auto-digested. "
                    f"You answered **{_n_user}** additional question(s). "
                    f"All stored as high-confidence insights."
                )
                with st.expander("View Modo AI answers", expanded=False):
                    for _qo in _iq:
                        _ans_text = _mres.get(_qo["question"])
                        if _ans_text:
                            st.markdown(f"**Q: {_qo['question']}**")
                            st.markdown(_ans_text[:600] + ("…" if len(_ans_text) > 600 else ""))
                            st.divider()
            else:
                st.success(
                    f"Interview complete — {_n_user} expert answers stored as high-confidence insights."
                )
            if st.button("Start New Interview", key="new_interview"):
                for _k2 in ["interview_questions", "interview_pending_qs",
                             "interview_modo_results"]:
                    st.session_state[_k2] = []  if _k2 != "interview_modo_results" else {}
                st.session_state["interview_idx"]          = 0
                st.session_state["interview_answers"]      = 0
                st.session_state["interview_modo_queried"] = False
                st.rerun()

        else:
            # User Q&A for questions Modo couldn't answer
            _q  = _pqs[_ii]
            _mres = st.session_state["interview_modo_results"]
            # Show Modo result for previously answered questions (context)
            if _ii == 0 and _mres:
                _n_auto = sum(1 for v in _mres.values() if v)
                if _n_auto:
                    st.info(
                        f"Modo AI answered {_n_auto} of {len(_iq)} questions automatically. "
                        f"Please answer the remaining {len(_pqs)}."
                    )
            st.progress(_ii / max(len(_pqs), 1), text=f"Question {_ii + 1} of {len(_pqs)}")
            st.markdown(f"**[{_q['topic']}]** {_q['question']}")
            st.caption(f"*Why this matters: {_q.get('why_asking', '')}*")
            _ans = st.text_area(
                "Your answer:", key=f"interview_ans_{_ii}", height=120,
                placeholder="Share what you know from experience…",
            )
            _col_submit, _col_skip = st.columns([2, 1])
            with _col_submit:
                if st.button("Submit Answer", key=f"interview_submit_{_ii}", type="primary"):
                    if _ans.strip():
                        try:
                            _store_interview_answer(_q["question"], _ans.strip(), _q["topic"])
                            st.session_state["interview_idx"]     += 1
                            st.session_state["interview_answers"] += 1
                            st.rerun()
                        except Exception as _e:
                            st.error(f"Failed to store answer: {_e}")
                    else:
                        st.warning("Please enter an answer before submitting.")
            with _col_skip:
                if st.button("Skip", key=f"interview_skip_{_ii}"):
                    st.session_state["interview_idx"] += 1
                    st.rerun()

    # Session ID management
    if "strat_session_id" not in st.session_state:
        st.session_state["strat_session_id"] = str(uuid.uuid4())
    if "strat_history" not in st.session_state:
        st.session_state["strat_history"] = []

    # Offer resume only when current chat is empty
    if not st.session_state["strat_history"]:
        _recent_sessions = _list_recent_sessions()
        if not _recent_sessions.empty:
            with st.expander("Resume a previous conversation?", expanded=False):
                for _, _srow in _recent_sessions.iterrows():
                    _sess_label = (
                        f"{_srow['updated_at'].strftime('%Y-%m-%d %H:%M')} — "
                        f"{int(_srow['msg_count'])} messages"
                    )
                    if st.button(_sess_label, key=f"resume_{_srow['session_id']}"):
                        st.session_state["strat_session_id"] = _srow["session_id"]
                        st.session_state["strat_history"] = _load_session(_srow["session_id"])
                        st.rerun()

    for msg in st.session_state["strat_history"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    user_input = st.chat_input("Ask about GB market fundamentals, system price, ancillary services…")
    if user_input:
        st.session_state["strat_history"].append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.markdown(user_input)

        with st.chat_message("assistant"):
            with st.spinner("Analysing…"):
                api_messages = [
                    {"role": m["role"], "content": m["content"]}
                    for m in st.session_state["strat_history"]
                ]
                try:
                    reply, updated, tool_events = _run_agent_turn(
                        api_messages, _build_strategist_system(user_input),
                        _STRATEGIST_TOOLS, _dispatch_strategist,
                    )
                except Exception as _agent_err:
                    reply = f"⚠️ API error: {_agent_err}. Please try again."
                    tool_events = []

            st.markdown(reply)

            if tool_events:
                with st.expander(f"Tools used ({len(tool_events)})", expanded=False):
                    for ev in tool_events:
                        st.caption(f"**{ev['tool']}** → {ev['result'][:120]}…")

        st.session_state["strat_history"].append({"role": "assistant", "content": reply})

        # Persist session to DB
        try:
            _save_session(st.session_state["strat_session_id"], st.session_state["strat_history"])
        except Exception:
            pass

        # Extract structured expert insights (stored in intl_market.gb_expert_insights)
        try:
            from services.gb_knowledge.expert_memory import extract_gb_insights
            n_insights = extract_gb_insights(user_input, reply, _ANTHROPIC_KEY)
            if n_insights:
                st.toast(f"Stored {n_insights} expert insight(s)")
        except Exception:
            pass

    if st.session_state["strat_history"] and st.button("Clear chat", key="clear_strat"):
        st.session_state["strat_history"] = []
        st.session_state["strat_session_id"] = str(uuid.uuid4())
        st.rerun()

# ---- Quant Agent -----------------------------------------------------------
with tab_quant:
    st.header("Quant — BESS Investment Economics")
    st.caption("Grounded on Modo index data · Parametric IRR model")

    if "quant_history" not in st.session_state:
        st.session_state["quant_history"] = []

    for msg in st.session_state["quant_history"]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    user_input_q = st.chat_input("Ask about BESS revenues, leaderboard, or IRR estimates…")
    if user_input_q:
        st.session_state["quant_history"].append({"role": "user", "content": user_input_q})
        with st.chat_message("user"):
            st.markdown(user_input_q)

        with st.chat_message("assistant"):
            with st.spinner("Calculating…"):
                api_messages_q = [
                    {"role": m["role"], "content": m["content"]}
                    for m in st.session_state["quant_history"]
                ]
                try:
                    reply_q, _, tool_events_q = _run_agent_turn(
                        api_messages_q, _build_quant_system(),
                        _QUANT_TOOLS, _dispatch_quant,
                    )
                except Exception as _agent_err:
                    reply_q = f"⚠️ API error: {_agent_err}. Please try again."
                    tool_events_q = []

            st.markdown(reply_q)

            if tool_events_q:
                with st.expander(f"Tools used ({len(tool_events_q)})", expanded=False):
                    for ev in tool_events_q:
                        st.caption(f"**{ev['tool']}** → {ev['result'][:120]}…")

        st.session_state["quant_history"].append({"role": "assistant", "content": reply_q})

        suggestions_q = _extract_memories(user_input_q, reply_q)
        for sug in suggestions_q:
            _save_memory(sug["category"], sug["subject"], sug["content"], source="auto")
        if suggestions_q:
            st.toast(f"Saved {len(suggestions_q)} memory item(s)")

    if st.session_state["quant_history"] and st.button("Clear chat", key="clear_quant"):
        st.session_state["quant_history"] = []
        st.rerun()

# ---- Data Management -------------------------------------------------------
with tab_mgmt:
    st.header("Data Management")

    col1, col2 = st.columns(2)

    with col1:
        st.subheader("Table Coverage")
        if st.button("Refresh counts"):
            _table_counts.clear()
        counts_df = _table_counts()
        st.dataframe(counts_df, use_container_width=True, hide_index=True)

    with col2:
        st.subheader("Run Backfill")
        st.caption("Fetches Modo data and upserts into DB.")
        bf_start = st.date_input("Backfill from", value=date.today() - timedelta(days=30), key="bf_start")
        bf_end   = st.date_input("Backfill to",   value=date.today(), key="bf_end")
        if st.button("Run Backfill", type="primary"):
            with st.spinner("Fetching from Modo Energy API…"):
                result = _run_ingestion_job(bf_start, bf_end, trigger="manual")
            if result["status"] == "success":
                st.success(f"Backfill complete in {result['duration']:.1f}s")
            else:
                st.error(f"Backfill failed: {result['error']}")
            if result.get("log"):
                st.code(result["log"])
            _table_counts.clear()
            _get_ingestion_logs.clear()
            st.rerun()

    st.divider()
    st.subheader("Knowledge Base Ingest")
    kb_col1, kb_col2 = st.columns(2)
    with kb_col1:
        kb_only = st.multiselect(
            "Sources (all if empty)",
            ["elexon", "entso_e", "timera", "modo", "meteologica", "modo_ai"],
            key="kb_only",
        )
    with kb_col2:
        if st.button("Run Knowledge Ingest", type="secondary"):
            with st.spinner("Fetching knowledge from all sources…"):
                kr = _run_knowledge_ingest_job(
                    only=kb_only or None, trigger="manual"
                )
            if kr["status"] == "success":
                st.success(f"{kr['total']} new docs in {kr['duration']:.1f}s")
                if kr.get("results"):
                    st.json(kr["results"])
            else:
                st.error(f"Knowledge ingest failed: {kr['error']}")
            _knowledge_doc_counts.clear()

    st.divider()
    st.subheader("Modo AI Distillation")
    st.caption(
        "Logs into app.modoenergy.com and asks 8 standard GB BESS market questions. "
        "Answers are stored as `modo_ai` knowledge docs and used by the Strategist agent. "
        "Runs automatically at 04:00 SGT nightly."
    )
    if st.button("Run Modo AI Distillation Now", type="secondary", key="modo_ai_btn"):
        with st.spinner("Opening Modo app and querying AI agent (≈2–3 min)…"):
            try:
                from services.gb_knowledge.modo_ai import ModoAIConnector
                from services.gb_knowledge.base import get_db_conn, ensure_table
                _ma_conn = get_db_conn()
                ensure_table(_ma_conn)
                _ma_n = ModoAIConnector().run(_ma_conn)
                _ma_conn.close()
                if _ma_n == 0:
                    st.warning("Modo AI distillation complete — 0 new docs inserted. Check debug screenshots below.")
                else:
                    st.success(f"Modo AI distillation complete — {_ma_n} new docs inserted.")
            except Exception as _ma_e:
                st.error(f"Modo AI distillation failed: {_ma_e}")
        _knowledge_doc_counts.clear()

    # Debug screenshots — written by the login/ask flow to /tmp
    import os as _os
    _debug_shots = sorted([
        f"/tmp/modo_{s}.png"
        for s in ["01_after_nav", "02_no_email", "03_after_email_submit",
                  "04_no_pass", "05_after_submit"]
        if _os.path.exists(f"/tmp/modo_{s}.png")
    ])
    if _debug_shots:
        with st.expander("Debug screenshots (last run)", expanded=False):
            for _shot_path in _debug_shots:
                st.caption(_shot_path)
                st.image(_shot_path)

    st.divider()
    st.subheader("Expert Memory — KB Digestion")
    st.caption(
        "Reads all undigested KB docs and uses Claude to extract durable market insights "
        "into the expert memory pool. These insights are injected into the Strategist's "
        "context at query time, so it 'knows' what it has read. Runs automatically at "
        "03:45 SGT nightly (after KB ingest)."
    )
    if st.button("Digest KB into Expert Memory", key="digest_kb_btn"):
        with st.spinner("Extracting insights from KB docs (1–2 min)…"):
            try:
                from services.gb_knowledge.expert_memory import digest_kb_docs
                _dk_n = digest_kb_docs(_ANTHROPIC_KEY, limit=200)
                st.success(f"Extracted {_dk_n} new insights from KB docs.")
            except Exception as _dk_e:
                st.error(f"KB digest failed: {_dk_e}")

    st.divider()
    st.subheader("Upload Documents to Knowledge Base")
    st.caption(
        "Uploaded files are ingested into the GB knowledge base and become searchable "
        "by the Strategist agent. Supported: PDF, Excel, Word, PPTX, TXT."
    )

    up_tab1, up_tab2 = st.tabs(["Upload Files", "Fetch from URL"])

    with up_tab1:
        uploaded_files = st.file_uploader(
            "Choose files",
            type=["pdf", "xlsx", "xls", "docx", "doc", "pptx", "ppt", "txt"],
            accept_multiple_files=True,
            key="kb_upload",
        )
        if uploaded_files:
            if st.button("Ingest uploaded files", type="primary", key="kb_upload_btn"):
                results = []
                progress = st.progress(0)
                for i, f in enumerate(uploaded_files):
                    with st.spinner(f"Processing {f.name}…"):
                        res = _ingest_uploaded_file(f.name, f.read())
                        results.append((f.name, res))
                    progress.progress((i + 1) / len(uploaded_files))
                progress.empty()
                ok = [r for _, r in results if r["status"] == "success"]
                if ok:
                    st.success(f"Ingested {len(ok)} file(s) successfully.")
                    for _, r in results:
                        if r["status"] == "success":
                            st.caption(f"✓ {r['msg']}")
                for fname, r in results:
                    if r["status"] == "error":
                        st.error(f"✗ {fname}: {r['msg']}")
                _knowledge_doc_counts.clear()
                st.rerun()

    with up_tab2:
        st.caption(
            "Paste a direct article URL to fetch and ingest. "
            "Works for public pages. Paywalled content (e.g. Modo research) must be "
            "saved as PDF first and uploaded via the Files tab."
        )
        fetch_url = st.text_input("Article URL", placeholder="https://modoenergy.com/research/…", key="fetch_url")
        if st.button("Fetch and ingest", type="primary", key="kb_fetch_btn") and fetch_url:
            with st.spinner("Fetching…"):
                res = _ingest_url(fetch_url.strip())
            if res["status"] == "success":
                st.success(res["msg"])
                _knowledge_doc_counts.clear()
            else:
                st.error(res["msg"])

    # Show existing uploaded docs
    uploaded_docs = _query(
        "SELECT title, doc_type, fetched_at::date AS uploaded "
        "FROM intl_market.gb_knowledge_docs WHERE source = 'upload' "
        "ORDER BY fetched_at DESC LIMIT 20"
    )
    if not uploaded_docs.empty:
        with st.expander(f"Uploaded documents ({len(uploaded_docs)})", expanded=False):
            st.dataframe(uploaded_docs, use_container_width=True, hide_index=True)

    st.divider()
    st.subheader("Daily Market Report")
    st.caption(
        "A PDF report (top 10 BESS performers, revenue breakdown, market summary) "
        "is emailed every day at **06:00 SGT**. "
        "Use the button below to send a test report for any date."
    )
    rpt_col1, rpt_col2 = st.columns(2)
    with rpt_col1:
        rpt_date = st.date_input(
            "Report date", value=date.today() - timedelta(days=1), key="rpt_date"
        )
        rpt_email = st.text_input(
            "Send to (comma-separated)", value="chen_dpeng@hotmail.com", key="rpt_email"
        )
    with rpt_col2:
        _smtp_user_default = os.environ.get("SMTP_USER", "dipengchen@gmail.com")
        rpt_from = _smtp_user_default
        st.text(f"Send from: {rpt_from}")
        st.write("")
        send_rpt = st.button("Send Report Now", type="primary", key="send_rpt_btn")

    if send_rpt:
        smtp_user = os.environ.get("SMTP_USER", "")
        smtp_pass = os.environ.get("SMTP_PASSWORD", "")
        if not smtp_user or not smtp_pass:
            st.error(
                "SMTP credentials not configured. "
                "Set SMTP_USER and SMTP_PASSWORD in the container environment."
            )
        else:
            with st.spinner("Generating PDF and sending email…"):
                try:
                    import importlib.util, pathlib
                    _rpt_path = pathlib.Path(__file__).with_name("daily_report.py")
                    _spec = importlib.util.spec_from_file_location("daily_report_ui", _rpt_path)
                    _mod = importlib.util.module_from_spec(_spec)
                    _spec.loader.exec_module(_mod)
                    pdf_bytes, _ai_cmnt = _mod.generate_report_pdf(rpt_date)
                    _mod.send_daily_report_email(pdf_bytes, rpt_date, rpt_email,
                                                 from_email=rpt_from,
                                                 ai_commentary=_ai_cmnt)
                    st.success(
                        f"Report sent to {rpt_email} (from {rpt_from})  "
                        f"({len(pdf_bytes):,} bytes)"
                    )
                except Exception as _rpt_exc:
                    st.error(f"Report failed: {_rpt_exc}")

    # --- WeCom send ---
    _wecom_default = os.environ.get("WECOM_WEBHOOK_URL", "")
    wecom_col1, wecom_col2 = st.columns([3, 1])
    with wecom_col1:
        wecom_url = st.text_input(
            "WeCom webhook URL(s)",
            value=_wecom_default,
            type="password",
            key="wecom_url",
            help="企业微信群机器人 webhook URL — comma-separated for multiple groups",
        )
    with wecom_col2:
        st.write("")
        st.write("")
        send_wecom = st.button("Send to WeCom", key="send_wecom_btn")

    if send_wecom:
        if not wecom_url:
            st.error("WeCom webhook URL is required.")
        else:
            with st.spinner("Generating PDF and sending to WeCom…"):
                try:
                    import importlib.util, pathlib
                    _rpt_path2 = pathlib.Path(__file__).with_name("daily_report.py")
                    _spec2 = importlib.util.spec_from_file_location("daily_report_wc", _rpt_path2)
                    _mod2 = importlib.util.module_from_spec(_spec2)
                    _spec2.loader.exec_module(_mod2)
                    pdf_bytes2, _ai_cmnt2 = _mod2.generate_report_pdf(rpt_date)
                    _mod2.send_daily_report_wecom(
                        pdf_bytes2, rpt_date,
                        webhook_url=wecom_url,
                        ai_commentary=_ai_cmnt2,
                    )
                    st.success(
                        f"Report sent to WeCom group  ({len(pdf_bytes2):,} bytes)"
                    )
                except Exception as _wc_exc:
                    st.error(f"WeCom send failed: {_wc_exc}")

    st.divider()
    st.subheader("Pricing Batch")
    st.caption("Runs nightly at 04:30 SGT. Computes options value, PF dispatch, and OLS forecast for top-50 BESS assets. Results saved to intl_market.gb_pricing_results.")
    pb_col1, pb_col2, pb_col3 = st.columns(3)
    pb_from = pb_col1.date_input("From", value=date.today() - timedelta(days=7),  key="pb_from")
    pb_to   = pb_col2.date_input("To",   value=date.today() - timedelta(days=1),  key="pb_to")
    with pb_col3:
        st.write("")
        st.write("")
        run_pb = st.button("Run Pricing Batch", key="run_pb_btn")
    _pb_missing = _get_pricing_missing_dates(pb_from.isoformat(), pb_to.isoformat())
    if _pb_missing:
        st.warning(f"**{len(_pb_missing)} date(s) missing** from gb_pricing_results: "
                   + ", ".join(_pb_missing))
    if run_pb:
        pb_dates = [pb_from + timedelta(days=i) for i in range((pb_to - pb_from).days + 1)]
        total_processed = 0
        all_errors: list[str] = []
        pb_progress = st.progress(0, text=f"Running pricing batch for {len(pb_dates)} date(s)…")
        try:
            import importlib.util, pathlib
            _pb_path2 = pathlib.Path(__file__).with_name("pricing_batch.py")
            _pb_spec2 = importlib.util.spec_from_file_location("pricing_batch_ui", _pb_path2)
            _pb_mod2  = importlib.util.module_from_spec(_pb_spec2)
            _pb_spec2.loader.exec_module(_pb_mod2)
            for _pb_i, _pb_d in enumerate(pb_dates):
                pb_progress.progress(
                    (_pb_i) / len(pb_dates),
                    text=f"Processing {_pb_d.isoformat()} ({_pb_i + 1}/{len(pb_dates)})…",
                )
                _r = _pb_mod2.run_pricing_batch(_pb_d, _conn())
                total_processed += _r.get("processed", 0)
                all_errors.extend(_r.get("errors", []))
            pb_progress.progress(1.0, text="Done.")
            st.success(
                f"Batch complete: {total_processed} asset-days processed across {len(pb_dates)} date(s)"
            )
            if all_errors:
                st.warning(f"{len(all_errors)} errors: " + "; ".join(all_errors[:5]))
        except Exception as _pb_exc2:
            st.error(f"Pricing batch failed: {_pb_exc2}")

    st.divider()
    st.subheader("Fuel Mix Backfill")
    st.caption("Backfill GB generation fuel mix from NESO CKAN into intl_market.gb_fuel_mix.")
    fm_col1, fm_col2, fm_col3 = st.columns(3)
    fm_from = fm_col1.date_input("From", value=date.today() - timedelta(days=30), key="fm_from")
    fm_to   = fm_col2.date_input("To",   value=date.today() - timedelta(days=1),  key="fm_to")
    with fm_col3:
        st.write("")
        st.write("")
        run_fm = st.button("Run Fuel Mix Backfill", key="run_fm_btn")
    _fm_missing = _get_fuel_mix_missing_dates(fm_from.isoformat(), fm_to.isoformat())
    if _fm_missing:
        st.warning(f"**{len(_fm_missing)} date(s) missing** from gb_fuel_mix: "
                   + ", ".join(_fm_missing))
    if run_fm:
        with st.spinner("Fetching fuel mix from NESO CKAN…"):
            try:
                import importlib.util, pathlib
                _fm_path2 = pathlib.Path(__file__).with_name("fuel_mix_ingest.py")
                _fm_spec2 = importlib.util.spec_from_file_location("fuel_mix_ingest_ui", _fm_path2)
                _fm_mod2  = importlib.util.module_from_spec(_fm_spec2)
                _fm_spec2.loader.exec_module(_fm_mod2)
                n = _fm_mod2.ingest_fuel_mix_range(fm_from, fm_to, _conn())
                st.success(f"Fuel mix backfill complete: {n} rows upserted ({fm_from} → {fm_to})")
            except Exception as _fm_exc2:
                st.error(f"Fuel mix backfill failed: {_fm_exc2}")

    st.divider()
    st.subheader("Scheduled Downloads")

    # Scheduler status
    try:
        sched = _start_scheduler()
        jobs = sched.get_jobs()
        if jobs:
            job = jobs[0]
            next_run = job.next_run_time
            next_str = next_run.strftime("%Y-%m-%d %H:%M SGT") if next_run else "—"
            st.success(f"Scheduler running · Next run: **{next_str}** (daily 03:00 SGT)")
        else:
            st.warning("Scheduler has no active jobs.")
    except Exception as e:
        st.error(f"Scheduler error: {e}")

    st.caption("Recent ingestion runs (auto-refreshes every 30s)")
    if st.button("Refresh logs"):
        _get_ingestion_logs.clear()
    logs_df = _get_ingestion_logs(limit=20)
    if logs_df.empty:
        st.info("No ingestion runs recorded yet.")
    else:
        for _, row in logs_df.iterrows():
            status_icon = "✅" if row["status"] == "success" else "❌"
            with st.expander(
                f"{status_icon} {row['run_at_sgt']}  [{row['trigger']}]  "
                f"{row['date_from']} → {row['date_to']}  ({row['duration_seconds']}s)",
                expanded=(row["status"] == "error"),
            ):
                if row["status"] == "error" and row["error_msg"]:
                    st.error(row["error_msg"])
                if row["rows_ingested"]:
                    st.json(row["rows_ingested"])

    st.divider()
    st.subheader("Agent Memory")
    mems = _load_memories(_APP_KEY)
    if mems.empty:
        st.info("No memories saved yet.")
    else:
        for _, row in mems.iterrows():
            c1, c2 = st.columns([10, 1])
            with c1:
                st.markdown(f"**[{row['category']}]** {row['subject']}: {row['content']}")
                st.caption(f"{row['source']} · {row['created_at']}")
            with c2:
                if st.button("🗑", key=f"del_{row['id']}"):
                    _delete_memory(row["id"])
                    st.rerun()

    st.divider()
    st.subheader("Add Memory Manually")
    with st.form("add_mem"):
        cat = st.selectbox("Category",
                           ["market_view", "methodology", "asset_note", "investment_thesis", "red_flag"])
        subj = st.text_input("Subject (≤8 words)")
        cont = st.text_area("Content (one sentence)")
        if st.form_submit_button("Save"):
            _save_memory(cat, subj, cont, source="manual")
            st.success("Saved")
            st.rerun()
