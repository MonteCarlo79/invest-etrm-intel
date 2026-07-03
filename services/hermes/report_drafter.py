"""Multi-source deep report drafter for Hermes.

Queries multiple specialist market agents in parallel, combines their outputs
with any user-uploaded reference files and the user's own notes, then uses
Claude Opus to synthesise a structured, conference-ready report.

Usage:
    from services.hermes.report_drafter import draft_report
    md = draft_report(
        topic="中国现货电力市场 BESS 投资机遇",
        user_notes="重点分析蒙西、山东两省；关注2025年价差收窄原因",
        markets=["spot", "bess-map"],
        file_texts=[{"filename": "policy.pdf", "text": "..."}],
        api_key=api_key,
        pg_url=pg_url,
    )
"""
from __future__ import annotations

import io
import logging
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional

logger = logging.getLogger(__name__)

# ── Tailored sub-questions per market agent ───────────────────────────────────
_MARKET_QUESTIONS = {
    "spot": (
        "为以下报告主题提供关键数据分析（最近90天数据）：{topic}。"
        "请涵盖：各省DA/RT价差趋势、省间交易量变化、市场基本面（装机结构/负荷水平）、"
        "现货价格波动的主要驱动因素、以及对储能资产盈利的影响。"
        "使用 markdown 表格呈现多省对比数据。"
    ),
    "bess-map": (
        "为以下报告主题提供储能经济性分析：{topic}。"
        "请涵盖：主要省份BESS捕获率排名、IRR测算、当前装机规模、"
        "蒙西/山东/广东等核心市场的峰谷价差和套利空间、以及对新增投资的建议。"
        "使用 markdown 表格呈现省份对比数据。"
    ),
    "mengxi": (
        "为以下报告主题提供内蒙古BESS资产运营分析：{topic}。"
        "请涵盖：4个资产近期P&L表现、实际vs理论策略差距、RT价格趋势、"
        "以及运营优化建议。"
    ),
    "gb": (
        "For the following report topic, provide a Great Britain electricity market analysis: {topic}. "
        "Cover: system prices, BESS revenue streams (FFR/DCR/BM), recent market trends, "
        "and investment outlook for battery storage."
    ),
    "au": (
        "For the following report topic, provide an Australia NEM market analysis: {topic}. "
        "Cover: pool prices, FCAS revenues, BESS project performance, and investment outlook."
    ),
    "internet": (
        "搜索以下报告主题的最新政策文件、行业动态、权威报告和市场新闻：{topic}。"
        "重点关注最近6个月内发布的内容。"
    ),
}

_DEFAULT_QUESTION = (
    "为以下报告主题提供你最重要的数据、分析和市场洞察：{topic}。"
    "覆盖核心数据点、近期趋势、风险和机遇。使用 markdown 表格呈现关键对比数据。"
)

_SYNTHESIS_SYSTEM = """\
You are a senior energy market analyst and investment advisor drafting a comprehensive \
conference report. The report will be converted to presentation slides or a conference \
brief, so use clear section headers and concise bullet points for data.

Respond in the SAME LANGUAGE as the Author's Notes — Chinese (Simplified) if the notes \
are in Chinese, English if in English. Mix languages only if the author does so.

## IMPORTANT: Detect author intent from the Notes and adapt structure accordingly.

If the Author's Notes contain phrases like "如何定位" / "战略定位" / "我们应该怎么做" / \
"指导开发" / "投资布局" / "where should we invest" / "how to position" — the author needs \
a DEVELOPER/INVESTOR STRATEGY FRAMEWORK as the PRIMARY output. In that case, lead with:

## 开发商战略定位框架 / Developer Strategy Framework
  - Which provinces / markets offer the best risk-adjusted returns RIGHT NOW
  - Recommended asset scale, duration, and technology positioning
  - Timing considerations: policy windows, grid connection queues, subsidy cliffs
  - Competitive positioning vs. other developers in each target market

Then follow with the supporting evidence sections.

If the Notes are primarily analytical (no positioning language), use the standard structure:

## 执行摘要 / Executive Summary
  - 3-5 key takeaways (bullets)

## 市场背景 / Market Overview
  - Context, current state, size and trajectory

## 核心数据分析 / Key Data & Analysis
  - The most important quantitative findings from the market agents
  - Use markdown tables for comparisons

## 投资逻辑与机遇 / Investment Thesis & Opportunities
  - Based on data + author's strategic notes

## 风险与挑战 / Risks & Challenges
  - Quantified where possible

## 结论与建议 / Conclusions & Recommendations
  - Actionable, specific, prioritised — name specific provinces and asset profiles

Rules:
- NEVER invent numbers. Only state data that was provided in the context.
- Cite the data source in parentheses after figures (e.g. "(spot agent, Jun 2026)").
- If a market agent returned an error or no data, omit that section gracefully.
- Integrate the Author's Notes as the strategic narrative backbone of the entire report.
- Reference file content as supporting evidence; analyse and synthesise it, do not merely summarise it.
- Name specific policies (文号), provinces, IRR ranges, and price levels wherever available.
- Total length: ~2000-3000 words. Do NOT truncate the Conclusions section.
"""


def _extract_file_text(filename: str, file_bytes: bytes) -> str:
    """Extract readable text from uploaded file bytes. Returns up to 6000 chars."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    try:
        if ext in ("xlsx", "xlsm", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        rows.append("\t".join(cells))
                    if len(rows) >= 100:
                        break
                if rows:
                    parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
            return "\n\n".join(parts)[:6000]

        if ext == "pdf":
            import pdfplumber
            pages = []
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page in pdf.pages[:15]:
                    t = page.extract_text()
                    if t:
                        pages.append(t)
            return "\n\n".join(pages)[:6000]

        if ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            texts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    texts.append(f"[Slide {i}] " + " | ".join(slide_texts))
            return "\n".join(texts)[:6000]

        if ext in ("docx", "doc"):
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)[:6000]

        if ext == "txt":
            return file_bytes.decode("utf-8", errors="replace")[:6000]

        # Unknown binary — skip
        return ""
    except Exception as exc:
        logger.warning("File text extraction failed for %s: %s", filename, exc)
        return ""


def draft_report(
    topic: str,
    user_notes: str,
    markets: list[str],
    file_texts: list[dict],
    api_key: str,
    pg_url: str = "",
) -> str:
    """Compile a deep report from market agent data, uploaded files, and user notes.

    Args:
        topic:      Report title / subject.
        user_notes: User's outline, thoughts, and key points to include.
        markets:    List of market agent keys to query (e.g. ["spot", "bess-map"]).
        file_texts: List of {"filename": ..., "text": ...} dicts from uploaded files.
        api_key:    Anthropic API key.
        pg_url:     Postgres URL for market agents.

    Returns:
        Markdown-formatted report string.
    """
    import anthropic
    from services.hermes.market_agent_bridge import run_market_query as _bridge

    client = anthropic.Anthropic(api_key=api_key)

    # ── 1. Query each market agent in parallel ────────────────────────────────
    def _query_market(market: str) -> tuple[str, str]:
        question_tpl = _MARKET_QUESTIONS.get(market, _DEFAULT_QUESTION)
        question = question_tpl.format(topic=topic)
        try:
            result = _bridge(market=market, question=question, api_key=api_key, pg_url=pg_url)
            return market, result
        except Exception as exc:
            logger.error("report_drafter: market agent %s failed: %s", market, exc)
            return market, f"[Agent unavailable: {exc}]"

    agent_results: dict[str, str] = {}
    if markets:
        with ThreadPoolExecutor(max_workers=min(len(markets), 4)) as pool:
            futures = {pool.submit(_query_market, m): m for m in markets}
            for fut in as_completed(futures):
                market, result = fut.result()
                agent_results[market] = result
                logger.info("report_drafter: %s agent done (%d chars)", market, len(result))

    # ── 2. Assemble synthesis context ─────────────────────────────────────────
    context_sections: list[str] = []

    if user_notes:
        context_sections.append(f"### Author's Notes & Outline\n{user_notes}")

    for entry in file_texts:
        fname = entry.get("filename", "file")
        text = entry.get("text", "")
        if text:
            context_sections.append(f"### Reference File: {fname}\n{text}")

    for market, result in agent_results.items():
        if result and not result.startswith("[Agent unavailable"):
            context_sections.append(f"### Market Data: {market} agent\n{result}")

    if not context_sections:
        return f"⚠️ 无可用数据：市场代理未返回结果，且未找到参考文件。主题：{topic}"

    combined = "\n\n---\n\n".join(context_sections)

    # ── 3. Synthesise with Claude Opus ────────────────────────────────────────
    user_prompt = (
        f"**Report Topic:** {topic}\n\n"
        f"**Requested Markets / Data Sources:** {', '.join(markets) if markets else '(none)'}\n\n"
        f"{combined}"
    )

    resp = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=8192,
        system=_SYNTHESIS_SYSTEM,
        messages=[{"role": "user", "content": user_prompt}],
    )
    report_text = resp.content[0].text.strip()
    logger.info("report_drafter: synthesis complete (%d chars)", len(report_text))

    # Prepend a header noting sources used
    file_names = ", ".join(e["filename"] for e in file_texts if e.get("text"))
    sources_line = f"**数据来源：** {', '.join(markets)}"
    if file_names:
        sources_line += f" + 参考文件：{file_names}"
    return f"# {topic}\n\n{sources_line}\n\n{report_text}"
