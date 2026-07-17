"""
Manual Cap Comp / FR Market Data Entry
=======================================
Extracts per-province 容量补偿 and 调频市场 data from:
  - Free-form text messages
  - Uploaded PDF / Excel / TXT files (as bytes)
  - URLs (fetched via requests)

Entry points:
  extract_capcomp_from_text(text, api_key, pg_url, year)  → {cap_comp_rows, fr_rows, summary}
  extract_capcomp_from_file(filename, file_bytes, api_key, pg_url, year)  → same
  extract_capcomp_from_url(url, api_key, pg_url, year)  → same
  is_capcomp_file(filename)  → bool — quick check if a file is likely cap comp / FR data
"""
from __future__ import annotations

import io
import json
import logging
import re
from datetime import date, datetime
from typing import Optional

import anthropic

logger = logging.getLogger(__name__)

# ── Filename detection ─────────────────────────────────────────────────────────

_CAPCOMP_FILE_KEYWORDS = [
    "容量补偿", "容量电价", "调频辅助", "调频市场", "调频价格",
    "capcomp", "fr_market", "fr-market",
]


def is_capcomp_file(filename: str) -> bool:
    """Return True if filename suggests cap comp or FR market data."""
    fn = filename.lower()
    for kw in _CAPCOMP_FILE_KEYWORDS:
        if kw.lower() in fn:
            return True
    return False


# ── Text extraction from file bytes ───────────────────────────────────────────

def _text_from_pdf(file_bytes: bytes) -> str:
    """Extract plain text from PDF bytes using pdfplumber."""
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            pages = []
            for page in pdf.pages[:30]:  # cap at 30 pages
                t = page.extract_text() or ""
                pages.append(t)
            return "\n\n".join(pages)
    except Exception as exc:
        logger.warning("PDF text extraction failed: %s", exc)
        return ""


def _text_from_excel(file_bytes: bytes) -> str:
    """Convert Excel bytes to a text representation using pandas."""
    try:
        import pandas as pd
        xl = pd.ExcelFile(io.BytesIO(file_bytes))
        parts = []
        for sheet in xl.sheet_names[:5]:
            df = xl.parse(sheet, header=None)
            parts.append(f"[Sheet: {sheet}]\n{df.to_string(index=False, header=False)}")
        return "\n\n".join(parts)
    except Exception as exc:
        logger.warning("Excel text extraction failed: %s", exc)
        return ""


def _text_from_bytes(filename: str, file_bytes: bytes) -> str:
    """Dispatch to the right extractor based on file extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        return _text_from_pdf(file_bytes)
    if ext in ("xlsx", "xls", "xlsm"):
        return _text_from_excel(file_bytes)
    if ext in ("docx", "doc"):
        try:
            import docx as _docx
            doc = _docx.Document(io.BytesIO(file_bytes))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("DOCX text extraction failed: %s", exc)
            return ""
    if ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("PPTX text extraction failed: %s", exc)
            return ""
    # TXT / CSV / other plain text
    for enc in ("utf-8", "gbk", "gb2312"):
        try:
            return file_bytes.decode(enc)
        except Exception:
            pass
    return file_bytes.decode("utf-8", errors="replace")


# ── Claude extraction prompt ───────────────────────────────────────────────────

_SYSTEM = (
    "你是中国电力市场政策数据提取专家。"
    "从提供的文本中提取各省的储能容量补偿标准和调频市场价格数据。"
    "Respond ONLY with valid JSON, no other text."
)

_PROMPT = """\
目标年份参考：{year}

以下是待提取的文本：

{context}

请提取所有省份的容量补偿和调频市场数据，以JSON格式回答：
{{
  "cap_comp_rows": [
    {{
      "province": "省份名称（如广东、山西、内蒙古（蒙西）等）",
      "cap_comp_yuan_kw": <容量补偿标准，元/kW·年，数字>,
      "peak_duration_hours": <年最高净负荷峰值时长，小时，数字或null>,
      "effective_year": <生效年份，整数，如2026>,
      "source": "<来源说明，如文件名或政策名称>"
    }}
  ],
  "fr_rows": [
    {{
      "province": "省份名称",
      "fr_price_yuan_kw_h": <调频容量价格，元/kW/h，数字或null（若仅知资金池总量可填null）>,
      "fr_pool_billion_yuan": <全省调频总资金池，亿元（该月或该年），数字或null>,
      "effective_year_month": "<数据所属年月，格式YYYY-MM，如2026-04>",
      "source": "<来源说明>"
    }}
  ]
}}

规则：
- 只填写文本中明确出现的数值（含可合理推断的数值）
- 省份名称使用标准名称（如"蒙西"写为"内蒙古（蒙西）"，"蒙东"写为"内蒙古（蒙东）"）
- 如果找不到某类数据，对应数组返回空列表 []
- 容量补偿元/kW（不含时长系数）；调频价格元/kW/h（即元/千瓦·时）
- 调频资金池：若为月度结算数据填写当月金额（亿元），若为年度数据填写全年金额（亿元）
- effective_year_month：调频数据用月度粒度（如"2026-04"），若只知年份则写"YYYY-01"
"""


def _call_claude(context: str, api_key: str, year: int) -> Optional[dict]:
    """Call Claude Haiku to extract structured cap comp / FR data from text."""
    # Truncate context to avoid token limits
    context = context[:12000]
    prompt = _PROMPT.format(year=year, context=context)
    try:
        client = anthropic.Anthropic(api_key=api_key)
        resp = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=2048,
            system=_SYSTEM,
            messages=[{"role": "user", "content": prompt}],
        )
        text = resp.content[0].text.strip()
        # Strip markdown code fences if present
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            # Model returned prose + JSON — find the outermost {...} block
            match = re.search(r'\{[\s\S]*\}', text)
            if match:
                return json.loads(match.group())
            raise
    except Exception as exc:
        logger.error("Claude cap comp extraction failed: %s", exc)
        return None


# ── Row validation + normalisation ────────────────────────────────────────────

_PROVINCE_ALIASES = {
    "蒙西": "内蒙古（蒙西）",
    "内蒙西": "内蒙古（蒙西）",
    "蒙东": "内蒙古（蒙东）",
    "内蒙东": "内蒙古（蒙东）",
    "冀南": "河北",
    "冀": "冀北",
    "西北": None,  # too ambiguous, skip
}


def _normalise_province(name: str) -> Optional[str]:
    name = name.strip()
    if name in _PROVINCE_ALIASES:
        return _PROVINCE_ALIASES[name]
    return name if name else None


def _safe_float(val) -> Optional[float]:
    if val is None:
        return None
    try:
        return float(val)
    except (TypeError, ValueError):
        return None


def _safe_year(val, default: int) -> int:
    if val is None:
        return default
    if isinstance(val, int) and 2015 <= val <= 2035:
        return val
    m = re.search(r"\b(20\d{2})\b", str(val))
    if m:
        yr = int(m.group(1))
        if 2015 <= yr <= 2035:
            return yr
    return default


def _build_cap_comp_row(raw: dict, default_year: int, default_source: str) -> Optional[dict]:
    province = _normalise_province(raw.get("province", ""))
    if not province:
        return None
    cap_val = _safe_float(raw.get("cap_comp_yuan_kw"))
    if cap_val is None or cap_val <= 0:
        return None
    return {
        "province": province,
        "effective_date": date(_safe_year(raw.get("effective_year"), default_year), 1, 1),
        "cap_comp_yuan_kw": cap_val,
        "peak_duration_hours": _safe_float(raw.get("peak_duration_hours")),
        "source": str(raw.get("source") or default_source)[:500],
    }


def _parse_year_month_field(val, default_year: int) -> date:
    """Parse 'YYYY-MM' or 'YYYY' from effective_year_month field → date(year, month, 1)."""
    if val:
        m = re.match(r"(\d{4})-(\d{1,2})$", str(val).strip())
        if m:
            yr, mo = int(m.group(1)), int(m.group(2))
            if 2015 <= yr <= 2035 and 1 <= mo <= 12:
                return date(yr, mo, 1)
    return date(_safe_year(val, default_year), 1, 1)


def _build_fr_row(raw: dict, default_year: int, default_source: str) -> Optional[dict]:
    province = _normalise_province(raw.get("province", ""))
    if not province:
        return None
    fr_price = _safe_float(raw.get("fr_price_yuan_kw_h"))
    fr_pool  = _safe_float(raw.get("fr_pool_billion_yuan"))
    # Require at least one of unit price or total pool size
    if (fr_price is None or fr_price <= 0) and (fr_pool is None or fr_pool <= 0):
        return None
    # Support monthly granularity via effective_year_month ("YYYY-MM") or fallback to effective_year
    eff_ym = raw.get("effective_year_month") or raw.get("effective_year")
    return {
        "province": province,
        "effective_date": _parse_year_month_field(eff_ym, default_year),
        "fr_price_yuan_kw_h": fr_price if (fr_price and fr_price > 0) else None,
        "fr_pool_billion_yuan": fr_pool,
        "source": str(raw.get("source") or default_source)[:500],
    }


# ── Main entry points ──────────────────────────────────────────────────────────

def _process_extracted(
    data: Optional[dict],
    api_key: str,
    pg_url: str,
    year: int,
    source_tag: str,
) -> dict:
    """Validate extracted data, upsert to DB, return summary."""
    from services.hermes.capcomp_etl import upsert_cap_comp_rows, upsert_fr_rows

    result = {
        "cap_comp_upserted": 0,
        "fr_upserted": 0,
        "conflicts": 0,
        "errors": [],
        "cap_comp_provinces": [],
        "fr_provinces": [],
    }

    if not data:
        result["errors"].append("Claude未返回有效JSON")
        return result

    # ── Cap comp rows ──
    cap_rows = []
    for raw in (data.get("cap_comp_rows") or []):
        row = _build_cap_comp_row(raw, year, source_tag)
        if row:
            cap_rows.append(row)
        else:
            logger.debug("Skipped invalid cap_comp raw row: %s", raw)

    if cap_rows:
        res = upsert_cap_comp_rows(cap_rows, pg_url, source_tag)
        result["cap_comp_upserted"] += res["upserted"]
        result["conflicts"] += res["conflicts"]
        result["errors"].extend(res["errors"])
        result["cap_comp_provinces"] = [r["province"] for r in cap_rows]

    # ── FR rows ──
    fr_rows = []
    for raw in (data.get("fr_rows") or []):
        row = _build_fr_row(raw, year, source_tag)
        if row:
            fr_rows.append(row)
        else:
            logger.debug("Skipped invalid fr raw row: %s", raw)

    if fr_rows:
        res = upsert_fr_rows(fr_rows, pg_url, source_tag)
        result["fr_upserted"] += res["upserted"]
        result["conflicts"] += res["conflicts"]
        result["errors"].extend(res["errors"])
        result["fr_provinces"] = [r["province"] for r in fr_rows]

    return result


def extract_capcomp_from_text(
    text: str,
    api_key: str,
    pg_url: str,
    year: Optional[int] = None,
) -> dict:
    """
    Extract cap comp / FR market data from a free-form text message and upsert to DB.
    Returns summary dict.
    """
    year = year or datetime.now().year
    source_tag = f"manual_text_{year}"
    data = _call_claude(text, api_key, year)
    return _process_extracted(data, api_key, pg_url, year, source_tag)


def extract_capcomp_from_file(
    filename: str,
    file_bytes: bytes,
    api_key: str,
    pg_url: str,
    year: Optional[int] = None,
) -> dict:
    """
    Extract cap comp / FR market data from an uploaded file (PDF/Excel/TXT) and upsert.
    Returns summary dict.
    """
    year = year or datetime.now().year
    source_tag = f"manual_file:{filename[:80]}"
    text = _text_from_bytes(filename, file_bytes)
    if not text.strip():
        return {
            "cap_comp_upserted": 0, "fr_upserted": 0,
            "conflicts": 0, "errors": ["文件无法提取文本内容"],
            "cap_comp_provinces": [], "fr_provinces": [],
        }
    data = _call_claude(text, api_key, year)
    return _process_extracted(data, api_key, pg_url, year, source_tag)


def extract_capcomp_from_image(
    image_bytes: bytes,
    filename: str,
    api_key: str,
    pg_url: str,
    year: Optional[int] = None,
) -> dict:
    """
    Extract cap comp / FR market data from an image (screenshot or photo) using
    Claude vision to transcribe the text, then the standard extraction pipeline.
    Returns the same summary dict as extract_capcomp_from_file.
    """
    import base64
    year = year or datetime.now().year
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "jpeg"
    _mime = {"jpg": "image/jpeg", "jpeg": "image/jpeg",
             "png": "image/png", "webp": "image/webp", "gif": "image/gif"}
    media_type = _mime.get(ext, "image/jpeg")
    client = anthropic.Anthropic(api_key=api_key)
    try:
        resp = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=4000,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": base64.standard_b64encode(image_bytes).decode("utf-8"),
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "请完整转录这张图片中的所有文字内容，"
                            "保留所有省份名称、数字、单位和政策规则，逐字转录，不要总结。"
                        ),
                    },
                ],
            }],
        )
        extracted_text = resp.content[0].text.strip()
    except Exception as exc:
        return {
            "cap_comp_upserted": 0, "fr_upserted": 0, "conflicts": 0,
            "errors": [f"图片文字识别失败：{exc}"],
            "cap_comp_provinces": [], "fr_provinces": [],
        }
    if not extracted_text:
        return {
            "cap_comp_upserted": 0, "fr_upserted": 0, "conflicts": 0,
            "errors": ["图片中未识别到文字"],
            "cap_comp_provinces": [], "fr_provinces": [],
        }
    source_tag = f"manual_image:{filename[:80]}"
    data = _call_claude(extracted_text, api_key, year)
    return _process_extracted(data, api_key, pg_url, year, source_tag)


def extract_capcomp_from_url(
    url: str,
    api_key: str,
    pg_url: str,
    year: Optional[int] = None,
) -> dict:
    """
    Fetch a URL and extract cap comp / FR data from its content.
    Returns summary dict.
    """
    import requests
    year = year or datetime.now().year
    source_tag = f"manual_url:{url[:100]}"
    try:
        resp = requests.get(url, timeout=20, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        content_type = resp.headers.get("content-type", "")
        if "pdf" in content_type or url.lower().endswith(".pdf"):
            text = _text_from_pdf(resp.content)
        elif "html" in content_type:
            # Strip HTML tags
            text = re.sub(r"<[^>]+>", " ", resp.text)
            text = re.sub(r"\s+", " ", text)
        else:
            text = resp.text
    except Exception as exc:
        return {
            "cap_comp_upserted": 0, "fr_upserted": 0,
            "conflicts": 0, "errors": [f"URL获取失败：{exc}"],
            "cap_comp_provinces": [], "fr_provinces": [],
        }
    data = _call_claude(text, api_key, year)
    return _process_extracted(data, api_key, pg_url, year, source_tag)


def format_result_message(result: dict) -> str:
    """Format extraction result as a human-readable Feishu/Telegram message."""
    lines = []

    if result["cap_comp_upserted"] > 0:
        provs = "、".join(result["cap_comp_provinces"])
        lines.append(f"✅ 容量补偿入库：{result['cap_comp_upserted']} 条（{provs}）")

    if result["fr_upserted"] > 0:
        provs = "、".join(result["fr_provinces"])
        lines.append(f"✅ 调频市场入库：{result['fr_upserted']} 条（{provs}）")

    if result["conflicts"] > 0:
        lines.append(f"⚠️ 数据冲突（差异>5%）：{result['conflicts']} 条，请在 bess-map 容量补偿Tab确认")

    if not lines:
        lines.append("⚠️ 未找到可入库的容量补偿或调频价格数据")
        lines.append("支持格式：「广东 容量补偿 165元/kW 净负荷6小时」或上传政策PDF/Excel/TXT")

    if result["errors"]:
        lines.append(f"错误：{result['errors'][0]}")

    return "\n".join(lines)


# ── Gap fill dispatcher ───────────────────────────────────────────────────────

_GAP_FILL_PROMPTS = {
    "province_cap_comp": (
        "从以下文本中提取指定省份的储能容量补偿标准（元/kW·年）和年最高净负荷峰值时段（小时）。"
        '以JSON回答: {"cap_comp_yuan_kw": <数值>, "peak_duration_hours": <数值或null>}'
    ),
    "province_fr_market": (
        "从以下文本中提取指定省份的调频容量价格（元/kW·h）和全省调频资金池（亿元/年）。"
        '以JSON回答: {"fr_price_yuan_kw_h": <数值>, "fr_pool_yi_yuan": <数值或null>}'
    ),
    "province_installed_monthly": (
        "从以下文本中提取指定省份的储能装机容量（MW）。"
        '以JSON回答: {"installed_mw": <数值>}'
    ),
    "province_sysopfee_monthly": (
        "从以下文本中提取指定省份的系统运行费（元/kWh）。"
        '以JSON回答: {"fee_yuan_kwh": <数值>}'
    ),
}


def extract_from_file_for_gap(
    file_bytes: bytes,
    filename: str,
    fill_table: str,
    province: str,
    month: str,
    api_key: str,
) -> dict:
    """
    Extract the relevant value(s) from a file for a specific gap fill target.
    Supports PDF, Excel, JPG/PNG (vision), DOCX, PPTX, TXT.
    Returns:
        {"extracted": True, "values": {...field: value}, "summary": str}
        {"extracted": False, "error": str}
    """
    import anthropic as _ant

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    prompt_template = _GAP_FILL_PROMPTS.get(fill_table)
    if not prompt_template:
        return {"extracted": False, "error": f"Unsupported fill_table: {fill_table}"}

    client = _ant.Anthropic(api_key=api_key)

    # Image path: use Claude vision
    if ext in ("jpg", "jpeg", "png", "webp", "gif"):
        _mime = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                 "webp": "image/webp", "gif": "image/gif"}
        import base64 as _b64
        b64 = _b64.standard_b64encode(file_bytes).decode()
        vision_prompt = (
            f"图片中包含电力市场数据。省份：{province}，月份：{month}。\n"
            f"请先转录图片中的文字，然后{prompt_template}"
        )
        try:
            resp = client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=300,
                messages=[{"role": "user", "content": [
                    {"type": "image", "source": {"type": "base64",
                     "media_type": _mime.get(ext, "image/jpeg"), "data": b64}},
                    {"type": "text", "text": vision_prompt},
                ]}],
            )
            text = resp.content[0].text.strip()
        except Exception as exc:
            return {"extracted": False, "error": str(exc)}
    else:
        # Text-based extraction
        raw_text = _text_from_bytes(filename, file_bytes)
        if not raw_text.strip():
            return {"extracted": False, "error": "无法从文件中提取文本"}
        full_prompt = (
            f"省份：{province}，月份/年份：{month}。\n"
            f"{prompt_template}\n\n文本内容：\n{raw_text[:6000]}"
        )
        try:
            resp = client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=300,
                messages=[{"role": "user", "content": full_prompt}],
            )
            text = resp.content[0].text.strip()
        except Exception as exc:
            return {"extracted": False, "error": str(exc)}

    # Parse JSON from response
    import json as _json, re as _re
    match = _re.search(r'\{[^{}]+\}', text)
    if not match:
        return {"extracted": False, "error": f"AI未能返回JSON: {text[:200]}"}
    try:
        values = _json.loads(match.group())
        values["fill_table"] = fill_table
        values["fill_province"] = province
        values["fill_month"] = month
        summary = "\n".join(f"{k}: {v}" for k, v in values.items()
                            if k not in ("fill_table", "fill_province", "fill_month"))
        return {"extracted": True, "values": values, "summary": summary}
    except Exception as exc:
        return {"extracted": False, "error": f"JSON解析失败: {exc} — {text[:200]}"}
