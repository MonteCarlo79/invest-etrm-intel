from __future__ import annotations
import logging
import time
from threading import Lock
from typing import Callable, Optional
from urllib.parse import quote

import requests

logger = logging.getLogger(__name__)

_GRAPH_BASE = "https://graph.microsoft.com/v1.0"
_TOKEN_ENDPOINT = "https://login.microsoftonline.com/consumers/oauth2/v2.0/token"
_SCOPES = "Files.ReadWrite offline_access User.Read"
_TOKEN_TTL_MARGIN = 300  # refresh 5 min before expiry


class OneDriveClient:
    """Microsoft Graph API client for personal OneDrive.

    Refresh tokens rotate on each use. The new refresh token is stored
    in memory and logged at WARNING level so it can be persisted if the
    service restarts before the next automatic rotation.
    """

    def __init__(
        self,
        client_id: str,
        client_secret: str,
        refresh_token: str,
        tenant: str = "consumers",
        on_token_rotated: Optional[Callable[[str], None]] = None,
    ) -> None:
        self.client_id = client_id
        self.client_secret = client_secret
        self._refresh_token = refresh_token
        self._tenant = tenant
        self._access_token: Optional[str] = None
        self._expires_at: float = 0.0
        self._lock = Lock()
        self._on_token_rotated = on_token_rotated

    # ── Token management ─────────────────────────────────────────────────────

    def _refresh(self) -> None:
        resp = requests.post(
            f"https://login.microsoftonline.com/{self._tenant}/oauth2/v2.0/token",
            data={
                "client_id": self.client_id,
                "client_secret": self.client_secret,
                "grant_type": "refresh_token",
                "refresh_token": self._refresh_token,
                "scope": _SCOPES,
            },
            timeout=15,
        )
        resp.raise_for_status()
        data = resp.json()
        self._access_token = data["access_token"]
        self._expires_at = time.time() + data.get("expires_in", 3600) - _TOKEN_TTL_MARGIN
        if "refresh_token" in data:
            new_rt = data["refresh_token"]
            if new_rt != self._refresh_token:
                self._refresh_token = new_rt
                logger.info("OneDrive refresh token rotated — persisting to DB")
                if self._on_token_rotated:
                    try:
                        self._on_token_rotated(new_rt)
                    except Exception as exc:
                        logger.error("Failed to persist rotated OneDrive token: %s", exc)

    def _token(self) -> str:
        with self._lock:
            if not self._access_token or time.time() >= self._expires_at:
                self._refresh()
            return self._access_token  # type: ignore[return-value]

    def _headers(self) -> dict[str, str]:
        return {"Authorization": f"Bearer {self._token()}"}

    # ── Graph API helpers ─────────────────────────────────────────────────────

    def _get(self, url: str, params: Optional[dict] = None) -> dict:
        resp = requests.get(url, headers=self._headers(), params=params, timeout=20)
        resp.raise_for_status()
        return resp.json()

    # ── Public API ────────────────────────────────────────────────────────────

    def list_items(self, folder_path: str = "/") -> list[dict]:
        """List files and folders at *folder_path* (e.g. '/' or '/Documents')."""
        if folder_path in ("", "/"):
            url = f"{_GRAPH_BASE}/me/drive/root/children"
        else:
            path = folder_path.strip("/")
            url = f"{_GRAPH_BASE}/me/drive/root:/{path}:/children"
        data = self._get(url, params={"$select": "id,name,size,lastModifiedDateTime,file,folder"})
        return data.get("value", [])

    def search(self, query: str) -> list[dict]:
        """Full-text search across the user's OneDrive."""
        data = self._get(
            f"{_GRAPH_BASE}/me/drive/root/search(q='{query}')",
            params={"$select": "id,name,size,lastModifiedDateTime,parentReference,webUrl"},
        )
        return data.get("value", [])

    def read_file(self, item_id: str) -> bytes:
        """Download file content by item ID."""
        resp = requests.get(
            f"{_GRAPH_BASE}/me/drive/items/{item_id}/content",
            headers=self._headers(),
            timeout=60,
            allow_redirects=True,
        )
        resp.raise_for_status()
        return resp.content

    def read_file_by_path(self, file_path: str) -> bytes:
        """Download file content by OneDrive path (e.g. 'bess-platform/data/电站.xlsx').

        Uses the path-based Graph API endpoint directly — no search() call needed.
        """
        enc_path = "/".join(quote(seg, safe="") for seg in file_path.strip("/").split("/"))
        url = f"{_GRAPH_BASE}/me/drive/root:/{enc_path}:/content"
        resp = requests.get(url, headers=self._headers(), timeout=60, allow_redirects=True)
        resp.raise_for_status()
        return resp.content

    def read_file_text(self, item_id: str, encoding: str = "utf-8") -> str:
        """Download a text file and return its content as a string."""
        return self.read_file(item_id).decode(encoding, errors="replace")

    def get_item_metadata(self, item_id: str) -> dict:
        """Return metadata (name, size, etc.) for a drive item."""
        resp = requests.get(
            f"{_GRAPH_BASE}/me/drive/items/{item_id}",
            headers=self._headers(),
            params={"$select": "id,name,size,file"},
            timeout=20,
        )
        resp.raise_for_status()
        return resp.json()

    def read_file_smart(self, item_id: str, max_chars: int = 6000) -> str:
        """Download and parse a file into readable text based on its extension.

        Supports: .xlsx/.xls (spreadsheet), .pdf, .docx, .pptx, .csv,
        .txt/.md/.json/.xml/.html and other plain-text formats.
        Binary files that cannot be parsed return a descriptive error.
        """
        meta = self.get_item_metadata(item_id)
        name = meta.get("name", "")
        ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
        raw = self.read_file(item_id)

        try:
            if ext in ("xlsx", "xlsm"):
                import openpyxl, io
                wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(c.strip() for c in cells):
                            rows.append("\t".join(cells))
                        if len(rows) >= 200:
                            rows.append("… (rows truncated)")
                            break
                    if rows:
                        parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
                text = "\n\n".join(parts)

            elif ext == "xls":
                import xlrd, io
                wb = xlrd.open_workbook(file_contents=raw)
                parts = []
                for sheet_name in wb.sheet_names():
                    ws = wb.sheet_by_name(sheet_name)
                    rows = []
                    for i in range(min(ws.nrows, 200)):
                        cells = [str(ws.cell_value(i, j)) for j in range(ws.ncols)]
                        if any(c.strip() for c in cells):
                            rows.append("\t".join(cells))
                    if rows:
                        parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
                text = "\n\n".join(parts)

            elif ext == "pdf":
                import pdfplumber, io
                pages = []
                with pdfplumber.open(io.BytesIO(raw)) as pdf:
                    for page in pdf.pages[:20]:
                        t = page.extract_text()
                        if t:
                            pages.append(t)
                text = "\n\n".join(pages)

            elif ext == "docx":
                import docx, io
                doc = docx.Document(io.BytesIO(raw))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

            elif ext == "pptx":
                from pptx import Presentation
                import io
                prs = Presentation(io.BytesIO(raw))
                slides = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    if texts:
                        slides.append(f"Slide {i}: " + " | ".join(texts))
                text = "\n".join(slides)

            elif ext == "csv":
                text = raw.decode("utf-8-sig", errors="replace")

            elif ext in ("html", "htm"):
                from bs4 import BeautifulSoup
                text = BeautifulSoup(raw, "html.parser").get_text(separator="\n")

            else:
                # Plain text (txt, md, json, xml, py, etc.) or unknown
                text = raw.decode("utf-8", errors="replace")
                # If it looks like binary (high ratio of replacement chars), bail
                if text.count("\ufffd") > len(text) * 0.05:
                    return f"[Binary file — cannot display as text: {name}]"

        except Exception as exc:
            logger.warning("read_file_smart parse failed for %s: %s", name, exc)
            return f"[Failed to parse {name}: {exc}]"

        if len(text) > max_chars:
            text = text[:max_chars] + "\n… (truncated)"
        return f"[File: {name}]\n\n{text}"

    _SIMPLE_UPLOAD_MAX = 4 * 1024 * 1024  # 4 MB — Graph API limit for simple PUT

    def upload_file(
        self,
        folder_path: str,
        filename: str,
        content: bytes,
        conflict_behavior: str = "replace",
    ) -> dict:
        """Upload (or replace) a file at *folder_path/filename*.

        Uses a resumable upload session for files > 4 MB.
        """
        folder = folder_path.strip("/")
        enc_name = quote(filename, safe="")
        if folder:
            enc_folder = "/".join(quote(seg, safe="") for seg in folder.split("/"))
            item_path = f"{enc_folder}/{enc_name}"
        else:
            item_path = enc_name

        if len(content) <= self._SIMPLE_UPLOAD_MAX:
            # NOTE: @microsoft.graph.conflictBehavior must be a literal query param —
            # passing via requests `params` dict would percent-encode the '@' to '%40'.
            url = f"{_GRAPH_BASE}/me/drive/root:/{item_path}:/content?@microsoft.graph.conflictBehavior={conflict_behavior}"
            resp = requests.put(
                url,
                headers={
                    **self._headers(),
                    "Content-Type": "application/octet-stream",
                },
                data=content,
                timeout=60,
            )
            if not resp.ok:
                try:
                    detail = resp.json()
                except Exception:
                    detail = resp.text[:300]
                raise requests.HTTPError(
                    f"{resp.status_code} uploading {filename}: {detail}",
                    response=resp,
                )
            return resp.json()

        # Large file: use upload session
        session_url = f"{_GRAPH_BASE}/me/drive/root:/{item_path}:/createUploadSession"
        sess_resp = requests.post(
            session_url,
            headers={**self._headers(), "Content-Type": "application/json"},
            json={"item": {"@microsoft.graph.conflictBehavior": conflict_behavior}},
            timeout=30,
        )
        sess_resp.raise_for_status()
        upload_url = sess_resp.json()["uploadUrl"]

        chunk_size = 5 * 1024 * 1024  # 5 MB chunks (must be multiple of 320 KiB)
        total = len(content)
        result: dict = {}
        offset = 0
        while offset < total:
            chunk = content[offset: offset + chunk_size]
            end = offset + len(chunk) - 1
            chunk_resp = requests.put(
                upload_url,
                headers={
                    "Content-Range": f"bytes {offset}-{end}/{total}",
                    "Content-Type": "application/octet-stream",
                },
                data=chunk,
                timeout=120,
            )
            chunk_resp.raise_for_status()
            if chunk_resp.status_code in (200, 201):
                result = chunk_resp.json()
            offset += len(chunk)
        return result

    def create_folder(self, parent_path: str, folder_name: str) -> dict:
        """Create a folder under *parent_path*."""
        parent = parent_path.strip("/")
        if parent:
            url = f"{_GRAPH_BASE}/me/drive/root:/{parent}:/children"
        else:
            url = f"{_GRAPH_BASE}/me/drive/root/children"
        resp = requests.post(
            url,
            headers=self._headers(),
            json={"name": folder_name, "folder": {}, "@microsoft.graph.conflictBehavior": "fail"},
            timeout=10,
        )
        resp.raise_for_status()
        return resp.json()


# ── Shared client factory (used by knowledge_pool vault_reader/vault_writer) ──

import os

_SHARED_CLIENT: Optional["OneDriveClient"] = None
_SHARED_LOCK = Lock()


def set_shared_onedrive_client(client: Optional["OneDriveClient"]) -> None:
    """Register an externally-built client as the process-wide shared one.

    app.py builds the chat OneDrive client at startup; registering it here
    keeps exactly one client (one refresh-token lineage) per process — two
    clients sharing a rotating MSA refresh token kill each other on the
    first rotation.
    """
    global _SHARED_CLIENT
    with _SHARED_LOCK:
        _SHARED_CLIENT = client


def _load_setting(pg_url: str, key: str) -> str:
    """Read a value from hermes_settings. Returns '' on any error."""
    url = pg_url or os.environ.get("PGURL", "")
    if not url:
        return ""
    try:
        import psycopg2
        with psycopg2.connect(url, connect_timeout=5) as conn:
            with conn.cursor() as cur:
                cur.execute("SELECT value FROM hermes_settings WHERE key = %s", (key,))
                row = cur.fetchone()
                return row[0] if row else ""
    except Exception as exc:
        logger.debug("hermes_settings read failed (%s): %s", key, exc)
        return ""


def _save_setting(pg_url: str, key: str, value: str) -> None:
    url = pg_url or os.environ.get("PGURL", "")
    if not url:
        return
    try:
        import psycopg2
        with psycopg2.connect(url, connect_timeout=5) as conn:
            with conn.cursor() as cur:
                cur.execute(
                    """INSERT INTO hermes_settings (key, value, updated_at) VALUES (%s, %s, NOW())
                       ON CONFLICT (key) DO UPDATE SET value = EXCLUDED.value, updated_at = NOW()""",
                    (key, value),
                )
            conn.commit()
    except Exception as exc:
        logger.warning(
            "hermes_settings write failed (%s): %s"
            " — run: python scripts/auth_microsoft_mail.py to refresh OneDrive credentials",
            key, exc,
        )


def get_shared_onedrive_client(pg_url: str = "") -> Optional["OneDriveClient"]:
    """Lazily build a process-wide OneDriveClient from env + hermes_settings.

    Returns None if OneDrive is not configured (or on any error) — callers
    must treat knowledge I/O as optional.
    """
    global _SHARED_CLIENT
    with _SHARED_LOCK:
        if _SHARED_CLIENT is not None:
            return _SHARED_CLIENT
        client_id = os.environ.get("ONEDRIVE_CLIENT_ID", "")
        client_secret = os.environ.get("ONEDRIVE_CLIENT_SECRET", "")
        if not client_id or not client_secret:
            return None
        refresh_token = _load_setting(pg_url, "onedrive_refresh_token") or os.environ.get(
            "ONEDRIVE_REFRESH_TOKEN", ""
        )
        if not refresh_token:
            return None

        def _rotated(new_token: str) -> None:
            _save_setting(pg_url, "onedrive_refresh_token", new_token)

        try:
            _SHARED_CLIENT = OneDriveClient(
                client_id=client_id,
                client_secret=client_secret,
                refresh_token=refresh_token,
                on_token_rotated=_rotated,
            )
        except Exception as exc:
            logger.warning(
                "Shared OneDriveClient init failed: %s"
                " — run: python scripts/auth_microsoft_mail.py to refresh OneDrive credentials",
                exc,
            )
            return None
        return _SHARED_CLIENT
