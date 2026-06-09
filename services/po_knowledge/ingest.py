"""Poland knowledge base ingestion orchestrator.

Sources:
  - Local reports (PDF/Excel/PNG from data/market-fundamentals-po/)
  - PSE (Polish TSO) market publications
  - URE (regulatory authority) announcements

Usage:
    python -m services.po_knowledge.ingest
    python -m services.po_knowledge.ingest --only local_reports
"""
from __future__ import annotations

import argparse
import io
import logging
import os
import pathlib
import sys
from datetime import date

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

logger = logging.getLogger(__name__)

_LOCAL_REPORTS_DIR = pathlib.Path(__file__).parent.parent.parent / "data" / "market-fundamentals-po"


# ── Connectors ────────────────────────────────────────────────────────────────

class LocalReportsConnector:
    """Ingests all PDF, Excel, and image reports from the local market-fundamentals-po directory."""

    source = "local_reports"

    def fetch(self) -> list[dict]:
        docs = []
        if not _LOCAL_REPORTS_DIR.exists():
            logger.warning("[po_local_reports] Directory not found: %s", _LOCAL_REPORTS_DIR)
            return []

        for path in sorted(_LOCAL_REPORTS_DIR.rglob("*")):
            if not path.is_file():
                continue
            ext = path.suffix.lower()
            if ext not in (".pdf", ".xlsx", ".xls", ".pptx", ".ppt", ".txt", ".png", ".jpg", ".jpeg"):
                continue

            url_key = f"local://{path.relative_to(_LOCAL_REPORTS_DIR.parent.parent).as_posix()}"
            title = path.stem.replace("_", " ").replace("-", " ")

            # Skip image files — no text to extract (log for reference only)
            if ext in (".png", ".jpg", ".jpeg"):
                logger.debug("[po_local_reports] Skipping image: %s", path.name)
                docs.append({
                    "doc_type": "image",
                    "title": title[:250],
                    "url": url_key,
                    "published_date": _infer_date(path.stem),
                    "content": f"[Image file: {path.name}] — Polish market screenshot (WXWork capture). No text extractable.",
                })
                continue

            try:
                content = _extract_text(path, ext)
            except Exception as exc:
                logger.warning("[po_local_reports] Skipping %s: %s", path.name, exc)
                continue

            if not content or len(content.strip()) < 100:
                logger.debug("[po_local_reports] Skipping %s: too little text", path.name)
                continue

            docs.append({
                "doc_type": _doc_type(ext),
                "title": title[:250],
                "url": url_key,
                "published_date": _infer_date(path.stem),
                "content": content[:100_000],
            })
            logger.debug("[po_local_reports] Extracted %s (%d chars)", path.name, len(content))

        return docs


def _doc_type(ext: str) -> str:
    return {
        ".pdf": "pdf",
        ".xlsx": "excel",
        ".xls": "excel",
        ".pptx": "presentation",
        ".ppt": "presentation",
        ".txt": "text",
    }.get(ext, "document")


def _infer_date(stem: str) -> date | None:
    import re
    m = re.search(r"(\d{4})[-_](\d{2})[-_](\d{2})", stem)
    if m:
        try:
            return date(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        except ValueError:
            pass
    # Aurora naming: Q1_26, Q2_26, Apr26, Mar26
    m = re.search(r"[Qq](\d)_?(\d{2})", stem)
    if m:
        quarter, yr = int(m.group(1)), int(m.group(2)) + 2000
        month = (quarter - 1) * 3 + 1
        try:
            return date(yr, month, 1)
        except ValueError:
            pass
    months = {"jan": 1, "feb": 2, "mar": 3, "apr": 4, "may": 5, "jun": 6,
              "jul": 7, "aug": 8, "sep": 9, "oct": 10, "nov": 11, "dec": 12}
    m = re.search(r"([A-Za-z]{3})(\d{2})", stem)
    if m:
        mon = m.group(1).lower()
        yr = int(m.group(2)) + 2000
        if mon in months:
            try:
                return date(yr, months[mon], 1)
            except ValueError:
                pass
    return None


def _extract_text(path: pathlib.Path, ext: str) -> str:
    if ext == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = [p.extract_text() for p in reader.pages if p.extract_text()]
            return "\n\n".join(pages)
        except Exception:
            pass
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                return "\n\n".join((p.extract_text() or "") for p in pdf.pages)
        except Exception as exc:
            raise RuntimeError(f"PDF extraction failed: {exc}") from exc

    if ext in (".xlsx", ".xls"):
        import pandas as pd
        xl = pd.ExcelFile(str(path))
        parts = []
        for sheet in xl.sheet_names:
            try:
                df = xl.parse(sheet)
                if not df.empty:
                    parts.append(f"[Sheet: {sheet}]\n{df.to_string(index=False)}")
            except Exception:
                pass
        return "\n\n".join(parts)

    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text.strip() for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except Exception as exc:
            raise RuntimeError(f"PPTX extraction failed: {exc}") from exc

    return ""


class PSEPublicationsConnector:
    """Thin wrapper — delegates to the richer PSEBalancingReportsConnector in entso_scraper."""

    source = "pse_pl"

    def fetch(self) -> list[dict]:
        try:
            from services.po_knowledge.entso_scraper import PSEBalancingReportsConnector
            conn = PSEBalancingReportsConnector()
            docs = conn.fetch(max_pages=3)
            # Normalise source tag to pse_pl for backwards compatibility
            for d in docs:
                d["source"] = "pse_pl"
            return docs
        except Exception as exc:
            logger.warning("[pse_pl] fetch failed: %s", exc)
            return []


# ── Orchestrator ──────────────────────────────────────────────────────────────

def run_knowledge_ingest(only: list[str] | None = None, verbose: bool = True) -> dict[str, int]:
    from dotenv import load_dotenv
    load_dotenv(
        os.path.join(os.path.dirname(__file__), "..", "..", "config", ".env"),
        override=False,
    )

    import psycopg2
    from services.po_knowledge.config import MARKET_CONFIG

    conn = psycopg2.connect(os.environ["PGURL"], keepalives=1, keepalives_idle=30)
    prefix = MARKET_CONFIG.table_prefix

    with conn.cursor() as cur:
        cur.execute(f"""
            CREATE TABLE IF NOT EXISTS intl_market.{prefix}knowledge_docs (
                id              SERIAL PRIMARY KEY,
                source          TEXT NOT NULL,
                doc_type        TEXT NOT NULL,
                title           TEXT,
                url             TEXT UNIQUE,
                published_date  DATE,
                content         TEXT NOT NULL,
                fetched_at      TIMESTAMPTZ NOT NULL DEFAULT NOW(),
                search_vector   TSVECTOR GENERATED ALWAYS AS (
                    to_tsvector('english',
                        coalesce(title,'') || ' ' || left(content,100000))
                ) STORED
            )
        """)
        cur.execute(
            f"CREATE INDEX IF NOT EXISTS {prefix}knowledge_docs_fts "
            f"ON intl_market.{prefix}knowledge_docs USING GIN(search_vector)"
        )
    conn.commit()

    from services.po_knowledge.entso_scraper import (
        PSEGridReportsConnector,
        PSEAFRRDataConnector,
        TGEMarketReportsConnector,
        UREPublicationsConnector,
        ENTSOEPublicationsConnector,
    )

    connectors = [
        ("local_reports",      "Local Aurora reports (PDF/Excel)", LocalReportsConnector()),
        ("pse_pl",             "PSE balancing publications",        PSEPublicationsConnector()),
        ("pse_grid",           "PSE grid/transmission reports",     PSEGridReportsConnector()),
        ("pse_afrr",           "PSE aFRR/FCR tender results",       PSEAFRRDataConnector()),
        ("tge_reports",        "TGE market reports",                TGEMarketReportsConnector()),
        ("ure_regulatory",     "URE regulatory publications",       UREPublicationsConnector()),
        ("entsoe_publications","ENTSO-E news and publications",     ENTSOEPublicationsConnector()),
    ]

    results: dict[str, int] = {}
    for key, label, connector in connectors:
        if only and key not in only:
            continue
        if verbose:
            print(f"  [{key}] {label}…", end="", flush=True)
        try:
            n = _run_connector(connector, conn, prefix)
            results[key] = n
            if verbose:
                print(f" {n} new docs")
        except Exception as exc:
            results[key] = 0
            if verbose:
                print(f" ERROR: {exc}")
            logger.error("[po_ingest:%s] %s", key, exc)

    conn.close()
    return results


def _run_connector(connector, conn, prefix: str) -> int:
    n = 0
    for doc in connector.fetch():
        with conn.cursor() as cur:
            cur.execute(
                f"INSERT INTO intl_market.{prefix}knowledge_docs "
                "(source, doc_type, title, url, published_date, content) "
                "VALUES (%s, %s, %s, %s, %s, %s) "
                "ON CONFLICT (url) DO UPDATE SET "
                "  content=EXCLUDED.content, title=EXCLUDED.title, fetched_at=NOW()",
                (
                    getattr(connector, "source", "unknown"),
                    doc["doc_type"],
                    doc.get("title", ""),
                    doc.get("url"),
                    doc.get("published_date"),
                    doc["content"],
                ),
            )
            if cur.rowcount > 0:
                n += 1
    conn.commit()
    return n


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    parser = argparse.ArgumentParser()
    parser.add_argument("--only", help="Comma-separated connectors to run")
    args = parser.parse_args()
    only = args.only.split(",") if args.only else None
    results = run_knowledge_ingest(only=only)
    print("Results:", results)
