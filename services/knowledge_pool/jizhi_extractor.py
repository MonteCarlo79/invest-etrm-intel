# services/knowledge_pool/jizhi_extractor.py
"""
Structured extraction and persistence for 136号文 机制竞价 bid data.

Tables managed:
  staging.jizhi_bids         — completed bid results (province × year × batch × tech_type)
  staging.jizhi_bid_winners  — 中标清单 (optional sub-table)
  staging.jizhi_upcoming     — upcoming bid calendar

Public API:
  _extract_pptx_text(file_bytes) -> str
  ensure_tables(pg_url)
  extract_bids(text, api_key) -> list[dict]
  extract_upcoming(text, api_key) -> list[dict]
  save_bids(records, source_doc_id, pg_url) -> int
  save_upcoming(records, pg_url) -> int
"""
from __future__ import annotations
import logging

import psycopg2

logger = logging.getLogger(__name__)

# Amazon Nova Pro via APAC cross-region inference — plain text generation, no tool use
# (avoids ModelErrorException that occurs with tool use + large documents)
_BEDROCK_MODEL = "apac.amazon.nova-pro-v1:0"
_BEDROCK_REGION = "ap-southeast-1"
_MAX_RETRIES = 2


def _bedrock_generate(prompt: str) -> str:
    """Call Nova Pro on AWS Bedrock with plain text generation. Returns response text.

    Retries up to _MAX_RETRIES times on transient errors.
    """
    import boto3, time
    client = boto3.client("bedrock-runtime", region_name=_BEDROCK_REGION)
    last_exc: Exception | None = None
    for _attempt in range(_MAX_RETRIES + 1):
        try:
            response = client.converse(
                modelId=_BEDROCK_MODEL,
                messages=[{"role": "user", "content": [{"text": prompt}]}],
                inferenceConfig={"maxTokens": 4096},
            )
            content = response.get("output", {}).get("message", {}).get("content", [])
            for item in content:
                if "text" in item:
                    return item["text"]
            raise RuntimeError(f"Bedrock returned no text (stopReason={response.get('stopReason')!r})")
        except Exception as exc:
            last_exc = exc
            err_str = str(exc)
            if _attempt < _MAX_RETRIES and (
                "ThrottlingException" in err_str or "ServiceUnavailableException" in err_str
            ):
                time.sleep(2 ** _attempt)
                logger.warning("[jizhi] bedrock attempt %d failed, retrying: %s", _attempt + 1, exc)
                continue
            raise
    raise last_exc  # unreachable but satisfies type checker


def _parse_json_from_text(text: str, key: str) -> list[dict]:
    """Extract a JSON array from a fenced or bare JSON block in model output."""
    import json, re
    # Try ```json ... ``` fence first
    m = re.search(r"```(?:json)?\s*(\{.*?\}|\[.*?\])\s*```", text, re.DOTALL)
    if m:
        try:
            parsed = json.loads(m.group(1))
            if isinstance(parsed, list):
                return parsed
            if isinstance(parsed, dict):
                return parsed.get(key, [])
        except json.JSONDecodeError:
            pass
    # Try to find the first { or [ and parse from there
    for start_char, end_char in [('{', '}'), ('[', ']')]:
        idx = text.find(start_char)
        if idx == -1:
            continue
        # Find matching closing bracket
        depth = 0
        for i, ch in enumerate(text[idx:], start=idx):
            if ch == start_char:
                depth += 1
            elif ch == end_char:
                depth -= 1
                if depth == 0:
                    try:
                        parsed = json.loads(text[idx:i+1])
                        if isinstance(parsed, list):
                            return parsed
                        if isinstance(parsed, dict):
                            return parsed.get(key, [])
                    except json.JSONDecodeError:
                        pass
                    break
    logger.warning("[jizhi] could not parse JSON from model response: %s", text[:300])
    return []

_DDL = """
CREATE TABLE IF NOT EXISTS staging.jizhi_bids (
    id                  SERIAL PRIMARY KEY,
    province            TEXT NOT NULL,
    year                INT  NOT NULL,
    batch               TEXT NOT NULL,
    tech_type           TEXT NOT NULL,
    price_floor         NUMERIC,
    price_cap           NUMERIC,
    mechanism_type      TEXT,
    mechanism_value     NUMERIC,
    supply_demand_ratio NUMERIC,
    cleared_price       NUMERIC,
    cleared_volume_gwh  NUMERIC,
    bid_date            DATE,
    verified            BOOLEAN NOT NULL DEFAULT FALSE,
    source_doc_id       INT,
    notes               TEXT,
    created_at          TIMESTAMPTZ NOT NULL DEFAULT NOW(),
    UNIQUE (province, year, batch, tech_type)
);
CREATE TABLE IF NOT EXISTS staging.jizhi_bid_winners (
    id            SERIAL PRIMARY KEY,
    bid_id        INT NOT NULL REFERENCES staging.jizhi_bids(id) ON DELETE CASCADE,
    project_name  TEXT NOT NULL,
    operator      TEXT,
    capacity_mw   NUMERIC,
    cleared_price NUMERIC,
    tech_type     TEXT,
    created_at    TIMESTAMPTZ NOT NULL DEFAULT NOW()
);
CREATE INDEX IF NOT EXISTS idx_jizhi_winners_bid
    ON staging.jizhi_bid_winners(bid_id);
CREATE TABLE IF NOT EXISTS staging.jizhi_upcoming (
    id                   SERIAL PRIMARY KEY,
    province             TEXT NOT NULL,
    year                 INT  NOT NULL,
    batch                TEXT NOT NULL,
    tech_type            TEXT NOT NULL,
    price_floor          NUMERIC,
    price_cap            NUMERIC,
    target_volume_gwh    NUMERIC,
    supply_demand_ratio  NUMERIC,
    bid_open_date        DATE,
    bid_close_date       DATE,
    source_url           TEXT,
    announcement_date    DATE,
    verified             BOOLEAN NOT NULL DEFAULT FALSE,
    notes                TEXT,
    created_at           TIMESTAMPTZ NOT NULL DEFAULT NOW(),
    UNIQUE (province, year, batch, tech_type, bid_open_date)
);
"""

_BIDS_TOOL = {
    "name": "save_bid_results",
    "description": "Save extracted 机制竞价 completed bid results from the document.",
    "input_schema": {
        "type": "object",
        "properties": {
            "bids": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "province":            {"type": "string",  "description": "Province name in Chinese, e.g. 广东"},
                        "year":                {"type": "integer", "description": "Year of the bidding, e.g. 2025"},
                        "batch":               {"type": "string",  "description": "One of: 存量, 增量_2025-12, 增量_2026-12, 增量_2027-12"},
                        "tech_type":           {"type": "string",  "description": "One of: 陆风, 海风, 光伏, 水电"},
                        "price_floor":         {"type": "number",  "description": "Minimum bid price in 元/kWh"},
                        "price_cap":           {"type": "number",  "description": "Maximum bid price in 元/kWh"},
                        "mechanism_type":      {"type": "string",  "description": "One of: 电量, 比例, 小时数, 浮动. Use null if not specified."},
                        "mechanism_value":     {"type": "number",  "description": "Value in GWh (电量), % (比例), or hours (小时数). Leave null for 浮动 type or if not explicitly stated in the document — do NOT copy price_floor or any price field here."},
                        "supply_demand_ratio": {"type": "number",  "description": "Competitive supply-demand ratio (e.g. 1.35 means 135% subscribed vs capacity offered). Leave null if not stated in document — do NOT default to 1."},
                        "cleared_price":       {"type": "number",  "description": "Final cleared price in 元/kWh"},
                        "cleared_volume_gwh":  {"type": "number",  "description": "Total cleared volume in GWh"},
                        "bid_date":            {"type": "string",  "description": "Bid date as YYYY-MM-DD"},
                        "notes":               {"type": "string"},
                    },
                    "required": ["province", "year", "batch", "tech_type"],
                },
            }
        },
        "required": ["bids"],
    },
}

_BIDS_PROMPT = """\
Extract ALL 机制竞价 completed bid results from the document below.
Return ONLY a JSON object with key "bids" containing an array. No explanation, no markdown prose — just the JSON.

Normalisation rules:
- batch: 存量 = grid-connected before 2025-05-31; 增量_2025-12 = before 2025-12-31; 增量_2026-12 = before 2026-12-31; 增量_2027-12 = before 2027-12-31
- tech_type: 陆风 / 海风 / 光伏 / 水电  (map 风电/wind → 陆风 unless specifically 海风)
- prices in 元/kWh  (divide by 1000 if document uses 元/MWh)
- cleared_volume_gwh in GWh  (divide by 1000 if document uses TWh)
- bid_date as YYYY-MM-DD

Critical field rules:
- mechanism_type: use 电量/比例/小时数/浮动. Set to null if unknown.
- mechanism_value: ONLY populate if the document explicitly states a separate mechanism quantity (e.g. "1000小时数", "20%", "500GWh"). NEVER copy price_floor or any price here. Leave null for 浮动 type.
- supply_demand_ratio: ONLY populate if the document explicitly states a 供需比 or subscription ratio (e.g. "申报量是需求量的1.35倍"). Leave null if not stated — do NOT default to 1.
- cleared_volume_gwh: total cleared capacity volume in GWh. Leave null if not stated.

Each bid object must have: province (string), year (int), batch (string), tech_type (string).
Optional: price_floor, price_cap, mechanism_type, mechanism_value, supply_demand_ratio, cleared_price, cleared_volume_gwh, bid_date (YYYY-MM-DD), notes.

Example output format:
{{"bids": [{{"province": "广东", "year": 2025, "batch": "存量", "tech_type": "光伏", "cleared_price": 0.35, "cleared_volume_gwh": 500}}]}}

Document:
{text}"""


def _extract_pptx_text(file_bytes: bytes) -> str:
    """Extract all text from a PPTX file as a single string.

    Handles text shapes, tables, and charts with XML-parsed category labels
    so that province→price associations are preserved for Claude extraction.
    Does not require an API key — purely local, no vision calls.

    Returns empty string if python-pptx is not installed or the file is invalid.
    """
    try:
        from pptx import Presentation  # type: ignore
        from io import BytesIO
        from lxml import etree  # type: ignore
    except ImportError:
        logger.warning("[jizhi] python-pptx or lxml not installed; cannot extract PPTX")
        return ""

    _C = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    def _chart_text(shape) -> str:
        """Return structured text for a chart shape using XML parsing."""
        try:
            root = etree.fromstring(shape.chart._element.xml.encode())
            parts: list[str] = []
            for ser in root.iter(f"{{{_C}}}ser"):
                # Series name is the first <c:v> child (before cat/val)
                ser_name = ""
                for sv in ser.findall(f".//{{{_C}}}tx//{{{_C}}}v"):
                    ser_name = sv.text or ""
                    break

                # Category labels: try <c:cat> then <c:xVal>
                cats: list[str] = [
                    v.text for v in ser.findall(f".//{{{_C}}}cat//{{{_C}}}v")
                    if v.text
                ]
                if not cats:
                    cats = [
                        v.text for v in ser.findall(f".//{{{_C}}}xVal//{{{_C}}}v")
                        if v.text
                    ]

                # Values: try <c:val> then <c:yVal>
                vals: list[float | str] = []
                for v in ser.findall(f".//{{{_C}}}val//{{{_C}}}v"):
                    if v.text:
                        try:
                            vals.append(float(v.text))
                        except ValueError:
                            vals.append(v.text)
                if not vals:
                    for v in ser.findall(f".//{{{_C}}}yVal//{{{_C}}}v"):
                        if v.text:
                            try:
                                vals.append(float(v.text))
                            except ValueError:
                                vals.append(v.text)

                if not vals:
                    continue
                if cats and len(cats) == len(vals):
                    pairs = ", ".join(
                        f"{c}={round(v, 6) if isinstance(v, float) else v}"
                        for c, v in zip(cats, vals)
                    )
                    parts.append(f"  {ser_name}: {pairs}" if ser_name else f"  {pairs}")
                else:
                    parts.append(
                        f"  {ser_name}: {vals}" if ser_name else f"  {vals}"
                    )
            return "\n".join(parts)
        except Exception as exc:
            logger.debug("[jizhi] chart XML parse error: %s", exc)
            return ""

    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as exc:
        logger.error("[jizhi] PPTX open failed: %s", exc)
        return ""

    slide_texts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        for shape in slide.shapes:
            # Text shapes (excluding chart labels handled separately)
            if shape.has_text_frame and not shape.has_chart:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)

            # Tables
            elif shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append("TABLE:\n" + "\n".join(rows))

            # Charts
            elif shape.has_chart:
                try:
                    title = (
                        shape.chart.chart_title.text_frame.text.strip()
                        if shape.chart.has_title else ""
                    )
                except Exception:
                    title = ""
                chart_body = _chart_text(shape)
                if title or chart_body:
                    header = f"[Chart: {title}]" if title else "[Chart]"
                    parts.append(header + ("\n" + chart_body if chart_body else ""))

        if parts:
            slide_texts.append(f"=== Slide {i} ===\n" + "\n".join(parts))

    return "\n\n".join(slide_texts)


def ensure_tables(pg_url: str) -> None:
    """Create jizhi_bids, jizhi_bid_winners, jizhi_upcoming if they don't exist."""
    conn = psycopg2.connect(pg_url)
    try:
        with conn.cursor() as cur:
            cur.execute(_DDL)
        conn.commit()
        logger.info("[jizhi] tables ensured")
    finally:
        conn.close()


def extract_bids(text: str, api_key: str = "") -> list[dict]:
    """Extract structured bid results from document text via Claude tool-use.

    Tries AWS Bedrock first (no API key needed, uses IAM role).
    Falls back to direct Anthropic API if Bedrock fails and api_key is provided.
    Returns [] on failure or empty input.
    """
    if not text.strip():
        return []
    prompt = _BIDS_PROMPT.format(text=text[:12000])
    # Try Bedrock first (plain text generation — no tool use)
    _bedrock_exc: Exception | None = None
    try:
        raw = _bedrock_generate(prompt)
        bids = _parse_json_from_text(raw, "bids")
        if bids:
            return bids
        # Empty list is valid (document has no bid results)
        return []
    except Exception as exc:
        logger.warning("[jizhi] bedrock extract_bids failed: %s", exc)
        _bedrock_exc = exc
    # Fallback: direct Anthropic API
    if not api_key:
        if _bedrock_exc is not None:
            raise _bedrock_exc
        return []
    import anthropic
from shared.anthropic_client import make_client as _make_anthropic_client
    client = _make_anthropic_client(api_key)
    try:
        response = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=4096,
            tools=[_BIDS_TOOL],
            tool_choice={"type": "tool", "name": "save_bid_results"},
            messages=[{"role": "user", "content": prompt}],
        )
        for block in response.content:
            if block.type == "tool_use" and block.name == "save_bid_results":
                return block.input.get("bids", [])
    except Exception as exc:
        logger.error("[jizhi] extract_bids failed: %s", exc)
    return []


_UPCOMING_TOOL = {
    "name": "save_upcoming_bids",
    "description": "Save upcoming 机制竞价 bid announcements from a notice or web page.",
    "input_schema": {
        "type": "object",
        "properties": {
            "upcoming": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "province":            {"type": "string"},
                        "year":                {"type": "integer"},
                        "batch":               {"type": "string",  "description": "存量 / 增量_2025-12 / 增量_2026-12 / 增量_2027-12"},
                        "tech_type":           {"type": "string",  "description": "陆风 / 海风 / 光伏 / 水电"},
                        "price_floor":         {"type": "number",  "description": "元/kWh"},
                        "price_cap":           {"type": "number",  "description": "元/kWh"},
                        "target_volume_gwh":   {"type": "number",  "description": "Target volume in GWh"},
                        "supply_demand_ratio": {"type": "number"},
                        "bid_open_date":       {"type": "string",  "description": "YYYY-MM-DD"},
                        "bid_close_date":      {"type": "string",  "description": "YYYY-MM-DD"},
                        "source_url":          {"type": "string"},
                        "announcement_date":   {"type": "string",  "description": "YYYY-MM-DD"},
                        "notes":               {"type": "string"},
                    },
                    "required": ["province", "year", "batch", "tech_type"],
                },
            }
        },
        "required": ["upcoming"],
    },
}

_UPCOMING_PROMPT = """\
Extract upcoming 机制竞价 bid announcements from the text below.
Return ONLY a JSON object with key "upcoming" containing an array. No explanation — just the JSON.

Focus on: province, bidding dates, price range (元/kWh), target volume (GWh), supply-demand ratio.
batch values: 存量 / 增量_2025-12 / 增量_2026-12 / 增量_2027-12
tech_type values: 陆风 / 海风 / 光伏 / 水电
dates as YYYY-MM-DD

Each item must have: province (string), year (int), batch (string), tech_type (string).
Optional: price_floor, price_cap, target_volume_gwh, supply_demand_ratio, bid_open_date, bid_close_date, source_url, announcement_date, notes.

Example: {{"upcoming": [{{"province": "广东", "year": 2025, "batch": "增量_2025-12", "tech_type": "光伏", "bid_open_date": "2025-09-01"}}]}}

Text:
{text}"""


def extract_upcoming(text: str, api_key: str = "") -> list[dict]:
    """Extract upcoming bid announcements from text via Claude tool-use.

    Tries AWS Bedrock first. Falls back to direct Anthropic API.
    Returns [] on failure or empty input.
    """
    if not text.strip():
        return []
    prompt = _UPCOMING_PROMPT.format(text=text[:12000])
    # Try Bedrock first (plain text generation — no tool use)
    _bedrock_exc: Exception | None = None
    try:
        raw = _bedrock_generate(prompt)
        return _parse_json_from_text(raw, "upcoming")
    except Exception as exc:
        logger.warning("[jizhi] bedrock extract_upcoming failed: %s", exc)
        _bedrock_exc = exc
    # Fallback: direct Anthropic API
    if not api_key:
        if _bedrock_exc is not None:
            raise _bedrock_exc
        return []
    import anthropic
    client = _make_anthropic_client(api_key)
    try:
        response = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=2048,
            tools=[_UPCOMING_TOOL],
            tool_choice={"type": "tool", "name": "save_upcoming_bids"},
            messages=[{"role": "user", "content": prompt}],
        )
        for block in response.content:
            if block.type == "tool_use" and block.name == "save_upcoming_bids":
                return block.input.get("upcoming", [])
    except Exception as exc:
        logger.error("[jizhi] extract_upcoming failed: %s", exc)
    return []


def save_bids(records: list[dict], source_doc_id: int | None, pg_url: str) -> int:
    """Upsert bid records to staging.jizhi_bids.

    Verified rows (verified=TRUE) are never overwritten.
    Returns count of rows actually inserted or updated.
    """
    if not records or not pg_url:
        return 0
    conn = psycopg2.connect(pg_url)
    count = 0
    try:
        with conn.cursor() as cur:
            for r in records:
                cur.execute(
                    """
                    INSERT INTO staging.jizhi_bids
                        (province, year, batch, tech_type,
                         price_floor, price_cap, mechanism_type, mechanism_value,
                         supply_demand_ratio, cleared_price, cleared_volume_gwh,
                         bid_date, verified, source_doc_id, notes)
                    VALUES (%s,%s,%s,%s, %s,%s,%s,%s, %s,%s,%s, %s,FALSE,%s,%s)
                    ON CONFLICT (province, year, batch, tech_type) DO UPDATE SET
                        price_floor         = EXCLUDED.price_floor,
                        price_cap           = EXCLUDED.price_cap,
                        mechanism_type      = EXCLUDED.mechanism_type,
                        mechanism_value     = EXCLUDED.mechanism_value,
                        supply_demand_ratio = EXCLUDED.supply_demand_ratio,
                        cleared_price       = EXCLUDED.cleared_price,
                        cleared_volume_gwh  = EXCLUDED.cleared_volume_gwh,
                        bid_date            = EXCLUDED.bid_date,
                        source_doc_id       = EXCLUDED.source_doc_id,
                        notes               = EXCLUDED.notes
                    WHERE staging.jizhi_bids.verified = FALSE
                    RETURNING id
                    """,
                    (
                        r.get("province"), r.get("year"), r.get("batch"), r.get("tech_type"),
                        r.get("price_floor"), r.get("price_cap"),
                        r.get("mechanism_type"), r.get("mechanism_value"),
                        r.get("supply_demand_ratio"), r.get("cleared_price"),
                        r.get("cleared_volume_gwh"), r.get("bid_date") or None,
                        source_doc_id, r.get("notes"),
                    ),
                )
                if cur.fetchone():
                    count += 1
        conn.commit()
    except Exception as exc:
        conn.rollback()
        logger.error("[jizhi] save_bids failed: %s", exc)
    finally:
        conn.close()
    return count


def save_upcoming(records: list[dict], pg_url: str) -> int:
    """Upsert upcoming bid records to staging.jizhi_upcoming.

    Verified rows are never overwritten.
    Returns count of rows inserted or updated.
    """
    if not records or not pg_url:
        return 0
    conn = psycopg2.connect(pg_url)
    count = 0
    try:
        with conn.cursor() as cur:
            for r in records:
                cur.execute(
                    """
                    INSERT INTO staging.jizhi_upcoming
                        (province, year, batch, tech_type,
                         price_floor, price_cap, target_volume_gwh, supply_demand_ratio,
                         bid_open_date, bid_close_date, source_url, announcement_date, notes)
                    VALUES (%s,%s,%s,%s, %s,%s,%s,%s, %s,%s,%s,%s,%s)
                    ON CONFLICT (province, year, batch, tech_type, bid_open_date) DO UPDATE SET
                        price_floor         = EXCLUDED.price_floor,
                        price_cap           = EXCLUDED.price_cap,
                        target_volume_gwh   = EXCLUDED.target_volume_gwh,
                        supply_demand_ratio = EXCLUDED.supply_demand_ratio,
                        bid_close_date      = EXCLUDED.bid_close_date,
                        source_url          = EXCLUDED.source_url,
                        notes               = EXCLUDED.notes
                    WHERE staging.jizhi_upcoming.verified = FALSE
                    RETURNING id
                    """,
                    (
                        r.get("province"), r.get("year"), r.get("batch"), r.get("tech_type"),
                        r.get("price_floor"), r.get("price_cap"),
                        r.get("target_volume_gwh"), r.get("supply_demand_ratio"),
                        r.get("bid_open_date") or None, r.get("bid_close_date") or None,
                        r.get("source_url"), r.get("announcement_date") or None,
                        r.get("notes"),
                    ),
                )
                if cur.fetchone():
                    count += 1
        conn.commit()
    except Exception as exc:
        conn.rollback()
        logger.error("[jizhi] save_upcoming failed: %s", exc)
    finally:
        conn.close()
    return count
