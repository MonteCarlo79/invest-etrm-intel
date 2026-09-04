"""
Reference document knowledge base — market rules, annual reports, policy docs.

Separate from the daily spot report pipeline (staging.spot_report_* tables).
Tables: staging.spot_knowledge_docs, staging.spot_knowledge_chunks

Usage:
    from services.knowledge_pool.knowledge_docs import (
        init_knowledge_tables,
        register_and_ingest,
        search_reference_docs,
        list_knowledge_docs,
        delete_knowledge_doc,
        CATEGORY_LABELS,
    )
"""
from __future__ import annotations

import datetime as dt
import hashlib
import logging

logger = logging.getLogger(__name__)
import io
import re
from typing import Optional

from .db import get_conn


# ── File-type extraction ──────────────────────────────────────────────────────

_VISION_MIME: dict[str, str] = {
    "jpg":  "image/jpeg",
    "jpeg": "image/jpeg",
    "png":  "image/png",
    "gif":  "image/gif",
    "webp": "image/webp",
}


def _describe_image(
    image_bytes: bytes,
    mime_type: str,
    api_key: str,
    context: str = "",
) -> str:
    """
    Send an image to Claude vision and return a text description.
    Used for standalone image files and embedded images/charts in PPTX.
    """
    import base64
    import anthropic
    from shared.anthropic_client import make_client as _make_anthropic_client

    prompt = (
        "Describe this image in detail for text indexing. "
        "If it is a chart or graph, state the chart type, title, axis labels, "
        "units, data series names, and summarise the key trends or values. "
        "If it contains text or tables, transcribe them. "
        "Be thorough — the description will be stored as a searchable text chunk."
    )
    if context:
        prompt = f"Context: {context}\n\n" + prompt

    client = _make_anthropic_client(api_key)
    resp = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": mime_type,
                        "data": base64.standard_b64encode(image_bytes).decode(),
                    },
                },
                {"type": "text", "text": prompt},
            ],
        }],
    )
    return resp.content[0].text.strip()


def _extract_pages_pdf(file_bytes: bytes) -> list[tuple[int, str]]:
    """Return [(page_no, text), ...] from a PDF."""
    import pdfplumber
    pages = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            pages.append((i, page.extract_text() or ""))
    return pages


def _extract_pages_pptx(
    file_bytes: bytes,
    api_key: Optional[str] = None,
) -> list[tuple[int, str]]:
    """
    Return [(slide_no, text), ...] from a PPTX file.

    - Text shapes: extracted directly.
    - Chart shapes: data (title, series, categories) extracted as structured text.
    - Picture shapes: described via Claude vision if api_key is provided.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(file_bytes))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []

        for shape in slide.shapes:
            # ── Text frames ────────────────────────────────────────────────
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)

            # ── Charts — extract data as structured text ───────────────────
            elif shape.has_chart:
                chart = shape.chart
                try:
                    title = (
                        chart.chart_title.text_frame.text.strip()
                        if chart.has_title else "Chart"
                    )
                    parts.append(f"[Chart: {title}]")
                    # Extract category labels (e.g. province names on x-axis)
                    categories: list[str] = []
                    try:
                        xml = chart._element.xml
                        cat_blocks = re.findall(
                            r"<c:cat>.*?</c:cat>", xml, re.DOTALL
                        ) or re.findall(
                            r"<c:xVal>.*?</c:xVal>", xml, re.DOTALL
                        )
                        if cat_blocks:
                            categories = re.findall(
                                r"<c:v>([^<]+)</c:v>", cat_blocks[0]
                            )
                    except Exception:
                        pass
                    for series in chart.series:
                        try:
                            vals = list(series.values)
                            if categories and len(categories) == len(vals):
                                pairs = [
                                    f"{categories[j]}={vals[j]}"
                                    for j in range(len(vals))
                                    if vals[j] is not None
                                ]
                                parts.append(
                                    f"  Series '{series.name}': {', '.join(pairs)}"
                                )
                            else:
                                parts.append(
                                    f"  Series '{series.name}': {vals}"
                                )
                        except Exception:
                            parts.append(f"  Series '{series.name}'")
                except Exception:
                    parts.append("[Chart]")

            # ── Pictures — Claude vision description ───────────────────────
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and api_key:
                try:
                    img_bytes = shape.image.blob
                    mime = shape.image.content_type or "image/png"
                    desc = _describe_image(
                        img_bytes, mime, api_key,
                        context=f"Slide {i}",
                    )
                    parts.append(f"[Image on slide {i}]: {desc}")
                except Exception:
                    pass

        pages.append((i, "\n".join(parts)))
    return pages


def _extract_pages_image(
    file_bytes: bytes,
    filename: str,
    api_key: Optional[str] = None,
) -> list[tuple[int, str]]:
    """Describe a standalone image file via Claude vision."""
    ext = filename.rsplit(".", 1)[-1].lower()
    mime = _VISION_MIME.get(ext, "image/jpeg")
    if not api_key:
        return [(1, f"[Image: {filename} — set ANTHROPIC_API_KEY to enable vision description]")]
    desc = _describe_image(file_bytes, mime, api_key)
    return [(1, desc)]


def _extract_pages_html(file_bytes: bytes) -> list[tuple[int, str]]:
    """Strip HTML tags and split into pages of 100 lines each."""
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "html.parser")
    # Remove script/style noise
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    pages = []
    for i in range(0, max(len(lines), 1), 100):
        pages.append((i // 100 + 1, "\n".join(lines[i:i + 100])))
    return pages


def _extract_pages_txt(file_bytes: bytes) -> list[tuple[int, str]]:
    """Split plain text into pages of 100 lines each."""
    text = file_bytes.decode("utf-8", errors="replace")
    lines = text.splitlines()
    page_size = 100
    pages = []
    for i in range(0, max(len(lines), 1), page_size):
        block = "\n".join(lines[i:i + page_size]).strip()
        if block:
            pages.append((i // page_size + 1, block))
    return pages or [(1, "")]


def _extract_pages_docx(file_bytes: bytes) -> list[tuple[int, str]]:
    """Return paragraphs from a DOCX file, grouped into pages of 50 paragraphs."""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    page_size = 50
    pages = []
    for i in range(0, max(len(paras), 1), page_size):
        block = "\n".join(paras[i:i + page_size])
        if block:
            pages.append((i // page_size + 1, block))
    return pages or [(1, "")]


def _extract_pages_xlsx(file_bytes: bytes) -> list[tuple[int, str]]:
    """Return one 'page' per sheet, with each row as a tab-separated line."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    pages = []
    for sheet_no, ws in enumerate(wb.worksheets, start=1):
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                rows.append(line)
        if rows:
            pages.append((sheet_no, f"[Sheet: {ws.title}]\n" + "\n".join(rows)))
    wb.close()
    return pages or [(1, "")]


def _extract_pages_xls(file_bytes: bytes) -> list[tuple[int, str]]:
    """Return one 'page' per sheet from legacy .xls files."""
    import xlrd
    wb = xlrd.open_workbook(file_contents=file_bytes)
    pages = []
    for sheet_no in range(wb.nsheets):
        ws = wb.sheet_by_index(sheet_no)
        rows = []
        for r in range(ws.nrows):
            line = "\t".join(str(ws.cell_value(r, c)) for c in range(ws.ncols)).strip()
            if line:
                rows.append(line)
        if rows:
            pages.append((sheet_no + 1, f"[Sheet: {ws.name}]\n" + "\n".join(rows)))
    return pages or [(1, "")]


def _extract_pages(
    file_bytes: bytes,
    filename: str,
    api_key: Optional[str] = None,
) -> list[tuple[int, str]]:
    """Dispatch to the right extractor based on file extension."""
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext in ("ppt", "pptx"):
        return _extract_pages_pptx(file_bytes, api_key=api_key)
    if ext in ("htm", "html"):
        return _extract_pages_html(file_bytes)
    if ext == "txt":
        return _extract_pages_txt(file_bytes)
    if ext in ("doc", "docx"):
        return _extract_pages_docx(file_bytes)
    if ext == "xlsx":
        return _extract_pages_xlsx(file_bytes)
    if ext == "xls":
        return _extract_pages_xls(file_bytes)
    if ext in _VISION_MIME:
        return _extract_pages_image(file_bytes, filename, api_key=api_key)
    return _extract_pages_pdf(file_bytes)


# ── Category definitions ────────────────────────────────────────────────────

CATEGORIES: dict[str, list[str]] = {
    "market_rules": [
        "交易规则", "市场规则", "结算规则", "交易管理", "竞价规则", "报价规则",
        "交易细则", "市场运营规则", "现货交易规则",
        "market rule", "trading rule", "settlement rule", "bidding rule",
    ],
    "annual_report": [
        "年度报告", "年报", "运行年报", "运营报告", "年度运行", "年度总结",
        "年度回顾", "全年运行", "电力市场年度",
        "annual report", "annual operations", "annual review",
    ],
    "policy_doc": [
        "通知", "办法", "规定", "意见", "政策", "指导意见", "管理办法",
        "实施方案", "工作方案", "发改委", "能源局",
        "policy", "notice", "regulation", "directive", "guideline", "circular",
    ],
    "technical_spec": [
        "技术规范", "技术标准", "规程", "技术要求", "调度规程", "并网规范",
        "技术条件", "标准规范",
        "specification", "technical standard", "grid code", "technical requirement",
    ],
    "research_report": [
        "研究报告", "分析报告", "调研报告", "白皮书", "研究院", "研究所",
        "market analysis", "research report", "white paper",
    ],
}

CATEGORY_LABELS: dict[str, str] = {
    "market_rules":      "Market Rules",
    "annual_report":     "Annual Report",
    "policy_doc":        "Policy Document",
    "technical_spec":    "Technical Spec",
    "research_report":   "Research Report",
    "monthly_report":    "Exchange Monthly Report",
    "conversation_log":  "Conversation Log",
    "other":             "Other",
}

CATEGORY_LABELS_ZH: dict[str, str] = {
    "market_rules":      "交易规则",
    "annual_report":     "年度报告",
    "policy_doc":        "政策文件",
    "technical_spec":    "技术规范",
    "research_report":   "研究报告",
    "monthly_report":    "交易所月报",
    "conversation_log":  "对话记录",
    "other":             "其他",
}


# ── DB setup ─────────────────────────────────────────────────────────────────

_DDL = """
CREATE TABLE IF NOT EXISTS staging.spot_knowledge_docs (
    id              SERIAL PRIMARY KEY,
    file_name       TEXT NOT NULL,
    file_hash       TEXT UNIQUE NOT NULL,
    category        TEXT NOT NULL DEFAULT 'other',
    app             TEXT NOT NULL DEFAULT 'shared',
    title           TEXT,
    doc_year        INT,
    file_size_bytes INT,
    page_count      INT DEFAULT 0,
    ingest_status   TEXT NOT NULL DEFAULT 'pending',
    parse_error     TEXT,
    active          BOOLEAN DEFAULT TRUE,
    created_at      TIMESTAMPTZ DEFAULT NOW()
);

CREATE TABLE IF NOT EXISTS staging.spot_knowledge_chunks (
    id          SERIAL PRIMARY KEY,
    doc_id      INT NOT NULL REFERENCES staging.spot_knowledge_docs(id),
    page_no     INT,
    chunk_index INT NOT NULL,
    chunk_text  TEXT NOT NULL,
    UNIQUE(doc_id, chunk_index)
);

CREATE INDEX IF NOT EXISTS idx_skc_fts
    ON staging.spot_knowledge_chunks
    USING GIN(to_tsvector('simple', chunk_text));
"""


_TABLES_INITIALIZED = False


def init_knowledge_tables() -> None:
    """Create tables if they don't exist. Idempotent — skips DDL after first run."""
    global _TABLES_INITIALIZED
    if _TABLES_INITIALIZED:
        return
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(_DDL)
            # Migration: add app column for tables created before app-scoping
            cur.execute("""
                ALTER TABLE staging.spot_knowledge_docs
                ADD COLUMN IF NOT EXISTS app TEXT NOT NULL DEFAULT 'shared'
            """)
        conn.commit()
    _TABLES_INITIALIZED = True


# ── Hashing ──────────────────────────────────────────────────────────────────

def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


# ── Auto-categorization ───────────────────────────────────────────────────────

def _keyword_category(text: str) -> str:
    t = text.lower()
    for cat, keywords in CATEGORIES.items():
        for kw in keywords:
            if kw.lower() in t:
                return cat
    return "other"


def auto_categorize(
    filename: str,
    text_sample: str,
    api_key: Optional[str] = None,
) -> str:
    """
    Detect category from filename + first-page text.

    Step 1: keyword heuristic on filename + first 1000 chars of text.
    Step 2: if heuristic returns 'other' and api_key is set, ask Haiku.
    """
    combined = f"{filename}\n{text_sample[:1000]}"
    cat = _keyword_category(combined)
    if cat != "other" or not api_key:
        return cat

    # LLM fallback — Haiku is cheap and fast
    try:
        import anthropic
        client = _make_anthropic_client(api_key)
        resp = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=20,
            system=(
                "Classify this Chinese electricity market document into exactly one category. "
                "Reply with only the category key, nothing else. "
                "Categories: market_rules | annual_report | policy_doc | "
                "technical_spec | research_report | other"
            ),
            messages=[{
                "role": "user",
                "content": f"Filename: {filename}\n\nText sample:\n{text_sample[:800]}",
            }],
        )
        cat_llm = (resp.content[0].text or "other").strip().lower()
        if cat_llm in CATEGORIES or cat_llm == "other":
            return cat_llm
    except Exception:
        pass

    return "other"


# ── Text chunking ─────────────────────────────────────────────────────────────

def _chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> list[str]:
    text = text.replace("\x00", "").strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start += chunk_size - overlap
    return chunks


def _infer_title(filename: str, first_page_text: str) -> str:
    """Strip file extension and year suffixes for a clean title."""
    stem = re.sub(r"[\(\（]\d{4}[\)\）]", "", filename)
    stem = re.sub(r"\.\w+$", "", stem).strip()
    return stem or filename


# ── Registration + ingestion ──────────────────────────────────────────────────

def register_and_ingest(
    file_bytes: bytes,
    filename: str,
    category_override: Optional[str] = None,
    app: str = "shared",
    api_key: Optional[str] = None,
    synthesize: bool = True,
) -> tuple[int, bool, str]:
    """
    Register and ingest a document from raw bytes (e.g. from Streamlit uploader).

    Args:
        app: Scopes the document to an agent. One of 'shared', 'strategist', 'trader'.
             'shared' means all agents can search it.

    Returns:
        (doc_id, is_new, category)
        is_new=False means the file already existed (same SHA-256 hash).
    """
    init_knowledge_tables()

    file_hash = sha256_bytes(file_bytes)

    # Dedup check
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT id, category FROM staging.spot_knowledge_docs WHERE file_hash = %s",
                (file_hash,),
            )
            row = cur.fetchone()
    if row:
        return row[0], False, row[1]

    # Extract pages (PDF or PPTX)
    pages_text: list[tuple[int, str]] = []
    try:
        pages_text = _extract_pages(file_bytes, filename, api_key=api_key)
    except Exception as exc:
        # Register as failed so user gets feedback
        with get_conn() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    """
                    INSERT INTO staging.spot_knowledge_docs
                        (file_name, file_hash, category, app, file_size_bytes, ingest_status, parse_error)
                    VALUES (%s, %s, %s, %s, %s, 'failed', %s)
                    RETURNING id
                    """,
                    (filename, file_hash, category_override or "other",
                     app, len(file_bytes), str(exc)),
                )
                doc_id = cur.fetchone()[0]
            conn.commit()
        return doc_id, True, category_override or "other"

    first_page_text = pages_text[0][1] if pages_text else ""
    category = category_override or auto_categorize(filename, first_page_text, api_key)
    title = _infer_title(filename, first_page_text)

    # Register doc and write chunks in one transaction
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO staging.spot_knowledge_docs
                    (file_name, file_hash, category, app, title, file_size_bytes,
                     page_count, ingest_status)
                VALUES (%s, %s, %s, %s, %s, %s, %s, 'parsed')
                RETURNING id
                """,
                (filename, file_hash, category, app, title,
                 len(file_bytes), len(pages_text)),
            )
            doc_id = cur.fetchone()[0]

            chunk_index = 0
            inserts = []
            for page_no, text in pages_text:
                for chunk in _chunk_text(text):
                    inserts.append((doc_id, page_no, chunk_index, chunk))
                    chunk_index += 1

            if inserts:
                cur.executemany(
                    """
                    INSERT INTO staging.spot_knowledge_chunks
                        (doc_id, page_no, chunk_index, chunk_text)
                    VALUES (%s, %s, %s, %s)
                    ON CONFLICT (doc_id, chunk_index) DO NOTHING
                    """,
                    inserts,
                )
        conn.commit()

    # Embed chunks in background (non-blocking, best-effort)
    threading_mod = __import__("threading")
    threading_mod.Thread(
        target=_embed_chunks_for_doc,
        args=(doc_id,),
        daemon=True,
    ).start()

    # Phase 1 hook: synthesize in background (best-effort, non-blocking)
    # Disabled during bulk ingestion (synthesize=False) to avoid burst-limit 403s;
    # the ECS synthesis task handles those docs instead.
    if api_key and synthesize:
        try:
            import threading
            from .synthesis import synthesize_on_ingest
            threading.Thread(
                target=synthesize_on_ingest,
                args=(doc_id, api_key),
                daemon=True,
            ).start()
        except Exception:
            pass

    return doc_id, True, category


def _embed_chunks_for_doc(doc_id: int) -> None:
    """Embed all un-embedded chunks for a document. Runs in a background thread."""
    try:
        from .embeddings import embed_texts, vec_to_pg
    except ImportError:
        return  # fastembed not installed — skip silently

    try:
        with get_conn() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    "SELECT id, chunk_text FROM staging.spot_knowledge_chunks "
                    "WHERE doc_id = %s AND embedding IS NULL",
                    (doc_id,),
                )
                rows = cur.fetchall()

        if not rows:
            return

        ids = [r[0] for r in rows]
        texts = [r[1] for r in rows]
        BATCH = 64
        for i in range(0, len(texts), BATCH):
            batch_ids = ids[i:i + BATCH]
            batch_texts = texts[i:i + BATCH]
            vecs = embed_texts(batch_texts)
            updates = [
                (vec_to_pg(v), cid)
                for v, cid in zip(vecs, batch_ids)
                if v is not None
            ]
            if updates:
                with get_conn() as conn:
                    with conn.cursor() as cur:
                        cur.executemany(
                            "UPDATE staging.spot_knowledge_chunks "
                            "SET embedding = %s::vector WHERE id = %s",
                            updates,
                        )
                    conn.commit()

        logger.info("Embedded %d chunks for doc_id=%d", len(rows), doc_id)
    except Exception as exc:
        logger.error("Chunk embedding failed for doc_id=%d: %s", doc_id, exc)


# Browser-like headers for mp.weixin.qq.com — the bare bot UA gets served an
# anti-bot challenge page instead of the article (seen since the 2026-08-30
# NAT-EIP cutover). Mirrors services/hermes/news_screener._WECHAT_HEADERS.
_WECHAT_FETCH_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (iPhone; CPU iPhone OS 17_0 like Mac OS X) "
        "AppleWebKit/605.1.15 (KHTML, like Gecko) "
        "Version/17.0 Mobile/15E148 Safari/604.1"
    ),
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
    "Referer": "https://mp.weixin.qq.com/",
}

_BOT_CHALLENGE_MARKERS = (
    "环境异常", "完成验证后即可继续访问", "请输入验证码", "操作频繁", "访问过于频繁",
)


def _is_challenge_page(html: str) -> bool:
    """True if the HTML is an anti-bot verification page, not real content."""
    head = html[:3000]
    return any(m in head for m in _BOT_CHALLENGE_MARKERS)


def register_url(
    url: str,
    api_key: Optional[str] = None,
) -> tuple[int, bool, str]:
    """
    Fetch a public URL, extract its text, and register as a knowledge doc.

    Returns (doc_id, is_new, category).  is_new=False if the URL was already ingested.
    Raises on network / parse errors so the caller can show a user-facing message.
    Raises ValueError when the page is an anti-bot challenge (e.g. WeChat 环境异常)
    so junk verification text is never ingested as document content.
    """
    import requests
    from bs4 import BeautifulSoup

    is_wechat = "mp.weixin.qq.com" in url
    headers = dict(_WECHAT_FETCH_HEADERS) if is_wechat else {
        "User-Agent": "Mozilla/5.0 (compatible; SpotMarketBot/1.0)"
    }
    resp = requests.get(url, headers=headers, timeout=30)
    resp.raise_for_status()

    if _is_challenge_page(resp.text):
        raise ValueError(
            "页面被反爬拦截（显示环境异常/需要验证），未获取到文章正文。"
            "请在微信中打开文章，复制正文后直接粘贴给我。"
        )

    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "nav", "header", "footer", "aside"]):
        tag.decompose()

    if is_wechat:
        title_el = (
            soup.find("h1", id="activity-name")
            or soup.find("h2", class_="rich_media_title")
            or soup.find("title")
        )
        content_div = (
            soup.find("div", id="js_content")
            or soup.find("div", class_="rich_media_content")
        )
        text = content_div.get_text(separator="\n", strip=True) if content_div else ""
    else:
        title_el = soup.find("h1") or soup.find("title")
        text = soup.get_text(separator="\n", strip=True)

    title = title_el.get_text(strip=True)[:200] if title_el else url.split("/")[-1][:100]
    if not text.strip():
        raise ValueError(
            "未能从该页面提取到文章正文（可能遭遇反爬验证）。"
            "请复制文章正文后直接粘贴给我。"
        )

    # Use SHA-256 of the URL as the dedup key (allows re-fetch to update)
    url_hash = sha256_bytes(url.encode())

    init_knowledge_tables()

    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT id, category FROM staging.spot_knowledge_docs WHERE file_hash = %s",
                (url_hash,),
            )
            row = cur.fetchone()
    if row:
        return row[0], False, row[1]

    category = auto_categorize(url.split("/")[-1], text[:1000], api_key)

    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO staging.spot_knowledge_docs
                    (file_name, file_hash, category, app, title,
                     file_size_bytes, page_count, ingest_status)
                VALUES (%s, %s, %s, 'shared', %s, %s, 1, 'parsed')
                RETURNING id
                """,
                (url[:255], url_hash, category, title, len(text.encode())),
            )
            doc_id = cur.fetchone()[0]

            inserts = [(doc_id, 1, i, chunk) for i, chunk in enumerate(_chunk_text(text))]
            if inserts:
                cur.executemany(
                    """
                    INSERT INTO staging.spot_knowledge_chunks
                        (doc_id, page_no, chunk_index, chunk_text)
                    VALUES (%s, %s, %s, %s)
                    ON CONFLICT (doc_id, chunk_index) DO NOTHING
                    """,
                    inserts,
                )
        conn.commit()

    return doc_id, True, category


# ── Retrieval ─────────────────────────────────────────────────────────────────

def _has_cjk(text: str) -> bool:
    """Return True if text contains CJK (Chinese/Japanese/Korean) characters."""
    return bool(re.search(r'[\u4e00-\u9fff\u3400-\u4dbf]', text))


def _cjk_bigrams(text: str, max_terms: int = 12) -> list[str]:
    """Extract overlapping 2-character bigrams from all CJK runs in text.

    PostgreSQL 'simple' text search cannot tokenise Chinese (no word boundaries),
    so we break the query into bigrams and use ILIKE for each.  This gives good
    recall for Chinese knowledge-base search without requiring pg_trgm or a
    dedicated CJK dictionary.

    Example:
        "中长期合同情况" → ["中长", "长期", "期合", "合同", "同情", "情况"]
    """
    bigrams: list[str] = []
    seen: set[str] = set()
    for run in re.findall(r'[\u4e00-\u9fff\u3400-\u4dbf]+', text):
        if len(run) == 1:
            if run not in seen:
                bigrams.append(run)
                seen.add(run)
        else:
            for i in range(len(run) - 1):
                bg = run[i:i + 2]
                if bg not in seen:
                    bigrams.append(bg)
                    seen.add(bg)
    return bigrams[:max_terms]


def search_reference_docs(
    query: str,
    category: Optional[str] = None,
    app: Optional[str] = None,
    limit: int = 5,
    filename_contains: Optional[tuple[str, ...]] = None,
) -> list[dict]:
    """
    Full-text search over staging.spot_knowledge_chunks.

    For Latin/English queries (>4 chars): PostgreSQL FTS with 'simple' config.
    For CJK queries: bigram ILIKE search — each 2-char bigram is an ILIKE
        condition; rank = number of bigrams matched, so the most-relevant
        chunks float to the top.  This handles Chinese without needing a CJK
        PostgreSQL dictionary or pg_trgm extension.

    Args:
        app: When set, returns docs where app = :app OR app = 'shared'.
             Pass 'strategist' or 'trader' to exclude the other agent's
             private documents.  Omit (None) to search all docs.
        filename_contains: When set, AND a file-name condition into the WHERE
             clause: (d.file_name ILIKE %term1% OR d.file_name ILIKE %term2% ...).
             Opt-in — no behavior change for existing callers.

    Returns list of dicts:
        doc_id, file_name, category, app, page_no, chunk_text, rank
    """
    init_knowledge_tables()

    conditions = ["d.active = TRUE"]
    params: list = []

    if _has_cjk(query):
        # CJK query: bigram ILIKE OR — matches any chunk containing at least one bigram.
        # Rank = sum of matched bigrams (higher = more relevant).
        bigrams = _cjk_bigrams(query)
        if bigrams:
            ilike_conds = " OR ".join("c.chunk_text ILIKE %s" for _ in bigrams)
            conditions.append(f"({ilike_conds})")
            params.extend(f"%{bg}%" for bg in bigrams)
            # Rank expression: count of bigrams matched
            case_parts = " + ".join(
                "(CASE WHEN c.chunk_text ILIKE %s THEN 1 ELSE 0 END)"
                for _ in bigrams
            )
            rank_expr = f"({case_parts})::float"
            params.extend(f"%{bg}%" for bg in bigrams)
        else:
            # Fallback: full ILIKE on the raw query (single-char or punctuation-only)
            conditions.append("c.chunk_text ILIKE %s")
            params.append(f"%{query}%")
            rank_expr = "1.0::float"
    elif len(query) <= 4:
        conditions.append("c.chunk_text ILIKE %s")
        params.append(f"%{query}%")
        rank_expr = "1.0::float"
    else:
        # Latin/numeric query: PostgreSQL FTS
        conditions.append(
            "to_tsvector('simple', c.chunk_text) @@ plainto_tsquery('simple', %s)"
        )
        params.append(query)
        rank_expr = (
            "ts_rank(to_tsvector('simple', c.chunk_text), "
            "plainto_tsquery('simple', %s))"
        )
        params.append(query)

    if category:
        conditions.append("d.category = %s")
        params.append(category)

    if app:
        conditions.append("(d.app = %s OR d.app = 'shared')")
        params.append(app)

    if filename_contains:
        fn_conds = " OR ".join("d.file_name ILIKE %s" for _ in filename_contains)
        conditions.append(f"({fn_conds})")
        params.extend(f"%{t}%" for t in filename_contains)

    where = " AND ".join(conditions)
    sql = f"""
        SELECT d.id AS doc_id, d.file_name, d.category, d.app,
               c.page_no, c.chunk_text,
               {rank_expr} AS rank
        FROM staging.spot_knowledge_chunks c
        JOIN staging.spot_knowledge_docs d ON d.id = c.doc_id
        WHERE {where}
        ORDER BY rank DESC
        LIMIT %s
    """
    params.append(limit)

    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(sql, params)
            cols = [d[0] for d in cur.description]
            return [dict(zip(cols, r)) for r in cur.fetchall()]


def vector_search_reference_docs(
    query: str,
    app: Optional[str] = None,
    limit: int = 10,
) -> list[dict]:
    """
    Vector similarity search over staging.spot_knowledge_chunks.
    Returns chunks ordered by cosine similarity to the query embedding.
    Falls back to empty list if embeddings not available.
    """
    try:
        from .embeddings import embed_one, vec_to_pg
    except ImportError:
        return []

    qvec = embed_one(query)
    if qvec is None:
        return []

    qvec_str = vec_to_pg(qvec)
    conditions = ["d.active = TRUE", "c.embedding IS NOT NULL"]
    params: list = [qvec_str]

    if app:
        conditions.append("(d.app = %s OR d.app = 'shared')")
        params.append(app)

    where = " AND ".join(conditions)
    params.extend([qvec_str, limit])

    sql = f"""
        SELECT d.id AS doc_id, d.file_name, d.category, d.app,
               c.page_no, c.chunk_text,
               1 - (c.embedding <=> %s::vector) AS rank
        FROM staging.spot_knowledge_chunks c
        JOIN staging.spot_knowledge_docs d ON d.id = c.doc_id
        WHERE {where}
        ORDER BY c.embedding <=> %s::vector
        LIMIT %s
    """

    try:
        init_knowledge_tables()
        with get_conn() as conn:
            with conn.cursor() as cur:
                cur.execute(sql, params)
                cols = [d[0] for d in cur.description]
                return [dict(zip(cols, r)) for r in cur.fetchall()]
    except Exception as exc:
        logger.error("Vector search failed: %s", exc)
        return []


# ── Doc management ────────────────────────────────────────────────────────────

def list_knowledge_docs() -> list[dict]:
    """Return all active knowledge docs ordered by most recently added."""
    init_knowledge_tables()
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT id, file_name, category, title, page_count,
                       ingest_status, created_at
                FROM staging.spot_knowledge_docs
                WHERE active = TRUE
                ORDER BY created_at DESC
                """,
            )
            cols = [d[0] for d in cur.description]
            return [dict(zip(cols, r)) for r in cur.fetchall()]


def delete_knowledge_doc(doc_id: int) -> None:
    """Soft-delete a document and its chunks remain but are excluded from queries."""
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE staging.spot_knowledge_docs SET active=FALSE WHERE id=%s",
                (doc_id,),
            )
        conn.commit()


# ── Conversation logging ──────────────────────────────────────────────────────

def log_conversation_turn(user_msg: str, agent_reply: str) -> None:
    """
    Append a Q&A turn to today's conversation log document.

    One document per calendar day (keyed by filename); each turn is one chunk.
    Logs are searchable by the agent via search_reference_docs(category='conversation_log').
    """
    init_knowledge_tables()

    today = dt.date.today().isoformat()
    file_name = f"conversation_log_{today}.md"
    # Stable per-day hash — same day always resolves to the same doc
    file_hash = sha256_bytes(f"__conv_log__{today}".encode())

    with get_conn() as conn:
        with conn.cursor() as cur:
            # Ensure today's log doc exists (idempotent)
            cur.execute(
                """
                INSERT INTO staging.spot_knowledge_docs
                    (file_name, file_hash, category, title, ingest_status)
                VALUES (%s, %s, 'conversation_log', %s, 'parsed')
                ON CONFLICT (file_hash) DO NOTHING
                """,
                (file_name, file_hash, f"Agent Conversation Log {today}"),
            )
            cur.execute(
                "SELECT id FROM staging.spot_knowledge_docs WHERE file_hash = %s",
                (file_hash,),
            )
            doc_id = cur.fetchone()[0]

            # Next available chunk index for this doc
            cur.execute(
                "SELECT COALESCE(MAX(chunk_index), -1) + 1 "
                "FROM staging.spot_knowledge_chunks WHERE doc_id = %s",
                (doc_id,),
            )
            next_idx = cur.fetchone()[0]

            chunk_text = (
                f"[User]: {user_msg}\n\n"
                f"[Agent]: {agent_reply}"
            )
            cur.execute(
                """
                INSERT INTO staging.spot_knowledge_chunks
                    (doc_id, page_no, chunk_index, chunk_text)
                VALUES (%s, 1, %s, %s)
                ON CONFLICT (doc_id, chunk_index) DO NOTHING
                """,
                (doc_id, next_idx, chunk_text),
            )
        conn.commit()
