"""Export helpers for Investment Advisor conversations.

Supported formats:
  - PDF  (reportlab)
  - PPTX (python-pptx)
  - DOCX (python-docx)

Usage:
    from services.intl_market_common.export_helpers import export_pdf, export_pptx, export_docx
    pdf_bytes  = export_pdf(history, title, market_name)
    pptx_bytes = export_pptx(history, title, market_name)
    docx_bytes = export_docx(history, title, market_name)

Each function returns bytes that can be passed directly to st.download_button.
``history`` is a list of dicts with keys ``role`` ("user"/"assistant") and ``content`` (str).
"""
from __future__ import annotations

import io
from datetime import date


# ── PDF ────────────────────────────────────────────────────────────────────────

def export_pdf(history: list[dict], title: str, market_name: str) -> bytes:
    """Generate a PDF report from the conversation history."""
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        HRFlowable, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=2.5 * cm,
        rightMargin=2.5 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )

    styles = getSampleStyleSheet()
    # Custom styles
    title_style = ParagraphStyle(
        "AdvisoryTitle",
        parent=styles["Title"],
        fontSize=18,
        spaceAfter=6,
        textColor=colors.HexColor("#1a365d"),
    )
    subtitle_style = ParagraphStyle(
        "AdvisorySubtitle",
        parent=styles["Normal"],
        fontSize=10,
        spaceAfter=14,
        textColor=colors.HexColor("#4a5568"),
        alignment=TA_CENTER,
    )
    q_style = ParagraphStyle(
        "Question",
        parent=styles["Normal"],
        fontSize=11,
        fontName="Helvetica-Bold",
        spaceAfter=4,
        spaceBefore=14,
        textColor=colors.HexColor("#2c5282"),
        leftIndent=0,
    )
    a_style = ParagraphStyle(
        "Answer",
        parent=styles["Normal"],
        fontSize=10,
        leading=15,
        spaceAfter=6,
        leftIndent=12,
        textColor=colors.HexColor("#2d3748"),
    )

    story = []

    # Cover header
    story.append(Paragraph(title, title_style))
    story.append(Paragraph(
        f"{market_name} Investment Advisory · Generated {date.today().strftime('%B %d, %Y')}",
        subtitle_style,
    ))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#2c5282")))
    story.append(Spacer(1, 0.4 * cm))

    q_num = 0
    for msg in history:
        role = msg.get("role", "")
        content = (msg.get("content") or "").strip()
        if not content:
            continue

        # Escape special XML chars for reportlab
        safe = content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        if role == "user":
            q_num += 1
            story.append(Paragraph(f"Q{q_num}. {safe}", q_style))
        elif role == "assistant":
            # Split on double-newlines for paragraph breaks
            paras = [p.strip() for p in content.split("\n\n") if p.strip()]
            for p in paras:
                safe_p = p.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                # Convert markdown bold (**text**) → <b>text</b>
                import re
                safe_p = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", safe_p)
                story.append(Paragraph(safe_p, a_style))

    if not story:
        story.append(Paragraph("No conversation to export.", styles["Normal"]))

    doc.build(story)
    return buf.getvalue()


# ── PPTX ───────────────────────────────────────────────────────────────────────

def export_pptx(history: list[dict], title: str, market_name: str) -> bytes:
    """Generate a PowerPoint presentation from the conversation history.

    Layout: title slide + one slide per Q&A pair.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colour palette
    DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
    MID_BLUE  = RGBColor(0x2C, 0x52, 0x82)
    LIGHT_GREY = RGBColor(0xF7, 0xFA, 0xFC)
    TEXT_DARK = RGBColor(0x2D, 0x37, 0x48)

    def _set_bg(slide, color: RGBColor):
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(slide, text: str, left, top, width, height, font_size: int,
                      bold=False, color=TEXT_DARK, word_wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    # ── Title slide ───────────────────────────────────────────────
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, DARK_BLUE)

    _add_text_box(slide, market_name.upper(),
                  Inches(1), Inches(1.8), Inches(11), Inches(0.8),
                  font_size=14, bold=False, color=RGBColor(0xA0, 0xBC, 0xD8))
    _add_text_box(slide, title,
                  Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                  font_size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text_box(slide, f"Generated {date.today().strftime('%B %d, %Y')}",
                  Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                  font_size=11, color=RGBColor(0xA0, 0xBC, 0xD8))

    # ── Q&A slides ────────────────────────────────────────────────
    pairs: list[tuple[str, str]] = []
    current_q = ""
    for msg in history:
        role = msg.get("role", "")
        content = (msg.get("content") or "").strip()
        if role == "user":
            current_q = content
        elif role == "assistant" and current_q:
            pairs.append((current_q, content))
            current_q = ""

    for idx, (q, a) in enumerate(pairs, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, LIGHT_GREY)

        # Q number badge
        _add_text_box(slide, f"Q{idx}", Inches(0.5), Inches(0.35), Inches(0.7), Inches(0.55),
                      font_size=14, bold=True, color=MID_BLUE)
        # Question
        q_short = (q[:180] + "…") if len(q) > 180 else q
        _add_text_box(slide, q_short, Inches(1.2), Inches(0.3), Inches(11.5), Inches(0.9),
                      font_size=13, bold=True, color=DARK_BLUE)

        # Horizontal rule via thin rectangle
        from pptx.util import Emu
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0.5), Inches(1.25), Inches(12.3), Emu(18000),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = MID_BLUE
        line.line.fill.background()

        # Answer (truncate to fit slide)
        a_display = a[:900] + ("…" if len(a) > 900 else "")
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = a_display
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT_DARK

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── DOCX ───────────────────────────────────────────────────────────────────────

def export_docx(history: list[dict], title: str, market_name: str) -> bytes:
    """Generate a Word document from the conversation history."""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    # Title
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1A, 0x36, 0x5D)

    # Subtitle
    sub = doc.add_paragraph(
        f"{market_name} Investment Advisory  ·  Generated {date.today().strftime('%B %d, %Y')}"
    )
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].font.size = Pt(10)
    sub.runs[0].font.color.rgb = RGBColor(0x4A, 0x55, 0x68)

    doc.add_paragraph()  # spacer

    q_num = 0
    for msg in history:
        role = msg.get("role", "")
        content = (msg.get("content") or "").strip()
        if not content:
            continue

        if role == "user":
            q_num += 1
            h2 = doc.add_heading(f"Q{q_num}. {content}", level=2)
            for run in h2.runs:
                run.font.color.rgb = RGBColor(0x2C, 0x52, 0x82)

        elif role == "assistant":
            paras = [p.strip() for p in content.split("\n\n") if p.strip()]
            for para_text in paras:
                p = doc.add_paragraph(para_text)
                p.paragraph_format.space_before = Pt(0)
                p.paragraph_format.space_after = Pt(4)
            doc.add_paragraph()  # spacer between answers

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()
