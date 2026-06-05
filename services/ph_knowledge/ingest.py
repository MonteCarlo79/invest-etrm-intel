"""Philippines knowledge base ingestion orchestrator.

Sources:
  - Local reports (PDF/Excel/PPTX from data/market-fundamentals-ph/)
  - DOE Philippines news (doe.gov.ph)
  - IEMOP market bulletins (iemop.ph)

Usage:
    python -m services.ph_knowledge.ingest             # all sources
    python -m services.ph_knowledge.ingest --only local_reports
    python -m services.ph_knowledge.ingest --only doe_news
    python -m services.ph_knowledge.ingest --only iemop_notices
"""
from __future__ import annotations

import argparse
import io
import logging
import os
import pathlib
import sys
from datetime import date

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

logger = logging.getLogger(__name__)

# Path to the local market fundamentals directory (relative to repo root)
_LOCAL_REPORTS_DIR = pathlib.Path(__file__).parent.parent.parent / "data" / "market-fundamentals-ph"


# ── Connectors ────────────────────────────────────────────────────────────────

class LocalReportsConnector:
    """Ingests all PDF, Excel, and PPTX reports from the local market-fundamentals-ph directory."""

    source = "local_reports"

    def fetch(self) -> list[dict]:
        docs = []
        if not _LOCAL_REPORTS_DIR.exists():
            logger.warning("[local_reports] Directory not found: %s", _LOCAL_REPORTS_DIR)
            return []

        for path in sorted(_LOCAL_REPORTS_DIR.rglob("*")):
            if not path.is_file():
                continue
            ext = path.suffix.lower()
            if ext not in (".pdf", ".xlsx", ".xls", ".pptx", ".ppt", ".txt"):
                continue

            url_key = f"local://{path.relative_to(_LOCAL_REPORTS_DIR.parent.parent).as_posix()}"
            title = path.stem.replace("_", " ").replace("-", " ")

            try:
                content = _extract_text(path, ext)
            except Exception as exc:
                logger.warning("[local_reports] Skipping %s: %s", path.name, exc)
                continue

            if not content or len(content.strip()) < 100:
                logger.debug("[local_reports] Skipping %s: too little text", path.name)
                continue

            docs.append({
                "doc_type": _doc_type(ext),
                "title": title[:250],
                "url": url_key,
                "published_date": _infer_date(path.stem),
                "content": content[:100_000],
            })
            logger.debug("[local_reports] Extracted %s (%d chars)", path.name, len(content))

        return docs


def _doc_type(ext: str) -> str:
    return {
        ".pdf": "pdf",
        ".xlsx": "excel",
        ".xls": "excel",
        ".pptx": "presentation",
        ".ppt": "presentation",
        ".txt": "text",
    }.get(ext, "document")


def _infer_date(stem: str) -> date | None:
    """Try to extract a date from the filename stem (e.g. '2025-05-28_Report')."""
    import re
    m = re.search(r"(\d{4})[-_](\d{2})[-_](\d{2})", stem)
    if m:
        try:
            return date(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        except ValueError:
            pass
    m = re.search(r"(\d{4})[-_](\d{2})", stem)
    if m:
        try:
            return date(int(m.group(1)), int(m.group(2)), 1)
        except ValueError:
            pass
    return None


def _extract_text(path: pathlib.Path, ext: str) -> str:
    """Extract plain text from a file."""
    if ext == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
            return "\n\n".join(pages)
        except Exception:
            pass
        # Fallback: pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                return "\n\n".join(
                    (p.extract_text() or "") for p in pdf.pages
                )
        except Exception as exc:
            raise RuntimeError(f"PDF extraction failed: {exc}") from exc

    if ext in (".xlsx", ".xls"):
        import pandas as pd
        xl = pd.ExcelFile(str(path))
        parts = []
        for sheet in xl.sheet_names:
            try:
                df = xl.parse(sheet)
                if not df.empty:
                    parts.append(f"[Sheet: {sheet}]\n{df.to_string(index=False)}")
            except Exception:
                pass
        return "\n\n".join(parts)

    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            from pptx.util import Emu
            prs = Presentation(str(path))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides)
        except Exception as exc:
            raise RuntimeError(f"PPTX extraction failed: {exc}") from exc

    return ""


class DOENewsConnector:
    """Scrapes recent news from the Philippine Department of Energy website."""

    source = "doe_ph"

    _BASE_URL = "https://www.doe.gov.ph/category/news-advisories"

    def fetch(self) -> list[dict]:
        try:
            import requests
            from bs4 import BeautifulSoup
        except ImportError:
            logger.warning("[doe_ph] requests/beautifulsoup4 not installed")
            return []

        docs = []
        try:
            resp = requests.get(
                self._BASE_URL,
                headers={"User-Agent": "BESSPlatformBot/1.0 (investment research)"},
                timeout=30,
            )
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "html.parser")

            # Find article links
            articles = soup.find_all("article") or soup.find_all("h2", class_="entry-title")
            for article in articles[:20]:
                link = article.find("a") if article.name != "a" else article
                if not link:
                    continue
                href = link.get("href", "")
                title = link.get_text(" ", strip=True)
                if not href or not title:
                    continue

                try:
                    art_resp = requests.get(
                        href,
                        headers={"User-Agent": "BESSPlatformBot/1.0"},
                        timeout=15,
                    )
                    art_soup = BeautifulSoup(art_resp.text, "html.parser")
                    for tag in art_soup.find_all(["script", "style", "nav", "header", "footer"]):
                        tag.decompose()
                    content = "\n\n".join(
                        el.get_text(" ", strip=True)
                        for el in art_soup.find_all(["p", "h1", "h2", "h3", "li"])
                        if el.get_text(strip=True)
                    )
                    if not content.strip():
                        continue
                    docs.append({
                        "doc_type": "news",
                        "title": title[:250],
                        "url": href,
                        "published_date": date.today(),
                        "content": f"DOE Philippines — {title}\n\n{content[:5000]}",
                    })
                except Exception:
                    pass
        except Exception as exc:
            logger.warning("[doe_ph] Fetch failed: %s", exc)

        return docs


class IEMOPNoticesConnector:
    """Fetches IEMOP (market operator) public bulletins and market notices."""

    source = "iemop"

    _BULLETINS_URL = "https://www.iemop.ph/market-operations/market-bulletins/"

    def fetch(self) -> list[dict]:
        try:
            import requests
            from bs4 import BeautifulSoup
        except ImportError:
            logger.warning("[iemop] requests/beautifulsoup4 not installed")
            return []

        docs = []
        try:
            resp = requests.get(
                self._BULLETINS_URL,
                headers={"User-Agent": "BESSPlatformBot/1.0"},
                timeout=30,
            )
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "html.parser")

            for link in soup.find_all("a", href=True)[:30]:
                href = link["href"]
                text = link.get_text(" ", strip=True)
                if not text or len(text) < 10:
                    continue
                # Only ingest links that look like bulletins
                if not any(kw in text.lower() for kw in ["bulletin", "notice", "advisory", "market", "wesm"]):
                    continue
                if not href.startswith("http"):
                    href = "https://www.iemop.ph" + href if href.startswith("/") else href

                try:
                    art_resp = requests.get(href, headers={"User-Agent": "BESSPlatformBot/1.0"}, timeout=15)
                    art_soup = BeautifulSoup(art_resp.text, "html.parser")
                    for tag in art_soup.find_all(["script", "style"]):
                        tag.decompose()
                    content = art_soup.get_text(" ", strip=True)
                    if len(content.strip()) < 100:
                        continue
                    docs.append({
                        "doc_type": "market_notice",
                        "title": text[:250],
                        "url": href,
                        "published_date": date.today(),
                        "content": f"IEMOP — {text}\n\n{content[:5000]}",
                    })
                except Exception:
                    pass
        except Exception as exc:
            logger.warning("[iemop] Fetch failed: %s", exc)

        return docs


# ── Orchestrator ──────────────────────────────────────────────────────────────

def run_knowledge_ingest(only: list[str] | None = None, verbose: bool = True) -> dict[str, int]:
    from dotenv import load_dotenv
    load_dotenv(
        os.path.join(os.path.dirname(__file__), "..", "..", "config", ".env"),
        override=False,
    )

    import psycopg2
    from services.ph_knowledge.config import MARKET_CONFIG

    conn = psycopg2.connect(os.environ["PGURL"], keepalives=1, keepalives_idle=30)
    prefix = MARKET_CONFIG.table_prefix

    with conn.cursor() as cur:
        cur.execute(f"""
            CREATE TABLE IF NOT EXISTS intl_market.{prefix}knowledge_docs (
                id              SERIAL PRIMARY KEY,
                source          TEXT NOT NULL,
                doc_type        TEXT NOT NULL,
                title           TEXT,
                url             TEXT UNIQUE,
                published_date  DATE,
                content         TEXT NOT NULL,
                fetched_at      TIMESTAMPTZ NOT NULL DEFAULT NOW(),
                search_vector   TSVECTOR GENERATED ALWAYS AS (
                    to_tsvector('english',
                        coalesce(title,'') || ' ' || left(content,100000))
                ) STORED
            )
        """)
        cur.execute(
            f"CREATE INDEX IF NOT EXISTS {prefix}knowledge_docs_fts "
            f"ON intl_market.{prefix}knowledge_docs USING GIN(search_vector)"
        )
    conn.commit()

    from services.ph_knowledge.wesm_scraper import (
        WESMReportConnector,
        WESMMarketWatchConnector,
        WESMMarketAssessmentConnector,
        WESMRetailAssessmentConnector,
        WESMOverridingConstraintsConnector,
        IEMOPKnowledgeCenterConnector,
        WESMBESSStudyConnector,
        IEMOPMarketDataConnector,
    )
    connectors = [
        ("local_reports",                "Local market reports (PDF/Excel/PPTX)",    LocalReportsConnector()),
        ("doe_news",                     "DOE Philippines news",                      DOENewsConnector()),
        ("iemop_notices",                "IEMOP market bulletins",                    IEMOPNoticesConnector()),
        ("wesm_reports",                 "WESM/IEMOP market reports & advisories",   WESMReportConnector()),
        ("wesm_market_watch",            "WESM Market Watch (weekly)",               WESMMarketWatchConnector()),
        ("wesm_assessment",              "WESM Market Assessment reports",            WESMMarketAssessmentConnector()),
        ("wesm_retail_assessment",       "WESM Retail Market Assessment",            WESMRetailAssessmentConnector()),
        ("wesm_overriding_constraints",  "WESM Over-riding Constraints",             WESMOverridingConstraintsConnector()),
        ("iemop_knowledge_center",       "IEMOP Knowledge Center",                   IEMOPKnowledgeCenterConnector()),
        ("wesm_bess_study",              "WESM BESS Study Reports",                  WESMBESSStudyConnector()),
        ("iemop_market_data",            "IEMOP Market Data publications",           IEMOPMarketDataConnector()),
    ]

    results: dict[str, int] = {}
    for key, label, connector in connectors:
        if only and key not in only:
            continue
        if verbose:
            print(f"  [{key}] {label}…", end="", flush=True)
        try:
            n = _run_connector(connector, conn, prefix)
            results[key] = n
            if verbose:
                print(f" {n} new docs")
        except Exception as exc:
            results[key] = 0
            if verbose:
                print(f" ERROR: {exc}")
            logger.error("[ph_ingest:%s] %s", key, exc)

    conn.close()
    return results


def _run_connector(connector, conn, prefix: str) -> int:
    n = 0
    for doc in connector.fetch():
        with conn.cursor() as cur:
            cur.execute(
                f"INSERT INTO intl_market.{prefix}knowledge_docs "
                "(source, doc_type, title, url, published_date, content) "
                "VALUES (%s, %s, %s, %s, %s, %s) "
                "ON CONFLICT (url) DO UPDATE SET "
                "  content=EXCLUDED.content, title=EXCLUDED.title, fetched_at=NOW()",
                (
                    getattr(connector, "source", "unknown"),
                    doc["doc_type"],
                    doc.get("title", ""),
                    doc.get("url"),
                    doc.get("published_date"),
                    doc["content"],
                ),
            )
            if cur.rowcount > 0:
                n += 1
    conn.commit()
    return n


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    parser = argparse.ArgumentParser()
    parser.add_argument("--only", help="Comma-separated connectors to run")
    args = parser.parse_args()
    only = args.only.split(",") if args.only else None
    results = run_knowledge_ingest(only=only)
    print("Results:", results)
